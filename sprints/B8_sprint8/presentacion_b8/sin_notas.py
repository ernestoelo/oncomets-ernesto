#!/usr/bin/env python
"""Copia del deck B8 SIN notas del presentador, para mandarle a Sebastián.

La pidió el 6-ago para mirar los mapas de calor y los parches sin el guion encima.

Es CIRUGÍA DE ZIP, no python-pptx ([[pptx-quitar-notas-y-respaldo]]): se reescribe el
paquete OOXML sacando `ppt/notesSlides/*`, la `<Relationship>` de tipo notesSlide de cada
`ppt/slides/_rels/slideN.xml.rels`, y el `<Override>` correspondiente de
`[Content_Types].xml`. El `notesMaster` se conserva, que es lo que hace un deck sin notas.

Por qué no python-pptx: su round-trip reserializa el paquete entero, y este deck se
construye sobre un template que EMBEBE sus fuentes ([[deck-template-fuentes-embebidas]]).
Con la cirugía de zip todo lo que no son notas queda byte-idéntico, y el script lo cuenta
al terminar para que se vea.

Las notas quedan ELIMINADAS del archivo, no vacías: no se recuperan abriéndolo.

Uso (después de correr generate_b8_deck.py):
  PYTHONPATH=/media/administrador/Storage1/sdonoso/clam_testing2/.pylibs \
  /home/sdonoso/miniconda3/envs/clam_latest/bin/python \
    sprints/B8_sprint8/presentacion_b8/sin_notas.py
"""
import os
import re
import zipfile

AQUI = os.path.dirname(os.path.abspath(__file__))
SRC = os.path.join(AQUI, "CLAM_Sprint8.pptx")
DST = os.path.join(AQUI, "CLAM_Sprint8_sin_notas.pptx")

REL = re.compile(rb'<Relationship[^>]*Type="[^"]*/notesSlide"[^>]*/>')
OVR = re.compile(rb'<Override[^>]*PartName="/ppt/notesSlides/[^"]*"[^>]*/>')
RELS = re.compile(r"ppt/slides/_rels/slide\d+\.xml\.rels")


def main():
    # El original se REGENERA con generate_b8_deck.py, así que no hace falta respaldarlo:
    # la fuente de verdad es el generador, que sí está versionado. Los .pptx no lo están.
    zin = zipfile.ZipFile(SRC)
    quitadas, iguales, tocadas = 0, 0, []
    with zipfile.ZipFile(DST, "w", zipfile.ZIP_DEFLATED) as zout:
        for it in zin.infolist():
            if it.filename.startswith("ppt/notesSlides/") and it.filename.endswith(".xml"):
                quitadas += 1
                continue
            data = orig = zin.read(it.filename)
            if it.filename == "[Content_Types].xml":
                data = OVR.sub(b"", data)
            elif RELS.fullmatch(it.filename):
                data = REL.sub(b"", data)
            if data == orig:
                iguales += 1
            else:
                tocadas.append(it.filename)
            zout.writestr(it, data)
    zin.close()

    print("notesSlides eliminados:", quitadas)
    print("partes reescritas:", len(tocadas), "· partes byte-idénticas:", iguales)

    from pptx import Presentation
    p = Presentation(DST)
    con = sum(1 for s in p.slides
              if s.has_notes_slide and s.notes_slide.notes_text_frame.text.strip())
    print("Guardado:", DST)
    print("verificación: %d láminas, %d con notas" % (len(p.slides), con))
    assert con == 0, "quedaron notas en la copia"


if __name__ == "__main__":
    main()
