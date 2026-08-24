#!/usr/bin/env python
"""generate_b8_deck.py — deck del sprint B8, con DOS EJES.

Eje 1, SI-MIL (Kapse et al., CVPR 2024): la tarea de investigación cerrada. Las ECUACIONES,
la FIGURA ORIGINAL del modelo (Fig. 2, pág. 4) y el formato de deck del proyecto.
Insumo: `simil_explicacion_matematica.md`; números: `simil_estudio.md`.

Eje 2, la medición de atención contra las marcas del patólogo
(`../atencion_vs_patologo/{prereg.md,resultados.md}`): el resultado que reordenó las cuatro
familias del objetivo de mitosis. Nada se re-mide acá: el experimento cerró el 2-ago y este
archivo solo lo presenta.

Reparto del 3-ago (pedido de Ernesto): SI-MIL se COMPACTA A LA MITAD, de 12 láminas de
contenido a 6, y NO se borra, porque fue una de las tareas de investigación. Entra la
sección de mitosis en registro muy pedagógico, con sus dos figuras. Método de compactado,
el mismo del recorte del 31-jul: se fusionan pares y lo que sale de la lámina se cuenta
hablando, con el guion REESCRITO, no pegado.

Las seis fusiones de SI-MIL: qué propone + la figura del paper · las dos entradas +
ecuación 1 · ecuación 2 + qué implica para nuestro modelo · el puente Top-K + la rama
interpretable · las ecuaciones 3 a 10 · qué reportan y el contraste + qué costaría y qué
preguntar.

Recorte previo del 31-jul: de 19 a 14 láminas, retirando el ejemplo numérico del orden.

Sección de cierre (4-ago): el grid E×S (`../grid_expertos_slots/resultados.md`, encargo 3 del
B8) entra DESPUÉS de los dos ejes y NO como tercer eje, por decisión de Ernesto. Cero Δ contra
CLAM por brazo, que es lo que el pre-registro §6 prohibió por diseño, y el veredicto es H_nula
y se cuenta como tal.

REDISEÑO del 4-ago (doce puntos pedidos por Ernesto tras revisar el deck armado). Sigue en 20
láminas, con la estructura cambiada en las dos puntas y el cuerpo aligerado:
  - Abre con OBJETIVOS DEL SPRINT (molde de recapitulación del B7) en vez del mapa del
    recorrido. El objetivo 1, el escalado de slots, es el único sin lámina: se cuenta en el
    guion, y de paso deja dicho uno de los mensajes que hay que llevar a la reunión.
  - Cierra con los TRES PAPERS de la rama de mitosis (heredan el lugar de la lámina de las
    cuatro familias, cuyo reordenamiento pasó al guion; sin las letras A/B/C/D) y una lámina
    nueva de OBJETIVOS PROPUESTOS. La lámina del determinismo se retiró: lo que quedaba en
    pie de ella es una línea de comparabilidad en el guion del grid.
  - Dos láminas se volvieron figura donde antes había un panel de texto: el estadístico
    (`escala_auc`) y, sobre todo, el nulo por traslación, que era la lámina que Ernesto
    dijo no entender (retirada el 19-ago, con sus figuras). La de los mapas pasa a la región anotada
    sola, al doble de lado, con la procedencia al pie (`pie_lineas`).
  - Transversales: títulos minimalistas, y cada guion abre con un punteo de 3 a 5 renglones
    de una línea antes de la prosa hablada, que sigue siendo prosa
    ([[notas-presentador-guion-didactico]]).
  - Corrección que atraviesa el deck: los checkpoints primarios NO vieron la lámina anotada
    «en entrenamiento» (está en validación), que es lo que sostienen el prereg y el
    resultados. La formulación anterior, «nunca vieron esta lámina», se pasaba de eso.

RECORTE del 5-ago: de 20 a 16 láminas, fusionando pares y agrandando lo que queda.

REORDENAMIENTO del 6-ago, pedido por SEBASTIÁN tras ver el deck, y ejecutado el 7-ago. Queda
en 14 láminas y la audiencia pasa a ser Benjamín, la semana del 11-ago:
  - El orden es GRID → ATENCIÓN → SI-MIL. Invierte la decisión del 3-ago (lo cerrado antes que
    lo vivo) y jubila el encuadre del 4-ago (el grid como cierre); las dos eran de encuadre
    nuestro, ésta viene del supervisor y no se re-decide.
  - Se van dos láminas: OBJETIVOS DEL SPRINT y los TRES PAPERS. El molde de la primera no se
    pierde, lo hereda «Objetivos propuestos»; el razonamiento de la segunda baja al guion de
    esa misma lámina, porque es la justificación del objetivo propuesto 2.
  - Las cuatro láminas de la medición pasan a describir el ESTADÍSTICO con precisión: qué es
    (U de Mann-Whitney normalizada), contra qué se mide cada grupo, y con cuánta precisión
    (el IC de Hanley-McNeil, que es una incertidumbre distinta de la sd entre checkpoints).
  - `build()` deja de ser un bloque único: una función por lámina y el orden en un solo lugar.

RECORTE del 19-ago, pedido por Ernesto en `correcciones.txt`. De 27 láminas a 15, y el
criterio es la reunión pendiente con Sebastián: lo que no sirva para llegar a ella, sale.
  - Se van, por orden explícito: la lámina de controles, la que fechaba el trabajo, las SEIS
    de regiones de escaneo, la de estado de la herramienta, la de la cola de cómputo, la de
    patrones y la de próximos pasos. La cola de cómputo no se menciona en ninguna parte, y
    tampoco se fecha lo que se hizo un día concreto.
  - Las dos que medían los factores del techo (la curva del filtro y el cruce) se FUNDEN en
    una sola de RESULTADOS, con la forma que pidió: los mapas de calor primero y el número
    al lado. Se conserva el 13 de 26 y la tabla de los dos factores; se va la meseta de
    tolerancia, que era método. Resuelto con Ernesto antes de borrar: el texto admitía
    leerse como «borrar las dos», y eso se habría llevado el único resultado cuantitativo.
  - Entra la FIGURA 1 del paper de HoVer-NeXt en lugar de la lámina de estado.
  - El eje de la lámina de resultados es el encadenamiento que pidió Sebastián: la atención
    de CLAM recorta y el detector trabaja adentro. No son dos herramientas por separado.
  - Barrido de referencias cruzadas después de cortar: la portada prometía «tres cosas» y
    ahora son cuatro; dos guiones apuntaban a la lámina de controles y ya no la nombran.
  - **Pedido posterior, el mismo día**: una lámina de CONTROL con la herramienta SOLA sobre la
    misma lámina, sin nada del modelo de atención encima. Queda de cierre y es la que hace
    legible a la de resultados: sin un brazo sin recorte, «13 de 26 con el recorte» no se
    compara con nada. No hizo falta correr nada — la corrida fue sobre la lámina entera y el
    recorte se aplicó después, sobre la salida. Deck en **16**.

Reglas que gobiernan el archivo:
  - Se construye SOBRE el template válido (Deep-LLM-V), nunca con Presentation() a secas:
    el template EMBEBE sus fuentes y ese es el motivo real de que un deck "no parezca el
    template" ([[deck-template-fuentes-embebidas]]).
  - TODO nativo salvo las IMÁGENES: las dos figuras de papers y las nuestras sobre tejido
    real (mapas de atención, la cadena, los detalles a resolución nativa). El criterio no es
    de dónde salió la imagen sino si el objeto es dibujable: tabla, gráfico o diagrama van
    nativos siempre (CLAUDE.md, ADDENDUM 3-ago).
  - Gramática de diagrama de Deep-LLM-V ([[deck-gramatica-diagrama-deep-llm-v]]).
  - Tipografía: TODO Barlow, forzado con forzar_barlow() sobre theme + master + layouts.
  - Sin rayas «—»/«–», sin la palabra «palanca», sin la expresión «al revés».

Uso:
  PYTHONPATH=/media/administrador/Storage1/sdonoso/clam_testing2/.pylibs \
  /home/sdonoso/miniconda3/envs/clam_latest/bin/python \
    sprints/B8_sprint8/presentacion_b8/generate_b8_deck.py
"""
import os

from lxml import etree
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR, MSO_SHAPE_TYPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn
from PIL import Image

REPO = "/media/administrador/Storage1/sdonoso/clam_testing2/oncomets-ernesto"
PRES = os.path.join(REPO, "papers/presentations")
ASSETS_BRAND = os.path.join(PRES, "assets_branding")
OUT_DIR = os.path.join(REPO, "sprints/B8_sprint8/presentacion_b8")
ASSETS = os.path.join(OUT_DIR, "assets")
DST = os.path.join(OUT_DIR, "CLAM_Sprint8.pptx")
DST_LEGACY = os.path.join(OUT_DIR, "CLAM_Sprint8_SIMIL.pptx")   # nombre previo, monográfico

TEMPLATE = os.path.join(REPO, "sprints/B7_sprint7/Modelo OncoMets Spatial V1 Deep-LLM-V.pptx")
# Solo la portada de marca. La lámina de TÍTULO del template se retira el 22-ago (pedido de
# Ernesto): los objetivos del sprint ocupan la posición 2, con el molde de la ex-lámina de
# cierre. Lo que esa lámina llevaba (el nombre del sprint y la fecha) pasa a la portada.
TPL_KEEP = (0,)            # portada de marca, nativa a 13.333
FECHA_REUNION = "24 de agosto de 2026"

# --- figura del paper (única imagen del deck), recortada de la página 4 a 400 DPI ---
FIG2_FULL = os.path.join(ASSETS, "simil_fig2_full.png")   # los tres paneles
FIG2_A = os.path.join(ASSETS, "simil_fig2_a.png")         # (a) SI-MIL overview
FIG2_B = os.path.join(ASSETS, "simil_fig2_b.png")         # (b) Conventional MIL branch
FIG2_C = os.path.join(ASSETS, "simil_fig2_c.png")         # (c) Self-Interpretable branch

# --- figuras de la medición de atención (producción NUESTRA, no figura de paper) ---
# Recortadas de las originales por prep_assets_atencion.py: las de archivo traen títulos de
# matplotlib y un hueco de ~390 px entre las dos regiones de escaneo del .bif.
FIG_MAPAS = os.path.join(ASSETS, "atencion_dos_regiones.png")      # grilla 2x2 (ar 1.1827)
# Preparado para el rediseño de la lámina 12: solo la región anotada, atención | marcas, que
# es el doble de lado y no arrastra la fila que se lee como imagen repetida. ar = 2.407, así
# que la caja de la lámina hay que recalcularla al adoptarlo (add_image_fit usa `h = w / ar`).
FIG_MAPAS_ANOTADA = os.path.join(ASSETS, "atencion_region_anotada.png")
FIG_MITOSIS = os.path.join(ASSETS, "mitosis_region_anotada.png")   # con el recuadro de foco
FIG_ZOOM = os.path.join(ASSETS, "mitosis_zoom.png")                # el detalle del recuadro

# --- figuras de la lámina de resultados del 19-ago (producción NUESTRA) ---
# Preparadas por prep_assets_hovernext.py. La cadena son tres paneles de la MISMA región en
# tres estados (atención · recorte · detecciones), montados con 26 px de aire: ar = 3.694.
# El detalle son cuatro recortes a resolución nativa sobre mitosis acertadas: ar = 4.069.
FIG_CADENA = os.path.join(ASSETS, "cadena_clam_hovernext.png")
FIG_HN_ZOOM = os.path.join(ASSETS, "hovernext_zoom.png")
# El brazo de CONTROL: la herramienta sola sobre la lámina entera, sus dos regiones de
# escaneo una al lado de la otra con las 177 detecciones. ar = 2.471.
FIG_HN_SOLO = os.path.join(ASSETS, "hovernext_solo.png")
# Figura EXTERNA: la Figura 1 del paper de HoVer-NeXt, extraída por
# prep_assets_paper_hovernext.py. ar = 1.621. Va como imagen por ser de un paper.
FIG_HN_PAPER = os.path.join(ASSETS, "hovernext_paper_fig1.png")

# --- láminas de contacto de la galería de mitosis (encargo 2), copiadas a `assets/` por
# prep_assets_hovernext.py. Son FOTOGRAFÍAS de tejido, o sea la excepción legítima a «todo
# nativo» del deck: un recorte de lámina no se dibuja con shapes.
FIG_GAL_SIN_MARCA = os.path.join(ASSETS, "galeria_sin_marca.png")   # 1852x1708, ar 1.084
FIG_GAL_FALLADAS = os.path.join(ASSETS, "galeria_falladas.png")     # 928x354,  ar 2.621

# --- procedencia de la lámina anotada, para la lámina de los mapas ---
# Sebastián preguntó de dónde salió la lámina y con qué checkpoints se midió. Las tres líneas
# de abajo son la respuesta, verificadas el 4-ago contra los CSV de labels de `environ/`, el
# `offset_129741.json` de `anotaciones_patologo/` y los `splits_*_bool.csv` de cada corrida.
PROVENANCIA = [
    "Lámina 129741, cohorte privada. H&E Ventana .bif a 20× (0,465 µm/px), bajo wsi/129741/. "
    "4799 parches de 256 px, 163 bajo alguna marca.",
    "Marcas del patólogo, 61 polígonos exportados de QuPath: 26 mitosis · 14 núcleos de alto "
    "grado · 6 necrosis · 5 células inmunes · 5 tumor · 2 tejido adiposo · 1 estroma · 2 fondo.",
    "Sus etiquetas en nuestros CSV: tasa mitótica score_3 · pleomorfismo nuclear score_2 · "
    "diferenciación tubular score_3 · grado general 3 · grado nuclear del CDIS alto.",
]

# --- números de la medición, verificados contra auc_por_checkpoint.csv (4 ckpt primarios,
# cabeza de la clase verdadera). El orden es el de la escalera que se dibuja en la lámina.
#
# La cuarta columna es el SEMIANCHO del intervalo de confianza al 95 % del AUC, calculado el
# 6-ago con la fórmula de Hanley y McNeil sobre `con_region/auc_por_checkpoint.csv`
# (universo `lamina`, rol `primario`): 1,96 × EE, con EE = √([A(1−A) + (n₊−1)(Q₁−A²) +
# (n₋−1)(Q₂−A²)] / (n₊·n₋)), Q₁ = A/(2−A), Q₂ = 2A²/(1+A), n₊ = parches del grupo y
# n₋ = 4799 − n₊. Es la precisión que da ESE tamaño de grupo, y es una incertidumbre
# DISTINTA de la sd entre checkpoints (que mide cambiar de modelo, no cambiar de marcas).
# Por eso el bigote de estroma (n = 12) mide el triple que el de tejido adiposo (n = 27).
ESCALERA = [
    ("Mitosis", 28, 0.890, 0.080, True),
    ("Núcleos de alto grado", 13, 0.828, 0.138, True),
    ("Tumor", 48, 0.826, 0.072, False),
    ("Necrosis", 18, 0.748, 0.131, False),
    ("Estroma", 12, 0.537, 0.167, False),
    ("Linfocitos", 23, 0.322, 0.095, False),
    ("Tejido adiposo", 27, 0.154, 0.050, False),
]

# --- la cuenta que lleva de 26 marcas a 28 parches, calculada el 6-ago sobre la geometría
# real (geojson + coords del h5, con el offset dx = 3829 ya adoptado). La regla de mapeo es
# la de `alinear_anotaciones_qupath.py:252`: un parche cuenta si su cuadrado de 256 px se
# solapa con la caja envolvente del polígono.
#   26 polígonos de mitosis → 36 pares (polígono, parche) → 28 parches distintos
#   16 polígonos caen enteros dentro de un parche · 10 se reparten entre dos  → +10
#   21 parches tienen una sola mitosis · 6 tienen dos · 1 tiene tres          →  −8
CUENTA_26_28 = [("26", "marcas del patólogo"),
                ("+10", "cruzan un borde y ocupan dos parches"),
                ("−8", "siete parches tienen más de una marca"),
                ("28", "parches con mitosis")]

# --- números del grid E×S, verificados contra grid_expertos_slots/resultados.md ---
# Contraste primario del pre-registro §6: (brazo que recorta S) menos (brazo que recorta E),
# pareado por fold, dentro de cada peldaño de igual capacidad total E·S. Signo positivo =
# a favor de recortar slots. Los signos por fold salen de la tabla del §4.
PELDANOS = [
    ("E·S = 270", "30×9  contra  27×10", +0.022, 0.063, "+-+-+"),
    ("E·S = 210", "30×7  contra  21×10", -0.014, 0.034, "--+-+"),
    ("E·S = 150", "30×5  contra  15×10", -0.002, 0.048, "--++-"),
]
# AUC medio por brazo (§6 del resultados.md). Rótulo, AUC, capacidad total.
RAMA_S = [("30×10", 0.825, "300"), ("30×9", 0.792, "270"), ("30×7", 0.797, "210"),
          ("30×5", 0.802, "150"), ("30×3", 0.786, "90")]
RAMA_E = [("30×10", 0.825, "300"), ("27×10", 0.770, "270"), ("21×10", 0.812, "210"),
          ("15×10", 0.804, "150")]

# --- objetivos DEL sprint, para la lámina 2 ---
# Reemplazan a los «propuestos» del 6-ago, que cerraban el deck. Pedido de Ernesto
# (22-ago): los objetivos van al principio, con el molde de esa lámina de cierre, y la
# lista incluye tanto los cuatro que se propusieron al abrir el sprint como los dos que
# aparecieron durante — investigar e implementar un detector de mitosis a nivel de núcleo,
# y encadenarlo con la atención. El estado es el REAL a la fecha de la reunión.
#
# La geometría es la de la lámina que se fue, sin re-derivar: fila de 19 pt sobre 7,75" de
# ancho, `row_tops = 1.10 + i * 0.68`, `row_h = 0.62`, marcador de estado en `x = 8.98`.
#
# Una línea cada uno: a 19 pt sobre 7,75" un ítem de dos renglones hace que las filas se
# toquen entre sí. Con seis filas el pie de la última cae en 1.10 + 5*0.68 + 0.62 = 5,12",
# que entra sobre los 5,63" de alto útil.
# Medidos uno por uno con `text_w` a 19 pt bold: los seis entran en UN renglón sobre los
# 7,75" de la fila (el más largo, el 5, mide 7,21"). El rasterizado del 22-ago mostró qué
# pasa si no: a dos renglones el ítem desborda su caja de 0,62" y las filas se tocan.
OBJETIVOS_SPRINT = [
    ("1. Escalar a la tarea la medición de capacidad con expertos.", True),
    ("2. Decidir si conviene recortar expertos o unidades.", True),
    ("3. Medir si la atención cae sobre las mitosis del patólogo.", True),
    ("4. Elegir y justificar la dirección de la rama de mitosis.", True),
    ("5. Investigar e implementar un detector de mitosis por núcleo.", False),
    ("6. Encadenar la atención con el detector y cruzar las marcas.", False),
]

# --- los 4 checkpoints primarios de la medición de atención ---
# Ernesto (5-ago): «5 folds, fold 0» no dice con qué se entrenó el modelo. Estas filas
# son el dataset de ENTRENAMIENTO de cada uno, contado el 5-ago sobre los splits reales:
#   environ/splits/grado_histologico_mitotic_rate_100/splits_0_bool.csv        → 120/15/18 (153)
#   environ/splits/grado_histologico_mitotic_rate_combined_100/splits_0_bool.csv → 783/98/97 (978)
#   environ/splits_5fold/grado_histologico_mitotic_rate_combined_100/          → 934 en total;
#       fold 0 → 746 train, fold 2 → 749 train. La lámina 129741 está en `val` en los dos,
#       y en `train` en los folds 1, 3 y 4 — por eso de la corrida de cinco aparecen SOLO dos.
# Métricas: atencion_vs_patologo/resultados.md §3.
CKPTS_TABLA = [
    ["Privado · 120 láminas", "score_2   ✗", "0,645", "0,840"],
    ["Privado + TCGA · 783 láminas", "score_3   ✓", "0,224", "0,878"],
    ["Privado + TCGA · 746 láminas · fold 0", "score_2   ✗", "0,712", "0,926"],
    ["Privado + TCGA · 749 láminas · fold 2", "score_2   ✗", "0,524", "0,917"],
]

# ---- paleta Deep-LLM-V (medida sobre el template) ----
ONCO_DARK = RGBColor(0x3E, 0x68, 0x77)    # bloque de proceso
ONCO_CONN = RGBColor(0x38, 0x62, 0x71)    # conector / borde
ONCO_PANEL = RGBColor(0xCD, 0xDF, 0xE1)   # panel contenedor / operador
ONCO_DATA = RGBColor(0xB7, 0xB7, 0xB7)    # bloque de dato
ONCO_INK = RGBColor(0x0E, 0x28, 0x41)     # borde / fondo oscuro

TEAL_TITLE = ONCO_DARK
TEAL_SQ = ONCO_CONN
TEAL_DIV = ONCO_CONN
LAV_TITLE = RGBColor(0xFF, 0xFF, 0xFF)
TEAL_SUB = ONCO_PANEL
TEAL_CARD = ONCO_PANEL
TEAL_CARD2 = RGBColor(0xE9, 0xF1, 0xF2)   # tinte del claro, para banding
GRIS_BODY = RGBColor(0x59, 0x59, 0x59)
GRIS_TXT = RGBColor(0x55, 0x55, 0x55)
INK = RGBColor(0x22, 0x22, 0x22)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BLACK = RGBColor(0x00, 0x00, 0x00)
PROG_BG = ONCO_DATA

# Los dos caminos de datos de SI-MIL vienen rotulados por COLOR en la figura original
# (naranja = features profundas, verde = PathExpert). El deck no puede inventar dos
# colores nuevos, pero tampoco puede pintar los dos caminos igual: se distinguen por
# peso de tono dentro de la familia del template, y la figura del paper queda como la
# referencia de color cuando se la proyecta.
VIA_PROF = ONCO_DARK          # camino profundo (naranja en la figura)
VIA_INTERP = ONCO_PANEL       # camino interpretable (verde en la figura)

F_TITLE = "Barlow"
F_BODY = "Barlow"
F_MONO = "Barlow"
# Barlow no trae griegas, ℝ, ∈, → ni ⊗: esos glifos caen al fallback del sistema.
# Se rasterizó la alternativa (declarar Cambria Math, que el template también embebe) y
# se ve PEOR: sus griegas son serif finas y contrastan con el Barlow que las rodea,
# mientras que el fallback sans casa bien. Se deja todo declarado Barlow, que además es
# lo que pidió Sebastián. Mismo criterio que el B7 tomó para «→».

# --- geometría de trabajo (se escala x1.3333 al final) ---
SW, SH = 10.0, 5.625
LOGO = os.path.join(ASSETS_BRAND, "logo_header.png")
CHECK_VERDE = os.path.join(ASSETS_BRAND, "check_verde.png")

# --- cabecera OncoMets (láminas técnicas del template) ---
ONCO_LOGO = os.path.join(REPO, "sprints/B7_sprint7/presentacion_b7/assets/logo_oncomets.png")
ONCO_TITLE = RGBColor(0x3E, 0x68, 0x77)
ONCO_LINE = RGBColor(0x3D, 0x68, 0x76)
ONCO_LOGO_L, ONCO_LOGO_T, ONCO_LOGO_W, ONCO_LOGO_H = 0.5625, 0.3278, 1.1033, 0.6630
ONCO_TIT_L, ONCO_TIT_W = 2.0558, 6.8799
ONCO_TIT_BASE = 0.9908
ONCO_TIT_T, ONCO_TIT_H = 0.2950, ONCO_TIT_BASE - 0.2950
ONCO_TIT_SZ = 18.75
ONCO_LINE_T, ONCO_LINE_H = 1.0658, 0.1050
ONCO_BAND = ONCO_LINE_T + ONCO_LINE_H     # 1.1708

# Área de contenido: se maqueta ya en su sitio definitivo para que reflow_onco() no tenga
# que comprimir nada (reancla el bloque a CONTENT_TOP_NEW y escala si se pasa del pie).
TOP = 1.24
BOT = 5.46

_uid = [2000]


# ============================================================================
# Helpers base
# ============================================================================
def _blank(prs):
    for lay in prs.slide_layouts:
        if (lay.name or "").lower().strip() == "blank":
            return lay
    return prs.slide_layouts[6]


def base_from_template():
    """Abre el template válido y le borra las láminas salvo TPL_KEEP.

    Conserva ppt/fonts/*.fntdata + <p:embeddedFontLst> (Barlow, Cambria Math), el theme y
    el master. Borrar una slide no tiene API en python-pptx: hay que sacarla del sldIdLst
    Y soltar la relación, si no queda huérfana en el paquete."""
    prs = Presentation(TEMPLATE)
    keep_ids = {id(prs.slides[i]._element) for i in TPL_KEEP}
    lst = prs.slides._sldIdLst
    for i, sld in enumerate(list(lst)):
        if i in TPL_KEEP:
            continue
        prs.part.drop_rel(sld.get(
            "{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id"))
        lst.remove(sld)
    return prs, keep_ids


def new_slide(prs):
    """Lámina en blanco SIN los placeholders del layout (el BLANK del template arrastra
    DATE/FOOTER/SLIDE_NUMBER y aparecen como cuadros vacíos)."""
    s = prs.slides.add_slide(_blank(prs))
    for ph in list(s.placeholders):
        ph._element.getparent().remove(ph._element)
    return s


def _add_runs(p, text, size, bold, color, font=F_BODY):
    """Mini-markup de sub/superíndices REALES (baseline OOXML, no caracteres Unicode):
    `_x` / `_(xx)` = subíndice ; `^x` / `^(xx)` = superíndice.

    Hace falta porque Unicode no tiene subíndice para casi ninguna letra (no hay ᵢ de
    molde para `w`, `j`, `ij`…), y este deck es de ecuaciones. Portado de
    generate_clam_mammoth_pptx.py con Barlow en vez de Carlito. Para un `_`/`^` literal,
    escaparlo con backslash en el string fuente."""
    def emit(s, base=None):
        if not s:
            return
        r = p.add_run(); r.text = s
        r.font.name = font; r.font.bold = bold; r.font.color.rgb = color
        r.font.size = Pt(size * 0.74 if base is not None else size)
        if base is not None:
            r._r.get_or_add_rPr().set("baseline", str(base))
    i, n, buf = 0, len(text), ""
    while i < n:
        c = text[i]
        if c == "\\" and i + 1 < n and text[i + 1] in "_^\\":
            buf += text[i + 1]; i += 2; continue
        if c in "_^" and i + 1 < n:
            base = -25000 if c == "_" else 30000
            nxt = text[i + 1]
            if nxt == "(":
                j = text.find(")", i + 2)
                if j != -1:
                    emit(buf); buf = ""
                    emit(text[i + 2:j], base)
                    i = j + 1; continue
                buf += c; i += 1; continue
            emit(buf); buf = ""
            emit(nxt, base)
            i += 2; continue
        buf += c; i += 1
    emit(buf)


def _set_runs(tf, lines, anchor=MSO_ANCHOR.TOP, markup=False):
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    for i, ln in enumerate(lines):
        txt, sz, bold, col = ln[0], ln[1], ln[2], ln[3]
        font = ln[4] if len(ln) > 4 else F_BODY
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = ln[5] if len(ln) > 5 else PP_ALIGN.LEFT
        p.space_after = Pt(4)
        if markup:
            _add_runs(p, txt, sz, bold, col, font)
        else:
            r = p.add_run(); r.text = txt
            r.font.size = Pt(sz); r.font.bold = bold
            r.font.name = font; r.font.color.rgb = col


def add_textbox(slide, l, t, w, h, lines, anchor=MSO_ANCHOR.TOP, markup=False):
    tb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    _set_runs(tb.text_frame, lines, anchor=anchor, markup=markup)
    return tb


def notes(slide, text):
    """Guion HABLADO en prosa (sin etiquetas de fase, sin nº de job ni nombres)."""
    slide.notes_slide.notes_text_frame.text = text


def _rect(slide, l, t, w, h, color, line=None):
    sp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(l), Inches(t), Inches(w), Inches(h))
    sp.fill.solid(); sp.fill.fore_color.rgb = color
    if line is None:
        sp.line.fill.background()
    else:
        sp.line.color.rgb = line; sp.line.width = Pt(1)
    sp.shadow.inherit = False
    return sp


def header_oncomets(slide, title, size=22):
    """Cabecera técnica del template: logo OncoMets + título Barlow bold + línea teal."""
    pic = slide.shapes.add_picture(ONCO_LOGO, Inches(ONCO_LOGO_L), Inches(ONCO_LOGO_T),
                                   Inches(ONCO_LOGO_W), Inches(ONCO_LOGO_H))
    pic.name = "ONCOHDR_logo"
    if title:
        tb = slide.shapes.add_textbox(Inches(ONCO_TIT_L), Inches(ONCO_TIT_T),
                                      Inches(ONCO_TIT_W), Inches(ONCO_TIT_H))
        tb.name = "ONCOHDR_title"
        _set_runs(tb.text_frame, [(title, min(size, ONCO_TIT_SZ), True, ONCO_TITLE, F_BODY)],
                  anchor=MSO_ANCHOR.BOTTOM)
    _rect(slide, 0.0, ONCO_LINE_T, SW, ONCO_LINE_H, ONCO_LINE).name = "ONCOHDR_line"


def content(prs, title, size=26):
    s = new_slide(prs)
    header_oncomets(s, title, size=size)
    return s


def divider(prs, title, subtitle):
    s = new_slide(prs)
    s.background.fill.solid(); s.background.fill.fore_color.rgb = TEAL_DIV
    s.shapes.add_picture(LOGO, Inches(0.42), Inches(0.36), height=Inches(0.62))
    add_textbox(s, 0.8, 1.85, SW - 1.6, 1.45,
                [(title, 44, True, LAV_TITLE, F_TITLE, PP_ALIGN.CENTER)], anchor=MSO_ANCHOR.MIDDLE)
    add_textbox(s, 0.8, 3.42, SW - 1.6, 0.7,
                [(subtitle, 18, False, TEAL_SUB, F_BODY, PP_ALIGN.CENTER)], anchor=MSO_ANCHOR.TOP)
    return s


def add_image_fit(slide, path, l, t, w, h, align="center"):
    iw, ih = Image.open(path).size
    ar = iw / ih; box_ar = w / h
    if ar > box_ar:
        nw = w; nh = w / ar
    else:
        nh = h; nw = h * ar
    nl = l + (w - nw) / 2
    nt = t + (h - nh) / 2 if align != "top" else t
    slide.shapes.add_picture(path, Inches(nl), Inches(nt), Inches(nw), Inches(nh))


def add_card(slide, l, t, w, h, idx, text, size=14):
    fill = TEAL_CARD if idx % 2 == 0 else TEAL_CARD2
    sp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(l), Inches(t), Inches(w), Inches(h))
    sp.fill.solid(); sp.fill.fore_color.rgb = fill
    sp.line.color.rgb = TEAL_SQ; sp.line.width = Pt(1.25); sp.shadow.inherit = False
    cd = 0.44
    circ = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(l + 0.16), Inches(t + (h - cd) / 2),
                                  Inches(cd), Inches(cd))
    circ.fill.solid(); circ.fill.fore_color.rgb = ONCO_DARK
    circ.line.fill.background(); circ.shadow.inherit = False
    _set_runs(circ.text_frame, [(str(idx), 15, True, WHITE, F_TITLE, PP_ALIGN.CENTER)],
              anchor=MSO_ANCHOR.MIDDLE)
    tb = slide.shapes.add_textbox(Inches(l + 0.74), Inches(t), Inches(w - 0.88), Inches(h))
    _set_runs(tb.text_frame, [(text, size, True, INK, F_BODY)], anchor=MSO_ANCHOR.MIDDLE)


def caption(slide, l, t, w, text, size=11, col=GRIS_TXT, align=PP_ALIGN.CENTER, bold=False):
    add_textbox(slide, l, t, w, 0.4, [(text, size, bold, col, F_BODY, align)])


def status_done(slide, cx, cy, size=0.42):
    slide.shapes.add_picture(CHECK_VERDE, Inches(cx - size / 2), Inches(cy - size / 2),
                             Inches(size), Inches(size))


def status_progress(slide, cx, cy, w=1.5, h=0.44, texto="Pendiente"):
    sp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(cx - w / 2), Inches(cy - h / 2),
                                Inches(w), Inches(h))
    sp.fill.solid(); sp.fill.fore_color.rgb = PROG_BG
    sp.line.color.rgb = ONCO_CONN; sp.line.width = Pt(1.25); sp.shadow.inherit = False
    _set_runs(sp.text_frame, [(texto, 11, True, BLACK, F_BODY, PP_ALIGN.CENTER)],
              anchor=MSO_ANCHOR.MIDDLE)


# Las dos piezas de la barra de remate de cada lámina, para que `auditar` pueda ver si
# algún objeto la cruza. Los cuatro cruces del 18-ago pasaron la auditoría porque cada
# objeto estaba DENTRO de su caja: el defecto no es de una caja, es de dos que se tocan.
#
# Se guardan los SHAPES y no la altura: `reflow_onco` corre después y reancla (y a veces
# comprime) el cuerpo entero, así que el 4,85 de acá no es donde la barra termina. Leer la
# posición real al auditar es lo único que sobrevive a esa pasada.
_TAKEAWAY = {}


def takeaway_bar(slide, text, t=4.85, col=TEAL_TITLE, size=14):
    linea = _rect(slide, 0.35, t, SW - 0.7, 0.02, TEAL_SQ)
    tb = add_textbox(slide, 0.35, t + 0.08, SW - 0.7, 0.62,
                     [(text, size, True, col, F_BODY, PP_ALIGN.CENTER)],
                     anchor=MSO_ANCHOR.MIDDLE)
    _TAKEAWAY[id(slide._element)] = (linea, tb)


# ============================================================================
# Gramática de diagrama de Deep-LLM-V
# ============================================================================
def _proc(slide, l, t, w, h, text, dim=None, size=11, col=None):
    """Bloque de proceso: rounded-rect #3E6877 con Barlow bold BLANCO."""
    sp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(l), Inches(t),
                                Inches(w), Inches(h))
    sp.fill.solid(); sp.fill.fore_color.rgb = col or ONCO_DARK
    sp.line.fill.background(); sp.shadow.inherit = False
    lines = [(text, size, True, WHITE, F_BODY, PP_ALIGN.CENTER)]
    if dim:
        lines.append((dim, size - 2, False, WHITE, F_BODY, PP_ALIGN.CENTER))
    _set_runs(sp.text_frame, lines, anchor=MSO_ANCHOR.MIDDLE, markup=True)
    for p in sp.text_frame.paragraphs:
        p.space_after = Pt(0)
    return sp


def _proc_claro(slide, l, t, w, h, text, size=11, dim=None):
    """Bloque de proceso en el tono CLARO (#CDDFE1 con texto teal): el detalle interno."""
    sp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(l), Inches(t),
                                Inches(w), Inches(h))
    sp.fill.solid(); sp.fill.fore_color.rgb = ONCO_PANEL
    sp.line.fill.background(); sp.shadow.inherit = False
    lines = [(text, size, True, ONCO_DARK, F_BODY, PP_ALIGN.CENTER)]
    if dim:
        lines.append((dim, size - 2, False, ONCO_DARK, F_BODY, PP_ALIGN.CENTER))
    _set_runs(sp.text_frame, lines, anchor=MSO_ANCHOR.MIDDLE, markup=True)
    for p in sp.text_frame.paragraphs:
        p.space_after = Pt(0)
    return sp


def _dato(slide, l, t, w, h, text, size=9.5):
    """Bloque de dato (forma del tensor): rect #B7B7B7 con Barlow NEGRO."""
    sp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(l), Inches(t), Inches(w), Inches(h))
    sp.fill.solid(); sp.fill.fore_color.rgb = ONCO_DATA
    sp.line.fill.background(); sp.shadow.inherit = False
    _set_runs(sp.text_frame, [(text, size, False, BLACK, F_BODY, PP_ALIGN.CENTER)],
              anchor=MSO_ANCHOR.MIDDLE, markup=True)
    return sp


def _grupo(slide, l, t, w, h, fill=None):
    """Panel contenedor que agrupa una región."""
    sp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(l), Inches(t),
                                Inches(w), Inches(h))
    sp.fill.solid(); sp.fill.fore_color.rgb = fill or ONCO_PANEL
    sp.line.color.rgb = fill or ONCO_PANEL; sp.line.width = Pt(0.6)
    sp.shadow.inherit = False
    return sp


def _conn(slide, x0, y0, x1, y1, arrow=True):
    """Conector recto #386271 de 2.37 pt, con punta opcional."""
    ln = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(x0), Inches(y0),
                                    Inches(x1), Inches(y1))
    ln.line.color.rgb = ONCO_CONN; ln.line.width = Pt(2.37)
    ln.shadow.inherit = False
    if arrow:
        lnpr = ln.line._get_or_add_ln()
        lnpr.append(lnpr.makeelement(qn('a:tailEnd'),
                                     {'type': 'triangle', 'w': 'med', 'len': 'med'}))
    return ln


def _conn_dash(slide, x0, y0, x1, y1):
    """Línea de expansión punteada: ata un bloque con su detalle ampliado."""
    ln = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(x0), Inches(y0),
                                    Inches(x1), Inches(y1))
    ln.line.color.rgb = ONCO_CONN; ln.line.width = Pt(1.0)
    ln.shadow.inherit = False
    lnpr = ln.line._get_or_add_ln()
    lnpr.append(lnpr.makeelement(qn('a:prstDash'), {'val': 'dash'}))
    return ln


def _dim(slide, l, t, w, text, size=9.5, align=PP_ALIGN.CENTER, col=None):
    """Etiqueta de forma del tensor suelta, pegada al bloque."""
    add_textbox(slide, l, t, w, 0.26,
                [(text, size, False, col or ONCO_INK, F_BODY, align)],
                anchor=MSO_ANCHOR.MIDDLE, markup=True)


def _oper(slide, cx, cy, sym="+", d=0.34):
    """Operador: óvalo #CDDFE1 con borde #0E2841."""
    sp = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(cx - d / 2), Inches(cy - d / 2),
                                Inches(d), Inches(d))
    sp.fill.solid(); sp.fill.fore_color.rgb = ONCO_PANEL
    sp.line.color.rgb = ONCO_INK; sp.line.width = Pt(1.0); sp.shadow.inherit = False
    _set_runs(sp.text_frame, [(sym, 11, True, ONCO_INK, F_BODY, PP_ALIGN.CENTER)],
              anchor=MSO_ANCHOR.MIDDLE)
    return sp


def _rot_label(slide, l, t, w, h, text, size=9, col=None):
    """Rótulo vertical al costado de un panel."""
    tb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    _set_runs(tb.text_frame, [(text, size, True, col or ONCO_DARK, F_BODY, PP_ALIGN.CENTER)],
              anchor=MSO_ANCHOR.MIDDLE)
    tb.rotation = 270
    return tb


def eq(slide, l, t, w, texto, num=None, size=15, h=0.52, fill=None, col=None):
    """Ecuación en su propio panel claro, con el número del paper a la derecha.

    Las ecuaciones son el objeto de este deck, así que se tratan como los bloques de
    arquitectura del template: caja propia, cuerpo grande, nada compitiendo al lado. El
    número va afuera del panel, chico, como lo pone el paper."""
    _grupo(slide, l, t, w, h, fill=fill or TEAL_CARD2)
    add_textbox(slide, l + 0.16, t, w - (0.62 if num else 0.32), h,
                [(texto, size, False, col or ONCO_INK, F_BODY, PP_ALIGN.CENTER)],
                anchor=MSO_ANCHOR.MIDDLE, markup=True)
    if num:
        add_textbox(slide, l + w - 0.58, t, 0.44, h,
                    [(num, 11, False, GRIS_BODY, F_BODY, PP_ALIGN.RIGHT)],
                    anchor=MSO_ANCHOR.MIDDLE)


def simple_table(slide, l, t, w, headers, rows, col_fracs, row_h=0.32, fs=9.5,
                 destacar=None, markup=False):
    """Tabla nativa: header teal + banding claro. `destacar` = índice de fila (0-based
    sobre `rows`) que se pinta con el celeste sólido, para la fila que importa."""
    ncol = len(headers); nrow = len(rows) + 1
    tbl = slide.shapes.add_table(nrow, ncol, Inches(l), Inches(t), Inches(w),
                                 Inches(row_h * nrow)).table
    for ci, fr in enumerate(col_fracs):
        tbl.columns[ci].width = Inches(w * fr)
    tbl.first_row = False; tbl.horz_banding = False
    for ri, row in enumerate([headers] + list(rows)):
        for ci, txt in enumerate(row):
            cell = tbl.cell(ri, ci)
            cell.margin_left = Inches(0.06); cell.margin_right = Inches(0.04)
            cell.margin_top = Inches(0.015); cell.margin_bottom = Inches(0.015)
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            if ri == 0:
                cell.fill.solid(); cell.fill.fore_color.rgb = TEAL_SQ
                col, bold = WHITE, True
            elif destacar is not None and ri - 1 == destacar:
                cell.fill.solid(); cell.fill.fore_color.rgb = TEAL_CARD
                col, bold = ONCO_DARK, True
            else:
                cell.fill.solid(); cell.fill.fore_color.rgb = TEAL_CARD2 if ri % 2 else TEAL_CARD
                col, bold = INK, (ci == 0)
            tf = cell.text_frame; tf.word_wrap = True
            p = tf.paragraphs[0]; p.alignment = PP_ALIGN.LEFT
            if markup:
                _add_runs(p, txt, fs, bold, col, F_BODY)
            else:
                r = p.add_run(); r.text = txt
                r.font.size = Pt(fs); r.font.bold = bold
                r.font.name = F_BODY; r.font.color.rgb = col
    return tbl


# ---------------------------------------------------------------------------
# Medición de texto real
# ---------------------------------------------------------------------------
# El chequeo programático de conformidad da «todo limpio» con texto que desborda su caja
# ([[deck-qa-puntos-ciegos-chequeo]]), porque nadie mide el texto. Acá se mide de verdad,
# con los mismos TTF de Barlow que están instalados bajo containment, y los paneles se
# dimensionan solos. Es el equivalente al auto-dimensionado de `render_table`.
BARLOW_DIR = "/media/administrador/Storage1/sdonoso/clam_testing2/fonts/barlow"
_FCACHE = {}
_MEDIDA = 40          # se mide a 40 px y se divide: da precisión sin depender del hinting


def _face(bold):
    from PIL import ImageFont
    key = bool(bold)
    if key not in _FCACHE:
        nombre = "Barlow-Bold.ttf" if bold else "Barlow-Regular.ttf"
        _FCACHE[key] = ImageFont.truetype(os.path.join(BARLOW_DIR, nombre), _MEDIDA)
    return _FCACHE[key]


def text_w(txt, size, bold=False):
    """Ancho del texto en PULGADAS del espacio de trabajo (lámina de 10 in, 72 pt/in).

    Las griegas y los símbolos que Barlow no tiene caen al fallback y miden distinto; se
    aproximan con el ancho de una «n», que es del orden correcto y sobreestima poco."""
    f = _face(bold)
    limpio = "".join(c if f.getbbox(c) else "n" for c in txt)
    return f.getlength(limpio) / _MEDIDA * size / 72.0


def wrap_lines(txt, ancho, size, bold=False):
    """Cuántas líneas ocupa `txt` dentro de `ancho` pulgadas (wrap por palabra)."""
    if not txt.strip():
        return 1
    n, actual = 1, ""
    for palabra in txt.split(" "):
        cand = (actual + " " + palabra).strip()
        if actual and text_w(cand, size, bold) > ancho:
            n += 1; actual = palabra
        else:
            actual = cand
    return n


def _alto_bloque(lineas, ancho, size, bold=False, space_after=3):
    """Alto en pulgadas de un bloque de párrafos ya envueltos."""
    total = 0.0
    for ln in lineas:
        total += wrap_lines(ln, ancho, size, bold) * size * 1.22 / 72.0 + space_after / 72.0
    return total


def panel(slide, l, t, w, h, title, tcol, lines, border, fill=TEAL_CARD2, tsize=14.5,
          bsize=12, markup=False):
    """Panel con título + líneas de cuerpo.

    `h=None` calcula el alto midiendo el texto, que es lo que evita que la última línea
    quede fuera de la caja. Con `h` explícito, el valor se respeta salvo que el texto no
    entre: en ese caso gana el texto, porque una caja corta y un renglón afuera se ve
    peor que una caja un poco más alta."""
    ancho_txt = w - 0.36
    necesario = (0.12 + _alto_bloque([title], ancho_txt, tsize, True)
                 + _alto_bloque(lines, ancho_txt, bsize) + 0.16)
    h = necesario if h is None else max(h, necesario)
    sp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(l), Inches(t), Inches(w), Inches(h))
    sp.fill.solid(); sp.fill.fore_color.rgb = fill
    sp.line.color.rgb = border; sp.line.width = Pt(1.5); sp.shadow.inherit = False
    tb = slide.shapes.add_textbox(Inches(l + 0.18), Inches(t + 0.08), Inches(ancho_txt), Inches(h - 0.12))
    runs = [(title, tsize, True, tcol, F_TITLE)] + [(ln, bsize, False, INK, F_BODY) for ln in lines]
    _set_runs(tb.text_frame, runs, markup=markup)
    for p in tb.text_frame.paragraphs:
        p.space_after = Pt(3)
    return sp


# ============================================================================
# Figuras nativas de la sección de mitosis
# ============================================================================
# Una escalera de siete AUC pide un gráfico y no una lista ([[deck-contenido-visual-no-bullets]]).
# Se dibuja con la gramática del template en vez de add_chart: un gráfico de python-pptx
# arrastra la paleta de Office y hay que repintarlo entero, y acá el objeto es una barra por
# fila, que es exactamente el arquetipo «dato» del template.
def barras_ranking(slide, l, t, w, h, datos, x_eje=2.62, ancho_eje=5.40, fs=10.5):
    """Barras horizontales de AUC con la marca del azar y el bigote del IC 95 %.

    `datos` = (nombre, n, auc, ic, interes), con `ic` = semiancho del intervalo. El bigote
    es lo que convierte la escalera en una lámina honesta: sin él, siete barras del mismo
    grosor sugieren siete números de la misma calidad, y no lo son — la precisión la fija
    el n de cada grupo, que va de 12 a 48 parches."""
    fila = h / len(datos)
    alto_barra = min(0.26, fila * 0.62)
    x0 = l + x_eje
    for i, (nombre, n, auc, ic, interes) in enumerate(datos):
        cy = t + i * fila + fila / 2
        add_textbox(slide, l, cy - 0.15, x_eje - 0.14, 0.30,
                    [("%s  (n = %d)" % (nombre, n), fs, interes, ONCO_DARK if interes else INK,
                      F_BODY, PP_ALIGN.RIGHT)], anchor=MSO_ANCHOR.MIDDLE)
        _rect(slide, x0, cy - alto_barra / 2, ancho_eje * auc, alto_barra,
              ONCO_DARK if interes else ONCO_PANEL)
        # bigote: tallo fino de auc−ic a auc+ic, con tope en los dos extremos. Va en el
        # azul oscuro del template para que se lea tanto sobre la barra como sobre el fondo.
        xa = x0 + ancho_eje * max(0.0, auc - ic)
        xb = x0 + ancho_eje * min(1.0, auc + ic)
        _rect(slide, xa, cy - 0.012, xb - xa, 0.024, ONCO_INK)
        for xc in (xa, xb):
            _rect(slide, xc - 0.012, cy - 0.085, 0.024, 0.170, ONCO_INK)
        # El valor va en COLUMNA FIJA al final del eje, no pegado a la punta del bigote.
        # Pegado se rompía con los grupos bajo el azar: el rótulo de linfocitos (0,322)
        # caía justo donde pasa la línea punteada del 0,5 y quedaba tachado. Los dos
        # objetos eran válidos y estaban en su caja, así que ningún chequeo automático lo
        # veía; salió del rasterizado ([[deck-qa-puntos-ciegos-chequeo]]). De paso los siete
        # números quedan alineados, que es como se comparan.
        add_textbox(slide, x0 + ancho_eje + 0.10, cy - 0.15, 0.80, 0.30,
                    [(_num(auc), fs, interes, ONCO_DARK if interes else GRIS_BODY, F_BODY)],
                    anchor=MSO_ANCHOR.MIDDLE)
    # el azar, que es la única referencia que este estadístico tiene
    x_nulo = x0 + ancho_eje * 0.5
    ln = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(x_nulo), Inches(t - 0.16),
                                    Inches(x_nulo), Inches(t + h + 0.02))
    ln.line.color.rgb = ONCO_INK; ln.line.width = Pt(1.25)
    lnpr = ln.line._get_or_add_ln()
    lnpr.append(lnpr.makeelement(qn('a:prstDash'), {'val': 'dash'}))
    ln.shadow.inherit = False
    add_textbox(slide, x_nulo - 0.80, t - 0.42, 1.60, 0.26,
                [("azar = 0,5", 9.5, True, ONCO_INK, F_BODY, PP_ALIGN.CENTER)],
                anchor=MSO_ANCHOR.MIDDLE)


def cadena_cuenta(slide, l, t, w, items, h=0.42, fs=14, fs_pie=8):
    """Una cuenta con el resultado ya hecho, en bloques encadenados.

    Es el arquetipo que la convención pide cuando un bullet dice un número que sale de otros
    ([[deck-contenido-visual-no-bullets]]): «26 marcas dan 28 parches» no se explica en una
    frase subordinada, se dibuja. Los extremos son bloques de proceso (los dos números que
    alguien puede cruzar y creer que se contradicen) y el medio son los dos ajustes."""
    n = len(items)
    cw = (w - 0.30 * (n - 1)) / n
    for i, (num, pie) in enumerate(items):
        x = l + i * (cw + 0.30)
        extremo = i in (0, n - 1)
        _rect(slide, x, t, cw, h, ONCO_DARK if extremo else ONCO_PANEL)
        add_textbox(slide, x, t, cw, h,
                    [(num, fs, True, WHITE if extremo else ONCO_INK, F_BODY, PP_ALIGN.CENTER)],
                    anchor=MSO_ANCHOR.MIDDLE)
        add_textbox(slide, x, t + h + 0.04, cw, 0.34,
                    [(l_, fs_pie, False, GRIS_BODY, F_BODY, PP_ALIGN.CENTER)
                     for l_ in pie.split("\n")])
        if i < n - 1:
            _conn(slide, x + cw + 0.04, t + h / 2, x + cw + 0.26, t + h / 2)


def cinta_ranking(slide, l, t, w, marcados, n=34, h=0.40, fs=9, ejes=True):
    """Los N parches de la lámina ordenados por atención, con los anotados resaltados.

    Es la figura del ESTADÍSTICO, que es justamente lo que no dejó ninguna imagen cuando se
    hizo la medición y por eso no se retuvo ([[hallazgo-necesita-forma-presentable]]). Cada
    celda es un parche; el orden es de más a menos atención."""
    cw = w / n
    for i in range(n):
        col = ONCO_DARK if i in marcados else ONCO_PANEL
        sp = _rect(slide, l + i * cw, t, cw * 0.86, h, col)
        sp.line.color.rgb = WHITE; sp.line.width = Pt(0.5)
    # El eje se dibuja una sola vez cuando hay dos cintas apiladas: es el mismo orden en las
    # dos, y repetirlo debajo de cada una lo convierte en ruido.
    if not ejes:
        return
    _conn(slide, l, t + h + 0.20, l + w, t + h + 0.20)
    add_textbox(slide, l, t + h + 0.26, 2.60, 0.26,
                [("más atención", fs, True, ONCO_DARK, F_BODY)])
    add_textbox(slide, l + w - 2.60, t + h + 0.26, 2.60, 0.26,
                [("menos atención", fs, True, GRIS_BODY, F_BODY, PP_ALIGN.RIGHT)])


# ============================================================================
# Figuras nativas de la sección del grid E×S
# ============================================================================
def _num(v, dec=3, signo=False):
    """Número con coma decimal y, si se pide, el signo delante (menos tipográfico)."""
    s = ("%+.*f" % (dec, v)) if signo else ("%.*f" % (dec, v))
    return s.replace("-", "−").replace(".", ",")


def barras_divergentes(slide, l, t, w, h, filas, esc=0.10, fs=10.5, destacar=None,
                       izq="a favor de recortar expertos",
                       der="a favor de recortar unidades",
                       metrica="diferencia de AUC entre los dos recortes, partición por "
                               "partición"):
    """Δ pareado por peldaño alrededor del cero, con la desviación como bigote.

    Esta figura NO está para leer la magnitud del Δ. Está para que se vea que el bigote cruza
    el cero en los tres peldaños y que la media se cambia de lado entre ellos: eso es lo que
    separa un resultado nulo de un «casi», y es justo lo que una tabla de tres números deja
    sin mostrar. Por eso el cero es una línea con peso y la barra de la media se dibuja
    ENCIMA del bigote, no al lado.

    Encabezado del 22-ago (pedido de Ernesto): el defecto era que la figura no decía QUÉ
    magnitud se compara, solo quién ganaba. Ahora arriba va la métrica y debajo las dos
    direcciones flanqueando el cero. Y se retiran del dibujo las dos leyendas que se leían
    como jerga suelta: la de escala (`± 0,1 AUC`) y la de los cuadros por partición. Los
    cuadros se siguen dibujando, sin rótulo, y cómo se leen queda en el guion, que es
    exactamente lo que pidió."""
    lab_w, val_w, chip_w = 2.50, 1.42, 0.96
    x0 = l + lab_w + 0.12
    ancho = w - lab_w - val_w - chip_w - 0.42
    semi = ancho / 2
    xc = x0 + semi
    fila = h / len(filas)

    # la banda del peldaño que se resalta va PRIMERO: si se dibujara después taparía el
    # cero, el bigote y los cuadros de esa misma fila
    if destacar is not None:
        _grupo(slide, l, t + destacar * fila + 0.02, w, fila - 0.04, fill=TEAL_CARD2)
    # qué magnitud se compara, arriba de todo y centrada sobre el cero. La caja se ensancha
    # más allá del área de barras (que mide 3,98" y le queda corta al rótulo): arriba de la
    # figura no hay nada que pisar, y centrada sobre el cero es donde se lee como su eje.
    hw = min(w, ancho + 2.30) / 2
    add_textbox(slide, xc - hw, t - 0.62, hw * 2, 0.26,
                [(metrica, 10, True, INK, F_BODY, PP_ALIGN.CENTER)], anchor=MSO_ANCHOR.MIDDLE)
    # y debajo, las dos direcciones, cada una de su lado del cero
    add_textbox(slide, xc - 2.70, t - 0.34, 2.58, 0.26,
                [("← " + izq, 9.5, True, GRIS_BODY, F_BODY, PP_ALIGN.RIGHT)],
                anchor=MSO_ANCHOR.MIDDLE)
    add_textbox(slide, xc + 0.12, t - 0.34, 2.58, 0.26,
                [(der + " →", 9.5, True, ONCO_DARK, F_BODY)], anchor=MSO_ANCHOR.MIDDLE)
    ln = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(xc), Inches(t - 0.04),
                                    Inches(xc), Inches(t + h + 0.04))
    ln.line.color.rgb = ONCO_INK; ln.line.width = Pt(1.5); ln.shadow.inherit = False

    # El rótulo y su subtítulo ocupan `off` arriba y abajo del centro de la fila. Con filas
    # bajas (la lámina fusionada del 5-ago las dejó en 0,39") los 0,26 fijos hacían que el
    # subtítulo de una fila pisara el rótulo de la siguiente.
    off = max(0.17, min(0.26, fila * 0.46))
    for i, (rot, sub, media, sd, signos) in enumerate(filas):
        cy = t + i * fila + fila / 2
        add_textbox(slide, l, cy - off, lab_w, off,
                    [(rot, fs, True, ONCO_DARK, F_BODY, PP_ALIGN.RIGHT)], anchor=MSO_ANCHOR.MIDDLE)
        add_textbox(slide, l, cy + 0.01, lab_w, off - 0.02,
                    [(sub, fs - 1.5, False, GRIS_BODY, F_BODY, PP_ALIGN.RIGHT)],
                    anchor=MSO_ANCHOR.MIDDLE)
        # el bigote primero: la barra de la media va encima
        xa, xb = xc + (media - sd) / esc * semi, xc + (media + sd) / esc * semi
        _rect(slide, xa, cy - 0.015, xb - xa, 0.03, ONCO_DATA)
        for x in (xa, xb):
            _rect(slide, x - 0.01, cy - 0.11, 0.02, 0.22, ONCO_DATA)
        # la media, con ancho mínimo para que un Δ de dos milésimas no desaparezca
        dx = media / esc * semi
        ancho_bar = max(abs(dx), 0.03)
        _rect(slide, xc if dx >= 0 else xc - ancho_bar, cy - 0.10, ancho_bar, 0.20, ONCO_DARK)
        add_textbox(slide, l + w - chip_w - val_w - 0.10, cy - 0.14, val_w, 0.28,
                    [("%s ± %s" % (_num(media, 3, True), _num(sd)), fs - 0.5, True, INK,
                      F_BODY, PP_ALIGN.RIGHT)], anchor=MSO_ANCHOR.MIDDLE)
        # un cuadro por fold: relleno = a favor de recortar slots
        cw, gap = 0.16, 0.04
        xs = l + w - chip_w
        for j, sg in enumerate(signos):
            sp = _rect(slide, xs + j * (cw + gap), cy - cw / 2, cw, cw,
                       ONCO_DARK if sg == "+" else TEAL_CARD2)
            sp.line.color.rgb = ONCO_CONN; sp.line.width = Pt(0.75)



def escalera_capacidad(slide, l, t, w, h, puntos, lo=0.75, hi=0.84, nota=None, fs=9.5):
    """Una rama del grid: AUC medio por brazo, con la línea que une los topes de las barras.

    El eje arranca en 0,75 y no en cero, y va rotulado en la lámina: entre el mejor y el peor
    brazo hay 0,055 de AUC, y a escala completa los ocho brazos serían la misma barra. Lo que
    la figura tiene que dejar ver es la FORMA (un escalón y después una meseta en una rama,
    una curva que ni siquiera es monótona en la otra), no la altura de cada barra."""
    paso = w / len(puntos)
    ancho_barra = min(0.58, paso * 0.50)
    base = t + h
    topes = []
    for i, (rot, auc, cap) in enumerate(puntos):
        cx = l + paso * (i + 0.5)
        alto = max(0.06, (auc - lo) / (hi - lo) * h)
        _rect(slide, cx - ancho_barra / 2, base - alto, ancho_barra, alto,
              ONCO_PANEL if i == 0 else ONCO_DARK)
        # el valor va DENTRO de la barra, no encima: la línea de tendencia llega al tope de
        # cada barra desde el tope del vecino, así que a los costados del rótulo pasa mucho
        # más arriba y lo cruza. Adentro no hay nada que la línea pueda tapar.
        add_textbox(slide, cx - 0.52, base - alto + 0.05, 1.04, 0.24,
                    [(_num(auc), fs, True, ONCO_DARK if i == 0 else WHITE, F_BODY,
                      PP_ALIGN.CENTER)], anchor=MSO_ANCHOR.MIDDLE)
        add_textbox(slide, cx - 0.52, base + 0.04, 1.04, 0.24,
                    [(rot, fs, True, INK if i else GRIS_BODY, F_BODY, PP_ALIGN.CENTER)])
        add_textbox(slide, cx - 0.52, base + 0.25, 1.04, 0.22,
                    [(cap, fs - 1.0, False, GRIS_BODY, F_BODY, PP_ALIGN.CENTER)])
        topes.append((cx, base - alto))
    for (x0, y0), (x1, y1) in zip(topes, topes[1:]):
        ln = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(x0), Inches(y0),
                                        Inches(x1), Inches(y1))
        ln.line.color.rgb = ONCO_INK; ln.line.width = Pt(1.25); ln.shadow.inherit = False
    if nota:
        i, texto = nota
        cx = l + paso * (i + 0.5)
        add_textbox(slide, cx - 1.10, base + 0.48, 2.20, 0.24,
                    [(texto, fs - 1.0, True, ONCO_DARK, F_BODY, PP_ALIGN.CENTER)])


# ============================================================================
# Figuras y pies del rediseño del 4-ago
# ============================================================================
def pie_lineas(slide, l, t, w, lineas, size=8, col=GRIS_BODY, space_after=1.5):
    """Pie de varias líneas en UN solo textbox, con el interlineado al mínimo.

    `caption` gasta 0,4" por renglón y con tres renglones se come el alto que necesita la
    figura de arriba. Acá el bloque mide exactamente lo que mide su texto, y devuelve ese
    alto para que quien lo llama pueda cerrar la lámina debajo."""
    h = _alto_bloque(lineas, w - 0.22, size, space_after=space_after)
    tb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    _set_runs(tb.text_frame, [(ln, size, False, col, F_BODY) for ln in lineas])
    for p in tb.text_frame.paragraphs:
        p.space_after = Pt(space_after)
    return h


def escala_auc(slide, l, t, w, valor, alto=0.16, fs=9.5, pie=None):
    """El estadístico puesto sobre su recorrido de 0 a 1, con el azar marcado.

    Un número suelto dentro de un renglón de texto no dice dónde cae. Sobre la escala, la
    distancia hasta el azar se lee sin que nadie la explique, que es de lo que se trata
    volver visual un panel de texto ([[deck-contenido-visual-no-bullets]])."""
    yb = t + 0.26
    _rect(slide, l, yb, w, alto, ONCO_PANEL)
    x_az = l + w * 0.5
    ln = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(x_az), Inches(yb - 0.07),
                                    Inches(x_az), Inches(yb + alto + 0.07))
    ln.line.color.rgb = ONCO_INK; ln.line.width = Pt(1.25); ln.shadow.inherit = False
    lnpr = ln.line._get_or_add_ln()
    lnpr.append(lnpr.makeelement(qn('a:prstDash'), {'val': 'dash'}))
    x_v = l + w * valor
    _rect(slide, x_v - 0.035, yb - 0.08, 0.07, alto + 0.16, ONCO_DARK)
    add_textbox(slide, x_v - 0.80, t - 0.03, 1.60, 0.28,
                [(_num(valor), 16, True, ONCO_DARK, F_BODY, PP_ALIGN.CENTER)],
                anchor=MSO_ANCHOR.MIDDLE)
    ye = yb + alto + 0.05
    add_textbox(slide, l - 0.24, ye, 0.48, 0.22,
                [("0", fs, False, GRIS_BODY, F_BODY, PP_ALIGN.CENTER)])
    add_textbox(slide, x_az - 0.80, ye, 1.60, 0.22,
                [("azar = 0,5", fs, True, ONCO_INK, F_BODY, PP_ALIGN.CENTER)])
    add_textbox(slide, l + w - 0.24, ye, 0.48, 0.22,
                [("1", fs, False, GRIS_BODY, F_BODY, PP_ALIGN.CENTER)])
    if pie:
        add_textbox(slide, l, ye + 0.26, w, 0.24,
                    [(pie, 11, False, INK, F_BODY, PP_ALIGN.CENTER)])
    return ye + (0.50 if pie else 0.22) - t


# ============================================================================
# Figuras nativas del rediseño del 5-ago (SI-MIL pedagógico)
# ============================================================================
def _mezcla(c0, c1, f):
    """Interpola dos colores de la paleta. Sirve para dar INTENSIDAD sin inventar tonos
    nuevos: todos los intermedios caen sobre la recta que une dos colores del template."""
    return RGBColor(*(int(round(a + (b - a) * f)) for a, b in
                      zip((c0[0], c0[1], c0[2]), (c1[0], c1[1], c1[2]))))


def barras_esquema(slide, l, t, w, h, alturas, col, destacados=(), col_dest=None):
    """Tira de barras verticales de altura dada, como ESQUEMA de un vector de pesos.

    No es un dato medido y no lleva eje ni valores: lo que tiene que leerse es la FORMA
    del vector (parejo antes de la compuerta, casi todo en cero después). Por eso las
    alturas se escriben a mano y son las mismas en cada regenerada."""
    n = len(alturas)
    paso = w / n
    ancho = paso * 0.62
    base = t + h
    for i, a in enumerate(alturas):
        alto = max(0.012, a * h)
        c = (col_dest or ONCO_DARK) if i in destacados else col
        _rect(slide, l + i * paso + (paso - ancho) / 2, base - alto, ancho, alto, c)
    return base


# Vector de pesos ANTES y DESPUÉS de la compuerta de la ecuación 5. Escritos a mano y
# fijos: son un esquema del mecanismo del paper, no una medición nuestra.
BETA_ANTES = [.52, .61, .48, .70, .55, .66, .44, .58, .74, .50, .63, .47, .68, .54, .60,
              .45, .72, .56, .49, .64, .53, .67, .46, .59, .71, .51, .62, .57, .43, .65]
BETA_DESPUES = [.04, .06, .03, .92, .05, .07, .03, .04, .98, .05, .08, .03, .06, .04, .05,
                .03, .95, .06, .04, .07, .05, .88, .03, .05, .09, .04, .06, .05, .03, .07]
BETA_VIVOS = (3, 8, 16, 21)


def grilla_contribuciones(slide, l, t, w, h, filas, cols, patron, celda_max=0.20):
    """La ecuación 9 dibujada: una celda por (parche, medición), con SIGNO.

    El punto de la ecuación es que la predicción se descompone en una suma de términos
    legibles, y que cada término tiene dirección. Por eso el color codifica el signo (teal
    = empuja hacia la clase, gris = en contra) y el tamaño codifica cuánto. Los valores son
    un esquema del mecanismo, no números del paper ni nuestros: no se rotula ninguno."""
    lab_w = 1.16
    x0 = l + lab_w
    paso_x = (w - lab_w) / len(cols)
    paso_y = h / len(filas)
    for j, c in enumerate(cols):
        add_textbox(slide, x0 + j * paso_x, t - 0.30, paso_x, 0.28,
                    [(c, 8.5, True, GRIS_BODY, F_BODY, PP_ALIGN.CENTER)],
                    anchor=MSO_ANCHOR.BOTTOM)
    for i, f in enumerate(filas):
        cy = t + i * paso_y + paso_y / 2
        add_textbox(slide, l, cy - 0.13, lab_w - 0.10, 0.26,
                    [(f, 9, False, INK, F_BODY, PP_ALIGN.RIGHT)], anchor=MSO_ANCHOR.MIDDLE)
        for j in range(len(cols)):
            v = patron[i][j]
            cx = x0 + j * paso_x + paso_x / 2
            if v == 0:
                # una contribución nula se DIBUJA, chiquita y casi blanca: dejar el hueco
                # vacío hace que la grilla se lea rota en vez de leerse como un cero
                _rect(slide, cx - 0.03, cy - 0.03, 0.06, 0.06,
                      _mezcla(TEAL_CARD2, ONCO_DATA, 0.18))
                continue
            lado = 0.07 + abs(v) / 3.0 * (celda_max - 0.07)
            col = _mezcla(TEAL_CARD2, ONCO_DARK if v > 0 else ONCO_DATA,
                          0.35 + 0.65 * abs(v) / 3.0)
            _rect(slide, cx - lado / 2, cy - lado / 2, lado, lado, col)
    return t + h


def destilacion(slide, l, t, w, h):
    """La ecuación 10 en tres bloques: las dos ramas se entrenan a la vez y la interpretable
    persigue a la profunda. Sin ese tercer término serían dos modelos en paralelo."""
    bw = (w - 0.30) / 2
    _proc(slide, l, t, bw, h, "Rama profunda", dim="su error", size=10)
    _proc_claro(slide, l + bw + 0.30, t, bw, h, "Rama interpretable", dim="su error", size=10)
    _conn_dash(slide, l + bw, t + h * 0.5, l + bw + 0.30, t + h * 0.5)
    add_textbox(slide, l, t + h + 0.04, w, 0.24,
                [("y un tercer término la empuja a parecerse a la profunda   (10)",
                  9, True, ONCO_DARK, F_BODY, PP_ALIGN.CENTER)])
    return t + h + 0.28


# ============================================================================
# Reflow, tipografía y re-base al tamaño del template
# ============================================================================
CONTENT_TOP_OLD = 0.86
CONTENT_TOP_NEW = ONCO_BAND + 0.05        # 1.221
SAFE_BOTTOM = SH - 0.14


def _scale_block(el, f):
    """Escala un shape completo por f: geometría, cuerpo tipográfico y métrica de tabla.
    El alto real de una tabla lo manda el texto de sus filas, no el alto del shape, así que
    hay que bajar también el `sz` y el alto de fila para que encoja de verdad."""
    for off in el.iter(qn("a:off")):
        off.set("x", str(int(round(int(off.get("x")) * f))))
        off.set("y", str(int(round(int(off.get("y")) * f))))
    for ext in el.iter(qn("a:ext")):
        ext.set("cx", str(int(round(int(ext.get("cx")) * f))))
        ext.set("cy", str(int(round(int(ext.get("cy")) * f))))
    for tag in ("a:rPr", "a:defRPr", "a:endParaRPr"):
        for rpr in el.iter(qn(tag)):
            if rpr.get("sz"):
                rpr.set("sz", str(max(100, int(round(int(rpr.get("sz")) * f)))))
    for gc in el.iter(qn("a:gridCol")):
        if gc.get("w"):
            gc.set("w", str(int(round(int(gc.get("w")) * f))))
    for tr in el.iter(qn("a:tr")):
        if tr.get("h"):
            tr.set("h", str(int(round(int(tr.get("h")) * f))))


def reflow_onco(prs, skip=()):
    """Reancla el contenido de las láminas con cabecera OncoMets bajo la banda teal, y lo
    comprime si no entra. Como acá se maqueta ya desde TOP=1.24, el ajuste típico es nulo:
    la función queda como red de seguridad."""
    for slide in prs.slides:
        if id(slide._element) in skip:
            continue
        cuerpo = [sh for sh in slide.shapes if not (sh.name or "").startswith("ONCOHDR_")]
        if len(cuerpo) == len(list(slide.shapes)):
            continue
        tops = [Emu(sh.top).inches for sh in cuerpo]
        bots = [Emu(sh.top).inches + Emu(sh.height).inches for sh in cuerpo]
        if not tops:
            continue
        top0, bot0 = min(tops), max(bots)
        f = 1.0
        if bot0 + (CONTENT_TOP_NEW - top0) > SAFE_BOTTOM and bot0 > top0:
            f = max(0.80, (SAFE_BOTTOM - CONTENT_TOP_NEW) / (bot0 - top0))
        if f < 1.0:
            for sh in cuerpo:
                _scale_block(sh._element, f)
        desplaz = CONTENT_TOP_NEW - top0 * f
        for sh in cuerpo:
            sh.top = Inches(Emu(sh.top).inches + desplaz)


def forzar_barlow(prs, fuente=F_BODY):
    """Deja Barlow como única tipografía del archivo (pedido de Sebastián).

    Los runs propios ya salen en Barlow, pero quedan tres focos fuera de alcance directo:
    el `endParaRPr`/`buFont` de las láminas heredadas, el fontScheme del theme (que en
    este template es el de Office, o sea Arial) y los `defRPr` del master y los layouts.
    Se normaliza el XML entero."""
    A = "{http://schemas.openxmlformats.org/drawingml/2006/main}"
    tags = tuple(A + t for t in ("latin", "ea", "cs", "sym", "buFont"))

    def normaliza(root):
        tocados = 0
        for el in root.iter():
            if el.tag in tags and el.get("typeface") != fuente:
                el.set("typeface", fuente)
                tocados += 1
        return tocados

    cola = [prs.part] + [s.part for s in prs.slides]
    cola += [s.notes_slide.part for s in prs.slides if s.has_notes_slide]
    vistos, partes = set(), []
    while cola:
        p = cola.pop()
        if id(p) in vistos:
            continue
        vistos.add(id(p))
        partes.append(p)
        for rel in p.rels.values():
            if not rel.is_external and rel.reltype.endswith(
                    ("slideLayout", "slideMaster", "notesMaster", "theme")):
                cola.append(rel.target_part)
    n = 0
    for p in partes:
        el = getattr(p, "_element", None)
        if el is not None:
            n += normaliza(el)
            continue
        if str(p.partname).startswith("/ppt/theme/"):
            raiz = etree.fromstring(p.blob)
            tocados = 0
            for cual in ("majorFont", "minorFont"):
                for fs in raiz.iter(A + cual):
                    for lat in fs.findall(A + "latin"):
                        if lat.get("typeface") != fuente:
                            lat.set("typeface", fuente)
                            tocados += 1
            if tocados:
                p._blob = etree.tostring(raiz, xml_declaration=True,
                                         encoding="UTF-8", standalone=True)
                n += tocados
    print("  tipografía: %d referencias forzadas a %s" % (n, fuente))


def scale_deck_to_1610(prs, k=13.333 / 10.0, skip=()):
    """Escala el deck entero por k y fija el tamaño de lámina en 13.333 x 7.5. Misma
    relación de aspecto que 10 x 5.625, así que la escala es uniforme y no deforma.
    `skip` = láminas heredadas, que ya están a 13.333."""
    for slide in prs.slides:
        if id(slide._element) in skip:
            continue
        tree = slide.shapes._spTree
        for off in tree.iter(qn("a:off")):
            off.set("x", str(int(round(int(off.get("x")) * k))))
            off.set("y", str(int(round(int(off.get("y")) * k))))
        for ext in tree.iter(qn("a:ext")):
            ext.set("cx", str(int(round(int(ext.get("cx")) * k))))
            ext.set("cy", str(int(round(int(ext.get("cy")) * k))))
        for tag in ("a:rPr", "a:defRPr", "a:endParaRPr"):
            for rpr in tree.iter(qn(tag)):
                if rpr.get("sz"):
                    rpr.set("sz", str(max(100, int(round(int(rpr.get("sz")) * k)))))
        for gc in tree.iter(qn("a:gridCol")):
            if gc.get("w"):
                gc.set("w", str(int(round(int(gc.get("w")) * k))))
        for tr in tree.iter(qn("a:tr")):
            if tr.get("h"):
                tr.set("h", str(int(round(int(tr.get("h")) * k))))
        for ln in tree.iter(qn("a:ln")):
            if ln.get("w"):
                ln.set("w", str(int(round(int(ln.get("w")) * k))))
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)


def _borde_inferior(problemas, idx, sh, t, h, b_ef, bar_t, propias):
    """Marca el shape que CRUZA la barra de remate.

    Las dos piezas de la barra se excluyen por identidad y no por posición: `reflow_onco`
    le sube el texto por encima de su propia línea, así que por posición la barra se
    marcaría a sí misma en casi toda lámina. Y terminar JUSTO sobre la línea no es
    cruzarla, por eso el margen va hacia abajo."""
    if bar_t is None or id(sh._element) in propias:
        return
    if t < bar_t - 0.01 and b_ef > bar_t + 0.02:
        problemas.append("s%02d  cruza la barra de remate: %s baja hasta %.2f\" "
                         "(la barra va en %.2f)" % (idx, sh.shape_type, b_ef, bar_t))


def auditar(prs, skip=()):
    """Chequeo de defectos que el ojo ve y un chequeo ingenuo no: texto que no entra en su
    caja, shapes fuera del lienzo, cuerpos por debajo del mínimo del template (7 pt) y
    objetos que cruzan la barra de remate.

    Se corre ANTES de escalar, o sea en el espacio de 10 x 5.625. No reemplaza mirar las
    láminas ([[deck-qa-puntos-ciegos-chequeo]]), pero caza la clase de defecto que más
    apareció acá: la última línea de un panel quedando afuera."""
    problemas = []
    for idx, slide in enumerate(prs.slides, start=1):
        if id(slide._element) in skip:
            continue
        par = _TAKEAWAY.get(id(slide._element))
        bar_t, propias = None, ()
        if par is not None:
            linea, tb = par
            bar_t = Emu(linea.top).inches
            propias = (id(linea._element), id(tb._element))
        for sh in slide.shapes:
            try:
                l, t = Emu(sh.left).inches, Emu(sh.top).inches
                w, h = Emu(sh.width).inches, Emu(sh.height).inches
            except TypeError:
                continue
            # Un shape rotado 270° reporta su bbox SIN rotar: da falso positivo de límites.
            if not getattr(sh, "rotation", 0):
                if l < -0.02 or t < -0.02 or l + w > SW + 0.02 or t + h > SH + 0.02:
                    problemas.append("s%02d  fuera del lienzo: %s (%.2f, %.2f, %.2f x %.2f)"
                                     % (idx, sh.shape_type, l, t, w, h))
            if not sh.has_text_frame:
                _borde_inferior(problemas, idx, sh, t, h, t + h, bar_t, propias)
                continue
            alto, chico = 0.0, None
            for p in sh.text_frame.paragraphs:
                if not p.runs:
                    continue
                sz = max((r.font.size.pt for r in p.runs if r.font.size), default=12)
                bold = any(r.font.bold for r in p.runs)
                txt = "".join(r.text for r in p.runs)
                # los runs de sub/superíndice vienen a 0.74x: no son un cuerpo chico real
                if sz >= 6 and (chico is None or sz < chico):
                    if not any(r._r.get_or_add_rPr().get("baseline") for r in p.runs):
                        chico = sz
                alto += wrap_lines(txt, max(w - 0.14, 0.2), sz, bold) * sz * 1.22 / 72.0
            # Una caja de texto no se ve: lo que se ve es el texto adentro. Medir la caja
            # marcaba captions de una línea sobre una caja de 0,4" que no cruzan nada, y un
            # chequeo que grita de más se termina ignorando, que es justo el modo de falla
            # que este chequeo vino a tapar. Un panel pintado sí se ve entero.
            b_ef = t + h
            if sh.shape_type == MSO_SHAPE_TYPE.TEXT_BOX and alto < h:
                anc = sh.text_frame.vertical_anchor
                if anc == MSO_ANCHOR.MIDDLE:
                    b_ef = t + (h + alto) / 2
                elif anc != MSO_ANCHOR.BOTTOM:
                    b_ef = t + alto
            _borde_inferior(problemas, idx, sh, t, h, b_ef, bar_t, propias)
            if alto > h + 0.06:
                problemas.append("s%02d  texto que no entra: sobra %.2f\" en «%s…»"
                                 % (idx, alto - h, sh.text_frame.text[:44].replace("\n", " ")))
            if chico is not None and chico < 7.0:
                problemas.append("s%02d  cuerpo por debajo del mínimo: %.1f pt" % (idx, chico))
    if problemas:
        print("  AUDITORÍA: %d avisos" % len(problemas))
        for p in problemas:
            print("   ·", p)
    else:
        print("  AUDITORÍA: sin avisos")
    return problemas


def _set_solo_run(par, texto):
    """Deja el párrafo con un solo run, conservando el formato del primero. El texto del
    template viene partido en varios runs, así que escribir solo runs[0] dejaría la cola
    del original pegada detrás."""
    runs = par.runs
    runs[0].text = texto
    for r in runs[1:]:
        r._r.getparent().remove(r._r)


def retitular_portada(prs):
    """Ajusta la lámina de apertura heredada del template sin redibujarla.

    Desde el 22-ago la lámina de TÍTULO del template ya no viaja en el deck (`TPL_KEEP`),
    así que el sprint y la fecha, que iban ahí, se rotulan sobre la propia portada. Va como
    caption al pie y no como shape del template: la portada es nativa a 13.333 y está en
    `keep_ids`, o sea que `scale_deck_to_1610` NO la escala, y por eso las coordenadas de
    este rótulo son las finales y no las de trabajo."""
    portada = prs.slides[0]
    for sh in portada.shapes:
        if not sh.has_text_frame:
            continue
        txt = sh.text_frame.text.strip()
        if txt.startswith("Care in people"):
            # Marcador sin reemplazar que viene del propio template.
            for par in list(sh.text_frame.paragraphs)[1:]:
                par._p.getparent().remove(par._p)
        elif txt.startswith("OncoMETS is an AI-powered platform"):
            # El párrafo arranca tan abajo que se sale por el borde inferior.
            sh.top = Inches(5.62)
            sh.left = Inches(0.28)
            sh.width = Inches(6.30)
    tb = portada.shapes.add_textbox(Inches(6.90), Inches(6.06), Inches(6.10), Inches(0.62))
    _set_runs(tb.text_frame, [
        ("OncoMets · Sprint 8", 16, True, ONCO_DARK, F_BODY, PP_ALIGN.RIGHT),
        (FECHA_REUNION, 12, False, GRIS_BODY, F_BODY, PP_ALIGN.RIGHT)],
        anchor=MSO_ANCHOR.BOTTOM)
    return portada


# ============================================================================
# Build
# ============================================================================
# ============================================================================
# Las láminas: una función por lámina
# ============================================================================
# Refactor del 7-ago. `build()` era un bloque único de mil líneas, así que cambiar
# el orden del deck significaba mover mil líneas y confiar en que la indentación
# sobreviviera. Ahora cada lámina es una función y `build()` es la lista de
# llamadas: el orden se lee de un vistazo y el próximo reordenamiento es mover
# renglones de una línea. El cuerpo de cada bloque quedó donde estaba, a cuatro
# espacios, porque ya lo estaba dentro de `build()`.
def lam_portada(prs):
    """Portada de marca, heredada del template."""
    s = retitular_portada(prs)
    notes(s, "Traigo cuatro cosas, y las puse en ese orden a propósito.\n"
             "\n"
             "Empiezo por un encargo que quedó del sprint pasado y ya cerró: si al modelo con "
             "expertos le sobra capacidad por dentro, y de qué lado conviene recortarla. Da un "
             "resultado negativo y lo dejo cerrado ahí mismo.\n"
             "\n"
             "Sigue una medición nuestra: si el modelo mira donde el patólogo marcó las mitosis. "
             "Cambió hacia dónde apunta el trabajo, y me voy a detener en cómo se mide, porque el "
             "número hay que poder defenderlo.\n"
             "\n"
             "Después, la revisión de uno de los papers encargados, que propone un modelo que se "
             "explica solo. Es una lectura terminada y la cuento comprimida.\n"
             "\n"
             "Y al final lo que más quiero discutir: un segundo modelo, que trabaja a nivel de "
             "núcleo, corriendo sobre el mapa de atención del nuestro. Es el encadenamiento que "
             "quedó pedido en la reunión anterior, y ya tiene mapas, número y cuatro láminas.")


def lam_objetivos_sprint(prs):
    """Los objetivos del sprint, en la posición 2.

    Reemplaza a la lámina de título del template. Hereda el molde EXACTO de la lámina de
    cierre que el 22-ago se retiró (`lam_objetivos_propuestos`), que a su vez replicaba el de
    la lámina «Objetivos del sprint» de las presentaciones anteriores: título 32 pt, lista
    numerada en infinitivo, marcador de estado a la derecha, y nada más.

    Lo único que cambia respecto de aquella es el marcador: ahí eran seis «Propuesto»
    iguales, y acá el estado es el real de cada objetivo a la fecha de la reunión."""
    s = content(prs, "Objetivos del sprint", size=32)
    row_tops = [1.10 + i * 0.68 for i in range(len(OBJETIVOS_SPRINT))]
    row_h = 0.62
    for (item, hecho), rt in zip(OBJETIVOS_SPRINT, row_tops):
        add_textbox(s, 0.35, rt, 7.75, row_h, [(item, 19, True, GRIS_BODY, F_BODY)],
                    anchor=MSO_ANCHOR.MIDDLE)
        if hecho:
            status_done(s, 8.98, rt + row_h / 2)
        else:
            status_progress(s, 8.98, rt + row_h / 2, texto="En curso")
    notes(s, "Los seis objetivos del sprint, y en qué quedó cada uno.\n"
             "Los cuatro primeros estaban propuestos desde el arranque y están cerrados.\n"
             "Los dos últimos aparecieron durante el sprint y siguen en curso.\n"
             "\n"
             "Los dos primeros son el encargo del modelo con expertos: escalar a la tarea completa "
             "la medición de capacidad, que veníamos arrastrando de siete láminas, y decidir de qué "
             "lado conviene recortar. Los dos cerrados, y el primero es la lámina que sigue.\n"
             "\n"
             "El tercero es la medición de atención contra las marcas del patólogo, y el cuarto, "
             "elegir y justificar hacia dónde sigue la rama de mitosis. También cerrados, y ocupan "
             "el grueso de lo que traigo.\n"
             "\n"
             "Los dos últimos no estaban en el plan de marzo: aparecieron durante el sprint, porque "
             "el tercero salió como salió. Investigar e implementar un modelo de detección de "
             "mitosis por núcleo, y encadenarlo con la atención para cruzar sus detecciones contra "
             "las marcas. Están en curso: el modelo ya corrió y ya hay número, pero medido sobre una "
             "sola lámina, y el barrido de las doce anotadas sigue en la cola.")


def lam_grid(prs):
    # ---- El grid E×S, en una sola lámina ----
    # Fusión pedida por Ernesto el 5-ago: las dos láminas del grid pasan a una, con los tres
    # diagramas y sin los renglones de prosa que los acompañaban. El veredicto sube al primer
    # renglón, que es donde estaba la barra de remate, y las dos lecturas de forma (el escalón
    # con meseta de una rama, la curva no monótona de la otra) bajan al guion.
    #
    # Las dos reglas del pre-registro siguen gobernando:
    #   - El veredicto es H_nula y se cuenta como tal. El +0,022 del primer peldaño es
    #     justamente lo que NO alcanza; presentarlo como hallazgo contradiría el prereg.
    #   - CERO Δ contra CLAM por brazo. El prereg §6 lo prohibió por diseño para no disparar
    #     ocho veces sobre el eje ya cerrado del Hallazgo 12, y encima sobre la tarea del dato
    #     abierto. Por eso CLAM no aparece.
    # Recorte del 22-ago (`correcciones.txt`): se van las dos líneas de prosa que Ernesto
    # marcó. La del veredicto («el signo se da vuelta…») porque la figura, ya rotulada, lo
    # muestra sola y el guion lo dice; y la coletilla del eje desde 0,72, que describía un
    # detalle de dibujo y no un dato. El alto que dejan libre se lo lleva el diagrama
    # divergente, que es lo que costaba leer.
    s = content(prs, "¿Recortar expertos o slots?")
    add_textbox(s, 0.36, TOP, 9.28, 0.22, [
        ("Presencia de carcinoma ductal in situ · 862 láminas, 730 con presencia y 132 sin · "
         "5 particiones.", 9, False, GRIS_BODY, F_BODY, PP_ALIGN.CENTER)])
    barras_divergentes(s, 0.36, TOP + 0.94, 9.28, 1.62, PELDANOS, destacar=0)
    for x, titulo in ((0.36, "Sacando slots, con los 30 expertos fijos"),
                      (5.22, "Sacando expertos, con los 10 slots fijos")):
        add_textbox(s, x, TOP + 2.76, 4.42, 0.24,
                    [(titulo, 11.5, True, ONCO_DARK, F_BODY, PP_ALIGN.CENTER)])
    escalera_capacidad(s, 0.36, TOP + 3.04, 4.42, 0.86, RAMA_S, lo=0.72, hi=0.845)
    escalera_capacidad(s, 5.22, TOP + 3.04, 4.42, 0.86, RAMA_E, lo=0.72, hi=0.845)
    notes(s, "Un encargo que quedó del sprint pasado y que este sprint cerró.\n"
             "El modelo tiene 30 expertos y 10 unidades por experto: 300 en total.\n"
             "A igual capacidad total, ¿conviene recortar por un lado o por el otro?\n"
             "Arriba, la diferencia pareada; abajo, qué cuesta ir sacando capacidad.\n"
             "\n"
             "Empiezo por acá porque es lo más corto. El sprint pasado medimos cómo se reparte el "
             "peso entre esas trescientas unidades sobre siete láminas; lo escalamos a las mil "
             "ciento setenta y seis de prueba y el número aguanta: se ocupan unas ciento sesenta "
             "unidades y los treinta expertos se usan los treinta. De ahí salía que, si sobraba "
             "capacidad, sobraba del lado de las unidades. La figura de arriba lo pone a prueba, a "
             "igual capacidad total y sobre las mismas particiones.\n"
             "\n"
             "Cómo se lee. Arriba está escrita la magnitud: la diferencia de área bajo la curva "
             "entre los dos recortes, partición por partición. La línea vertical del medio es el "
             "cero, y la escala llega a una décima de área bajo la curva hacia cada lado. La barra "
             "oscura es el promedio, y sale del cero hacia el lado que gana: a la derecha si "
             "conviene recortar unidades, a la izquierda si conviene recortar expertos. La línea "
             "gris que la atraviesa es cuánto se mueve entre las cinco particiones. Y los cinco "
             "cuadraditos son las particiones: relleno oscuro cuando ésa votó a favor de recortar "
             "unidades, claro cuando votó al revés.\n"
             "\n"
             "Con eso los tres pares se leen solos. El primero va a favor de recortar unidades, dos "
             "centésimas; el segundo se da vuelta; el tercero es prácticamente cero, y en los tres "
             "la línea gris cruza el cero. Si la dirección importara, el signo sería el mismo en los "
             "tres.\n"
             "\n"
             "Abajo, qué cuesta ir sacando capacidad. Del lado de las unidades hay un escalón y "
             "después una meseta, sin ningún quiebre donde el total cae sobre la ocupación medida; "
             "del lado de los expertos ni siquiera baja de forma ordenada, y el peor caso de la "
             "tanda es el recorte más chico. Con eso cierro el encargo.")


def lam_pregunta_medible(prs):
    # ---- La pregunta y el estadístico, en una sola lámina ----
    # Fusión pedida por Ernesto el 5-ago: la lámina de la pregunta se quedaba en bullets (dos
    # paneles de hipótesis y una fila de fichas) y la del estadístico tenía el dibujo. Se
    # juntan, los dos paneles bajan a UNA línea, y lo que ocupa la lámina son las dos cintas
    # de parches reordenados, que es lo que hay que ver. Las hipótesis se cuentan hablando.
    # Los rótulos son los del pre-registro §2 y NO se asignan por cuál ganó: la primaria es
    # la del patólogo, o sea la que el resultado terminó refutando.
    # Recorte del 22-ago (`correcciones.txt`). Cuatro cosas salen de la lámina y bajan al
    # guion, y las cuatro por el mismo motivo: eran metodología escrita, no lo que hay que
    # mirar. La cita del patólogo, porque el título nuevo ya la enuncia como pregunta. El
    # nombre del estadístico («la U de Mann-Whitney normalizada»), que Ernesto marcó como
    # jerga sin anclaje. La línea de las dos lecturas pre-registradas. Y las tres tarjetas.
    # Con eso la lámina queda con lo único que sí hay que ver, que son las dos cintas, y se
    # las agranda de 0,42" a 0,78" de alto ([[deck-contenido-visual-no-bullets]]).
    s = content(prs, "¿La atención de CLAM cae sobre los núcleos de la mitosis?")
    add_textbox(s, 0.36, TOP, 9.28, 0.56, [
        ("Se ordenan los 4799 parches por atención y se mira dónde cayeron los 28 que "
         "contienen una mitosis marcada.", 12, True, INK, F_BODY, PP_ALIGN.CENTER),
        ("El estadístico es el AUC de ranking. Cada grupo de tejido tiene el suyo; el de "
         "las cintas es el de mitosis.", 9.5, False, GRIS_BODY, F_BODY, PP_ALIGN.CENTER)])
    # El eje «más atención / menos atención» va una sola vez, debajo de las dos cintas: el
    # orden es el mismo y repetirlo bajo cada una lo vuelve ruido.
    # Los rótulos de las cintas eran `_rot_label` a 270°, y se PISABAN entre sí: un rótulo
    # rotado ocupa a lo alto lo que mide de ancho (1,60"), y las dos cintas están a 0,56" una
    # de otra. El bbox que reporta el shape es el de antes de rotar, así que la auditoría no
    # podía verlo ([[deck-qa-puntos-ciegos-chequeo]]) y salió del rasterizado. Ahora van
    # horizontales en la calle de la izquierda, que para eso se corrieron las cintas.
    ct1, ct2, CH = TOP + 0.96, TOP + 1.92, 0.78
    for ct, txt, col in ((ct1, "SI FUERA\nAZAR", GRIS_BODY),
                         (ct2, "LO\nOBSERVADO", ONCO_DARK)):
        # 0,92" de caja: a 9 pt «OBSERVADO» mide 0,68" y con los 0,80" de antes quedaba
        # 0,66" útiles, así que el rasterizado lo partía en «OBSERVAD / O».
        add_textbox(s, 0.28, ct, 0.92, CH,
                    [(ln, 9, True, col, F_BODY, PP_ALIGN.RIGHT) for ln in txt.split("\n")],
                    anchor=MSO_ANCHOR.MIDDLE)
    cinta_ranking(s, 1.24, ct1, 7.10, marcados={2, 7, 11, 16, 21, 27, 31}, h=CH, ejes=False)
    add_textbox(s, 8.46, ct1, 1.20, CH,
                [("0,5", 19, True, GRIS_BODY, F_BODY, PP_ALIGN.CENTER)],
                anchor=MSO_ANCHOR.MIDDLE)
    cinta_ranking(s, 1.24, ct2, 7.10, marcados={0, 1, 2, 4, 5, 7, 9}, h=CH)
    add_textbox(s, 8.46, ct2, 1.20, CH,
                [("0,89", 19, True, ONCO_DARK, F_BODY, PP_ALIGN.CENTER)],
                anchor=MSO_ANCHOR.MIDDLE)
    # La cuenta de pares, que es la definición del estadístico escrita con los números de
    # esta lámina (`atencion_vs_anotaciones.py:139`: n_neg = n_total − n_pos = 4799 − 28).
    # Iba comprimida en una fórmula de una línea y Ernesto marcó que no se entendía: pasa a
    # dos renglones de lenguaje corriente, que es como se cuenta hablando.
    add_textbox(s, 0.36, TOP + 3.44, 9.28, 0.72, [
        ("Se compara cada uno de los 28 parches con mitosis contra cada uno de los 4771 "
         "restantes.", 13, True, INK, F_BODY, PP_ALIGN.CENTER),
        ("En el 89 % de esas comparaciones, el parche con mitosis recibió más atención.",
         13, True, TEAL_TITLE, F_BODY, PP_ALIGN.CENTER)])
    notes(s, "Lo que medimos no es un mapa de calor: es un número, y sale de un ranking.\n"
             "Se ordenan los 4799 parches por atención y se mira dónde cayeron los 28 de mitosis.\n"
             "Si la atención no supiera nada daría 0,5; observamos 0,89.\n"
             "\n"
             "Esto empieza con una frase del patólogo: en mitosis los núcleos son finos y dispersos, "
             "y quizá esos parches no reciben atención suficiente.\n"
             "\n"
             "Esa frase se puede medir, y por eso la traigo. No es una opinión sobre el modelo: es "
             "una afirmación sobre dónde cae la atención. Y la podemos contrastar porque tenemos las "
             "marcas del patólogo, o sea que sabemos en qué parches hay una mitosis. Sin ellas solo "
             "tendríamos un mapa de calor, y de un mapa de calor uno no puede decir si el rojo está "
             "donde corresponde. Las marcas convierten una impresión en un número.\n"
             "\n"
             "Antes de correr nada dejamos escritas las dos respuestas posibles: que los parches "
             "marcados no rankearan mejor que el azar, y el problema estaría en cómo el modelo los "
             "combina; o que rankearan alto, y entonces lo que se pierde está antes. Quedó en pie la "
             "segunda.\n"
             "\n"
             "Cómo se mide. Se ordenan los cuatro mil setecientos noventa y nueve parches por la "
             "atención que recibieron. Cada casilla de las dos cintas es un parche, y los oscuros "
             "son los veintiocho con una mitosis marcada. Si la atención no supiera nada de mitosis "
             "estarían repartidos por toda la fila, que es la cinta de arriba y da cero coma cinco; "
             "si se concentrara ahí se amontonarían a la izquierda, que es la de abajo.\n"
             "\n"
             "El número sale de una cuenta simple. Comparo cada uno de esos veintiocho parches "
             "contra cada uno de los cuatro mil setecientos setenta y uno restantes y miro cuál "
             "recibió más atención: en el ochenta y nueve por ciento de las comparaciones ganó el "
             "parche con mitosis. Es un estadístico clásico para comparar dos grupos sin suponer "
             "nada sobre sus distribuciones; solo mira el orden, así que sirve entre grupos de "
             "tamaños muy distintos, y su referencia es cero coma cinco siempre.")


def lam_mapas(prs):
    # ---- Atención y marcas del patólogo ----
    # Rediseño del 4-ago, y es el cambio que Ernesto pidió con más énfasis. Se proyecta SOLO
    # la región anotada, que ocupa el doble de lado: la grilla 2x2 anterior arrastraba la
    # fila de la región sin marcas, cuyo panel de anotaciones es tejido pelado y se leía como
    # la misma imagen dos veces. Con esa fila fuera, el panel que explicaba las dos regiones
    # ya no explica nada de lo que se ve, así que se va entero al guion: una leyenda no
    # arregla una figura confusa, y la parte confusa acá no aportaba
    # ([[deck-qa-puntos-ciegos-chequeo]]).
    s = content(prs, "Atención y marcas del patólogo")
    # Caja con la relación de aspecto EXACTA del asset (1502 x 624 = 2,407), para que
    # add_image_fit no deje aire a los costados y los rótulos caigan sobre su panel. El ancho
    # tiene tope: con más de 8,10 la figura no deja lugar para las tres líneas del pie.
    MAP_L, MAP_W = 0.95, 8.10
    MAP_H = MAP_W / 2.407
    MAP_GAP = 22 / 1502.0 * MAP_W        # el aire de 22 px del montaje, en pulgadas
    col_w = (MAP_W - MAP_GAP) / 2
    add_image_fit(s, FIG_MAPAS_ANOTADA, MAP_L, TOP + 0.28, MAP_W, MAP_H, align="top")
    for i, rot in enumerate(("Atención del modelo", "Marcas del patólogo")):
        cx = MAP_L + i * (col_w + MAP_GAP) + col_w / 2
        caption(s, cx - 1.40, TOP, 2.80, rot, size=12, col=TEAL_TITLE, bold=True)
    # La procedencia al pie, que es lo que preguntó Sebastián: de qué lámina salió esto y con
    # qué se midió. Va en un solo textbox porque tres `caption` gastarían 1,2" de alto.
    pie_lineas(s, 0.36, TOP + 0.28 + MAP_H + 0.10, 9.28, PROVENANCIA, size=8)
    notes(s, "La región anotada, con la atención del modelo y las marcas del patólogo.\n"
             "A la izquierda pone el color el modelo; a la derecha, las marcas las puso el patólogo.\n"
             "Al pie, de dónde salió la lámina y qué tiene marcado.\n"
             "\n"
             "A la izquierda está la atención del modelo sobre el tejido: el rojo es mucha atención, "
             "el azul poca. A la derecha, las marcas del patólogo, cada color un tipo de tejido.\n"
             "\n"
             "Al pie está de dónde sale todo esto. Es una lámina de nuestra cohorte privada, "
             "escaneada a veinte aumentos, de la que salen cuatro mil setecientos noventa y nueve "
             "parches de doscientos cincuenta y seis píxeles de lado. El patólogo dibujó sesenta y "
             "un polígonos en QuPath: veintiséis de mitosis, catorce de núcleos de alto grado, seis "
             "de necrosis, cinco de células inmunes, cinco de tumor, dos de tejido adiposo, uno de "
             "estroma y dos de fondo. Debajo de alguno de ellos caen ciento sesenta y tres parches, "
             "y de ésos, veintiocho contienen una mitosis.\n"
             "\n"
             "En la última línea están las etiquetas de esta lámina en nuestros datos, las de "
             "reporte, que en los CSV se escriben como score dos y score tres: tasa mitótica alta, "
             "pleomorfismo nuclear intermedio, diferenciación tubular alta, grado general tres. La primera es coherente con veintiséis mitosis marcadas, y "
             "es la que le vamos a pedir al modelo dentro de dos láminas.\n"
             "\n"
             "Sobre si las marcas coinciden con lo que mira el modelo, mirando de a dos uno intuye "
             "que caen en zonas calientes. Pero intuir no es medir, y por eso el resultado es el "
             "número de la lámina anterior y no esta imagen. Lo que agrega la imagen es que el "
             "número no vive de una casualidad de una esquina.")


def lam_escalera(prs):
    # ---- El resultado: la escalera de los siete grupos ----
    # Recorte del 22-ago (`correcciones.txt`). Se van las DOS tarjetas de incertidumbre y el
    # renglón que nombraba el intervalo al 95 %: Ernesto marcó que ni «el bigote es el
    # intervalo al 95 %» ni el caso del estroma se entendían en la lámina. Lo que queda es
    # qué se está mostrando y cómo se midió, dicho sin nombrar el estadístico. Las dos
    # incertidumbres bajan al guion, donde siguen enteras ([[auc-atencion-dos-incertidumbres]]).
    # El alto liberado (0,74" de las tarjetas) se lo lleva la escalera, que es el objeto.
    s = content(prs, "El resultado, grupo por grupo")
    barras_ranking(s, 0.36, TOP + 0.40, 9.28, 3.00, ESCALERA)
    add_textbox(s, 0.36, TOP + 3.54, 9.28, 0.60, [
        ("Cada barra es el AUC de la atención sobre los parches que el patólogo marcó de esa "
         "clase, medidos contra todo el resto de la lámina.", 11.5, True, INK, F_BODY,
         PP_ALIGN.CENTER),
        ("La línea sobre cada barra es la precisión que da el número de parches marcados: "
         "cuantos menos hay, más larga.", 9.5, False, GRIS_BODY, F_BODY, PP_ALIGN.CENTER)])
    notes(s, "Una barra por clase de tejido marcada, con el mismo estadístico de recién.\n"
             "Mitosis queda arriba de todo, y por encima incluso de tumor.\n"
             "La escalera baja hasta la grasa, que el modelo evita.\n"
             "La línea sobre cada barra dice cuánta precisión da el número de parches de ese grupo.\n"
             "\n"
             "Cada barra es una de las clases marcadas, y su valor es el estadístico de recién "
             "calculado con los parches de esa clase: la probabilidad de que un parche de ese tejido "
             "reciba más atención que uno cualquiera del resto de la lámina. Cada clase se mide "
             "contra todo lo demás, nunca contra otra clase.\n"
             "\n"
             "Mitosis da cero coma ochocientos noventa, muy lejos del azar y por encima de tumor, "
             "que son regiones grandes y bien delimitadas, mucho más fáciles de acertar que "
             "veintiocho parches sueltos. Pero lo que convence es la escalera completa. Abajo de "
             "todo está el tejido adiposo, en cero coma ciento cincuenta y cuatro: tan por debajo "
             "del azar significa que el modelo lo evita, no que lo ignore. Arriba, tumor y núcleos "
             "de alto grado alrededor de cero coma ochocientos veinticinco. Ese orden no lo diseñó "
             "nadie, salió de medir.\n"
             "\n"
             "Ahora la línea sobre cada barra, porque sin ella siete barras del mismo grosor parecen "
             "siete números de la misma calidad. Hay dos incertidumbres. Una es de modelo: si repito "
             "la medición con otros checkpoints, con las marcas fijas, el valor de mitosis se mueve "
             "unas cuatro centésimas. La otra es la dibujada, y es la grande: cuánto se movería si "
             "el patólogo hubiera marcado otras mitosis en vez de éstas. Depende de cuántos parches "
             "tiene el grupo, y en mitosis vale unas ocho centésimas, el doble que la primera.\n"
             "\n"
             "El caso que obliga a mirarla es el estroma: con doce parches su línea va de cero coma "
             "treinta y siete a cero coma setenta, así que esta lámina no distingue estroma evitado "
             "de estroma atendido. Mitosis aguanta, porque su línea no se acerca al azar. Los "
             "núcleos de alto grado tampoco los presento como resultado: su línea es la segunda más "
             "larga.")


def lam_28_parches(prs):
    # ---- Dónde caen los 28 parches ----
    s = content(prs, "Los 28 parches de mitosis")
    # Mismo criterio que la lámina anterior: cada caja con la relación de aspecto exacta de
    # su figura, y las dos a la MISMA altura, para que los pies queden en una sola línea.
    # Pedido 6 de Ernesto (6-ago): acá va la cuenta que lleva de 26 marcas a 28 parches. Va
    # ARRIBA, pegada al título, porque el título dice 28 y la pregunta nace ahí. El costo es
    # que las dos figuras bajan de 3,34 a 2,68 de alto (−20 %); si se prefiere el tamaño
    # anterior, lo que hay que mover es esta banda, no las figuras.
    #
    # El neto de +2 es coincidencia de ESTA lámina y no se presenta como regla: son dos
    # efectos grandes que casi se cancelan (`atencion_vs_patologo/resultados.md` §1.a).
    cadena_cuenta(s, 0.36, TOP, 9.28, CUENTA_26_28, h=0.32, fs=13)
    FIG_H = 2.68
    reg_w = FIG_H * (800 / 676.0)
    zoom_w = FIG_H * (453 / 552.0)
    FIG_T = TOP + 0.74
    add_image_fit(s, FIG_MITOSIS, 0.36, FIG_T, reg_w, FIG_H, align="top")
    add_image_fit(s, FIG_ZOOM, 0.36 + reg_w + 0.15, FIG_T, zoom_w, FIG_H, align="top")
    # Una línea: a dos renglones el segundo cae bajo la barra de remate y queda tachado.
    caption(s, 0.36, FIG_T + FIG_H + 0.04, reg_w,
            "la región anotada, con los 28 parches en blanco", size=9.5)
    caption(s, 0.36 + reg_w + 0.15, FIG_T + FIG_H + 0.04, zoom_w, "el detalle del recuadro",
            size=9.5)
    # El panel de dos renglones baja a UNA línea, y sin caja: un panel alrededor de una sola
    # frase es una caja alrededor de nada, y acá la lámina la manda la figura.
    COL_L = 0.36 + reg_w + zoom_w + 0.46
    COL_W = 9.64 - COL_L
    add_textbox(s, COL_L, FIG_T + 0.10, COL_W, 1.10,
                [("Los parches blancos caen sobre el rojo, no sobre el borde ni sobre el "
                  "azul.", 13.5, True, ONCO_DARK, F_BODY)], anchor=MSO_ANCHOR.MIDDLE)
    # La escala, que es la causa de fondo de la cuenta de arriba Y el puente a la lámina
    # siguiente: lo que el modelo comprime a un vector es un parche que es casi todo otra cosa.
    add_textbox(s, COL_L, FIG_T + 1.34, COL_W, 1.20, [
        ("Una mitosis ocupa entre el 2 % y el 4 % del parche.", 12, True, TEAL_TITLE, F_BODY),
        ("36 px de lado contra los 256 del parche. De ahí que 10 de las 26 caigan sobre un "
         "borde, y que el vector con el que el modelo representa el parche resuma, sobre "
         "todo, tejido que no es la mitosis.", 9.5, False, GRIS_BODY, F_BODY)])
    takeaway_bar(s, "Y aun así, este modelo predice mal esta lámina.", t=TOP + 3.78,
                 size=13)
    notes(s, "Arriba, de dónde salen 28 parches si el patólogo dibujó 26 marcas.\n"
             "En blanco, los 28 parches que contienen una mitosis marcada.\n"
             "Debajo, en colores, dónde puso la atención el modelo.\n"
             "Los blancos caen sobre el rojo.\n"
             "\n"
             "Despejo primero la cuenta de arriba, porque los dos números conviven en todo el "
             "trabajo y parecen contradecirse. El patólogo dibujó veintiséis marcas y yo vengo "
             "hablando de veintiocho parches: son unidades distintas. Diez marcas caen sobre el "
             "borde entre dos parches y ocupan los dos; y siete parches tienen más de una mitosis "
             "adentro.\n"
             "\n"
             "La causa de fondo es de escala. Una marca de mitosis mide treinta y seis píxeles de "
             "lado y el parche mide doscientos cincuenta y seis: la mitosis ocupa entre el dos y el "
             "cuatro por ciento del área del parche. Con ese tamaño, que diez de veintiséis caigan "
             "sobre un borde deja de sorprender.\n"
             "\n"
             "Y hay una implicación más seria. El modelo no ve el parche: ve un vector de quinientos "
             "doce números en el que el extractor lo resume entero. Si la mitosis ocupa dos de cada "
             "cien píxeles, su aporte a ese resumen es muy chico y el resto lo llena el tejido de "
             "alrededor. El parche puede estar bien elegido y la evidencia que justifica elegirlo "
             "puede haberse diluido en el camino.\n"
             "\n"
             "Abajo, el mapa de atención con los veintiocho parches de mitosis en blanco encima. El "
             "color lo pone el modelo y los cuadrados los pone el patólogo, así que la pregunta es "
             "si los blancos caen sobre el rojo: no están en el borde ni repartidos, están sobre el "
             "corazón de la mancha.\n"
             "\n"
             "Y este mismo modelo se equivoca al clasificar esta lámina.")


def lam_mira_responde(prs):
    # ---- Mira bien y responde mal ----
    # Tabla rehecha el 5-ago: Ernesto no entendía a qué modelo correspondía cada fila, ni
    # por qué de la corrida de cinco particiones aparecían dos y no cinco. La columna ahora
    # nombra el DATASET DE ENTRENAMIENTO con su tamaño (contado sobre los splits reales, ver
    # CKPTS_TABLA) y el pie dice por qué son esas dos particiones: son las únicas donde esta
    # lámina cayó en validación.
    s = content(prs, "La atención acierta y la clasificación falla")
    # Pedido del 22-ago: decir CUÁLES son las 4 clases de la tarea, que el deck daba por
    # sabidas. Verificadas contra `environ/csv/dataset_grado_histologico_tasa_mitotica_label.csv`.
    add_textbox(s, 0.36, TOP, 9.28, 0.78, [
        ("Cuatro modelos de tasa mitótica. La lámina cayó en validación en los cuatro: "
         "ninguno la tuvo en entrenamiento.", 12, True, GRIS_BODY, F_BODY, PP_ALIGN.CENTER),
        ("Las 4 clases de la tarea: score_1 tasa baja (636 láminas) · score_2 intermedia "
         "(287) · score_3 alta (254) · no_identificado, que el informe no consigna (693).",
         9.5, False, GRIS_BODY, F_BODY, PP_ALIGN.CENTER),
        ("La etiqueta verdadera es score_3, tasa mitótica alta, coherente con las 26 marcas "
         "de mitosis.", 11, False, GRIS_BODY, F_BODY, PP_ALIGN.CENTER)])
    simple_table(s, 0.50, TOP + 0.88, 9.00,
                 ["Entrenado con", "Qué respondió", "Confianza en score_2",
                  "AUC de atención"],
                 CKPTS_TABLA,
                 col_fracs=[0.34, 0.20, 0.23, 0.23], row_h=0.34, fs=11, destacar=2)
    # Los dos paneles se funden en UNA lectura: la tabla ya dice quién falla y quién mira
    # mejor, así que enumerarlo al lado era leerla en voz alta. Lo que hay que agregar es la
    # consecuencia, y eso es una frase. El resto (los 8 modelos que sí la vieron, el cambio
    # de diagnóstico) baja al guion.
    add_textbox(s, 0.36, TOP + 2.92, 9.28, 0.96, [
        ("El modelo que mejor localiza las mitosis es el que peor clasifica la lámina.",
         13.5, True, TEAL_TITLE, F_BODY, PP_ALIGN.CENTER),
        ("El problema no está en elegir los parches: está en la información que sobrevive "
         "cuando el parche se resume en un vector.", 13.5, True, TEAL_TITLE, F_BODY,
         PP_ALIGN.CENTER)], anchor=MSO_ANCHOR.MIDDLE)
    notes(s, "Cuatro modelos de tasa mitótica, ninguno la tuvo en entrenamiento.\n"
             "La última columna es cuánto mira cada uno las mitosis: todos miran bien.\n"
             "La fila destacada es el que mejor mira, y es el que responde peor.\n"
             "\n"
             "Primero, qué tarea es ésta. Tasa mitótica es uno de los tres componentes del grado "
             "histológico, y en nuestros datos tiene cuatro etiquetas, que en la tabla aparecen "
             "como score uno, dos y tres más no identificado: tasa baja, seiscientas "
             "treinta y seis láminas; intermedia, doscientas ochenta y siete; alta, doscientas "
             "cincuenta y cuatro; y no identificado, seiscientas noventa y tres, que son aquellas "
             "cuyo informe no consigna el dato. Esta lámina es de tasa alta.\n"
             "\n"
             "En la tabla hay cuatro modelos entrenados sobre esa tarea, con distinta cantidad de "
             "datos, y ninguno vio esta lámina en entrenamiento. Tres de los cuatro responden tasa "
             "intermedia: se equivocan hacia abajo.\n"
             "\n"
             "La fila destacada hace la lámina. Ese modelo mira mejor que los otros tres, con cero "
             "coma novecientos veintiséis en la última columna, y se equivoca con más convicción, "
             "con un setenta y uno por ciento de confianza en la respuesta equivocada.\n"
             "\n"
             "Suena a paradoja y no lo es. El modelo trabaja en dos pasos: primero elige qué parches "
             "importan, después combina sus vectores en una única respuesta. La última columna mide "
             "el primer paso, y está bien resuelto; falla el segundo, y falla antes de combinar "
             "nada, porque la evidencia ya venía diluida en cada parche. No hay que enseñarle a "
             "mirar, hay que conseguir que lo que mira sobreviva a la compresión.")


def lam_simil_propone(prs):
    # ---- Qué propone + la figura del paper ----
    # Rediseño del 5-ago (pedido de Ernesto): los bullets al mínimo para que la figura del
    # paper crezca todo lo que da. La figura es más ancha que alta (ar 3,195), así que en la
    # caja anterior de 9,56 x 2,52 la limitaba el ALTO y se dibujaba a 8,05 x 2,52. Ahora la
    # limita el ancho de la lámina y sale a 9,60 x 3,00, un 19 % más de lado. El bloque de
    # tres líneas de arriba baja a una frase y la barra de remate se retira: su texto ERA esa
    # frase, y repetirla costaba media pulgada de figura.
    s = content(prs, "SI-MIL: qué propone")
    add_textbox(s, 0.36, TOP, 9.28, 0.50, [
        ("Que la explicación SEA la predicción: predecir con una combinación lineal de "
         "mediciones de núcleos que tienen nombre de patología, y dejar la red profunda como "
         "el maestro que enseña dónde mirar.", 13.5, True, TEAL_TITLE, F_BODY,
         PP_ALIGN.CENTER)], anchor=MSO_ANCHOR.MIDDLE)
    FIG_H_S3 = 9.60 / 3.195
    add_image_fit(s, FIG2_FULL, 0.20, TOP + 0.56, 9.60, FIG_H_S3, align="top")
    pie_lineas(s, 0.36, TOP + 0.56 + FIG_H_S3 + 0.08, 9.28, [
        "Taming Deep MIL for Self-Interpretability in Gigapixel Histopathology · Kapse et al., "
        "CVPR 2024.",
        "(a) el recorrido completo, de la lámina a las dos predicciones · (b) y (c) el detalle "
        "de cada rama.",
    ], size=9)
    notes(s, "La ambición del paper: que la explicación SEA la predicción.\n"
             "Dos caminos que salen de la misma lámina: uno mide, el otro selecciona.\n"
             "La rama profunda no llega a producción: se descarta entera.\n"
             "\n"
             "Paso al paper encargado, que llega al mismo problema por otro camino: nosotros miramos "
             "un modelo ya entrenado para ver dónde puso la atención, y ellos proponen construirlo "
             "para que no haya que ir a mirarlo.\n"
             "\n"
             "La idea es que la explicación sea la predicción, no un anexo. La red profunda deja de "
             "predecir y su único producto pasa a ser una selección de veinte parches, los más "
             "atendidos. La predicción se calcula aparte, con una combinación lineal de mediciones "
             "que un patólogo puede leer, y en producción la rama profunda se descarta entera.\n"
             "\n"
             "Esas mediciones son doscientas cuarenta y seis y no salen de la nada. Antes pasa por "
             "el parche un segmentador de núcleos, HoVer-Net, que separa cada núcleo de sus vecinos "
             "y le pone una de cinco clases: epitelial neoplásico, conectivo, inflamatorio, "
             "necrótico y epitelial no neoplásico. Sobre ese mapa se calculan las mediciones. "
             "Doscientas cinco son morfométricas: diez propiedades por núcleo, como el área, la "
             "excentricidad o la solidez; cuatro estadísticos que resumen cada propiedad en el "
             "parche, la media, la desviación, la asimetría y la curtosis; y las cinco clases. Diez "
             "por cuatro por cinco, doscientas, más cinco conteos, uno por clase. Las cuarenta y una "
             "restantes describen cómo están distribuidos en el espacio: medidas de grafo, tratando "
             "a las células como nodos, y medidas de mezcla de tipos celulares.\n"
             "\n"
             "Cada una tiene nombre legible: una es la asimetría de la solidez de los núcleos "
             "neoplásicos, otra la mezcla de células conectivas dentro de la región neoplásica. Por "
             "eso pueden decir que el reporte del modelo es la predicción misma.\n"
             "\n"
             "En la figura, los dos caminos que bajan de la lámina describen los mismos parches de "
             "dos maneras, y la caja del medio es la bisagra: toma la atención del camino profundo y "
             "con ella elige los veinte parches sobre el otro camino.")


def lam_simil_resultados(prs):
    # ---- Resultados y costo de adopción ----
    # Rediseño del 4-ago: la tabla es el objeto de la lámina, así que se agranda a todo el
    # ancho y se queda sola. Los tres paneles del inventario de costos y las dos tarjetas de
    # preguntas salieron: su contenido ya estaba entero en el guion, y acá competían con la
    # única fila que hay que leer ([[deck-contenido-visual-no-bullets]]).
    s = content(prs, "Resultados y costo de adopción")
    caption(s, 0.36, TOP, 9.28,
            "Su Tabla 2: el mismo método aplicado sobre distintos modelos de base, en "
            "TCGA-BRCA", size=10.5, col=GRIS_BODY, bold=True)
    simple_table(s, 0.36, TOP + 0.34, 9.28,
                 ["Modelo de base", "Profundo  Acc / AUC", "Con SI-MIL  Acc / AUC"],
                 [["ABMIL", "0,937 / 0,974", "0,944 / 0,968"],
                  ["CLAM (el nuestro)", "0,937 / 0,972", "0,925 / 0,957"],
                  ["TransMIL", "0,934 / 0,936", "0,929 / 0,933"]],
                 col_fracs=[0.32, 0.34, 0.34], row_h=0.50, fs=13, destacar=1)
    add_textbox(s, 0.36, TOP + 2.54, 9.28, 0.62, [
        ("Sobre nuestro modelo de base bajan las dos: la exactitud y el área bajo la curva. "
         "El titular de que no hay compromiso entre rendimiento e interpretabilidad se "
         "sostiene sobre el modelo más simple.", 13, True, TEAL_TITLE, F_BODY,
         PP_ALIGN.CENTER)], anchor=MSO_ANCHOR.MIDDLE)
    takeaway_bar(s, "Publican el dataset ya procesado, y la cohorte pública de mama está "
                    "adentro: verificar el cruce es una tarde.", t=TOP + 3.34, size=13)
    notes(s, "La tabla compara el método sobre distintos modelos de base, y uno es el nuestro.\n"
             "Sobre el nuestro, la fila destacada, las dos métricas bajan.\n"
             "Publican el dataset ya procesado, y la cohorte de mama está adentro.\n"
             "\n"
             "Esta tabla nos interesa más que las otras porque no compara el método contra otros "
             "métodos, sino aplicado sobre distintos modelos de base, y uno de ellos es el nuestro.\n"
             "\n"
             "Sobre el primero, el modelo de atención más simple, la exactitud sube un poco al "
             "agregarle la rama interpretable. Sobre el segundo, que es el nuestro y es la fila "
             "destacada, bajan las dos: la exactitud algo más de un punto y el área bajo la curva un "
             "punto y medio. Lo digo sin ánimo de desacreditar el trabajo: el titular de que no hay "
             "compromiso entre rendimiento e interpretabilidad se sostiene sobre el primer modelo, y "
             "la fila que nos correspondería es la que baja.\n"
             "\n"
             "La comparación con lo nuestro es nítida sin ser una competencia. La interpretación de "
             "ellos aparece durante el entrenamiento, sobre mediciones que tienen nombre desde "
             "antes; la nuestra aparece después, sobre modelos ya congelados. Por eso el costo de "
             "equivocarse es distinto: allá empeora el modelo, acá queda mal la descripción.\n"
             "\n"
             "Si quisiéramos probarlo acá, el costo más duro es de escala física: el segmentador de "
             "núcleos está entrenado a cuarenta aumentos y solo a cuarenta, tanto que ellos mismos "
             "filtraron sus datasets para quedarse con láminas de esa magnificación, y nuestras "
             "cohortes están a escalas distintas. El costo de cómputo se desarma solo: publican el "
             "dataset ya procesado y la cohorte pública de mama está adentro.\n"
             "\n"
             "Y dejo una pregunta de fondo. Cuando le mostraron los reportes a un patólogo, algo más "
             "de un cuarto de las mediciones que el modelo declara importantes le resultaron no "
             "relevantes. Me parece muy honesto que lo publiquen, y la pregunta es si ese número es "
             "aceptable para nuestro estándar clínico.")


# ============================================================================
# El orden del deck
# ============================================================================
# Datos de la sección del 19-ago
# ============================================================================
# Sale de sprints/B8_sprint8/hovernext_129741/cruce_marcas.md. Se escribe a mano acá para
# que la lámina no dependa de leer los CSV en tiempo de build.
#
# El recorte del 19-ago se llevó los datos de las regiones de escaneo (alcance, perfiles,
# sensibilidad, probe de rotación), la curva del techo y su barrido de tolerancia, y el
# estado de la cola de cómputo. Están en el historial: no se re-escriben acá.

# cruce_marcas.md §3: los dos factores del techo y su intersección. La unidad son MARCAS
# (26), NO los 28 parches con marca: una marca puede caer sobre dos parches vecinos, así
# que las dos cuentas no se encadenan.
# La escalera de lo que cuesta cada brazo, en área revisada. El área sale de la geometría
# real del h5: parche de 256 px a 0,465 µm/px = 0,0142 mm². Las marcas recuperadas de
# techo_conjunto.csv (brazo CLAM), y el 13 de la lámina entera verificado también contra las
# 82 detecciones de la región anotada sola: da lo mismo, porque las 26 marcas caen todas ahí.
ESCALERA_AREA = [["La lámina entera, sin recorte", "4799", "68,0 mm²", "177", "13 de 26"],
                 ["Solo la región anotada", "2496", "35,4 mm²", "82", "13 de 26"],
                 ["El 12 % más atendido por CLAM", "300", "4,3 mm²", "48", "11 de 26"]]

# Simplificada el 22-ago de 5 columnas a 3 (`correcciones.txt`: «no se entienden las
# columnas»). Se van «Ambas» —que era la misma cuenta que la última— y «Cota mín( )», una
# cota que el propio pie declaraba floja. Quedan las dos que sostienen la lectura: cuántas
# marcas caen dentro del recorte, y de ésas cuántas además detecta HoVer-NeXt.
# Verificadas contra `cruce_marcas/techo_conjunto.csv`, brazo CLAM, columnas
# `en_mascara` y `ambas`.
# ---------------------------------------------------------------------------
# Datos de la Fase A (21-ago), que son las cuatro láminas nuevas del 22-ago
# ---------------------------------------------------------------------------
# A2.bis §2 — la comparación a CARGA FIJA: cuántos objetos pide cada brazo para llegar al
# mismo recall. A K fijo la comparación es tramposa, porque una máscara más grande gana por
# ser más grande ([[carga-fija-no-k-fijo]]). Verificada contra
# `escalera_brazos/escalera_brazos.csv`, interpolando entre los K del barrido.
#
# NO va la fila de 26/26 que traía el plan con «2496» en las cuatro columnas. Ese 2496 es el
# CHEQUEO DE SANIDAD (la región entera), no la carga mínima para llegar a 26/26: el CSV
# muestra que Mammoth llega a 26/26 en K=750 y CLAM en K=1392. Poner 2496 como «carga para
# llegar a 26/26» sería falso en la propia unidad de la tabla. El tope del detector, que es
# lo que esa fila venía a mostrar, queda dicho en el bloque de conclusión y en el pie.
ESCALERA_BRAZOS = [["8 de 26", "71 parches", "70", "33", "80", "no aplica"],
                   ["13 de 26", "130", "95", "84", "126", "82 núcleos"],
                   ["19 de 26", "300", "217", "241", "251", "imposible"],
                   ["22 de 26", "450", "300", "382", "372", "imposible"]]

# A2 §3 — marcas dentro del top-K por atención, confinado a la región anotada, denominador
# 26 marcas. Verificado contra `cruce_marcas_gate/techo_conjunto.csv` (columna `en_mascara`,
# brazo Mammoth) y `cruce_marcas/techo_conjunto.csv` (mismo brazo, tarea CDIS).
GATE_TOPK = [["4,0 %", "1 de 26", "14 de 26"],
             ["7,6 %", "8 de 26", "18 de 26"],
             ["12,0 %", "11 de 26", "22 de 26"],
             ["20,0 %", "17 de 26", "23 de 26"]]

# A0 — las 13 acreditadas contra las 13 falladas. El grupo acreditado es el CONTROL POSITIVO
# que hace legible al fallado (patrón P3). Verificado contra
# `a0_falladas/marcas_vs_instancias.csv`: las dos mitades salen indistinguibles en todo
# menos en la etiqueta que salió de la cabeza de clase.
A0_FILAS = [["Núcleo segmentado sobre la marca", "13 de 13", "13 de 13"],
            ["Distancia al centroide, mediana", "1,9 µm", "2,1 µm"],
            ["Clase que le puso el detector", "mitosis ×13",
             "epithelial-cell ×12, neutrophil ×1"],
            ["Instancias dentro de 15 µm, mediana", "9", "9"]]

CRUCE_CONJUNTA = [["4,0 % de la región", "12 de 26", "8"],
                  ["7,6 %", "15 de 26", "10"],
                  ["12,0 %", "19 de 26", "11"],
                  ["la región entera", "26 de 26", "13"]]


# ============================================================================
# Láminas del 19-ago
# ============================================================================
# Ernesto recortó el deck el 19-ago (`correcciones.txt`): se fueron la sección entera de
# regiones de escaneo, la lámina que fechaba el trabajo, la de la cola de cómputo, la de
# patrones y la de próximos pasos. De las dos que medían los factores del techo queda UNA,
# de resultados, con la forma que pidió: los mapas primero y el número al lado.

def lam_hovernext_paper(prs):
    # ---- La herramienta, contada con la figura del paper ----
    # Ernesto pidió la figura del diagrama original en lugar de la lámina de estado. La
    # figura ES el contenido, así que se le da todo el alto disponible y lo que la acompaña
    # es una columna angosta: lo que hicimos nosotros con ella, no lo que dice el paper.
    s = content(prs, "HoVer-NeXt y la clase de mitosis")
    # La figura ES el contenido, así que se lleva todo el alto disponible y no se le pone
    # rótulo arriba: el título de la lámina ya dice qué es (CLAUDE.md, ADDENDUM 19-jul).
    # Recorte del 22-ago: Ernesto pidió sacar los dos paneles y ampliar la figura, cuyo
    # recorrido quiere oír paso por paso en el guion.
    #
    # El plan proponía llevarla a los 9,28" de ancho de la caja útil, y eso NO entra: con la
    # razón 1,621 del asset, 9,28" de ancho piden 5,72" de alto y el cuerpo mide 4,22". Acá
    # manda el ALTO, no el ancho. Con la barra de remate puesta el techo son ~3,2" de alto
    # (5,25" de ancho), o sea un 7 % más que antes, que no es «ampliar». Así que la barra
    # también se retira, que además es lo que hace el propio template en sus láminas de
    # arquitectura: sin subtítulo y sin barra (CLAUDE.md, ADDENDUM 19-jul). Con eso la
    # figura pasa de 4,90" a 6,10" de ancho, un 24 % más de lado.
    FIG_H = 3.76
    FIG_W = FIG_H * 1.621                      # razón exacta del asset (2301 x 1419)
    FIG_L = 0.35 + (9.28 - FIG_W) / 2          # centrada en la caja útil
    add_image_fit(s, FIG_HN_PAPER, FIG_L, TOP + 0.04, FIG_W, FIG_H, align="top")
    caption(s, 0.35, TOP + 0.04 + FIG_H + 0.06, 9.28,
            "Figura 1 de Baumann et al., MIDL 2024 · reproducida del paper", size=8)
    notes(s, "Un segmentador de núcleos con una clase de mitosis entre las suyas.\n"
             "La figura del paper: el teselado, la red, y el cosido de las teselas.\n"
             "Encaja sin reescalar: trabaja a media micra por píxel, como nuestras láminas privadas.\n"
             "\n"
             "Vuelvo a nuestra línea. La idea viene de la reunión anterior: si el modelo mira donde "
             "hay mitosis y aun así falla, quizá lo que falta no es atención sino detalle, y "
             "conviene un instrumento a nivel de núcleo.\n"
             "\n"
             "La figura es del paper y muestra las tres piezas. Arriba, el flujo completo: la lámina "
             "se parte en teselas, cada tesela pasa por la red, y un segundo componente cose las "
             "salidas en una sola imagen del tamaño de la lámina. Abajo a la izquierda, la red, que "
             "es un codificador y dos decodificadores que lo comparten. Y abajo a la derecha, el "
             "cosido: las teselas se procesan con solapamiento, así que hay que resolver los núcleos "
             "del borde para no contarlos dos veces ni partirlos.\n"
             "\n"
             "Las dimensiones. La entrada es una tesela de tres canales, rojo, verde y azul, "
             "muestreada a media micra por píxel. El codificador es una ConvNeXt versión dos, tamaño "
             "tiny, preentrenada en imágenes naturales. El primer decodificador es el de instancia y "
             "saca cinco canales: dos son mapas de regresión de distancia, que permiten separar un "
             "núcleo pegado a otro, y los otros tres son una clasificación por píxel en fondo, "
             "interior del núcleo y borde del núcleo. El segundo es el de clase y saca ocho canales: "
             "fondo más las siete clases de núcleo con las que se entrenó, entre ellas mitosis. En "
             "total trece canales por píxel, y el mapa final sale de combinar las dos salidas: el de "
             "instancia dice dónde empieza y termina cada núcleo, y el de clase dice qué es.\n"
             "\n"
             "Corrijo algo que quedó dando vueltas de la última vez. Las doscientas cuarenta y seis "
             "mediciones no son de este modelo: son del paper anterior, que usa HoVer-Net, sin la "
             "equis, como front-end de segmentación. Éste no produce ningún vector de "
             "características por parche, produce el mapa de núcleos que describí.\n"
             "\n"
             "Por qué elegimos éste. La escala: está pensado para media micra por píxel y nuestras "
             "láminas privadas están casi exactamente ahí, así que no hay que reescalar nada, que es "
             "donde aparecen los problemas silenciosos. Y entre sus clases tiene una de mitosis.")


def lam_resultados(prs):
    # ---- La lámina de resultados: la cadena que pidió Sebastián, y el número ----
    # Funde las dos láminas que antes iban separadas (el techo del filtro y el cruce) en una
    # sola, y le cambia la forma: lo que se ve son los mapas, y el número los acompaña. El
    # eje es el encadenamiento — la atención recorta, el detector trabaja adentro — porque
    # eso es lo que Sebastián pidió integrar, y no dos herramientas contadas por separado.
    #
    # Tres salvedades van en el CUERPO y no solo en el guion, porque son las que más fácil
    # se pierden: la unidad de la tabla son MARCAS (26) y no los parches (28) de la lámina 7,
    # la cota del mínimo es floja, y no hay ninguna precisión calculada — las marcas son
    # positivos parciales, así que una detección sin marca no es un error.
    s = content(prs, "Resultados: recortar y detectar")

    # ---- izquierda: la cadena, y el detalle a resolución nativa ----
    # Las dos figuras se apilan y hay que dejar sitio para el pie: los altos salen de las
    # razones EXACTAS de los assets (3.694 la cadena, 4.069 el detalle), y el ancho es lo
    # que se ajusta para que la pila entre completa arriba de la barra de remate.
    xl, wl = 0.35, 4.62
    H_CAD, H_ZOOM = wl / 3.694, wl / 4.069
    GAP = 26 / 6205.0 * wl                     # el aire del montaje, en pulgadas
    pw = (wl - 2 * GAP) / 3
    # Los tres rótulos tienen que entrar en UN renglón de 1,53": el tercio es angosto y
    # `caption` no avisa cuando parte una palabra, solo la tapa con la figura de abajo.
    for i, rot in enumerate(("Atención de CLAM", "El 12 % más atendido",
                             "HoVer-NeXt detecta")):
        caption(s, xl + i * (pw + GAP), TOP + 0.00, pw, rot, size=9, col=TEAL_TITLE,
                bold=True)
    add_image_fit(s, FIG_CADENA, xl, TOP + 0.26, wl, H_CAD, align="top")
    caption(s, xl, TOP + 0.28 + H_CAD, wl,
            "Cuatro de las 13 que coinciden, a resolución nativa", size=9,
            col=TEAL_TITLE, bold=True)
    add_image_fit(s, FIG_HN_ZOOM, xl, TOP + 0.56 + H_CAD, wl, H_ZOOM, align="top")
    pie_lineas(s, xl, TOP + 0.62 + H_CAD + H_ZOOM, wl, [
        "Amarillo, la detección; blanco, la marca del patólogo. La barra son 25 µm.",
        "En el 12 % recortado caen 48 de las 82 detecciones de la región. Eso cuenta "
        "detecciones; la tabla de al lado cuenta marcas."], size=8)

    # ---- derecha: el número y los dos factores ----
    xr, wr = 5.22, 4.43
    _grupo(s, xr, TOP + 0.02, wr, 0.78, fill=TEAL_CARD)
    add_textbox(s, xr, TOP + 0.02, wr, 0.78, [
        ("13 de las 26 marcas   =   50,0 %", 17, True, ONCO_DARK, F_BODY, PP_ALIGN.CENTER),
        ("Emparejamiento uno a uno: una detección no puede acreditar dos marcas",
         9, False, GRIS_BODY, F_BODY, PP_ALIGN.CENTER)], anchor=MSO_ANCHOR.MIDDLE)
    add_textbox(s, xr, TOP + 0.90, wr, 0.26,
                [("Los dos factores, y su intersección", 12, True, ONCO_DARK, F_BODY)])
    simple_table(s, xr, TOP + 1.22, wr,
                 ["Recorte por atención", "Marcas dentro del recorte", "Y además detectadas"],
                 CRUCE_CONJUNTA, col_fracs=[0.34, 0.36, 0.30],
                 row_h=0.34, fs=10, destacar=2)
    pie_lineas(s, xr, TOP + 2.98, wr, [
        "El porcentaje es la fracción de los 2496 parches de la región anotada que quedan "
        "encendidos: el 12 % son 4,25 mm².",
        "Unidad: marcas del patólogo (26), no los 28 parches con marca.",
        "Sin precisión, y es deliberado: las marcas son positivos parciales."], size=8)

    takeaway_bar(s, "Desde el 7,6 % de la región el cuello ya no es el recorte, es el "
                    "detector")
    notes(s, "La atención de CLAM recorta la región, y el detector de núcleos trabaja adentro.\n"
             "Recupera 13 de las 26 marcas del patólogo: la mitad.\n"
             "Desde el 7,6 % de la región el cuello deja de ser el recorte y pasa a ser el detector.\n"
             "Las 13 que faltan no fallan por poco: no hay ninguna detección de mitosis cerca.\n"
             "\n"
             "Esto es lo que quedó pedido en la reunión anterior: integrar el modelo de atención al "
             "camino de detección.\n"
             "\n"
             "La cadena tiene tres pasos, y cada uno lo hace una herramienta distinta. Primero "
             "nuestro modelo mira la lámina entera y reparte atención entre sus parches. Segundo, el "
             "recorte: nos quedamos con el doce por ciento más atendido y apagamos el resto. Ese "
             "porcentaje ordena toda la lámina: la región anotada tiene dos mil cuatrocientos "
             "noventa y seis parches, así que el doce por ciento son trescientos parches encendidos, "
             "y en superficie, cuatro coma veinticinco milímetros cuadrados. Tercero, el detector de "
             "núcleos trabaja adentro de ese recorte, y es un modelo ajeno que no sabe nada del "
             "nuestro ni de nuestra tarea.\n"
             "\n"
             "Lo que produce es un inventario de núcleos: en esta lámina segmentó más de doscientas "
             "treinta y ocho mil células y a ciento setenta y siete les puso mitosis. Cada marca "
             "amarilla es un núcleo, uno solo; las blancas son las veintiséis mitosis que dibujó el "
             "patólogo. Dos fuentes independientes sobre la misma imagen, y toda la lámina consiste "
             "en compararlas.\n"
             "\n"
             "Los tres paneles de la izquierda son la misma región en tres estados: la atención, el "
             "recorte, y encima lo que encontró el detector. Las detecciones se agrupan sobre el "
             "recorte, o sea que la atención y el detector miran el mismo tejido sin haberse "
             "consultado. Abajo, cuatro coincidencias a resolución nativa: no están de acuerdo sobre "
             "una zona, están de acuerdo sobre una célula.\n"
             "\n"
             "El número es trece de veintiséis, la mitad, con emparejamiento uno a uno: si uno le "
             "pregunta a cada marca cuál es la detección más cercana el número sube, pero son "
             "siempre las mismas trece contadas varias veces, y una detección no puede acreditar dos "
             "marcas.\n"
             "\n"
             "Las otras trece no fallan por poco: movimos la tolerancia en un rango de diez veces y "
             "el número se queda clavado en trece, y la detección de mitosis más cercana a una marca "
             "fallada está a ciento quince micras de mediana, unos seis núcleos. No apunta al lado, "
             "no hay nada ahí. Averiguamos por qué, y es la lámina dieciséis.\n"
             "\n"
             "La tabla junta los dos límites de la prueba, y es lo que cambió la conversación. Cada "
             "fila es un tamaño de recorte; una columna dice cuántas marcas entran, y la de al lado "
             "cuántas de ésas además se detectan. Veníamos discutiendo cuán grande convenía hacer el "
             "recorte, y a partir del siete coma seis por ciento de la región eso deja de ser el "
             "problema: por generoso que uno lo haga, el techo se queda en trece. El factor que "
             "manda pasó a ser la detección.\n"
             "\n"
             "Una cosa que este número no dice. La mitad no es la exhaustividad del detector en "
             "mitosis: el denominador son las marcas del patólogo, no las mitosis que hay en la "
             "lámina. Tampoco calculamos precisión, y es deliberado: el patólogo marca solo donde la "
             "evidencia es clara, así que una detección sin marca no es un error.")


def lam_hovernext_solo(prs):
    # ---- El control de la cadena: la herramienta SOLA, sin recorte ----
    # Pedido de Ernesto el 19-ago, después del recorte. Es la lámina que hace legible a la
    # anterior: sin un brazo sin máscara, «13 de 26 con el recorte» no se puede comparar con
    # nada. Y el dato existe sin correr nada nuevo, porque la corrida fue así de entrada — la
    # lámina completa — y el recorte se aplicó después, sobre la salida.
    #
    # La afirmación que la lámina NO hace, y que es la fácil de colar: nada sobre las 95
    # detecciones de la región sin anotar. No tienen marcas, así que no son ni aciertos ni
    # errores. Van dibujadas del mismo color que las otras justamente por eso.
    s = content(prs, "El detector sobre la lámina completa")
    FIG_L, FIG_W = 0.35, 5.30
    FIG_H = FIG_W / 2.471                      # razón exacta del asset (3400 x 1376)
    add_image_fit(s, FIG_HN_SOLO, FIG_L, TOP + 0.26, FIG_W, FIG_H, align="top")
    gap = 26 / 3400.0 * FIG_W
    pw = (FIG_W - gap) / 2
    for i, rot in enumerate(("Región sin marcas · 95", "Región anotada · 82")):
        caption(s, FIG_L + i * (pw + gap), TOP + 0.00, pw, rot, size=9, col=TEAL_TITLE,
                bold=True)
    pie_lineas(s, FIG_L, TOP + 0.34 + FIG_H, FIG_W, [
        "Las 177 detecciones de la lámina entera, sin máscara de ninguna clase. La corrida "
        "fue así: el recorte se aplica después, sobre la salida.",
        "De la región sin marcas, la de la izquierda, no se afirma nada: sus 95 detecciones "
        "no son ni aciertos ni errores, y por eso van del mismo color que las otras."],
        size=8.5)

    xr, wr = 5.90, 3.75
    _grupo(s, xr, TOP + 0.02, wr, 0.74, fill=TEAL_CARD)
    add_textbox(s, xr, TOP + 0.02, wr, 0.74, [
        ("El mismo 13 de 26, sobre 68 mm²", 15, True, ONCO_DARK, F_BODY, PP_ALIGN.CENTER),
        ("Recortar no cambió cuántas marcas se recuperan", 9.5, False, GRIS_BODY, F_BODY,
         PP_ALIGN.CENTER)], anchor=MSO_ANCHOR.MIDDLE)
    add_textbox(s, xr, TOP + 0.88, wr, 0.26,
                [("Lo que cuesta cada brazo", 12, True, ONCO_DARK, F_BODY)])
    simple_table(s, xr, TOP + 1.18, wr,
                 ["Qué se revisa", "Parches", "Área", "Det.", "Marcas"],
                 ESCALERA_AREA, col_fracs=[0.36, 0.15, 0.18, 0.12, 0.19],
                 row_h=0.36, fs=9.5, destacar=2)
    pie_lineas(s, xr, TOP + 2.72, wr, [
        "Área desde el h5: un parche de 256 px a 0,465 µm/px son 0,0142 mm².",
        "Del brazo sin recorte al del 12 %, la superficie a revisar baja 16 veces y las "
        "marcas recuperadas bajan de 13 a 11.",
        "Las 26 marcas caen todas en la región anotada, así que restringirse a ella no "
        "pierde ninguna: por eso las dos primeras filas dan lo mismo."], size=8)
    takeaway_bar(s, "El recorte reduce la superficie a revisar, no aumenta las marcas "
                    "recuperadas")
    notes(s, "El brazo de control: el detector solo, sin nada del modelo de atención encima.\n"
             "La lámina entera son 68 mm² y 177 detecciones.\n"
             "Recupera las mismas 13 de 26 que con el recorte puesto.\n"
             "De la región sin marcas no se afirma nada.\n"
             "\n"
             "Falta el brazo de control, y va después a propósito: si digo trece de veintiséis con "
             "el recorte puesto, la pregunta inmediata es comparado con qué.\n"
             "\n"
             "No hubo que correr nada nuevo, porque la corrida fue así desde el principio: la lámina "
             "completa, y el recorte aplicado después sobre la salida. Si hubiéramos recortado antes "
             "para ahorrar cómputo, este brazo no existiría.\n"
             "\n"
             "Se ven las ciento setenta y siete detecciones sobre la lámina entera. De las noventa y "
             "cinco de la izquierda no afirmo nada: ahí no hay marcas del patólogo, así que no son "
             "ni aciertos ni errores.\n"
             "\n"
             "La tabla es lo que quiero discutir. Revisando la lámina entera, sesenta y ocho "
             "milímetros cuadrados, se recuperan trece de las veintiséis marcas. Revisando solo la "
             "región anotada, la mitad de la superficie, se recuperan las mismas trece; no es una "
             "sorpresa sino una comprobación, porque las veintiséis marcas caen todas ahí. Y "
             "revisando el doce por ciento más atendido, cuatro coma tres milímetros cuadrados, se recuperan "
             "once.\n"
             "\n"
             "Entre el primer brazo y el último la superficie a revisar baja dieciséis veces y las "
             "marcas bajan de trece a once. Eso separa dos preguntas que conviene no mezclar. Si la "
             "pregunta es cuánta superficie hay que ponerle delante a un patólogo, el recorte "
             "responde bien; si la pregunta es cuántas mitosis se encuentran, no ayuda, porque el "
             "techo lo pone el detector.\n"
             "\n"
             "Sobre el costo: la lámina entera tardó dieciocho minutos, así que recortar para "
             "ahorrar cómputo hoy no hace falta. Recortar para achicar lo que hay que revisar sí "
             "tiene sentido, y los dos motivos no piden el mismo tamaño de recorte.\n"
             "\n"
             "Las cuatro láminas que siguen son las pruebas que le hicimos a esta cadena.")


# ============================================================================
# Las cuatro láminas nuevas del 22-ago: la Fase A
# ============================================================================
# Cerró entera el 21-ago (A0, A1, A2, A2.bis) y el deck del 19-ago no tenía nada de ella.
# Las tres primeras son TABLAS NATIVAS y la cuarta lleva las láminas de contacto, que son
# fotografías de tejido y por eso caen en la excepción legítima a «todo nativo».
#
# La necrosis (A4) NO entra, por decisión de Ernesto: el encargo pedía comparar la necrosis
# del detector contra las marcas del patólogo, y esa mitad exige la clase `dead`, que vive
# solo en el otro juego de pesos (B2, sin lanzar). Sin esa mitad, A4 contesta otra pregunta.


def lam_escalera_brazos(prs):
    """A2.bis — los brazos comparados a CARGA FIJA, que es la comparación justa."""
    s = content(prs, "CLAM, Mammoth y el detector a igual carga")
    add_textbox(s, 0.36, TOP, 9.28, 0.46, [
        ("Cuántos objetos hay que mirar para llegar al mismo número de marcas recuperadas.",
         12, True, INK, F_BODY, PP_ALIGN.CENTER),
        ("A recorte fijo la comparación premia a la máscara más grande por ser más grande; "
         "a carga fija, no.", 9.5, False, GRIS_BODY, F_BODY, PP_ALIGN.CENTER)])
    simple_table(s, 0.36, TOP + 0.58, 9.28,
                 ["Llegar a", "CLAM", "Mammoth", "CLAM ∩ Mammoth", "CLAM ∪ Mammoth",
                  "HoVer-NeXt solo"],
                 ESCALERA_BRAZOS,
                 col_fracs=[0.15, 0.16, 0.15, 0.19, 0.19, 0.16],
                 row_h=0.38, fs=10.5, destacar=1)
    add_textbox(s, 0.36, TOP + 2.60, 9.28, 0.56, [
        ("Hasta 13 marcas los dos caminos piden una carga parecida, pero de distinto "
         "tamaño: un parche mide 119 µm de lado y un núcleo es un punto ya localizado.",
         12.5, True, TEAL_TITLE, F_BODY, PP_ALIGN.CENTER)], anchor=MSO_ANCHOR.MIDDLE)
    pie_lineas(s, 0.36, TOP + 3.26, 9.28, [
        "Por encima de 13 el detector no es una opción a ninguna carga: está topado ahí y no "
        "pasa de ese número ni con la lámina entera. La atención sola llega a las 26.",
        "La unión de las dos máscaras nunca es el mejor brazo a carga fija. Lo que parecía "
        "ventaja era el tamaño de la máscara, no la calidad del orden.",
        "Mammoth no se suma a CLAM: lo reemplaza. Son brazos alternativos, y esto no dice "
        "nada sobre cuál clasifica mejor, solo sobre dónde pone cada uno la atención.",
    ], size=8)
    notes(s, "Cuántos objetos hay que mirar para llegar al mismo número de marcas.\n"
             "A recorte fijo la comparación premia a la máscara más grande por ser más grande.\n"
             "Los dos caminos piden una carga parecida hasta 13 marcas, y de ahí el detector se "
             "topa.\n"
             "\n"
             "Faltaba una comparación, y es la primera: el detector, ¿agrega algo sobre la atención "
             "sola?\n"
             "\n"
             "Antes, cómo está armada la tabla, porque la forma de comparar importa más que los "
             "números. A recorte fijo, dejando entrar la misma cantidad de parches en cada brazo, "
             "gana el que tenga la máscara más grande, y gana por tamaño y no por calidad. Así que "
             "damos vuelta la pregunta: fijamos el resultado y medimos el costo, así que cada fila "
             "es un nivel de marcas recuperadas y cada celda dice cuántos objetos hay que mirar para "
             "llegar ahí.\n"
             "\n"
             "Los brazos son cuatro de atención y uno de detección. Los dos primeros son nuestro "
             "modelo y la variante con expertos, que no se suman: son alternativos, porque la "
             "variante es el mismo modelo con la primera capa cambiada. Los dos del medio son la "
             "intersección y la unión de sus máscaras. Y el último es el detector solo, sin "
             "recorte. Su columna lleva dos marcadores distintos: no aplica cuando se le pide "
             "menos de trece, porque no ordena sus detecciones y no se le puede pedir una carga "
             "menor; e imposible por encima de trece, que es su techo.\n"
             "\n"
             "Tres cosas que solo aparecen así. La unión nunca es el mejor brazo: a recorte fijo "
             "parecía el mejor de todos, y acá queda al nivel del peor de los dos que la componen, "
             "porque su ventaja era el tamaño. La intersección es la más eficiente donde la carga es "
             "chica: donde los dos modelos coinciden, aciertan. Y la tercera: hasta trece marcas los "
             "dos caminos piden una carga parecida, ochenta y dos núcleos contra ochenta y cuatro "
             "parches, pero los objetos son de tamaño muy distinto, porque un parche mide ciento "
             "diecinueve micras de lado y un núcleo es un punto ya localizado. Por encima de trece "
             "el detector deja de ser una opción a cualquier carga, mientras que la atención sola "
             "llega a las veintiséis.\n"
             "\n"
             "Esta lámina es sobre dónde pone la atención cada modelo, no sobre cuál clasifica "
             "mejor.")


def lam_gate_invasivo(prs):
    """A2 — el encargo de Sebastián: la misma cadena con el checkpoint de invasivo."""
    s = content(prs, "El checkpoint de carcinoma invasivo")
    add_textbox(s, 0.36, TOP, 9.28, 0.28, [
        ("El encargo era repetir la cadena con el otro checkpoint. Las dos lecturas no "
         "coinciden, y ésa es la lámina.", 12, True, INK, F_BODY, PP_ALIGN.CENTER)])

    # izquierda: por AUC parecen equivalentes
    xl, wl = 0.36, 4.34
    _grupo(s, xl, TOP + 0.42, wl, 1.72, fill=TEAL_CARD2)
    add_textbox(s, xl + 0.16, TOP + 0.52, wl - 0.32, 0.26,
                [("Por AUC parecen equivalentes", 12.5, True, ONCO_DARK, F_BODY)])
    add_textbox(s, xl + 0.16, TOP + 0.84, wl - 0.32, 0.92, [
        ("gate de invasivo   0,865", 14, True, ONCO_DARK, F_BODY),
        ("intervalo al 95 %: 0,778 a 0,951", 9, False, GRIS_BODY, F_BODY),
        ("CDIS   0,919", 14, True, GRIS_BODY, F_BODY),
        ("intervalo al 95 %: 0,848 a 0,989", 9, False, GRIS_BODY, F_BODY)])
    add_textbox(s, xl + 0.16, TOP + 1.80, wl - 0.32, 0.26,
                [("Los dos intervalos se solapan de sobra.", 9.5, False, GRIS_BODY, F_BODY)])

    # derecha: por recorte, no
    xr, wr = 5.30, 4.34
    add_textbox(s, xr, TOP + 0.42, wr, 0.26,
                [("Por recorte, el gate es claramente peor", 12.5, True, TEAL_TITLE, F_BODY)])
    simple_table(s, xr, TOP + 0.74, wr,
                 ["Recorte", "gate de invasivo", "CDIS"],
                 GATE_TOPK, col_fracs=[0.28, 0.38, 0.34], row_h=0.34, fs=10.5, destacar=2)
    add_textbox(s, 0.36, TOP + 2.36, 9.28, 0.52, [
        ("Un AUC casi igual y un recorte muy distinto: el AUC resume todos los umbrales y "
         "el recorte es uno solo, de los extremos.", 12.5, True, TEAL_TITLE, F_BODY,
         PP_ALIGN.CENTER)], anchor=MSO_ANCHOR.MIDDLE)
    pie_lineas(s, 0.36, TOP + 2.98, 9.28, [
        "El brazo del gate es Mammoth: el CLAM plano de esa tarea no existe en disco y lo "
        "entrena un job que sigue en cola. Así que todavía no se separa si el efecto es de "
        "la tarea o del brazo.",
        "El chequeo de sanidad pasa: en la región entera los dos brazos convergen a 13 de "
        "26, que es el factor de detección solo.",
        "Los folds difieren a propósito. En cada tarea se toma aquel donde esta lámina cayó "
        "en prueba y no fue vista en entrenamiento.",
    ], size=8)
    notes(s, "El encargo era repetir la cadena con el otro checkpoint.\n"
             "Por AUC los dos parecen equivalentes: los intervalos se solapan de sobra.\n"
             "Por recorte no: el checkpoint del gate es claramente peor.\n"
             "\n"
             "El otro encargo era repetir la cadena con el checkpoint de carcinoma invasivo, en vez "
             "del de carcinoma in situ. No hubo que volver a correr el detector: la lámina se corrió "
             "entera y el recorte se aplica después. Las dos lecturas no coinciden, y por eso la "
             "lámina está partida en dos.\n"
             "\n"
             "A la izquierda, por área bajo la curva parecen equivalentes: cero coma ochocientos "
             "sesenta y cinco el gate contra cero coma novecientos diecinueve el otro, con "
             "intervalos que se solapan de sobra. Con este número solo uno concluiría que da igual "
             "cuál usar.\n"
             "\n"
             "A la derecha, la misma pregunta hecha por recorte, que es como se usa en la práctica. "
             "En el cuatro por ciento de la región el gate mete una marca de veintiséis y el otro "
             "catorce; en el siete coma seis por ciento, ocho contra dieciocho; en el doce, once "
             "contra veintidós. Para poner las mismas marcas delante del patólogo hay que darle "
             "bastante más superficie.\n"
             "\n"
             "Por qué discrepan está en la línea del medio. El área bajo la curva resume todos los "
             "umbrales a la vez; el recorte es un umbral solo, y de los extremos. El checkpoint del "
             "gate deja los parches con mitosis en el percentil ochenta y seis y el otro en el "
             "noventa y cinco: esa diferencia casi no mueve el área bajo la curva y decide entera la "
             "lista de los trescientos más atendidos.\n"
             "\n"
             "Dos salvedades del pie. El brazo del gate es la variante con expertos, no el modelo "
             "plano: el plano de esa tarea no existe en disco y lo entrena un trabajo que sigue en "
             "la cola, así que todavía no separamos si esto es propiedad de la tarea o del brazo. La "
             "otra es el chequeo de sanidad, y pasa: en la región entera los dos brazos convergen a "
             "trece de veintiséis.")


def lam_a0_falla_la_clase(prs):
    """A0 — por qué el número se queda en 13 de 26. Va ANTES de las imágenes."""
    s = content(prs, "Las 13 que se escapan sí estaban segmentadas")
    add_textbox(s, 0.36, TOP, 9.28, 0.46, [
        ("Las 13 marcas que el detector no acredita tienen un núcleo segmentado encima. "
         "Las 13.", 13, True, TEAL_TITLE, F_BODY, PP_ALIGN.CENTER),
        ("Ninguna se escapó por falta de segmentación: se escaparon por la clase.",
         10, False, GRIS_BODY, F_BODY, PP_ALIGN.CENTER)])
    simple_table(s, 0.86, TOP + 0.68, 8.28,
                 ["", "Acreditadas (13)", "Falladas (13)"],
                 A0_FILAS, col_fracs=[0.40, 0.26, 0.34], row_h=0.44, fs=11, destacar=2)
    add_textbox(s, 0.36, TOP + 2.90, 9.28, 0.56, [
        ("Falla la cabeza de clase, no la segmentación. El objeto ya está detectado y "
         "delineado: lo que hay que arreglar es la etiqueta.", 12.5, True, TEAL_TITLE,
         F_BODY, PP_ALIGN.CENTER)], anchor=MSO_ANCHOR.MIDDLE)
    pie_lineas(s, 0.36, TOP + 3.56, 9.28, [
        "Las dos mitades salen indistinguibles en todo menos en la etiqueta: mismo tejido, "
        "misma densidad de núcleos, misma puntería de la segmentación.",
        "12 de las 13 caen en epithelial-cell, que es la clase padre: en un carcinoma de "
        "mama una figura mitótica es una célula epitelial, así que pierde la distinción "
        "fina y no la confunde con otro tejido.",
        "La explicación de fondo: la clase de mitosis del detector se entrenó y validó solo "
        "en colon, y esta lámina es de mama.",
    ], size=8)
    notes(s, "Las 13 marcas que el detector no acredita tienen un núcleo segmentado encima. Las 13.\n"
             "No se escaparon por falta de segmentación: se escaparon por la clase.\n"
             "Las acreditadas son el control, y las dos mitades salen indistinguibles salvo en la "
             "etiqueta.\n"
             "\n"
             "Vuelvo a las trece que se escapan. Habíamos cerrado en que no hay ninguna detección de "
             "mitosis cerca, y nos detuvimos ahí sin preguntar ausencia de qué.\n"
             "\n"
             "El núcleo estaba. Las trece marcas falladas tienen una instancia segmentada encima, "
             "las trece, sin excepción. Y no cerca: encima, a dos coma un micras de mediana, sobre "
             "marcas cuyo lado mediano es dieciséis coma siete micras. No falló la segmentación, "
             "falló la etiqueta que le puso la cabeza de clasificación.\n"
             "\n"
             "La columna de la izquierda son las trece acreditadas, pasadas por el mismo "
             "procedimiento: si no las recuperara a ellas, el problema sería el procedimiento y no "
             "el grupo de estudio. Las dos mitades salen indistinguibles en todo, y la única "
             "variable que las separa es la clase.\n"
             "\n"
             "Adónde fueron a parar dice algo. Doce de las trece recibieron la clase célula "
             "epitelial, que es la clase padre: en un carcinoma de mama una figura mitótica es una "
             "célula epitelial, así que el modelo no la confunde con estroma ni con un linfocito, "
             "sino que pierde la distinción fina.\n"
             "\n"
             "Esto cambia el costo del arreglo. Lo caro, barrer la lámina y segmentar doscientas "
             "treinta y ocho mil células, ya está pagado y guardado: si lo que falla es la etiqueta, "
             "alcanza con una segunda etapa de clasificación sobre los objetos que ya existen.\n"
             "\n"
             "Y la explicación de fondo es de dominio: la clase de mitosis de este detector se "
             "entrenó y se validó solo en colon, mientras que segmentar y tipificar núcleos en mama "
             "sí tiene número propio y bueno. Lo que sabe hacer fuera de su tejido lo hace bien; lo "
             "que nunca se validó fuera de colon es lo que falla.")


def lam_galeria_164(prs):
    """A1 — en qué se fija el detector. Es el pedido literal de Ernesto."""
    s = content(prs, "En qué se fija el detector")
    # Las dos láminas de contacto arriba, una al lado de la otra, y el pie a TODO el ancho
    # abajo. Con el pie metido en la columna derecha la lámina quedaba sobrecargada y
    # `reflow_onco` comprimía runs por debajo del mínimo de 7 pt del template, que es el
    # síntoma de una lámina que no entra ([[deck-contenido-visual-no-bullets]]).
    FL, FH = 0.36, 3.00
    FW = FH * 1.084                            # razón exacta del asset (1852 x 1708)
    caption(s, FL, TOP, FW, "Las 164 detecciones sin marca del patólogo", size=9.5,
            col=TEAL_TITLE, bold=True)
    add_image_fit(s, FIG_GAL_SIN_MARCA, FL, TOP + 0.26, FW, FH, align="top")

    xr = FL + FW + 0.30
    wr = 9.64 - xr
    FH2 = wr / 2.621                           # razón exacta del asset (928 x 354)
    caption(s, xr, TOP, wr, "Las 13 marcas que se escaparon", size=9.5, col=TEAL_TITLE,
            bold=True)
    add_image_fit(s, FIG_GAL_FALLADAS, xr, TOP + 0.26, wr, FH2, align="top")
    caption(s, xr, TOP + 0.32 + FH2, wr, "centradas, sin ninguna detección cerca", size=8)

    # El panel se apoya sobre el pie, que va a TODO el ancho: su alto y el hueco de arriba
    # salen de lo que queda libre hasta `pie_lineas`, no de una constante. Con 0.98" y un
    # hueco de 0.34" el panel invadía el pie 0.35" y las dos cajas se leían superpuestas
    # (ninguna desborda la suya, así que `auditar` lo da por bueno).
    yb = TOP + 0.32 + FH2 + 0.22
    _grupo(s, xr, yb, wr, 0.72, fill=TEAL_CARD)
    add_textbox(s, xr + 0.12, yb, wr - 0.24, 0.72, [
        ("145 de las 164 son el mismo tipo de recorte", 13, True, ONCO_DARK, F_BODY,
         PP_ALIGN.CENTER),
        ("el 88 %: epitelio tumoral denso, núcleo hipercromático y condensado",
         9, False, GRIS_BODY, F_BODY, PP_ALIGN.CENTER),
        ("Se despega una familia de 15, con 65 % de fondo: tejido laxo o borde de lámina.",
         9, False, GRIS_BODY, F_BODY, PP_ALIGN.CENTER)], anchor=MSO_ANCHOR.MIDDLE)

    pie_lineas(s, 0.36, TOP + 3.56, 9.28, [
        "Las 164 NO son falsos positivos: el patólogo marca solo donde la evidencia es "
        "clara, así que una detección sin marca puede ser una mitosis real sin marcar. No se "
        "calcula precisión y ningún panel las pinta como error.",
        "El parecido es de píxeles, no semántico: dos recortes vecinos comparten color y "
        "textura, no necesariamente la misma entidad biológica.",
    ], size=8)
    notes(s, "Las 164 detecciones que no acreditan ninguna marca, y las 13 que se escaparon.\n"
             "145 de las 164 son el mismo tipo de recorte: el 88 %.\n"
             "No son falsos positivos, y no se calcula precisión.\n"
             "\n"
             "Cierro con las imágenes, que es lo que pediste ver. La lámina de contacto grande son "
             "las ciento sesenta y cuatro detecciones que no acreditaron ninguna marca; ésas más las "
             "trece acreditadas son las ciento setenta y siete. Cada recorte está a resolución "
             "nativa, en una ventana de sesenta micras de lado.\n"
             "\n"
             "Se parecen entre sí bastante más de lo que esperábamos. Agrupándolas por color y "
             "textura, ciento cuarenta y cinco de las ciento sesenta y cuatro, el ochenta y ocho por "
             "ciento, caen en un solo grupo: epitelio tumoral denso, recortes oscuros y saturados "
             "casi sin fondo, con un núcleo hipercromático y condensado en el centro. Es la forma "
             "que el detector busca, y la misma que tienen las trece que sí acertó. Se despega un "
             "grupo chico: quince recortes claros, con dos tercios de fondo, que son tejido laxo o "
             "borde de lámina.\n"
             "\n"
             "A la derecha están las trece falladas, centradas, cada una sin ninguna detección "
             "cerca. En varias se ve una figura mitótica de manual, y ése es el bloque que más "
             "informa.\n"
             "\n"
             "Y lo que no se puede decir de esta lámina, que acá importa más que en ninguna otra: "
             "las ciento sesenta y cuatro no son falsos positivos. El patólogo marca solo donde la "
             "evidencia es clara, no pretende marcarlo todo, así que una detección sin marca puede "
             "ser perfectamente una mitosis real sin marcar. Por eso no calculamos precisión y "
             "ningún panel las pinta como error. Y el parecido que describí es de píxeles, no "
             "semántico.\n"
             "\n"
             "Lo que me gustaría discutir es hacia dónde movemos el número del detector, y hay dos "
             "cosas baratas antes de tocar nada. Una es mirar si en esos trece la clase mitosis "
             "quedó segunda por poco, que se responde con lo que ya está guardado y decide entre "
             "recalibrar un umbral o reentrenar la cabeza. La otra: si el problema es que su clase "
             "de mitosis nunca vio mama, existe un conjunto público de mitosis en mama, con "
             "doscientos casos y licencia abierta. Y viene con una advertencia que nos toca: los "
             "detectores de mitosis se caen al cambiar de escáner.")


# ============================================================================
# Orden pedido por Sebastián el 6-ago, y no se re-decide acá: abre el grid de
# expertos y slots, sigue la medición de atención contra las marcas del patólogo,
# y SI-MIL queda al final. Invierte la decisión del 3-ago (lo cerrado antes que lo
# vivo) y jubila el encuadre del 4-ago (el grid como sección de cierre): las dos
# eran de encuadre nuestro, ésta viene del supervisor.
def build():
    prs, keep_ids = base_from_template()
    prs.slide_width = Inches(SW); prs.slide_height = Inches(SH)

    lam_portada(prs)
    lam_objetivos_sprint(prs)

    # ---- el grid abre ----
    lam_grid(prs)

    # ---- la medición de atención ----
    lam_pregunta_medible(prs)
    lam_mapas(prs)
    lam_escalera(prs)
    lam_28_parches(prs)
    lam_mira_responde(prs)

    # ---- SI-MIL, al final ----
    # Ernesto retiró las dos láminas del detalle de las ramas el 22-ago: quedan «qué
    # propone» y los resultados, que es lo que se discute. El detalle arquitectónico que
    # llevaban baja al guion.
    lam_simil_propone(prs)
    lam_simil_resultados(prs)

    # ---- el segundo hilo: la herramienta, y los resultados de cruzarla con la atención ----
    lam_hovernext_paper(prs)
    lam_resultados(prs)
    lam_hovernext_solo(prs)

    # ---- la Fase A, que cerró el 21-ago y el deck del 19 no tenía ----
    lam_escalera_brazos(prs)
    lam_gate_invasivo(prs)
    lam_a0_falla_la_clase(prs)
    lam_galeria_164(prs)

    # ---- cierre: reflow, auditoría, escala al tamaño del template, tipografía ----
    reflow_onco(prs, skip=keep_ids)
    auditar(prs, skip=keep_ids)
    scale_deck_to_1610(prs, skip=keep_ids)
    forzar_barlow(prs)
    os.makedirs(OUT_DIR, exist_ok=True)
    prs.save(DST)
    print("Guardado:", DST, "·", len(prs.slides), "slides · 13.333x7.5")
    # El deck dejó de ser monográfico de SI-MIL, así que el archivo con el nombre viejo
    # quedaría al lado del nuevo, desactualizado y sin dueño. Se retira acá y no a mano,
    # para que regenerar desde cero no lo resucite.
    if os.path.exists(DST_LEGACY):
        os.remove(DST_LEGACY)
        print("Retirado el nombre previo:", os.path.basename(DST_LEGACY))


if __name__ == "__main__":
    build()
