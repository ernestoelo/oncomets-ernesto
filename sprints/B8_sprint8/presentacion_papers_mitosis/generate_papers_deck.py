#!/usr/bin/env python
"""generate_papers_deck.py — deck de los CUATRO papers de la reunión del 12-ago-2026.

Diez láminas, una por paper más el encuadre, el cuadro comparativo, la recomendación y las
preguntas. Contenido verificado que sale de dos documentos y de ninguna otra parte:

  ../tareas_geometricas/hojas_reunion.md      hojas 0, 1 y 3 más el anexo de MIDOG
  ../papers_11_agosto/hojas_papers_nuevos.md  hojas 0, 5 y 6

CERO números nuevos: si una cifra no está en esas hojas, no entra acá.

El cuarteto es SIMÉTRICO, dos papers por tarea, y ese es el encuadre del deck:

  mitosis        PU learning (Zhao, MELBA 2022) · ZoomMIL (Thandiackal, ECCV 2022)
  grado nuclear  NPKC-MIL (Wang y Yuan, iScience 2024) · pleomorfismo (Mercan, npj BC 2022)

CellViT y MS-CLAM NO entran: sus hojas siguen siendo fichas válidas, pero no son el material
de esta reunión.

El molde de las cuatro láminas de paper es el mismo, para que se lean igual: figura original
de los autores a la izquierda con la procedencia al pie, tres paneles cortos a la derecha y
barra de remate. Es la excepción explícita a «todo nativo»: la figura de un paper va como
imagen; el cuadro de la lámina 7 y los diagramas de las láminas 2, 8 y 10 van nativos.

Dos cosas gobiernan la geometría y las dos se descubrieron ANTES de escribir código:

  · El ancho de la columna izquierda NO se puede fijar. Las cuatro figuras van de 1,27 a
    3,09 de relación de aspecto, así que se fija el ALTO de la caja (ALTO_FIG) y el ancho
    sale de alto x ar, topeado por COL_MIN para que la columna de bloques no se angoste.
    El pie va pegado al borde inferior REAL de la imagen dibujada, no al de la caja.
  · `auditar` NO ve dos defectos que este deck puede tener, distintos de los que ya lista
    [[deck-qa-puntos-ciegos-chequeo]]: una tabla nativa cuyo texto no entra (el `row_h` de
    `simple_table` es un mínimo y PowerPoint crece la fila por su cuenta) y una pila de
    `panel(h=None)` que se mete debajo de la `takeaway_bar` (no es «fuera del lienzo», así
    que no dispara aviso). Por eso están `chequeo_pila` y `alto_tabla`, que miden las dos
    cosas con los TTF reales y las imprimen.

Notas del presentador: punteo de tres a cinco renglones de una línea, línea en blanco, y
después el guion HABLADO en prosa corrida ([[notas-presentador-guion-didactico]]). Cargan el
peso explicativo a propósito: son el material con el que se estudian los cuatro papers.

Helpers copiados de `../presentacion_b8/generate_b8_deck.py` (solo los que este deck usa).

Uso:
    PYTHONPATH=/media/administrador/Storage1/sdonoso/clam_testing2/.pylibs \
    /home/sdonoso/miniconda3/envs/clam_latest/bin/python \
      sprints/B8_sprint8/presentacion_papers_mitosis/generate_papers_deck.py
"""
import os

from lxml import etree
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn
from PIL import Image

REPO = "/media/administrador/Storage1/sdonoso/clam_testing2/oncomets-ernesto"
OUT_DIR = os.path.join(REPO, "sprints/B8_sprint8/presentacion_papers_mitosis")
ASSETS = os.path.join(OUT_DIR, "assets")
DST = os.path.join(OUT_DIR, "Papers_Mitosis.pptx")

TEMPLATE = os.path.join(REPO, "sprints/B7_sprint7/Modelo OncoMets Spatial V1 Deep-LLM-V.pptx")
TPL_KEEP = (0, 1)          # portada de marca + lámina de título, nativas a 13.333

TITULO_PORTADA = "OncoMets · Cuatro papers"
FECHA_REUNION = "12 de agosto de 2026"

FIG_PU = os.path.join(ASSETS, "fig_pulearning.png")
FIG_ZOOM = os.path.join(ASSETS, "fig_zoommil.png")
FIG_NPKC = os.path.join(ASSETS, "fig_npkcmil.png")
FIG_PLEO = os.path.join(ASSETS, "fig_pleomorfismo.png")

# ---- paleta Deep-LLM-V (medida sobre el template) ----
ONCO_DARK = RGBColor(0x3E, 0x68, 0x77)    # bloque de proceso
ONCO_CONN = RGBColor(0x38, 0x62, 0x71)    # conector / borde
ONCO_PANEL = RGBColor(0xCD, 0xDF, 0xE1)   # panel contenedor / operador
ONCO_DATA = RGBColor(0xB7, 0xB7, 0xB7)    # bloque de dato
ONCO_INK = RGBColor(0x0E, 0x28, 0x41)     # borde / fondo oscuro

TEAL_TITLE = ONCO_DARK
TEAL_SQ = ONCO_CONN
TEAL_CARD = ONCO_PANEL
TEAL_CARD2 = RGBColor(0xE9, 0xF1, 0xF2)   # tinte del claro, para banding
GRIS_BODY = RGBColor(0x59, 0x59, 0x59)
GRIS_TXT = RGBColor(0x55, 0x55, 0x55)
INK = RGBColor(0x22, 0x22, 0x22)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BLACK = RGBColor(0x00, 0x00, 0x00)

F_TITLE = "Barlow"
F_BODY = "Barlow"
# Barlow no trae griegas, ℝ, ∈, → ni ⊗: esos glifos caen al fallback del sistema. Se
# rasterizó la alternativa (declarar Cambria Math, que el template también embebe) y se ve
# PEOR: sus griegas son serif finas y contrastan con el Barlow que las rodea. Acá el único
# glifo de esa familia es la «×» de las magnificaciones y la «µ» de µm/px.

# --- geometría de trabajo (se escala x1.3333 al final) ---
SW, SH = 10.0, 5.625

# --- cabecera OncoMets (láminas técnicas del template) ---
ONCO_LOGO = os.path.join(REPO, "sprints/B7_sprint7/presentacion_b7/assets/logo_oncomets.png")
ONCO_TITLE = RGBColor(0x3E, 0x68, 0x77)
ONCO_LINE = RGBColor(0x3D, 0x68, 0x76)
ONCO_LOGO_L, ONCO_LOGO_T, ONCO_LOGO_W, ONCO_LOGO_H = 0.5625, 0.3278, 1.1033, 0.6630
ONCO_TIT_L, ONCO_TIT_W = 2.0558, 6.8799
ONCO_TIT_BASE = 0.9908
ONCO_TIT_T, ONCO_TIT_H = 0.2950, ONCO_TIT_BASE - 0.2950
ONCO_TIT_SZ = 18.75
ONCO_LINE_T, ONCO_LINE_H = 1.0658, 0.1050
ONCO_BAND = ONCO_LINE_T + ONCO_LINE_H     # 1.1708

TOP = 1.24

# --- el molde de las cuatro láminas de paper ---
# La REGIÓN izquierda es constante en las cuatro (COL_IZQ) y la columna de bloques también
# (COL_DER): lo que varía es cuánto de la región izquierda ocupa la figura, que depende de
# su relación de aspecto. Así las cuatro láminas tienen la misma retícula, el pie siempre
# tiene el mismo ancho para envolver, y la figura queda tan grande como su forma permite.
LM, RM, GAP = 0.34, 0.34, 0.26
COL_DER = 4.10                             # columna de bloques, constante
COL_IZQ = SW - LM - RM - GAP - COL_DER     # región de la figura y su pie, constante
ALTO_FIG = 3.00                            # alto MÁXIMO de la figura
BAR_PAPER = 4.80                           # barra de remate de las láminas de paper
FIG_T0 = TOP + 0.06

_AVISOS = []


# ============================================================================
# Helpers (portados de ../presentacion_b8/generate_b8_deck.py)
# ============================================================================
def _blank(prs):
    for lay in prs.slide_layouts:
        if (lay.name or "").lower().strip() == "blank":
            return lay
    return prs.slide_layouts[6]


def base_from_template():
    """Abre el template válido y le borra las láminas salvo TPL_KEEP.

    Conserva ppt/fonts/*.fntdata + <p:embeddedFontLst> (Barlow, Cambria Math), el theme y
    el master. Borrar una slide no tiene API en python-pptx: hay que sacarla del sldIdLst
    Y soltar la relación, si no queda huérfana en el paquete."""
    prs = Presentation(TEMPLATE)
    keep_ids = {id(prs.slides[i]._element) for i in TPL_KEEP}
    lst = prs.slides._sldIdLst
    for i, sld in enumerate(list(lst)):
        if i in TPL_KEEP:
            continue
        prs.part.drop_rel(sld.get(
            "{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id"))
        lst.remove(sld)
    return prs, keep_ids


def new_slide(prs):
    """Lámina en blanco SIN los placeholders del layout (el BLANK del template arrastra
    DATE/FOOTER/SLIDE_NUMBER y aparecen como cuadros vacíos)."""
    s = prs.slides.add_slide(_blank(prs))
    for ph in list(s.placeholders):
        ph._element.getparent().remove(ph._element)
    return s


def _add_runs(p, text, size, bold, color, font=F_BODY):
    """Mini-markup de sub/superíndices REALES (baseline OOXML, no caracteres Unicode):
    `_x` / `_(xx)` = subíndice ; `^x` / `^(xx)` = superíndice. Para un `_`/`^` literal,
    escaparlo con backslash en el string fuente."""
    def emit(s, base=None):
        if not s:
            return
        r = p.add_run(); r.text = s
        r.font.name = font; r.font.bold = bold; r.font.color.rgb = color
        r.font.size = Pt(size * 0.74 if base is not None else size)
        if base is not None:
            r._r.get_or_add_rPr().set("baseline", str(base))
    i, n, buf = 0, len(text), ""
    while i < n:
        c = text[i]
        if c == "\\" and i + 1 < n and text[i + 1] in "_^\\":
            buf += text[i + 1]; i += 2; continue
        if c in "_^" and i + 1 < n:
            base = -25000 if c == "_" else 30000
            nxt = text[i + 1]
            if nxt == "(":
                j = text.find(")", i + 2)
                if j != -1:
                    emit(buf); buf = ""
                    emit(text[i + 2:j], base)
                    i = j + 1; continue
                buf += c; i += 1; continue
            emit(buf); buf = ""
            emit(nxt, base)
            i += 2; continue
        buf += c; i += 1
    emit(buf)


def _set_runs(tf, lines, anchor=MSO_ANCHOR.TOP, markup=False):
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    for i, ln in enumerate(lines):
        txt, sz, bold, col = ln[0], ln[1], ln[2], ln[3]
        font = ln[4] if len(ln) > 4 else F_BODY
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = ln[5] if len(ln) > 5 else PP_ALIGN.LEFT
        p.space_after = Pt(4)
        if markup:
            _add_runs(p, txt, sz, bold, col, font)
        else:
            r = p.add_run(); r.text = txt
            r.font.size = Pt(sz); r.font.bold = bold
            r.font.name = font; r.font.color.rgb = col


def add_textbox(slide, l, t, w, h, lines, anchor=MSO_ANCHOR.TOP, markup=False):
    tb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    _set_runs(tb.text_frame, lines, anchor=anchor, markup=markup)
    return tb


def notes(slide, text):
    """Guion HABLADO en prosa (sin etiquetas de fase, sin nº de job ni nombres)."""
    slide.notes_slide.notes_text_frame.text = text


def _rect(slide, l, t, w, h, color, line=None):
    sp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(l), Inches(t), Inches(w), Inches(h))
    sp.fill.solid(); sp.fill.fore_color.rgb = color
    if line is None:
        sp.line.fill.background()
    else:
        sp.line.color.rgb = line; sp.line.width = Pt(1)
    sp.shadow.inherit = False
    return sp


def header_oncomets(slide, title, size=22):
    """Cabecera técnica del template: logo OncoMets + título Barlow bold + línea teal."""
    pic = slide.shapes.add_picture(ONCO_LOGO, Inches(ONCO_LOGO_L), Inches(ONCO_LOGO_T),
                                   Inches(ONCO_LOGO_W), Inches(ONCO_LOGO_H))
    pic.name = "ONCOHDR_logo"
    if title:
        tb = slide.shapes.add_textbox(Inches(ONCO_TIT_L), Inches(ONCO_TIT_T),
                                      Inches(ONCO_TIT_W), Inches(ONCO_TIT_H))
        tb.name = "ONCOHDR_title"
        _set_runs(tb.text_frame, [(title, min(size, ONCO_TIT_SZ), True, ONCO_TITLE, F_BODY)],
                  anchor=MSO_ANCHOR.BOTTOM)
    _rect(slide, 0.0, ONCO_LINE_T, SW, ONCO_LINE_H, ONCO_LINE).name = "ONCOHDR_line"


def content(prs, title, size=26):
    s = new_slide(prs)
    header_oncomets(s, title, size=size)
    return s


def add_image_fit(slide, path, l, t, w, h, align="center"):
    iw, ih = Image.open(path).size
    ar = iw / ih; box_ar = w / h
    if ar > box_ar:
        nw = w; nh = w / ar
    else:
        nh = h; nw = h * ar
    nl = l + (w - nw) / 2
    nt = t + (h - nh) / 2 if align != "top" else t
    slide.shapes.add_picture(path, Inches(nl), Inches(nt), Inches(nw), Inches(nh))


def add_card(slide, l, t, w, h, idx, text, size=14):
    fill = TEAL_CARD if idx % 2 == 0 else TEAL_CARD2
    sp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(l), Inches(t), Inches(w), Inches(h))
    sp.fill.solid(); sp.fill.fore_color.rgb = fill
    sp.line.color.rgb = TEAL_SQ; sp.line.width = Pt(1.25); sp.shadow.inherit = False
    cd = 0.44
    circ = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(l + 0.16), Inches(t + (h - cd) / 2),
                                  Inches(cd), Inches(cd))
    circ.fill.solid(); circ.fill.fore_color.rgb = ONCO_DARK
    circ.line.fill.background(); circ.shadow.inherit = False
    _set_runs(circ.text_frame, [(str(idx), 15, True, WHITE, F_TITLE, PP_ALIGN.CENTER)],
              anchor=MSO_ANCHOR.MIDDLE)
    tb = slide.shapes.add_textbox(Inches(l + 0.74), Inches(t), Inches(w - 0.88), Inches(h))
    _set_runs(tb.text_frame, [(text, size, True, INK, F_BODY)], anchor=MSO_ANCHOR.MIDDLE)


def caption(slide, l, t, w, text, size=11, col=GRIS_TXT, align=PP_ALIGN.CENTER, bold=False):
    add_textbox(slide, l, t, w, 0.4, [(text, size, bold, col, F_BODY, align)])


def takeaway_bar(slide, text, t=4.85, col=TEAL_TITLE, size=14):
    _rect(slide, 0.35, t, SW - 0.7, 0.02, TEAL_SQ)
    add_textbox(slide, 0.35, t + 0.08, SW - 0.7, 0.62,
                [(text, size, True, col, F_BODY, PP_ALIGN.CENTER)], anchor=MSO_ANCHOR.MIDDLE)


# ---- gramática de diagrama de Deep-LLM-V ----
def _proc(slide, l, t, w, h, text, dim=None, size=11, col=None):
    """Bloque de proceso: rounded-rect #3E6877 con Barlow bold BLANCO."""
    sp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(l), Inches(t),
                                Inches(w), Inches(h))
    sp.fill.solid(); sp.fill.fore_color.rgb = col or ONCO_DARK
    sp.line.fill.background(); sp.shadow.inherit = False
    lines = [(text, size, True, WHITE, F_BODY, PP_ALIGN.CENTER)]
    if dim:
        lines.append((dim, size - 2, False, WHITE, F_BODY, PP_ALIGN.CENTER))
    _set_runs(sp.text_frame, lines, anchor=MSO_ANCHOR.MIDDLE, markup=True)
    for p in sp.text_frame.paragraphs:
        p.space_after = Pt(0)
    return sp


def _proc_claro(slide, l, t, w, h, text, size=11, dim=None):
    """Bloque de proceso en el tono CLARO (#CDDFE1 con texto teal): el detalle interno."""
    sp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(l), Inches(t),
                                Inches(w), Inches(h))
    sp.fill.solid(); sp.fill.fore_color.rgb = ONCO_PANEL
    sp.line.fill.background(); sp.shadow.inherit = False
    lines = [(text, size, True, ONCO_DARK, F_BODY, PP_ALIGN.CENTER)]
    if dim:
        lines.append((dim, size - 2, False, ONCO_DARK, F_BODY, PP_ALIGN.CENTER))
    _set_runs(sp.text_frame, lines, anchor=MSO_ANCHOR.MIDDLE, markup=True)
    for p in sp.text_frame.paragraphs:
        p.space_after = Pt(0)
    return sp


def _dato(slide, l, t, w, h, text, size=9.5):
    """Bloque de dato: rect #B7B7B7 con Barlow NEGRO."""
    sp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(l), Inches(t), Inches(w), Inches(h))
    sp.fill.solid(); sp.fill.fore_color.rgb = ONCO_DATA
    sp.line.fill.background(); sp.shadow.inherit = False
    _set_runs(sp.text_frame, [(text, size, False, BLACK, F_BODY, PP_ALIGN.CENTER)],
              anchor=MSO_ANCHOR.MIDDLE, markup=True)
    return sp


def _grupo(slide, l, t, w, h, fill=None):
    """Panel contenedor que agrupa una región."""
    sp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(l), Inches(t),
                                Inches(w), Inches(h))
    sp.fill.solid(); sp.fill.fore_color.rgb = fill or ONCO_PANEL
    sp.line.color.rgb = fill or ONCO_PANEL; sp.line.width = Pt(0.6)
    sp.shadow.inherit = False
    return sp


def _conn(slide, x0, y0, x1, y1, arrow=True):
    """Conector recto #386271 de 2.37 pt, con punta opcional."""
    ln = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(x0), Inches(y0),
                                    Inches(x1), Inches(y1))
    ln.line.color.rgb = ONCO_CONN; ln.line.width = Pt(2.37)
    ln.shadow.inherit = False
    if arrow:
        lnpr = ln.line._get_or_add_ln()
        lnpr.append(lnpr.makeelement(qn('a:tailEnd'),
                                     {'type': 'triangle', 'w': 'med', 'len': 'med'}))
    return ln


def simple_table(slide, l, t, w, headers, rows, col_fracs, row_h=0.32, fs=9.5,
                 destacar=None, markup=False):
    """Tabla nativa: header teal + banding claro. `destacar` = índice de fila (0-based
    sobre `rows`) que se pinta con el celeste sólido, para la fila que importa."""
    ncol = len(headers); nrow = len(rows) + 1
    tbl = slide.shapes.add_table(nrow, ncol, Inches(l), Inches(t), Inches(w),
                                 Inches(row_h * nrow)).table
    for ci, fr in enumerate(col_fracs):
        tbl.columns[ci].width = Inches(w * fr)
    tbl.first_row = False; tbl.horz_banding = False
    for ri, row in enumerate([headers] + list(rows)):
        for ci, txt in enumerate(row):
            cell = tbl.cell(ri, ci)
            cell.margin_left = Inches(0.06); cell.margin_right = Inches(0.04)
            cell.margin_top = Inches(0.015); cell.margin_bottom = Inches(0.015)
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            if ri == 0:
                cell.fill.solid(); cell.fill.fore_color.rgb = TEAL_SQ
                col, bold = WHITE, True
            elif destacar is not None and ri - 1 == destacar:
                cell.fill.solid(); cell.fill.fore_color.rgb = TEAL_CARD
                col, bold = ONCO_DARK, True
            else:
                cell.fill.solid(); cell.fill.fore_color.rgb = TEAL_CARD2 if ri % 2 else TEAL_CARD
                col, bold = INK, (ci == 0)
            tf = cell.text_frame; tf.word_wrap = True
            p = tf.paragraphs[0]; p.alignment = PP_ALIGN.LEFT
            if markup:
                _add_runs(p, txt, fs, bold, col, F_BODY)
            else:
                r = p.add_run(); r.text = txt
                r.font.size = Pt(fs); r.font.bold = bold
                r.font.name = F_BODY; r.font.color.rgb = col
    return tbl


# ---------------------------------------------------------------------------
# Medición de texto con los TTF reales de Barlow, instalados bajo containment. El chequeo
# programático de conformidad da «todo limpio» con texto que desborda su caja
# ([[deck-qa-puntos-ciegos-chequeo]]) porque nadie mide el texto. Acá se mide de verdad y
# los paneles se dimensionan solos.
BARLOW_DIR = "/media/administrador/Storage1/sdonoso/clam_testing2/fonts/barlow"
_FCACHE = {}
_MEDIDA = 40          # se mide a 40 px y se divide: da precisión sin depender del hinting


def _face(bold):
    from PIL import ImageFont
    key = bool(bold)
    if key not in _FCACHE:
        nombre = "Barlow-Bold.ttf" if bold else "Barlow-Regular.ttf"
        _FCACHE[key] = ImageFont.truetype(os.path.join(BARLOW_DIR, nombre), _MEDIDA)
    return _FCACHE[key]


def text_w(txt, size, bold=False):
    """Ancho del texto en PULGADAS del espacio de trabajo (lámina de 10 in, 72 pt/in).

    Los símbolos que Barlow no tiene caen al fallback y miden distinto; se aproximan con el
    ancho de una «n», que es del orden correcto y sobreestima poco."""
    f = _face(bold)
    limpio = "".join(c if f.getbbox(c) else "n" for c in txt)
    return f.getlength(limpio) / _MEDIDA * size / 72.0


def wrap_lines(txt, ancho, size, bold=False):
    """Cuántas líneas ocupa `txt` dentro de `ancho` pulgadas (wrap por palabra)."""
    if not txt.strip():
        return 1
    n, actual = 1, ""
    for palabra in txt.split(" "):
        cand = (actual + " " + palabra).strip()
        if actual and text_w(cand, size, bold) > ancho:
            n += 1; actual = palabra
        else:
            actual = cand
    return n


def _alto_bloque(lineas, ancho, size, bold=False, space_after=3):
    """Alto en pulgadas de un bloque de párrafos ya envueltos."""
    total = 0.0
    for ln in lineas:
        total += wrap_lines(ln, ancho, size, bold) * size * 1.22 / 72.0 + space_after / 72.0
    return total


def panel(slide, l, t, w, h, title, tcol, lines, border, fill=TEAL_CARD2, tsize=14.5,
          bsize=12, markup=False):
    """Panel con título + líneas de cuerpo.

    `h=None` calcula el alto midiendo el texto, que es lo que evita que la última línea
    quede fuera de la caja. Con `h` explícito, el valor se respeta salvo que el texto no
    entre: en ese caso gana el texto.

    El ancho con el que se MIDE no es el del textbox: PowerPoint le mete 0,1" de margen a
    cada lado, así que el texto envuelve 0,2" antes de lo que dice la caja. La versión del
    B8 medía sobre el ancho de la caja y por eso subestimaba el wrap en una línea."""
    ancho_txt = w - 0.36
    ancho_medida = ancho_txt - 0.20
    necesario = (0.12 + _alto_bloque([title], ancho_medida, tsize, True)
                 + _alto_bloque(lines, ancho_medida, bsize) + 0.16)
    h = necesario if h is None else max(h, necesario)
    sp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(l), Inches(t), Inches(w), Inches(h))
    sp.fill.solid(); sp.fill.fore_color.rgb = fill
    sp.line.color.rgb = border; sp.line.width = Pt(1.5); sp.shadow.inherit = False
    tb = slide.shapes.add_textbox(Inches(l + 0.18), Inches(t + 0.08), Inches(ancho_txt), Inches(h - 0.12))
    runs = [(title, tsize, True, tcol, F_TITLE)] + [(ln, bsize, False, INK, F_BODY) for ln in lines]
    _set_runs(tb.text_frame, runs, markup=markup)
    for p in tb.text_frame.paragraphs:
        p.space_after = Pt(3)
    return sp


def pie_lineas(slide, l, t, w, lineas, size=8, col=GRIS_BODY, space_after=1.5):
    """Pie de varias líneas en UN solo textbox, con el interlineado al mínimo.

    `caption` gasta 0,4" por renglón y con tres renglones se come el alto que necesita la
    figura de arriba. Acá el bloque mide exactamente lo que mide su texto, y devuelve ese
    alto para que quien lo llama pueda cerrar la lámina debajo."""
    h = _alto_bloque(lineas, w - 0.22, size, space_after=space_after)
    tb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    _set_runs(tb.text_frame, [(ln, size, False, col, F_BODY) for ln in lineas])
    for p in tb.text_frame.paragraphs:
        p.space_after = Pt(space_after)
    return h


CONTENT_TOP_NEW = ONCO_BAND + 0.05        # 1.221
SAFE_BOTTOM = SH - 0.14


def _scale_block(el, f):
    """Escala un shape completo por f: geometría, cuerpo tipográfico y métrica de tabla."""
    for off in el.iter(qn("a:off")):
        off.set("x", str(int(round(int(off.get("x")) * f))))
        off.set("y", str(int(round(int(off.get("y")) * f))))
    for ext in el.iter(qn("a:ext")):
        ext.set("cx", str(int(round(int(ext.get("cx")) * f))))
        ext.set("cy", str(int(round(int(ext.get("cy")) * f))))
    for tag in ("a:rPr", "a:defRPr", "a:endParaRPr"):
        for rpr in el.iter(qn(tag)):
            if rpr.get("sz"):
                rpr.set("sz", str(max(100, int(round(int(rpr.get("sz")) * f)))))
    for gc in el.iter(qn("a:gridCol")):
        if gc.get("w"):
            gc.set("w", str(int(round(int(gc.get("w")) * f))))
    for tr in el.iter(qn("a:tr")):
        if tr.get("h"):
            tr.set("h", str(int(round(int(tr.get("h")) * f))))


def reflow_onco(prs, skip=()):
    """Reancla el contenido de las láminas con cabecera OncoMets bajo la banda teal, y lo
    comprime si no entra. Como acá se maqueta ya desde TOP, el ajuste típico es nulo: la
    función queda como red de seguridad."""
    for slide in prs.slides:
        if id(slide._element) in skip:
            continue
        cuerpo = [sh for sh in slide.shapes if not (sh.name or "").startswith("ONCOHDR_")]
        if len(cuerpo) == len(list(slide.shapes)):
            continue
        tops = [Emu(sh.top).inches for sh in cuerpo]
        bots = [Emu(sh.top).inches + Emu(sh.height).inches for sh in cuerpo]
        if not tops:
            continue
        top0, bot0 = min(tops), max(bots)
        f = 1.0
        if bot0 + (CONTENT_TOP_NEW - top0) > SAFE_BOTTOM and bot0 > top0:
            f = max(0.80, (SAFE_BOTTOM - CONTENT_TOP_NEW) / (bot0 - top0))
        if f < 1.0:
            for sh in cuerpo:
                _scale_block(sh._element, f)
        desplaz = CONTENT_TOP_NEW - top0 * f
        for sh in cuerpo:
            sh.top = Inches(Emu(sh.top).inches + desplaz)


def forzar_barlow(prs, fuente=F_BODY):
    """Deja Barlow como única tipografía del archivo (pedido de Sebastián).

    Los runs propios ya salen en Barlow, pero quedan tres focos fuera de alcance directo:
    el `endParaRPr`/`buFont` de las láminas heredadas, el fontScheme del theme (que en este
    template es el de Office, o sea Arial) y los `defRPr` del master y los layouts."""
    A = "{http://schemas.openxmlformats.org/drawingml/2006/main}"
    tags = tuple(A + t for t in ("latin", "ea", "cs", "sym", "buFont"))

    def normaliza(root):
        tocados = 0
        for el in root.iter():
            if el.tag in tags and el.get("typeface") != fuente:
                el.set("typeface", fuente)
                tocados += 1
        return tocados

    cola = [prs.part] + [s.part for s in prs.slides]
    cola += [s.notes_slide.part for s in prs.slides if s.has_notes_slide]
    vistos, partes = set(), []
    while cola:
        p = cola.pop()
        if id(p) in vistos:
            continue
        vistos.add(id(p))
        partes.append(p)
        for rel in p.rels.values():
            if not rel.is_external and rel.reltype.endswith(
                    ("slideLayout", "slideMaster", "notesMaster", "theme")):
                cola.append(rel.target_part)
    n = 0
    for p in partes:
        el = getattr(p, "_element", None)
        if el is not None:
            n += normaliza(el)
            continue
        if str(p.partname).startswith("/ppt/theme/"):
            raiz = etree.fromstring(p.blob)
            tocados = 0
            for cual in ("majorFont", "minorFont"):
                for fs in raiz.iter(A + cual):
                    for lat in fs.findall(A + "latin"):
                        if lat.get("typeface") != fuente:
                            lat.set("typeface", fuente)
                            tocados += 1
            if tocados:
                p._blob = etree.tostring(raiz, xml_declaration=True,
                                         encoding="UTF-8", standalone=True)
                n += tocados
    print("  tipografía: %d referencias forzadas a %s" % (n, fuente))


def scale_deck_to_1610(prs, k=13.333 / 10.0, skip=()):
    """Escala el deck entero por k y fija el tamaño de lámina en 13.333 x 7.5. Misma
    relación de aspecto que 10 x 5.625, así que la escala es uniforme y no deforma."""
    for slide in prs.slides:
        if id(slide._element) in skip:
            continue
        tree = slide.shapes._spTree
        for off in tree.iter(qn("a:off")):
            off.set("x", str(int(round(int(off.get("x")) * k))))
            off.set("y", str(int(round(int(off.get("y")) * k))))
        for ext in tree.iter(qn("a:ext")):
            ext.set("cx", str(int(round(int(ext.get("cx")) * k))))
            ext.set("cy", str(int(round(int(ext.get("cy")) * k))))
        for tag in ("a:rPr", "a:defRPr", "a:endParaRPr"):
            for rpr in tree.iter(qn(tag)):
                if rpr.get("sz"):
                    rpr.set("sz", str(max(100, int(round(int(rpr.get("sz")) * k)))))
        for gc in tree.iter(qn("a:gridCol")):
            if gc.get("w"):
                gc.set("w", str(int(round(int(gc.get("w")) * k))))
        for tr in tree.iter(qn("a:tr")):
            if tr.get("h"):
                tr.set("h", str(int(round(int(tr.get("h")) * k))))
        for ln in tree.iter(qn("a:ln")):
            if ln.get("w"):
                ln.set("w", str(int(round(int(ln.get("w")) * k))))
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)


def auditar(prs, skip=()):
    """Chequeo de defectos que el ojo ve y un chequeo ingenuo no: texto que no entra en su
    caja, shapes fuera del lienzo y cuerpos por debajo del mínimo del template (7 pt).

    Se corre ANTES de escalar, o sea en el espacio de 10 x 5.625. No reemplaza mirar las
    láminas ([[deck-qa-puntos-ciegos-chequeo]]) y NO ve los dos defectos que este deck
    puede tener: para eso están `chequeo_pila` y `alto_tabla`."""
    problemas = []
    for idx, slide in enumerate(prs.slides, start=1):
        if id(slide._element) in skip:
            continue
        for sh in slide.shapes:
            try:
                l, t = Emu(sh.left).inches, Emu(sh.top).inches
                w, h = Emu(sh.width).inches, Emu(sh.height).inches
            except TypeError:
                continue
            # Un shape rotado 270° reporta su bbox SIN rotar: da falso positivo de límites.
            if not getattr(sh, "rotation", 0):
                if l < -0.02 or t < -0.02 or l + w > SW + 0.02 or t + h > SH + 0.02:
                    problemas.append("s%02d  fuera del lienzo: %s (%.2f, %.2f, %.2f x %.2f)"
                                     % (idx, sh.shape_type, l, t, w, h))
            if not sh.has_text_frame:
                continue
            alto, chico = 0.0, None
            for p in sh.text_frame.paragraphs:
                if not p.runs:
                    continue
                sz = max((r.font.size.pt for r in p.runs if r.font.size), default=12)
                bold = any(r.font.bold for r in p.runs)
                txt = "".join(r.text for r in p.runs)
                # los runs de sub/superíndice vienen a 0.74x: no son un cuerpo chico real
                if sz >= 6 and (chico is None or sz < chico):
                    if not any(r._r.get_or_add_rPr().get("baseline") for r in p.runs):
                        chico = sz
                alto += wrap_lines(txt, max(w - 0.14, 0.2), sz, bold) * sz * 1.22 / 72.0
            if alto > h + 0.06:
                problemas.append("s%02d  texto que no entra: sobra %.2f\" en «%s…»"
                                 % (idx, alto - h, sh.text_frame.text[:44].replace("\n", " ")))
            if chico is not None and chico < 7.0:
                problemas.append("s%02d  cuerpo por debajo del mínimo: %.1f pt" % (idx, chico))
    if problemas:
        print("  AUDITORÍA: %d avisos" % len(problemas))
        for p in problemas:
            print("   ·", p)
    else:
        print("  AUDITORÍA: sin avisos")
    return problemas


def _set_solo_run(par, texto):
    """Deja el párrafo con un solo run, conservando el formato del primero. El texto del
    template viene partido en varios runs, así que escribir solo runs[0] dejaría la cola
    del original pegada detrás."""
    runs = par.runs
    runs[0].text = texto
    for r in runs[1:]:
        r._r.getparent().remove(r._r)


def retitular_portada(prs, titulo, fecha):
    """Ajusta las dos láminas de apertura heredadas del template sin redibujarlas.

    El título y la fecha van por PARÁMETRO: la versión del B8 los traía hardcodeados y este
    deck no es un deck de sprint."""
    portada = prs.slides[0]
    for sh in portada.shapes:
        if not sh.has_text_frame:
            continue
        txt = sh.text_frame.text.strip()
        if txt.startswith("Care in people"):
            # Marcador sin reemplazar que viene del propio template.
            for par in list(sh.text_frame.paragraphs)[1:]:
                par._p.getparent().remove(par._p)
        elif txt.startswith("OncoMETS is an AI-powered platform"):
            # El párrafo arranca tan abajo que se sale por el borde inferior.
            sh.top = Inches(5.62)
            sh.left = Inches(0.28)
            sh.width = Inches(6.30)
    titulo_sl = prs.slides[1]
    for sh in titulo_sl.shapes:
        if not sh.has_text_frame:
            continue
        txt = sh.text_frame.text.strip()
        if txt == "OncoMets - Spatial":
            _set_solo_run(sh.text_frame.paragraphs[0], titulo)
        elif txt == "14/11/2025":
            _set_solo_run(sh.text_frame.paragraphs[0], fecha)
    return titulo_sl


# ============================================================================
# Los dos chequeos que `auditar` no hace
# ============================================================================
def chequeo_pila(nombre, borde_inferior, barra):
    """La pila de `panel(h=None)` puede meterse debajo de la barra de remate sin salirse
    del lienzo, que es justo lo que `auditar` no mira. Se imprime SIEMPRE el borde inferior
    del último panel para poder leerlo, y se avisa si toca la barra."""
    holgura = barra - borde_inferior
    marca = "OK " if holgura >= 0.04 else "PISA"
    print("  %-22s último panel hasta %.2f\"  ·  barra en %.2f\"  ·  holgura %+.2f\"  %s"
          % (nombre, borde_inferior, barra, holgura, marca))
    if holgura < 0.04:
        _AVISOS.append("%s: la pila de paneles llega a %.2f\" y la barra está en %.2f\""
                       % (nombre, borde_inferior, barra))


def alto_tabla(headers, rows, w, col_fracs, fs, row_h):
    """`row_h` de `simple_table` es un MÍNIMO: PowerPoint crece la fila si el texto no
    entra, así que el alto declarado del shape miente. Acá se estima el alto real midiendo
    cada celda con los TTF de Barlow."""
    total = 0.0
    for ri, fila in enumerate([headers] + list(rows)):
        alto = 0.0
        for ci, txt in enumerate(fila):
            ancho = w * col_fracs[ci] - 0.10
            n = wrap_lines(txt, ancho, fs, ri == 0 or ci == 0)
            alto = max(alto, n * fs * 1.22 / 72.0 + 0.03)
        total += max(row_h, alto)
    return total


# ============================================================================
# El molde de las cuatro láminas de paper
# ============================================================================
def lamina_paper(prs, titulo, png, pie, bloques, remate, nombre):
    """Figura del paper a la izquierda con la procedencia al pie, tres paneles cortos a la
    derecha y barra de remate. Igual en las cuatro para que se lean igual.

    El ancho de la figura NO se puede fijar: las cuatro van de 1,27 a 3,09 de relación de
    aspecto. Se fija el alto (ALTO_FIG) y el ancho sale de alto x ar, topeado por COL_IZQ.

    El pie ocupa siempre COL_IZQ entera, así que su alto NO depende de la figura y se puede
    medir primero: con eso el alto de la figura se recorta para que la pila figura + pie
    entre en la zona. Sin ese recorte el pie se mete DEBAJO de la barra de remate, que es un
    defecto que `auditar` no ve (el textbox no se sale del lienzo y su texto entra en su
    propia caja) y que apareció de verdad en la primera corrida."""
    s = content(prs, titulo)

    iw, ih = Image.open(png).size
    ar = iw / ih
    zona_h = (BAR_PAPER - 0.08) - FIG_T0
    pie_h = _alto_bloque(pie, COL_IZQ - 0.22, 8, space_after=1.5)

    fh = min(ALTO_FIG, zona_h - 0.10 - pie_h)
    fw = min(fh * ar, COL_IZQ)
    fh = fw / ar

    grupo_h = fh + 0.10 + pie_h
    fig_t = FIG_T0 + max(0.0, (zona_h - grupo_h) / 2.0)
    fig_l = LM + (COL_IZQ - fw) / 2.0        # la figura, centrada en su región

    add_image_fit(s, png, fig_l, fig_t, fw, fh, align="top")
    pie_t = fig_t + fh + 0.10
    pie_lineas(s, LM, pie_t, COL_IZQ, pie, size=8)
    chequeo_pila(nombre + " · pie", pie_t + pie_h, BAR_PAPER)

    rl = LM + COL_IZQ + GAP
    t = FIG_T0
    for tit, lineas in bloques:
        sp = panel(s, rl, t, COL_DER, None, tit, TEAL_TITLE, lineas, ONCO_CONN,
                   tsize=12.5, bsize=10.5)
        t += Emu(sp.height).inches + 0.12
    chequeo_pila(nombre + " · paneles", t - 0.12, BAR_PAPER)

    takeaway_bar(s, remate, t=BAR_PAPER, size=13)
    print("       figura %.2f x %.2f\" (ar %.2f) · pie %.2f\"" % (fw, fh, ar, pie_h))
    return s


# ============================================================================
# Las láminas: una función por lámina
# ============================================================================
def lam_portada(prs):
    """Portada de marca, heredada del template."""
    s = retitular_portada(prs, TITULO_PORTADA, FECHA_REUNION)
    notes(s, "Cuatro papers, dos por cada tarea del encargo.\n"
             "Dos salieron de nuestra búsqueda y dos llegaron de la reunión pasada.\n"
             "El orden lo decide una sola pregunta: qué anotación exige cada uno.\n"
             "\n"
             "Traigo los cuatro papers que quedaron para hoy. Dos son los que fiché por mi "
             "cuenta hace un par de semanas y dos son los que llegaron después, así que el "
             "material está mezclado a propósito y no separado por origen.\n"
             "\n"
             "Voy a ir uno por uno, con la figura de los autores a la vista, y después los "
             "pongo a los cuatro en un mismo cuadro para poder compararlos. Cierro con lo que "
             "propongo hacer primero y con las preguntas que necesito responder para poder "
             "empezar. Son preguntas que se contestan hoy, no cosas que haya que ir a medir.")
    return s


def lam_por_que(prs):
    """Por qué estos cuatro: dos tareas, dos papers cada una, y el criterio que los ordena."""
    s = content(prs, "Por qué estos cuatro")
    add_textbox(s, 0.36, 1.26, 9.28, 0.40, [
        ("Dos tareas, dos papers cada una. Los ordena una sola pregunta: qué anotación exige "
         "cada uno.", 13, True, TEAL_TITLE, F_BODY, PP_ALIGN.CENTER)],
        anchor=MSO_ANCHOR.MIDDLE)

    gt, gh, gw = 1.74, 2.72, 4.42
    for gl, tarea, papers, escala in (
        (0.42, "Mitosis",
         [("PU learning", "Zhao et al., MELBA 2022  ·  positivos parciales, los nuestros"),
          ("ZoomMIL", "Thandiackal et al., ECCV 2022  ·  solo la etiqueta de lámina")],
         "Se cuenta a 40× y nuestro privado está escaneado a 20×"),
        (5.16, "Grado nuclear",
         [("NPKC-MIL", "Wang y Yuan, iScience 2024  ·  solo la etiqueta de lámina"),
          ("Pleomorfismo nuclear", "Mercan et al., npj Breast Cancer 2022  ·  puntaje por región")],
         "Trabaja a 0,5 µm/px y nuestro privado está en 0,465"),
    ):
        _grupo(s, gl, gt, gw, gh)
        _proc(s, gl + (gw - 2.60) / 2, 1.86, 2.60, 0.40, tarea, size=13)
        for i, (nombre, detalle) in enumerate(papers):
            _proc_claro(s, gl + 0.22, 2.40 + i * 0.86, gw - 0.44, 0.74,
                        nombre, size=12, dim=detalle)
        _dato(s, gl + 0.22, 4.08, gw - 0.44, 0.30, escala, size=9.5)

    takeaway_bar(s, "La escala física juega en dirección contraria en cada rama, y por eso "
                    "conviene decidirlas por separado.", t=4.62, size=13)
    notes(s, "El encargo tenía dos tareas y mi búsqueda quedó volcada a una sola.\n"
             "Los dos papers que llegaron después llenan justo ese hueco.\n"
             "El cuarteto queda simétrico: dos por tarea.\n"
             "Los ordeno por la anotación que exigen, que es lo que descarta rápido.\n"
             "Y la escala física juega en dirección contraria en cada rama.\n"
             "\n"
             "La lista cambió en las últimas semanas, y el cambio ordena bastante la discusión.\n"
             "\n"
             "El encargo original tenía dos tareas, mitosis y grado nuclear. Cuando salí a "
             "buscar papers, la búsqueda se me fue casi entera a mitosis: de los cuatro que "
             "fiché, tres eran de esa rama, y grado nuclear se apoyaba en un paper de "
             "segmentación de células que ni siquiera tiene clase mitótica ni puntaje de "
             "pleomorfismo. Entraba de refilón y se notaba.\n"
             "\n"
             "Los dos papers que llegaron después llenan exactamente ese hueco. Los dos son de "
             "grado nuclear, así que la lista queda pareja, dos por tarea, y esa simetría es la "
             "forma que le doy hoy a la reunión.\n"
             "\n"
             "Para ordenarlos uso un solo criterio, que es el que descarta candidatos rápido: "
             "qué anotación exige cada método, comparada con la que efectivamente tenemos. Lo "
             "que tenemos hoy son marcas parciales de una sola lámina, y eso deja fuera de "
             "juego a métodos que sobre el papel se ven muy bien. Debajo de cada paper puse "
             "qué pide: uno pide marcas parciales, que es literalmente lo que tenemos, dos "
             "piden solamente la etiqueta de la lámina, que también tenemos, y el último pide "
             "un puntaje por región, que hoy no tenemos pero es más barato de conseguir que "
             "marcar objeto por objeto.\n"
             "\n"
             "Y hay una consecuencia que conviene decir temprano porque cambia el orden en que "
             "gastaría el esfuerzo. La escala física juega en dirección contraria en cada "
             "rama. La mitosis se cuenta a cuarenta aumentos y nuestro privado está escaneado "
             "a veinte, así que toda esa rama arrastra un reescalado con riesgo. El paper de "
             "pleomorfismo, en cambio, entrena y evalúa a medio micrón por píxel, que es "
             "prácticamente nuestro privado. La rama de grado nuclear no tiene ese problema, "
             "y eso pesa más que cualquier diferencia de método.")
    return s


def lam_pu_learning(prs):
    """PU learning, Zhao et al. Figura recompuesta a tres columnas por dos filas."""
    s = lamina_paper(
        prs, "PU learning para detección celular", FIG_PU,
        ["Zhao et al., Positive-unlabeled learning for cell detection in histopathology images "
         "with incomplete annotations · MELBA 2022, Fig. 1, recompuesta para proyectarla.",
         "Tres columnas (anotación completa, baseline y método propuesto) por dos filas. "
         "Quedan afuera la columna del competidor especializado y la fila que lo compara."],
        [("Qué propone", [
            "Lo no marcado deja de contarse como negativo.",
            "Reescribe medio término de la pérdida, no el modelo."]),
         ("Los números", [
             "F1 de detección de 0,470 a 0,507 sobre el baseline.",
             "Recupera el 70 % del hueco de la anotación incompleta."]),
         ("Lo que lo frena", [
             "Retienen el 73 % de las marcas; nosotros, 26 en 4799.",
             "El código es de 2018 y no corre en nuestra GPU."])],
        "Es el único de los cuatro cuyo régimen de anotación es el nuestro: no lo esquiva, "
        "lo usa.", "s04 PU learning")
    notes(s, "Es el único que usa marcas parciales en vez de esquivarlas.\n"
             "El problema está en media pérdida: castigar lo no marcado como si fuera negativo.\n"
             "Recupera el setenta por ciento de lo que se pierde por anotar de forma incompleta.\n"
             "La salvedad grande es que su régimen de prueba no es el nuestro.\n"
             "\n"
             "Arregla un error que cometeríamos el primer día, si nos pusiéramos a entrenar un "
             "detector con las marcas que ya tenemos.\n"
             "\n"
             "Un detector de células se entrena con una pérdida que tiene "
             "dos mitades. Una premia encontrar lo que está marcado, y la otra castiga "
             "encontrar cosas donde no hay nada marcado. Esa segunda mitad supone algo que casi "
             "nunca es cierto en patología: que todo lo que el patólogo no marcó es "
             "efectivamente negativo. Cuando el patólogo anota algunas figuras mitóticas y no "
             "todas, esa suposición es falsa, y el detector aprende activamente a no encontrar "
             "las mitosis que quedaron sin marcar. O sea que lo estás entrenando en contra de "
             "lo que querés.\n"
             "\n"
             "Lo que proponen es reescribir esa segunda mitad usando solamente dos cosas: los "
             "positivos marcados y el resto, que llaman no etiquetado en vez de negativo. Es un "
             "cambio de una de las dos mitades de la pérdida y nada más: no cambian la "
             "arquitectura, y los propios autores dicen que sirve para cualquier detector. El "
             "nombre del método viene de ahí, de aprender con positivos y no etiquetados.\n"
             "\n"
             "Los números que ven son de un conjunto público de mitosis en mama, con la "
             "anotación incompleta simulada a propósito para poder medir. El titular del paper "
             "es una mejora chica, pero el titular engaña, y esta es la forma correcta de leer "
             "la tabla. Esa mejora chica es contra otro método especializado. Contra el "
             "baseline, que es lo que haríamos nosotros entrenando de la manera normal, la "
             "diferencia es casi cuatro puntos de F1. Y lo que ordena todo es la fila del "
             "techo, que es entrenar con la anotación completa: entre entrenar mal con marcas "
             "parciales y entrenar con todo anotado hay un hueco, y el método recupera el "
             "setenta por ciento de ese hueco. En recall, que es encontrar lo que hay, recupera casi "
             "todo.\n"
             "\n"
             "Y hay un detalle del paper que me parece adoptable aunque no implementemos nunca "
             "el detector. Cuando la anotación es parcial, ellos eligen los hiperparámetros "
             "mirando el recall y no la precisión, con el argumento de que un supuesto falso "
             "positivo puede ser perfectamente una célula real que quedó sin marcar. Es el "
             "mismo razonamiento que habíamos escrito por nuestra cuenta, y acá viene con forma "
             "de procedimiento.\n"
             "\n"
             "Queda la parte incómoda, y prefiero decirla yo antes de que la pregunten. "
             "El régimen que ellos evalúan no es el nuestro. Ellos borran marcas hasta dejar "
             "una por región y terminan reteniendo tres cuartas partes de las anotaciones. "
             "Nosotros tenemos veintiséis marcas en cuatro mil setecientos noventa y nueve "
             "parches. Es otro orden de magnitud, el paper no evalúa ese caso y no hay forma "
             "honesta de extrapolar su curva hasta ahí. Y el código, aunque está publicado, "
             "está escrito contra versiones de dos mil dieciocho que no tienen soporte para "
             "nuestra placa: sirve como referencia de la pérdida, no como software que se clone "
             "y se corra.\n"
             "\n"
             "Sobre la figura: la columna de la izquierda es la anotación completa, o sea la "
             "verdad; la del medio es el detector entrenado de la forma normal, y la derecha es "
             "el método propuesto. La recompuse para que se lea proyectada, y eso quiere decir "
             "que saqué la columna del competidor especializado y la fila que lo compara contra "
             "él. Están en el paper, no las escondí, simplemente no entran en una lámina.")
    return s


def lam_zoommil(prs):
    """ZoomMIL, Thandiackal et al. Figura 1 entera."""
    s = lamina_paper(
        prs, "ZoomMIL: aprender dónde ampliar", FIG_ZOOM,
        ["Thandiackal et al., Differentiable Zooming for Multiple Instance Learning on "
         "Whole-Slide Images · ECCV 2022, Fig. 1.",
         "(a) el patólogo, (b) MIL de una sola escala, (c) MIL multiescala, (d) ZoomMIL. "
         "Cada esquema cruza los planos de 2,5× y de 10×; en (d) el plano de 10× solo se "
         "calcula donde el modelo decidió ampliar."],
        [("Qué propone", [
            "Mira todo en aumento bajo y aprende qué acercar.",
            "Se entrena con la etiqueta de lámina."]),
         ("Los números", [
             "En mama, 5,2 puntos de F1 sobre CLAM.",
             "Y con 12,8 veces menos cómputo."]),
         ("Lo que lo frena", [
             "En el privado no hay 40× al que ampliar.",
             "Con objetos chicos el zoom se degrada."])],
        "Para mitosis queda tercero. Donde sí sirve es en armonizar las escalas de las "
        "cohortes.", "s05 ZoomMIL")
    notes(s, "Imita lo que hace el patólogo: barrer en aumento bajo y acercarse donde importa.\n"
             "Lo entrena de punta a punta con la etiqueta de lámina, sin anotación nueva.\n"
             "En un conjunto de mama le saca cinco puntos a CLAM gastando doce veces menos.\n"
             "Para mitosis lo frena la cohorte, no el método: en el privado no hay 40×.\n"
             "\n"
             "Nadie mira una lámina entera a máximo aumento. Se barre en aumento bajo, se eligen "
             "unas pocas zonas y a esas se les hace zoom. De ahí sale este método, que le "
             "enseña al modelo a trabajar igual. El nuestro hoy no hace nada de eso: procesa "
             "todos los parches a la misma escala y después les pone atención.\n"
             "\n"
             "Lo que proponen es hacer exactamente lo del patólogo, pero de forma que se pueda "
             "entrenar. El modelo mira todo en aumento bajo, decide qué parches merecen mirarse "
             "de cerca, y vuelve a mirar solamente esos en el aumento siguiente. La palabra "
             "clave del título es derivable: consiguen que esa decisión de dónde ampliar tenga "
             "gradiente, así que se aprende sola durante el entrenamiento y no hace falta "
             "ninguna anotación nueva. Se entrena con la etiqueta de la lámina, que es la que "
             "ya tenemos. El bloque de atención que usan por dentro es de la misma familia que "
             "el nuestro, así que no es un mundo ajeno.\n"
             "\n"
             "Los números tienen una corrección que vale la pena, y es la razón por la que "
             "valía bajar el PDF en vez de quedarse con el resumen. Lo que circula es que le "
             "sacan cinco puntos a nuestro tipo de modelo, y es cierto, pero eso pasa en un "
             "conjunto de mama, no en el conjunto de metástasis que uno esperaría. Ahí, además, "
             "lo hace gastando un orden de magnitud menos de cómputo, que es el argumento "
             "fuerte del paper. En el otro conjunto queda a la par, no mejor.\n"
             "\n"
             "Y ese otro conjunto es el que dice lo que más me importa, y no está en una tabla "
             "sino en un párrafo. Explican que como las regiones metastásicas pueden ser "
             "extremadamente chicas, tuvieron que renunciar a empezar en aumento bajo y arrancar "
             "más cerca, y que aun así el rendimiento se resiente. O sea: cuando el objeto que "
             "hay que encontrar es chico, el mecanismo de zoom se degrada, y lo dicen ellos. "
             "Una micrometástasis mide cientos de micras. Una figura mitótica de las que marcó "
             "nuestro patólogo mide algo menos de diecisiete micras. Está un orden de magnitud "
             "por debajo del caso que ya les daba problemas.\n"
             "\n"
             "El otro freno no es del método, es de nuestra cohorte. Un modelo que aprende a "
             "hacer zoom no puede hacer zoom a una magnificación que el archivo no contiene, y "
             "nuestro privado está escaneado a veinte aumentos. Hay además un detalle del "
             "código que nos mordería a nosotros específicamente: su preprocesamiento lee el "
             "aumento nativo de una propiedad de un fabricante de escáneres, y si no la "
             "encuentra asume cuarenta aumentos con una advertencia. Nuestro privado es de otro "
             "fabricante y no expone esa propiedad, así que caería en ese camino y se lo "
             "trataría como si fuera el doble de fino de lo que es. Es el error silencioso de "
             "factor dos contra el que ya tenemos regla en el proyecto.\n"
             "\n"
             "Por eso para mitosis lo dejo tercero. Pero lo separo de mitosis y lo dejo vivo "
             "para otra cosa, porque para el problema de que nuestras cohortes están a escalas "
             "distintas sin avisar, un modelo multiescala parametrizado en micras por píxel es "
             "la vía natural, y esa infraestructura ya está escrita y sin lanzar.")
    return s


def lam_npkc(prs):
    """NPKC-MIL, Wang y Yuan. Núcleos como restricción de la pérdida."""
    s = lamina_paper(
        prs, "NPKC-MIL: núcleos dentro de la pérdida", FIG_NPKC,
        ["Wang y Yuan, Nuclei-level prior knowledge constrained MIL for breast histopathology "
         "WSI classification · iScience 2024, Fig. 1.",
         "A features del parche · B parches ordenados por atención · C pérdida de lámina · "
         "D de parche · E de núcleos · F la suma. La rama de núcleos sale del extremo de "
         "atención alta de B."],
        [("Qué propone", [
            "No toca el agregador: suma dos términos a la pérdida.",
            "Uno de parche y otro de núcleos, sobre un grafo."]),
         ("Los números", [
             "96,25 % de acierto contra 86,25 de CLAM, en binario.",
             "Sus ablaciones no cierran: +2,5 y +1,25 sueltas, +12,5 juntas."]),
         ("Lo que lo frena", [
             "Tarea binaria, sobre unas 80 láminas de test.",
             "Hereda entero el costo de segmentar núcleos."])],
        "La rama de núcleos entrena sobre los 8 parches de atención más alta: la idea "
        "que ya se propuso acá, publicada.", "s06 NPKC-MIL")
    notes(s, "Le agrega al modelo una dimensión que hoy no mira: el núcleo.\n"
             "No toca el agregador, suma dos términos a la pérdida.\n"
             "La rama de núcleos corre solo sobre los ocho parches de atención más alta.\n"
             "Sus números son altos, pero las ablaciones no cierran y hay que decirlo.\n"
             "\n"
             "Un modelo con atención mira la lámina y mira el parche, y nunca mira el núcleo. De "
             "esa crítica parte, y es fácil de aceptar: le falta una escala de análisis, y "
             "justo es la escala en la que el patólogo describe lo que ve.\n"
             "\n"
             "La propuesta es conservadora, y eso me gusta. No cambian el agregador, o sea que "
             "el modelo sigue siendo el nuestro. Lo que hacen es agregarle dos penalizaciones a "
             "la pérdida. Una es de parche: una red convolucional que clasifica parches "
             "sueltos. La otra es de núcleos: segmentan los núcleos, convierten cada núcleo en "
             "un nodo, conectan cada nodo con sus vecinos más cercanos y le corren encima una "
             "red convolucional de grafos. A cada núcleo le cuelgan dieciséis medidas hechas a "
             "mano, nueve de forma, como el eje mayor, el área o la excentricidad, y siete de "
             "textura. Entrenan con la suma de los tres términos, y el argumento declarado es "
             "la interpretabilidad: esas medidas tienen nombre.\n"
             "\n"
             "Ahora el detalle que hace que este paper valga la reunión, y está verificado en "
             "el texto. La rama de núcleos no corre sobre la lámina entera. Entrena solamente "
             "sobre los ocho parches de atención más alta. O sea que es exactamente la idea que "
             "ya habíamos discutido acá, la de correr el análisis de núcleos únicamente sobre "
             "los mejores parches que el modelo selecciona, publicada y con números. La única "
             "diferencia es de grado: acá se habló de veinte parches y el paper usa ocho.\n"
             "\n"
             "Los números son buenos, diez puntos de acierto sobre nuestro tipo de modelo. "
             "Pero al lado hay que leer la tabla de sus propias ablaciones, que es la que no "
             "cierra, y prefiero traerla dicha por nosotros antes de que la encuentre otro. "
             "Agregar solo la restricción de parche suma dos puntos y medio. Agregar solo la de "
             "núcleos suma poco más de un punto. Las dos juntas suman doce y medio. El paper no "
             "explica ese salto. Y con unas ochenta láminas de test, cada punto de acierto es "
             "menos de una lámina, así que la tabla entera se mueve con un puñado de casos.\n"
             "\n"
             "Además de eso pesa el encuadre. Su tarea es binaria, lámina sana "
             "contra lámina cancerosa, que es mucho más fácil que las nuestras, y de ahí no se "
             "deduce nada sobre puntuar pleomorfismo en tres clases. Y hereda entero el costo "
             "que ya habíamos pausado: sin segmentar núcleos no hay rama de núcleos. Lo bueno "
             "es que los ocho parches vuelven esa cuenta manejable, y esa es justamente la "
             "cuenta que conviene rehacer hoy.\n"
             "\n"
             "Un último dato de su conjunto de datos, porque nos toca: juntan tres fuentes "
             "distintas para llegar a las cuatrocientas setenta y seis láminas, y las tres "
             "están a escalas físicas distintas. Es nuestro problema de magnificación, sin "
             "tratamiento y sin siquiera una mención.")
    return s


def lam_pleomorfismo(prs):
    """Pleomorfismo nuclear, Mercan et al. El que ataca nuestra tarea de frente."""
    s = lamina_paper(
        prs, "Pleomorfismo nuclear como regresión", FIG_PLEO,
        ["Mercan et al., Deep learning for fully-automated nuclear pleomorphism scoring in "
         "breast cancer · npj Breast Cancer 2022, Fig. 1.",
         "La lámina de entrada, el tumor invasivo resaltado y el puntaje continuo sobre el "
         "tejido, con un detalle a gran aumento y la barra de color a la derecha. Entre "
         "panel y panel, las dos redes."],
        [("Qué propone", [
            "Regresión continua entre 1 y 3, no tres clases.",
            "Antes acota el análisis al tumor invasivo."]),
         ("De dónde sale la etiqueta", [
             "Del promedio de 10 patólogos por región.",
             "El desacuerdo es la señal del continuo."]),
         ("Lo que lo frena", [
             "Nuestro puntaje sale del informe, por lámina.",
             "El código no está publicado."])],
        "Trabaja a 0,5 µm/px, que es prácticamente nuestro privado: el único de los cuatro "
        "que no pide reescalar nada.", "s07 Pleomorfismo")
    notes(s, "Es nuestra tarea, con nuestro nombre, y a nuestra escala física.\n"
             "Puntúa como regresión continua en vez de clasificar en tres clases.\n"
             "La etiqueta es el promedio de diez patólogos por región.\n"
             "Agrega promediando, que es lo contrario de lo que pide el recuento mitótico.\n"
             "El código no está y la anotación que usan es mucho más cara que la nuestra.\n"
             "\n"
             "De los cuatro, este es el que nos toca de frente. Ataca la tarea que ya tenemos "
             "en nuestros datos con ese mismo nombre, y hoy está en cero coma setenta y siete "
             "de área bajo la curva a cinco particiones.\n"
             "\n"
             "Proponen dos cosas y la segunda es la que importa. La primera es acotar el "
             "análisis al tumor invasivo con un detector de células epiteliales que entrenan "
             "aparte y después dejan congelado, para no puntuar tejido que no corresponde. La "
             "segunda es puntuar el pleomorfismo como una regresión continua entre uno y tres, "
             "en vez de clasificarlo en tres clases separadas.\n"
             "\n"
             "Y acá está el corazón del paper, que es de dónde sacan la etiqueta para poder "
             "entrenar una regresión. No usan un consenso ni una votación por mayoría. Usan el "
             "promedio de los puntajes de diez patólogos de seis países sobre cada región. El "
             "argumento es que forzar una mayoría tira a la basura la información que vive en "
             "el desacuerdo, y que ese desacuerdo es justamente la evidencia de que el "
             "pleomorfismo es un continuo y no tres cajones. Por eso pueden entrenar una "
             "regresión: la referencia ya viene con decimales.\n"
             "\n"
             "Cómo pasan del parche a la lámina también vale decirlo, porque es lo contrario de "
             "la otra rama. Recorren la lámina con parches solapados, promedian dentro de cada "
             "bloque y después promedian los bloques. Promediar acá es lo correcto, porque el "
             "pleomorfismo es el aspecto predominante del tejido. El recuento mitótico es lo "
             "opuesto: es un máximo local, se cuenta en el punto más caliente. Las dos tareas "
             "del encargo piden operadores contrarios, y eso significa que no hay un solo "
             "diseño que sirva para las dos.\n"
             "\n"
             "Los números son fuertes. A nivel de región concuerdan con el panel mejor que ocho "
             "de sus diez integrantes, y su acuerdo promedio contra cada patólogo es el más alto "
             "de todo el panel. Lo honesto es leerlo junto con el techo: entre patólogos el "
             "acuerdo es de ese mismo orden, o sea que la tarea tiene un límite humano y el "
             "modelo ya está pegado a él. No es que sea sobrehumano, es que la tarea es "
             "ambigua.\n"
             "\n"
             "Lo que más me interesa es que trabaja a medio micrón por píxel. Nuestro privado "
             "está en cero coma cuatrocientos sesenta y cinco. Es el único de los cuatro donde "
             "la escala física juega a favor sin reescalar nada.\n"
             "\n"
             "Los frenos son concretos. El código no está publicado, y su primera "
             "etapa, el detector de epitelio, es interno del grupo y viene de otro paper suyo: "
             "reproducir la cadena entera no es clonar un repositorio. Lo que sí está público "
             "es el test, ciento dieciocho láminas y un evaluador oficial. Y la brecha de "
             "anotación es real: ellos tienen diez lecturas por región, nosotros tenemos un "
             "puntaje por lámina sacado del informe. Mucho más barato y mucho más grueso. Esa "
             "brecha es lo que habría que discutir antes de prometer nada.")
    return s


def lam_cuadro(prs):
    """Los cuatro en un cuadro. Tabla nativa."""
    s = content(prs, "Los cuatro, en un cuadro")
    headers = ["", "Anotación que exige", "¿La tenemos?", "Qué habilita", "Costo",
               "Lo que lo frena"]
    rows = [
        ["PU learning", "Marcas parciales de mitosis", "Sí, es la nuestra",
         "Contar mitosis, y de ahí el punto caliente", "Alto: detector nuevo",
         "26 marcas en una lámina no entrenan nada"],
        ["ZoomMIL", "Solo la etiqueta de lámina", "Sí, no usa marcas",
         "Contexto multiescala y armonizar µm/px", "Medio: pirámide y GPU",
         "En el privado no hay 40× al que ampliar"],
        ["NPKC-MIL", "Solo la etiqueta de lámina", "Sí, no usa marcas",
         "Medidas de núcleo con nombre dentro de la pérdida", "Alto: segmentar núcleos",
         "Tarea binaria y ablaciones que no cierran"],
        ["Pleomorfismo", "Puntaje por región, 10 lecturas", "No hoy: la tenemos por lámina",
         "Nuestra tarea de grado nuclear, de frente", "Alto: sin código publicado",
         "La brecha entre las dos anotaciones"],
    ]
    # Cuando la prosa compite con una tabla gana la tabla, así que el cuerpo y el alto de
    # fila se eligieron con `alto_tabla` para LLENAR la zona en vez de dejar blanco abajo:
    # a 9 pt la tabla medía 1,64" y sobraba pulgada y media hasta la barra de remate.
    fracs = [0.13, 0.17, 0.14, 0.20, 0.15, 0.21]
    W, FS, RH = 9.28, 11.5, 0.58
    alto = alto_tabla(headers, rows, W, fracs, FS, RH)
    print("  s08 cuadro           tabla estimada en %.2f\"  ·  desde 1,32 llega a %.2f\""
          % (alto, 1.32 + alto))
    if 1.32 + alto > 4.66:
        _AVISOS.append("s08: la tabla llega a %.2f\" y la barra está en 4,66" % (1.32 + alto))
    simple_table(s, 0.36, 1.32, W, headers, rows, fracs, row_h=RH, fs=FS)
    takeaway_bar(s, "Las dos columnas del medio son las que deciden: sin la anotación que "
                    "pide, un método no se puede ni empezar.", t=4.66, size=13)
    notes(s, "El cuadro está ordenado por la anotación que exige cada uno.\n"
             "Tres de los cuatro se pueden empezar con lo que ya tenemos.\n"
             "El cuarto pide una anotación nueva, pero de las baratas.\n"
             "La columna de costo es de montaje, no de cómputo por lámina.\n"
             "\n"
             "Lo ordené por la segunda columna, que es la que importa: qué anotación exige cada "
             "método.\n"
             "\n"
             "Leído así, la mitad izquierda separa el campo. El primero pide marcas parciales "
             "de mitosis, que es literalmente lo que tenemos, y por eso es el único que usa "
             "nuestras marcas en vez de esquivarlas. El segundo y el tercero piden solamente la "
             "etiqueta de la lámina, así que se pueden empezar mañana sin pedirle nada al "
             "patólogo. El cuarto pide algo que hoy no tenemos, un puntaje por región, pero es "
             "la anotación más barata de todas las que están en juego: es poner un número a un "
             "recuadro, no marcar objeto por objeto.\n"
             "\n"
             "La columna de qué habilita es la que responde para qué sirve cada uno, y no se "
             "superponen. Uno produce un conteo de mitosis, que es exactamente la cantidad que "
             "define el puntaje clínico. Otro produce contexto multiescala, que además es la "
             "vía natural para el problema de que nuestras cohortes están a escalas distintas. "
             "El tercero mete medidas de núcleo con nombre dentro del entrenamiento. Y el "
             "cuarto ataca nuestra tarea de grado nuclear directamente.\n"
             "\n"
             "La columna de costo se refiere a montar el método, no a lo que cuesta correrlo "
             "por lámina. Tres de los cuatro son altos, y por razones distintas: uno porque hay "
             "que escribir un andamiaje de detección que hoy no tenemos, otro porque arrastra "
             "la segmentación de núcleos, y el último porque el código no está publicado.\n"
             "\n"
             "Y la última columna es la que hay que tener a mano cuando alguien se entusiasme "
             "con cualquiera de las filas.")
    return s


def lam_recomendacion(prs):
    """La recomendación, en dos carriles."""
    s = content(prs, "Qué propongo hacer primero")
    gt, gh, gw = 1.32, 3.20, 4.42
    for gl, tarea, pasos, remate in (
        (0.42, "Mitosis",
         [("Paso 1 · Go / no-go barato",
           "Un detector público de mitosis sobre unas pocas láminas, medido contra las 26 marcas"),
          ("Paso 2 · Solo si el 1 pasa",
           "Ajuste fino con positivos parciales, y de ahí el conteo en el punto caliente")],
         "Primero TCGA, que ya está a la escala de los datos públicos"),
        (5.16, "Grado nuclear",
         [("Entrar por el paper de pleomorfismo",
           "Es el único que trabaja a nuestra escala física, y la anotación que pide es barata"),
          ("NPKC-MIL queda detrás",
           "Aporta la receta de correr núcleos sobre los mejores parches, con su costo atado")],
         "Sin reescalar nada: 0,5 µm/px contra nuestros 0,465"),
    ):
        _grupo(s, gl, gt, gw, gh)
        _proc(s, gl + (gw - 2.60) / 2, 1.44, 2.60, 0.40, tarea, size=13)
        for i, (nombre, detalle) in enumerate(pasos):
            _proc_claro(s, gl + 0.22, 2.00 + i * 1.14, gw - 0.44, 0.86,
                        nombre, size=11.5, dim=detalle)
        _conn(s, gl + gw / 2, 2.88, gl + gw / 2, 3.12)
        _dato(s, gl + 0.22, 4.10, gw - 0.44, 0.32, remate, size=9.5)
    takeaway_bar(s, "En las dos ramas lo barato va primero, y en mitosis el primer paso no "
                    "depende del paper que lo motivó.", t=4.66, size=13)
    notes(s, "Mitosis sigue igual: la familia de positivos parciales primero, en dos pasos.\n"
             "El primer paso es barato y se puede cerrar sin comprometer nada más.\n"
             "Grado nuclear entra por el paper de pleomorfismo, por la escala física.\n"
             "El de núcleos queda detrás, aportando una receta y arrastrando su costo.\n"
             "\n"
             "Las separo en dos carriles porque no compiten entre sí: se pueden avanzar en "
             "paralelo y con presupuestos muy distintos.\n"
             "\n"
             "En mitosis la recomendación no cambia respecto de lo que ya había dicho, y va en "
             "dos pasos. El primero es un go o no go barato, antes de gastar una hora de placa: "
             "agarrar un detector de mitosis ya entrenado con datos públicos, correrlo sobre "
             "unas pocas láminas nuestras y medir cuánto de lo que marcó el patólogo encuentra. "
             "Si no encuentra lo que el patólogo marcó, esa rama se cierra ahí y nos ahorramos "
             "todo lo demás. Es el mismo patrón que ya nos ahorró casi un día de cómputo en "
             "otro eje del proyecto.\n"
             "\n"
             "El segundo paso solo existe si el primero pasa: ajustar el detector con positivos "
             "parciales a medida que lleguen más láminas anotadas, y usar el conteo en el punto "
             "caliente como entrada de un clasificador chico que devuelva el puntaje. Ese "
             "conteo, dicho sea de paso, no lo entrega ninguno de los cuatro papers: lo ponemos "
             "nosotros.\n"
             "\n"
             "Hay un orden de cohortes que importa y conviene fijarlo antes de diseñar nada. "
             "Los datos públicos de mitosis se anotan a cuarenta aumentos, que es donde se "
             "cuenta clínicamente, y ahí TCGA calza sin reescalar. Nuestro privado está al "
             "doble de grueso, así que es un brazo aparte y con su propio riesgo. La lámina que "
             "está anotada es privada, así que la validación contra las marcas cae justo en el "
             "brazo difícil, y eso hay que tenerlo previsto y no descubrirlo después.\n"
             "\n"
             "En grado nuclear la entrada es por el paper de pleomorfismo, y el motivo es "
             "sobre todo la escala: es el único de los cuatro que trabaja donde ya estamos. "
             "Además la anotación que pide, un puntaje por región, es la más barata de "
             "conseguir de todas las que aparecieron hoy.\n"
             "\n"
             "El de núcleos lo dejo detrás, pero no lo descarto, porque aporta algo que ninguno "
             "de los otros aporta: la receta publicada de correr el análisis de núcleos "
             "solamente sobre los mejores parches. Lo que arrastra es el costo de la "
             "segmentación, que es la conversación que ya habíamos pausado.")
    return s


def lam_tres_cosas(prs):
    """Las tres cosas que hay que decir sí o sí."""
    s = content(prs, "Tres cosas que no quiero dejar sin decir")
    textos = [
        "El primer paso de mitosis no depende del paper de Zhao: necesita pesos públicos de "
        "un detector, y esos los tiene el ecosistema del desafío de mitosis. La prueba barata "
        "está desacoplada de la decisión cara.",
        "El régimen de anotación que ese paper evalúa no es el nuestro: ellos retienen el "
        "73 % de las marcas y nosotros tenemos 26 en 4799 parches. No hay forma de extrapolar "
        "su curva hasta ahí.",
        "La escala física juega en dirección contraria en cada rama: la mitosis se cuenta a "
        "40× y estamos a 20×, mientras que el paper de pleomorfismo trabaja a 0,5 µm/px, que "
        "es prácticamente nuestro privado.",
    ]
    for i, txt in enumerate(textos):
        add_card(s, 0.60, 1.42 + i * 1.14, 8.80, 0.94, i + 1, txt, size=13)
    takeaway_bar(s, "Las tres son salvedades, no objeciones: acotan lo que se puede prometer "
                    "de cada rama.", t=4.86, size=13)
    notes(s, "Lo barato de mitosis se puede hacer aunque el paper caro no se apruebe.\n"
             "Su régimen de anotación es mucho más denso que el nuestro y hay que decirlo.\n"
             "La escala física empuja en sentidos contrarios según la tarea.\n"
             "Las tres son salvedades, no motivos para no hacer nada.\n"
             "\n"
             "Estas tres se pierden si la reunión se va por las ramas, y las tres cambian lo que "
             "puedo prometer.\n"
             "\n"
             "La primera es la mejor propiedad del plan de mitosis y quiero que quede clara. El "
             "primer paso, la prueba barata, no depende del paper que la motivó. Lo que "
             "necesita son los pesos de un detector de mitosis ya entrenado, y esos los tiene "
             "el ecosistema de un desafío público que existe justamente para eso. O sea que "
             "podemos saber si la rama tiene sentido antes de comprometernos a escribir nada. "
             "La prueba barata está desacoplada de la decisión cara.\n"
             "\n"
             "La segunda es la salvedad más importante de todo lo que traje. Ese paper evalúa "
             "un escenario donde el anotador dejó tres de cada cuatro marcas puestas. Nosotros "
             "estamos en veintiséis marcas sobre casi cinco mil parches. No es la misma "
             "situación con menos datos, es un régimen distinto, el paper no lo evalúa y "
             "cualquier extrapolación sería inventada. El planteo del método sigue siendo el "
             "correcto para nosotros, pero sus números no se trasladan.\n"
             "\n"
             "Y la tercera es la que ordena las prioridades. La escala física empuja en "
             "sentidos contrarios según la tarea. La mitosis se cuenta a cuarenta aumentos y "
             "nuestro privado está a veinte, así que toda esa rama arrastra un reescalado con "
             "riesgo desde el primer día. El paper de pleomorfismo trabaja a medio micrón por "
             "píxel, que es prácticamente nuestro privado, y ahí no hay nada que reescalar. Si "
             "hubiera que elegir una sola rama para gastar el trimestre, este es el argumento "
             "más sólido que tengo, y no depende de qué método sea mejor.")
    return s


def lam_preguntas(prs):
    """Las preguntas que deciden."""
    s = content(prs, "Las preguntas que deciden")
    gt, gh, gw = 1.38, 2.92, 4.42
    for gl, tarea, preguntas in (
        (0.42, "Mitosis",
         [("¿Vienen más láminas anotadas?",
           "Con 26 marcas en una lámina no se entrena; con varias decenas de láminas, sí, "
           "y pueden ser parciales"),
          ("¿Quién es «GDT»?",
           "Saber de quién es la anotación decide si sirve como referencia y si se puede pedir más")]),
        (5.16, "Grado nuclear",
         [("¿Se puede conseguir puntaje por región?",
           "Es lo que separa nuestra etiqueta de informe de la que usa el paper de pleomorfismo"),
          ("¿Hay GPU para segmentar núcleos?",
           "Solo sobre los mejores parches de cada lámina, que es lo que lo vuelve viable")]),
    ):
        _grupo(s, gl, gt, gw, gh)
        _proc(s, gl + (gw - 2.60) / 2, 1.50, 2.60, 0.40, tarea, size=13)
        for i, (preg, porque) in enumerate(preguntas):
            _proc_claro(s, gl + 0.22, 2.06 + i * 1.10, gw - 0.44, 0.94,
                        preg, size=11.5, dim=porque)
    takeaway_bar(s, "Las cuatro se responden hoy, sin medir nada, y de ellas depende qué se "
                    "empieza primero.", t=4.46, size=13)
    notes(s, "Cuatro preguntas, dos por rama, y ninguna necesita que vayamos a medir.\n"
             "En mitosis deciden si la rama existe: sin más láminas anotadas, no existe.\n"
             "En grado nuclear deciden por cuál de los dos papers se entra.\n"
             "Con las respuestas de hoy puedo dejar cerrado qué se hace primero.\n"
             "\n"
             "Cierro con las preguntas, porque son las que convierten todo esto en trabajo.\n"
             "\n"
             "En mitosis la primera es la que decide si la rama existe. Con veintiséis marcas "
             "en una sola lámina no se entrena nada. Con varias decenas de láminas sí, y lo "
             "importante es que el método admite que esas marcas sean parciales. Eso convierte "
             "un pedido vago, necesitamos más anotaciones, en un pedido con forma: marcas por "
             "punto sobre figuras mitóticas, en un número concreto de láminas, y está bien que "
             "estén incompletas. Es mucho más fácil de conceder que pedir una anotación "
             "exhaustiva.\n"
             "\n"
             "La segunda parece menor y no lo es. En el archivo de anotaciones que tenemos "
             "aparecen unas iniciales que no sé de quién son. Saber quién marcó decide dos "
             "cosas: si eso sirve como referencia para medir, y a quién habría que pedirle las "
             "láminas que faltan.\n"
             "\n"
             "En grado nuclear la primera pregunta es la que separa nuestra etiqueta de la de "
             "ellos. Hoy nuestro puntaje sale del informe y es por lámina. Ellos tienen un "
             "puntaje por región, promediado entre varios lectores. Sin al menos algo por "
             "región no podemos reproducir lo que hacen, y quiero saber si eso es pedible o si "
             "hay que diseñar alrededor de la etiqueta gruesa.\n"
             "\n"
             "Y la última es de presupuesto. Segmentar núcleos sobre láminas enteras ya "
             "sabemos que no cierra por costo. Sobre los mejores parches de cada lámina la "
             "cuenta es otra, y es la que vuelve viable la idea que el paper de núcleos ya "
             "publicó. Necesito saber si hay lugar en la placa para eso.\n"
             "\n"
             "Las cuatro se pueden responder hoy mismo, sin medir nada, y con esas respuestas "
             "puedo dejar cerrado qué empiezo primero.")
    return s


# ============================================================================
# El orden del deck
# ============================================================================
def build():
    prs, keep_ids = base_from_template()
    prs.slide_width = Inches(SW); prs.slide_height = Inches(SH)

    lam_portada(prs)
    lam_por_que(prs)

    # ---- un paper por lámina, mismo molde en las cuatro ----
    lam_pu_learning(prs)
    lam_zoommil(prs)
    lam_npkc(prs)
    lam_pleomorfismo(prs)

    # ---- los cuatro juntos, y qué hacer con ellos ----
    lam_cuadro(prs)
    lam_recomendacion(prs)
    lam_tres_cosas(prs)
    lam_preguntas(prs)

    # ---- cierre: reflow, auditoría, escala al tamaño del template, tipografía ----
    reflow_onco(prs, skip=keep_ids)
    auditar(prs, skip=keep_ids)
    scale_deck_to_1610(prs, skip=keep_ids)
    forzar_barlow(prs)
    os.makedirs(OUT_DIR, exist_ok=True)
    prs.save(DST)
    print("Guardado:", DST, "·", len(prs.slides), "slides · 13.333x7.5")
    if _AVISOS:
        print("  PENDIENTES DE MAQUETADO: %d" % len(_AVISOS))
        for a in _AVISOS:
            print("   ·", a)


if __name__ == "__main__":
    build()
