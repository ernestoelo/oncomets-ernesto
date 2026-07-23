#!/usr/bin/env python
"""generate_b6_deck.py — deck de la reunión del VIERNES (B6): mammoth / interpretabilidad OBJ-A.

Formato = deck B4 (10x5.625 in, Barlow, paleta teal) — corrige el formato que Benjamín
rechazó del B5 (13.333, Carlito). Contenido = mecanismo mammoth (slides B5 re-formateadas)
+ interpretabilidad OBJ-A (material nuevo del run 30-jun). Eje = ENTENDIMIENTO (no rendimiento).

- Slides de OBJETIVOS/contenido → NATIVAS en Barlow a tamaños B4 (el fix de Benjamín).
- Diagramas de arquitectura → reusados de los .pptx standalone, ESCALADOS x0.75 a 10x5.625
  (mismo aspect ratio; conservan Carlito, convención aceptada para diagramas). Editables.
- Slides de mecanismo (interior/tensor, keep_slots) → NATIVAS con matrices+dims Y código real
  de mammoth.py en paralelo (pedido de Ernesto para la reunión).
- Figuras de interpretabilidad (heatmap/top-k/cross-slide) y del paper → como imagen (excepción).
- Notas del presentador → guion HABLADO en prosa (sin etiquetas de fase, sin nº de job/nombres),
  integrando las explicaciones del estudio (5 bloques de estudio_reunion_viernes.md).

Uso:
  PYTHONPATH=/media/administrador/Storage1/sdonoso/clam_testing2/.pylibs \
  /home/sdonoso/miniconda3/envs/clam_latest/bin/python \
    sprints/B6_sprint6/presentacion_viernes/generate_b6_deck.py
"""
import copy
import io
import os

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn
from PIL import Image

REPO = "/media/administrador/Storage1/sdonoso/clam_testing2/oncomets-ernesto"
PRES = os.path.join(REPO, "papers/presentations")
ASSETS_BRAND = os.path.join(PRES, "assets_branding")
PAPER_FIGS = os.path.join(ASSETS_BRAND, "paper_figs")
OBJA = os.path.join(REPO, "sprints/B5_sprint5/mammoth_entendimiento/interpretabilidad")
OUT_DIR = os.path.join(REPO, "sprints/B7_sprint7/presentacion_b7")
DST = os.path.join(OUT_DIR, "CLAM_Sprint7.pptx")
# Portada branded de Sebastián (extraída de Plantilla.pptx s00). El deck se construye a
# 10x5.625 (paleta/fuentes = las de Plantilla) y se ESCALA x1.3333 a 13.333x7.5 al final
# (misma relación 16:9) → re-base sobre el template de Sebastián. Ver scale_deck_to_1610().
PORTADA_PLANTILLA = os.path.join(OUT_DIR, "assets/portada_plantilla.jpg")

# --- BASE del deck = "Modelo OncoMets Spatial V1 Deep-LLM-V.pptx" (template VÁLIDO) ---
# Ernesto fijó este archivo como el template a respetar (19-jul), no Plantilla.pptx.
# Sus s02-s16 son las mismas láminas técnicas que Plantilla s04-s18, con la cabecera
# OncoMets IDÉNTICA al píxel (mismos nombres de shape: Google Shape;115;p13 / 197;p29 /
# 198;p29) → la geometría medida vale para los dos.
#
# Por qué se construye SOBRE el archivo y no con Presentation():
# el template EMBEBE sus fuentes (ppt/fonts/*.fntdata + <p:embeddedFontLst>). Barlow no
# está en el servidor (fc-list → 0) ni necesariamente en la máquina donde se presenta: un
# deck generado con el template default de python-pptx (cero fuentes embebidas) hace que
# PowerPoint SUSTITUYA la tipografía y se vea fuera de template AUNQUE las cabeceras estén
# bien. Ese era el síntoma real. Abrir el template y borrarle las láminas preserva el
# paquete de fuentes (verificado por round-trip), el theme y el master.
#
# OJO — este template embebe SOLO Barlow + Cambria Math. NO trae "Barlow ExtraBold" ni
# "Consolas" (Plantilla sí). Por eso F_TITLE usa Barlow bold: ver la nota en F_TITLE.
TEMPLATE = os.path.join(REPO, "sprints/B7_sprint7/Modelo OncoMets Spatial V1 Deep-LLM-V.pptx")
# Láminas del template que se CONSERVAN tal cual (nativas, ya a 13.333): portada de marca
# y lámina de título. Reemplazan a la portada JPG que se usaba antes.
TPL_KEEP = (0, 1)
FECHA_REUNION = "22/07/2026"   # miércoles

# --- diagramas standalone (13.333x7.5, sin header) que reusamos escalados ---
DIAG_MAM = os.path.join(PRES, "Diagrama_CLAM_mammoth.pptx")     # s0 pipeline · s1 interior/tensor
DIAG_FUSED = os.path.join(PRES, "Diagrama_mammoth_fused.pptx")  # s0 flujo oficial · s1 keep_slots
# --- figuras ---
FIG1_MAM = os.path.join(PAPER_FIGS, "mammoth_fig1_overview.png")  # t-SNE + barras (paper Fig 1)
FIG3_MAM = os.path.join(PAPER_FIGS, "mammoth_fig3_routing.png")   # ruteo→fenotipo (paper Fig 3)
FIG2_ARCH = os.path.join(PAPER_FIGS, "mammoth_fig2_arch.png")     # arquitectura oficial (Fig 2, 4375x1758)
A14Q = os.path.join(OBJA, "TCGA-E2-A14Q-01Z-00-DX1_cdis_f0")
HEATMAP_MONTAGE = os.path.join(A14Q, "heatmap_montage.png")
TOPK_SUBSET = os.path.join(A14Q, "topk_subset_6experts.png")
XSLIDE_E8 = os.path.join(OBJA, "cross_slide/expert_08_crossslide.png")
XSLIDE_E26 = os.path.join(OBJA, "cross_slide/expert_26_crossslide.png")

# ---- paleta Deep-LLM-V (medida sobre el template válido) ----
# El template pinta el contenido técnico con cuatro colores y nada más:
#   #3E6877  bloque de proceso   (rounded-rect relleno, texto Barlow bold BLANCO encima)
#   #386271  conector / borde    (líneas de 2.37pt)
#   #CDDFE1  panel contenedor y operadores (+, ×)
#   #B7B7B7  bloque de dato      (rect relleno, texto Barlow regular NEGRO encima)
#   #0E2841  borde/fondo oscuro
# Hasta el 19-jul el cuerpo del deck seguía usando la paleta de B4 (#217589, #31859C,
# #DDEAEE, #F3F8F9 y una familia naranja #E2723B/#B4521E que NO existe en Deep-LLM-V):
# se había migrado la cabecera pero no el cuerpo. Acá se migra el cuerpo. Los NOMBRES de
# las constantes se conservan a propósito — los usa build() en ~700 líneas — y lo que
# cambia es el valor, para que el cuerpo hable el mismo idioma que la cabecera.
ONCO_DARK  = RGBColor(0x3E, 0x68, 0x77)   # bloque de proceso
ONCO_CONN  = RGBColor(0x38, 0x62, 0x71)   # conector / borde
ONCO_PANEL = RGBColor(0xCD, 0xDF, 0xE1)   # panel contenedor / operador
ONCO_DATA  = RGBColor(0xB7, 0xB7, 0xB7)   # bloque de dato
ONCO_INK   = RGBColor(0x0E, 0x28, 0x41)   # borde/fondo oscuro

TEAL_TITLE = ONCO_DARK                    # antes #217589
TEAL_SQ    = ONCO_CONN                    # antes #31859C
BAR_GRIS   = RGBColor(0xF2, 0xF2, 0xF2)   # barra del header B4 (ya no se usa)
TEAL_DIV   = ONCO_CONN                    # antes #2E7E8F
LAV_TITLE  = RGBColor(0xFF, 0xFF, 0xFF)   # antes #CDD6F4 (lavanda B4, fuera de paleta)
TEAL_SUB   = ONCO_PANEL                   # antes #B8D4D9
TEAL_CARD  = ONCO_PANEL                   # antes #DDEAEE
# El template no define un segundo claro. Para el banding de tablas y las tarjetas
# alternas se usa un TINTE de #CDDFE1 (misma familia de tono, no un color nuevo): un
# blanco puro desaparecería sobre el fondo blanco de la lámina.
TEAL_CARD2 = RGBColor(0xE9, 0xF1, 0xF2)   # antes #F3F8F9
CODE_BG    = ONCO_INK                     # panel de código (antes #1E2A2E)
CODE_FG    = ONCO_PANEL                   # texto de código
CODE_CMT   = ONCO_DATA                    # comentarios de código
GRIS_BODY  = RGBColor(0x59, 0x59, 0x59)   # prosa secundaria (neutro, no es marca)
GRIS_TXT   = RGBColor(0x55, 0x55, 0x55)
# Deep-LLM-V no tiene familia cálida. Lo que en B4 era el naranja de "acento / lo nuevo"
# pasa al primario del template; el contraste ahora lo da oscuro-vs-claro, no cálido-vs-frío.
ORA_T      = ONCO_DARK                    # antes #B4521E
ORA_ACC    = ONCO_DARK                    # antes #E2723B
INK        = RGBColor(0x22, 0x22, 0x22)
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
BLACK      = RGBColor(0x00, 0x00, 0x00)   # texto sobre bloque de dato (archetipo B7B7B7)

# Varias láminas emparejan un panel "pregunta abierta / lo nuevo" con uno "respuesta /
# base", que en B4 se distinguían por temperatura (naranja vs teal). Deep-LLM-V no tiene
# familia cálida, así que el par se distingue por PESO de tono: el destacado va relleno
# con el celeste sólido del template, el neutro con el tinte claro. El gris #B7B7B7 NO
# sirve acá — lee como "deshabilitado", no como "prestá atención".
PANEL_ACC  = ONCO_PANEL                   # panel destacado
PANEL_NEU  = TEAL_CARD2                   # panel neutro

# El template válido embebe SOLO Barlow (regular/bold/italic) + Cambria Math. Antes
# F_TITLE era "Barlow ExtraBold": esa fuente NO viaja en el paquete, así que PowerPoint la
# sustituía. Barlow bold sí viaja, y es exactamente lo que usan los títulos técnicos del
# propio template (25pt Barlow bold #3E6877) → más fiel Y sin riesgo de sustitución.
F_TITLE = "Barlow"
F_BODY  = "Barlow"
# Consolas tampoco viaja en este template, pero se instala con Office en Windows, que es
# donde se presenta. Se conserva para los paneles de código; si alguna vez falla, el
# reemplazo seguro es Barlow (se pierde el ancho fijo).
F_MONO  = "Consolas"

# --- geometría B4 (10 x 5.625), extraída del volcado real ---
SW, SH = 10.0, 5.625
HDR = 0.785
LOGO = os.path.join(ASSETS_BRAND, "logo_header.png")
PORTADA = os.path.join(ASSETS_BRAND, "portada_fullbleed.jpg")
CHECK_VERDE = os.path.join(ASSETS_BRAND, "check_verde.png")  # ticket "hecho" (reusado del deck B3)
# Pill "En progreso" = el arquetipo "bloque de dato" del template (Google Shape;391):
# relleno #B7B7B7 con texto Barlow NEGRO. Antes era naranja muy claro (#FDEFE6), fuera
# de paleta. El gris neutro lee como "pendiente" sin inventar un color.
PROG_BG   = ONCO_DATA

# --- cabecera OncoMets (láminas TÉCNICAS de Plantilla.pptx s04-s18) ---
# El logo se extrajo del blipFill del freeform id=4 de Plantilla s05 (258x155 RGBA).
ONCO_LOGO  = os.path.join(OUT_DIR, "assets/logo_oncomets.png")
ONCO_TITLE = RGBColor(0x3E, 0x68, 0x77)  # título Barlow bold (medido en Plantilla)
ONCO_LINE  = RGBColor(0x3D, 0x68, 0x76)  # línea horizontal bajo el encabezado

# Geometría LITERAL de la cabecera técnica, medida sobre Plantilla s05/s06 (13.333-space)
# y convertida a 10-space dividiendo por 1.3333. Vuelve a 13.333 exacta al escalar al final.
#   logo   L=0.750 T=0.437 W=1.471 H=0.884   (freeform con blipFill en Plantilla)
#   título L=2.741 T=0.850 W=9.173 H=0.471   Barlow bold 25pt #3E6877
#   línea  L=0.000 T=1.421 W=13.334 H=0.140  (grupo en Plantilla, aquí un rect plano)
ONCO_LOGO_L, ONCO_LOGO_T, ONCO_LOGO_W, ONCO_LOGO_H = 0.5625, 0.3278, 1.1033, 0.6630
ONCO_TIT_L, ONCO_TIT_W = 2.0558, 6.8799
# El template da la caja del título en T=0.6375 H=0.3533 (base = 0.9908), dimensionada para
# UNA línea. Varios títulos nuestros son más largos que "Patch Encoder" y caen a dos: con esa
# caja la 2ª línea desbordaba hacia abajo y la cortaba la línea teal (y=1.0658). Se conserva
# la MISMA base y se extiende la caja hacia ARRIBA con anclaje inferior: los títulos de una
# línea caen exactamente donde el template los pone, y los de dos crecen hacia el aire libre
# de la banda (el logo llega hasta x=1.67 y el título arranca en 2.06 → no se tocan).
ONCO_TIT_BASE = 0.9908
ONCO_TIT_T, ONCO_TIT_H = 0.2950, ONCO_TIT_BASE - 0.2950
ONCO_TIT_SZ = 18.75
ONCO_LINE_T, ONCO_LINE_H = 1.0658, 0.1050
ONCO_BAND = ONCO_LINE_T + ONCO_LINE_H   # 1.1708 · donde puede empezar el contenido

# --- assets sección magnificación multi-escala (B6, reunión Sebastián) ---
MSCROP = os.path.join(ASSETS_BRAND, "multiscale_crop")
MS_FINE = os.path.join(MSCROP, "fine_112um.png")   # crop real TCGA-BRCA campo 112µm (~20×)
# El crop de contexto va en la copia con el recuadro RECOLOREADO a la paleta del template
# (ver recolor_crop_box.py): el original de B6 lo trae en el naranja de B4 y chocaba con el
# esquema nativo de la misma lámina, que marca ese mismo campo en teal.
MS_CTX = os.path.join(OUT_DIR, "assets/context_512um_onco.png")

# --- assets sección comparación pareada (B7: atención CLAM vs mammoth + Q1) ---
INTERP = os.path.join(REPO, "results/b7_mammoth_interp/interpretabilidad")
# La lámina CDIS positiva: el contraste de entropía más fuerte de las 7 (0.642 -> 0.927),
# y es la tarea donde la variante también mide mejor.
ATT_SBS = os.path.join(INTERP, "carcinoma_ductal_insitu_presente_ci_reform",
                       "TCGA-D8-A1XB-01Z-00-DX2", "attention_side_by_side.png")
Q1_JSON = os.path.join(REPO, "sprints/B7_sprint7/respuesta_q1_expertos_slots.json")

_uid = [2000]


# ============================================================================
# Helpers base
# ============================================================================
def _blank(prs):
    for lay in prs.slide_layouts:
        if (lay.name or "").lower().strip() == "blank":
            return lay
    return prs.slide_layouts[6]


def base_from_template():
    """Abre el template válido y le borra las láminas salvo TPL_KEEP.

    Lo que se conserva y es el motivo de hacerlo así: ppt/fonts/*.fntdata + el
    <p:embeddedFontLst> de presentation.xml (Barlow, Cambria Math), el theme y el slide
    master. Borrar una slide en python-pptx no tiene API: hay que sacarla del sldIdLst Y
    soltar la relación, si no queda huérfana en el paquete.

    Devuelve (prs, keep_ids), donde keep_ids identifica las láminas heredadas. Son las
    únicas que YA están a 13.333, así que scale_deck_to_1610() debe saltearlas.
    """
    prs = Presentation(TEMPLATE)
    # Los elementos a conservar se capturan ANTES de tocar el sldIdLst: después de borrar,
    # los índices se corren y ya no se puede resolver TPL_KEEP.
    keep_ids = {id(prs.slides[i]._element) for i in TPL_KEEP}
    lst = prs.slides._sldIdLst
    for i, sld in enumerate(list(lst)):
        if i in TPL_KEEP:
            continue
        prs.part.drop_rel(sld.get(
            "{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id"))
        lst.remove(sld)
    return prs, keep_ids


def new_slide(prs):
    """Lámina en blanco SIN los placeholders del layout.

    El layout BLANK del template arrastra DATE/FOOTER/SLIDE_NUMBER; add_slide los clona en
    la lámina y aparecen como cuadros vacíos. El deck dibuja todo a mano, así que se
    eliminan. (Con el template default de python-pptx esto era un no-op.)"""
    s = prs.slides.add_slide(_blank(prs))
    for ph in list(s.placeholders):
        ph._element.getparent().remove(ph._element)
    return s


def _set_runs(tf, lines, anchor=MSO_ANCHOR.TOP):
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    for i, ln in enumerate(lines):
        txt, sz, bold, col = ln[0], ln[1], ln[2], ln[3]
        font = ln[4] if len(ln) > 4 else F_BODY
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = ln[5] if len(ln) > 5 else PP_ALIGN.LEFT
        p.space_after = Pt(4)
        r = p.add_run(); r.text = txt
        r.font.size = Pt(sz); r.font.bold = bold; r.font.name = font; r.font.color.rgb = col


def add_textbox(slide, l, t, w, h, lines, anchor=MSO_ANCHOR.TOP):
    tb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    _set_runs(tb.text_frame, lines, anchor=anchor)
    return tb


def notes(slide, text):
    """Guion HABLADO en prosa (sin etiquetas de fase, sin nº de job ni nombres)."""
    slide.notes_slide.notes_text_frame.text = text


def _rect(slide, l, t, w, h, color, line=None):
    sp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(l), Inches(t), Inches(w), Inches(h))
    sp.fill.solid(); sp.fill.fore_color.rgb = color
    if line is None:
        sp.line.fill.background()
    else:
        sp.line.color.rgb = line; sp.line.width = Pt(1)
    sp.shadow.inherit = False
    return sp


def logo_mark(slide, size=0.785, logo_h=0.645):
    _rect(slide, 0, 0, size, size, TEAL_SQ)
    slide.shapes.add_picture(LOGO, Inches(0.065), Inches(0.045), height=Inches(logo_h))


def header(slide, title, bar=True, size=26):
    """Header B4: barra gris + cuadrado teal + logo + título Barlow ExtraBold teal."""
    if bar:
        _rect(slide, 0.0, 0.0, SW, HDR, BAR_GRIS)
    logo_mark(slide)
    if title:
        tb = slide.shapes.add_textbox(Inches(0.99), Inches(0.10), Inches(8.9), Inches(0.62))
        _set_runs(tb.text_frame, [(title, size, True, TEAL_TITLE, F_TITLE)], anchor=MSO_ANCHOR.MIDDLE)


def header_oncomets(slide, title, size=22):
    """Header OncoMets = el de las láminas TÉCNICAS de Plantilla.pptx (s04-s18).

    Plantilla usa DOS cabeceras según el tipo de lámina: la Environ (cuadro teal +
    barra gris) para lo administrativo/agenda, y ésta para el contenido técnico
    (s04 "Patch Encoder" en adelante). Identidad medida sobre Plantilla s05:
      logo OncoMets, título Barlow bold #3E6877, línea horizontal #3D6876.

    Geometría LITERAL de Plantilla s05/s06, convertida a 10-space (÷1.3333). Antes se
    usaba una versión compactada a la banda HDR (0.785) para no re-maquetar el contenido;
    ahora se replica la banda real (termina en y=1.171) y el contenido de las técnicas se
    baja en bloque con reflow_onco(), que además comprime las láminas que se pasarían.
    """
    pic = slide.shapes.add_picture(ONCO_LOGO, Inches(ONCO_LOGO_L), Inches(ONCO_LOGO_T),
                                   Inches(ONCO_LOGO_W), Inches(ONCO_LOGO_H))
    pic.name = "ONCOHDR_logo"
    if title:
        # 18.75 = 25pt tras el x1.3333 = EXACTAMENTE el de los títulos técnicos de Plantilla.
        # Es tope, no tamaño fijo: los títulos cortos pueden venir más chicos por diseño.
        tb = slide.shapes.add_textbox(Inches(ONCO_TIT_L), Inches(ONCO_TIT_T),
                                      Inches(ONCO_TIT_W), Inches(ONCO_TIT_H))
        tb.name = "ONCOHDR_title"
        _set_runs(tb.text_frame, [(title, min(size, ONCO_TIT_SZ), True, ONCO_TITLE, F_BODY)],
                  anchor=MSO_ANCHOR.BOTTOM)
    _rect(slide, 0.0, ONCO_LINE_T, SW, ONCO_LINE_H, ONCO_LINE).name = "ONCOHDR_line"


def add_image_fit(slide, path, l, t, w, h, align="center"):
    iw, ih = Image.open(path).size
    ar = iw / ih; box_ar = w / h
    if ar > box_ar:
        nw = w; nh = w / ar
    else:
        nh = h; nw = h * ar
    nl = l + (w - nw) / 2
    nt = t + (h - nh) / 2 if align != "top" else t
    slide.shapes.add_picture(path, Inches(nl), Inches(nt), Inches(nw), Inches(nh))


def add_card(slide, l, t, w, h, idx, text, size=14):
    """Tarjeta numerada: óvalo #3E6877 con el número en blanco + texto Barlow."""
    fill = TEAL_CARD if idx % 2 == 0 else TEAL_CARD2
    sp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(l), Inches(t), Inches(w), Inches(h))
    sp.fill.solid(); sp.fill.fore_color.rgb = fill
    sp.line.color.rgb = TEAL_SQ; sp.line.width = Pt(1.25); sp.shadow.inherit = False
    cd = 0.44
    circ = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(l + 0.16), Inches(t + (h - cd) / 2),
                                  Inches(cd), Inches(cd))
    circ.fill.solid(); circ.fill.fore_color.rgb = ORA_ACC
    circ.line.fill.background(); circ.shadow.inherit = False
    _set_runs(circ.text_frame, [(str(idx), 15, True, WHITE, F_TITLE, PP_ALIGN.CENTER)],
              anchor=MSO_ANCHOR.MIDDLE)
    tb = slide.shapes.add_textbox(Inches(l + 0.74), Inches(t), Inches(w - 0.88), Inches(h))
    _set_runs(tb.text_frame, [(text, size, True, INK, F_BODY)], anchor=MSO_ANCHOR.MIDDLE)


def caption(slide, l, t, w, text, size=11, col=GRIS_TXT, align=PP_ALIGN.CENTER, bold=False):
    add_textbox(slide, l, t, w, 0.4, [(text, size, bold, col, F_BODY, align)])


def status_done(slide, cx, cy, size=0.42):
    """Ticket verde 'hecho/cerrado' (icono reusado del deck B3), centrado en (cx, cy)."""
    slide.shapes.add_picture(CHECK_VERDE, Inches(cx - size / 2), Inches(cy - size / 2),
                             Inches(size), Inches(size))


def status_progress(slide, cx, cy, w=1.5, h=0.44):
    """Pill 'En progreso': el arquetipo gris de dato del template, centrada en (cx, cy)."""
    sp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(cx - w / 2), Inches(cy - h / 2),
                                Inches(w), Inches(h))
    sp.fill.solid(); sp.fill.fore_color.rgb = PROG_BG
    sp.line.color.rgb = ONCO_CONN; sp.line.width = Pt(1.25); sp.shadow.inherit = False
    # Texto NEGRO, no teal: sobre el gris #B7B7B7 el teal oscuro queda en ~2.9:1 de
    # contraste, insuficiente a 11pt. El negro es además lo que usa el template encima
    # de sus bloques grises.
    _set_runs(sp.text_frame, [("En progreso", 11, True, BLACK, F_BODY, PP_ALIGN.CENTER)],
              anchor=MSO_ANCHOR.MIDDLE)


def takeaway_bar(slide, text, t=4.85, col=TEAL_TITLE, size=14):
    """Barra de remate al pie (línea teal + frase centrada)."""
    _rect(slide, 0.35, t, SW - 0.7, 0.02, TEAL_SQ)
    add_textbox(slide, 0.35, t + 0.08, SW - 0.7, 0.62,
                [(text, size, True, col, F_BODY, PP_ALIGN.CENTER)], anchor=MSO_ANCHOR.MIDDLE)


def dim_pipeline(slide, blocks, t, h=0.62, bw=1.06, gap=0.25, arrow_sz=15):
    """Pipeline al pie: fila de bloques (variable + dimensión) conectados por flechas →.
    `blocks` = lista de (variable, dimensión). Centrado horizontalmente.

    Gramática de Deep-LLM-V (arquetipo Google Shape;395): el bloque de proceso es un
    rounded-rect RELLENO de #3E6877 con el texto Barlow bold BLANCO encima. Hasta el
    19-jul esto estaba invertido — bloque claro con texto teal — que es el negativo del
    molde; se veía como otra plantilla aunque los colores fueran de la familia.
    """
    n = len(blocks)
    total = n * bw + (n - 1) * gap
    x = (SW - total) / 2
    for i, (var, dim) in enumerate(blocks):
        sp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(t), Inches(bw), Inches(h))
        sp.fill.solid(); sp.fill.fore_color.rgb = ONCO_DARK
        sp.line.fill.background(); sp.shadow.inherit = False
        lines = [(var, 12, True, WHITE, F_TITLE, PP_ALIGN.CENTER)]
        if dim:
            # La dimensión es el "dato" del bloque: mismo blanco, un punto más chica y sin
            # negrita, para que la variable siga siendo lo que se lee primero.
            lines.append((dim, 9.5, False, WHITE, F_MONO, PP_ALIGN.CENTER))
        _set_runs(sp.text_frame, lines, anchor=MSO_ANCHOR.MIDDLE)
        for p in sp.text_frame.paragraphs:
            p.space_after = Pt(0)
        if i < n - 1:
            ar = slide.shapes.add_textbox(Inches(x + bw), Inches(t), Inches(gap), Inches(h))
            _set_runs(ar.text_frame, [("→", arrow_sz, True, TEAL_SQ, F_BODY, PP_ALIGN.CENTER)],
                      anchor=MSO_ANCHOR.MIDDLE)
        x += bw + gap


# ============================================================================
# Gramática de diagrama de Deep-LLM-V
# ============================================================================
# Arquetipos medidos sobre el template (nombres de shape entre paréntesis):
#   proceso   (Google Shape;395)  rounded-rect #3E6877 · Barlow 12 bold BLANCO · centrado
#   dato      (Google Shape;391)  rect #B7B7B7 · Barlow 12 regular NEGRO
#   panel     (Google Shape;323)  rounded-rect #CDDFE1, borde del mismo color, agrupa
#   operador  (Google Shape;374)  óvalo #CDDFE1, borde #0E2841, símbolo (+, ×)
#   conector  (Google Shape;399)  línea #386271 de 2.37pt, con o sin punta

def _proc(slide, l, t, w, h, text, dim=None, size=11):
    """Bloque de proceso: el arquetipo 395 del template."""
    sp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(l), Inches(t),
                                Inches(w), Inches(h))
    sp.fill.solid(); sp.fill.fore_color.rgb = ONCO_DARK
    sp.line.fill.background(); sp.shadow.inherit = False
    lines = [(text, size, True, WHITE, F_BODY, PP_ALIGN.CENTER)]
    if dim:
        lines.append((dim, size - 2, False, WHITE, F_BODY, PP_ALIGN.CENTER))
    _set_runs(sp.text_frame, lines, anchor=MSO_ANCHOR.MIDDLE)
    for p in sp.text_frame.paragraphs:
        p.space_after = Pt(0)
    return sp


def _dato(slide, l, t, w, h, text, size=9.5):
    """Bloque de dato (forma del tensor): el arquetipo 391 del template."""
    sp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(l), Inches(t), Inches(w), Inches(h))
    sp.fill.solid(); sp.fill.fore_color.rgb = ONCO_DATA
    sp.line.fill.background(); sp.shadow.inherit = False
    _set_runs(sp.text_frame, [(text, size, False, BLACK, F_BODY, PP_ALIGN.CENTER)],
              anchor=MSO_ANCHOR.MIDDLE)
    return sp


def _grupo(slide, l, t, w, h):
    """Panel contenedor que agrupa una región: el arquetipo 323 del template."""
    sp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(l), Inches(t),
                                Inches(w), Inches(h))
    sp.fill.solid(); sp.fill.fore_color.rgb = ONCO_PANEL
    sp.line.color.rgb = ONCO_PANEL; sp.line.width = Pt(0.6)
    sp.shadow.inherit = False
    return sp


def _conn(slide, x0, y0, x1, y1, arrow=True):
    """Conector recto #386271 2.37pt, con punta opcional (arquetipo 399)."""
    ln = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(x0), Inches(y0),
                                    Inches(x1), Inches(y1))
    ln.line.color.rgb = ONCO_CONN; ln.line.width = Pt(2.37)
    ln.shadow.inherit = False
    if arrow:
        # python-pptx no expone las puntas de flecha; se inyecta el tailEnd en el XML.
        lnpr = ln.line._get_or_add_ln()
        end = lnpr.makeelement(qn('a:tailEnd'), {'type': 'triangle', 'w': 'med', 'len': 'med'})
        lnpr.append(end)
    return ln


def _proc_claro(slide, l, t, w, h, text, size=11, dim=None):
    """Bloque de proceso en el TONO CLARO del template (#CDDFE1 con texto teal bold).

    El template no usa un solo tono para sus bloques: sus láminas de arquitectura (s11-s16)
    alternan el relleno oscuro #3E6877 con este claro, y ese contraste es lo que da
    jerarquía al diagrama — lo oscuro es el camino principal, lo claro el detalle interno.
    Hasta ahora el deck solo tenía el oscuro (_proc) y los diagramas salían todos planos.
    """
    sp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(l), Inches(t),
                                Inches(w), Inches(h))
    sp.fill.solid(); sp.fill.fore_color.rgb = ONCO_PANEL
    sp.line.fill.background(); sp.shadow.inherit = False
    lines = [(text, size, True, ONCO_DARK, F_BODY, PP_ALIGN.CENTER)]
    if dim:
        lines.append((dim, size - 2, False, ONCO_DARK, F_BODY, PP_ALIGN.CENTER))
    _set_runs(sp.text_frame, lines, anchor=MSO_ANCHOR.MIDDLE)
    for p in sp.text_frame.paragraphs:
        p.space_after = Pt(0)
    return sp


def _dim(slide, l, t, w, text, size=9.5, align=PP_ALIGN.CENTER, col=None):
    """Etiqueta de forma del tensor SUELTA, al lado del bloque: «[N × 512]».

    Es como rotula el template (s13-s16): la dimensión va como texto libre pegado al
    bloque, no dentro de una caja gris. El bloque gris (_dato) sigue existiendo para
    cuando la forma ES el objeto del diagrama, pero para anotar un flujo esto pesa menos
    y deja los bloques más grandes.
    """
    add_textbox(slide, l, t, w, 0.26,
                [(text, size, False, col or ONCO_INK, F_BODY, align)],
                anchor=MSO_ANCHOR.MIDDLE)


def _oper(slide, cx, cy, sym="+", d=0.34):
    """Operador: óvalo #CDDFE1 con borde #0E2841 (arquetipo 374), centrado en (cx, cy)."""
    sp = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(cx - d / 2), Inches(cy - d / 2),
                                Inches(d), Inches(d))
    sp.fill.solid(); sp.fill.fore_color.rgb = ONCO_PANEL
    sp.line.color.rgb = ONCO_INK; sp.line.width = Pt(1.0); sp.shadow.inherit = False
    _set_runs(sp.text_frame, [(sym, 11, True, ONCO_INK, F_BODY, PP_ALIGN.CENTER)],
              anchor=MSO_ANCHOR.MIDDLE)
    return sp


def _conn_dash(slide, x0, y0, x1, y1):
    """Línea de expansión punteada: une un bloque con su detalle ampliado.

    Es el recurso con el que el template abre el «Transformer Block» en su lámina s12 —
    el bloque chico del flujo y su interior ampliado quedan atados por dos punteadas, así
    el zoom se lee como zoom y no como otro paso del pipeline."""
    ln = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(x0), Inches(y0),
                                    Inches(x1), Inches(y1))
    ln.line.color.rgb = ONCO_CONN; ln.line.width = Pt(1.0)
    ln.shadow.inherit = False
    lnpr = ln.line._get_or_add_ln()
    lnpr.append(lnpr.makeelement(qn('a:prstDash'), {'val': 'dash'}))
    return ln


def _rot_label(slide, l, t, w, h, text, size=9, col=None):
    """Rótulo vertical al costado de un panel (el «Transformer Block» de s13-s15)."""
    tb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    _set_runs(tb.text_frame, [(text, size, True, col or ONCO_DARK, F_BODY, PP_ALIGN.CENTER)],
              anchor=MSO_ANCHOR.MIDDLE)
    tb.rotation = 270
    return tb


def ratio_bar(slide, l, t, w, h, frac, label, valor, de, nota, col=None):
    """Barra de proporción: cuánto de un presupuesto se usa de verdad.

    Para Q1 el mensaje es una FRACCIÓN (30 de 30 contra 159 de 300), y una fracción se lee
    de un vistazo como barra y no como celda de tabla. El riel va en el claro del template
    y lo ocupado en el oscuro; el número grande queda encima para que sobreviva a la
    proyección."""
    col = col or ONCO_DARK
    add_textbox(slide, l, t, w, 0.30, [(label, 13, True, TEAL_TITLE, F_TITLE)])
    # El valor y el «de N» van en UN párrafo con dos runs: se leen como una sola cifra,
    # con el total en cuerpo chico apoyado en la línea de base del número grande.
    tb = slide.shapes.add_textbox(Inches(l), Inches(t + 0.32), Inches(w), Inches(0.52))
    tf = tb.text_frame; tf.word_wrap = True; tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]; p.space_after = Pt(0)
    for txt, sz, bold, cl in ((valor, 27, True, col), (f"  de {de}", 13, False, GRIS_BODY)):
        r = p.add_run(); r.text = txt
        r.font.size = Pt(sz); r.font.bold = bold; r.font.name = F_TITLE; r.font.color.rgb = cl
    riel_t = t + 0.92
    _rect(slide, l, riel_t, w, h, ONCO_PANEL)
    if frac > 0:
        _rect(slide, l, riel_t, w * frac, h, col)
    add_textbox(slide, l, riel_t + h + 0.06, w, 0.30, [(nota, 10.5, False, GRIS_BODY, F_BODY)])


def scale_axis(slide, l, t, w, marcas, h=0.10):
    """Eje de escala física (µm) con tramos marcados encima y debajo.

    `marcas` = lista de (x0_frac, x1_frac, texto, arriba, color). Sirve para mostrar de un
    vistazo que un objeto entra en el parche y el otro no: es el argumento entero de la
    lámina, y como dibujo no necesita bullets."""
    _rect(slide, l, t, w, h, ONCO_PANEL)
    for x0f, x1f, txt, arriba, col in marcas:
        x0, x1 = l + w * x0f, l + w * x1f
        _rect(slide, x0, t, max(x1 - x0, 0.03), h, col)
        ty = t - 0.42 if arriba else t + h + 0.06
        # Padding generoso a los lados: la caja del rótulo tiene que poder ser MÁS ancha
        # que su tramo, si no un tramo corto parte el texto en dos líneas. El rótulo va
        # centrado, así que dos cajas vecinas pueden solaparse sin que el texto se toque.
        add_textbox(slide, x0 - 0.80, ty, (x1 - x0) + 1.60, 0.38,
                    [(txt, 10, True, col, F_BODY, PP_ALIGN.CENTER)],
                    anchor=MSO_ANCHOR.BOTTOM if arriba else MSO_ANCHOR.TOP)


def code_panel(slide, l, t, w, h, lines, title=None):
    """Panel de código nativo: fondo oscuro + texto monoespaciado (Consolas).
    `lines` = lista de strings; los que empiezan con '#' se pintan como comentario."""
    _rect(slide, l, t, w, h, CODE_BG)
    tb = slide.shapes.add_textbox(Inches(l + 0.12), Inches(t + 0.08), Inches(w - 0.24), Inches(h - 0.16))
    tf = tb.text_frame; tf.word_wrap = True; tf.vertical_anchor = MSO_ANCHOR.TOP
    runs = []
    if title:
        runs.append((title, 10, True, CODE_CMT, F_MONO))
    for ln in lines:
        col = CODE_CMT if ln.strip().startswith("#") else CODE_FG
        runs.append((ln if ln else " ", 9.5, False, col, F_MONO))
    _set_runs(tf, runs)
    for p in tf.paragraphs:
        p.space_after = Pt(1)
    return tb


def pipeline_mammoth(prs):
    """Lámina 'Dónde entra MAMMOTH en el pipeline', NATIVA con la gramática Deep-LLM-V.

    Sustituye al diagrama que se traía de B4 con copy_diagram_scaled(). Aquel venía en
    Carlito, con 129 runs bajo 10pt (la mayoría a 6) y su propia paleta: a escala de
    proyección no se leía y se veía como otra plantilla.

    Se dejaron FUERA los paneles de fórmulas que traía el original. No es una pérdida:
    la lámina siguiente ('Por dentro de MAMMOTH: dimensiones y código') ya lleva esa
    matemática en una tabla nativa y legible. El trabajo de ESTA lámina es ubicar el
    bloque en el pipeline, y para eso el flujo horizontal + las formas del tensor
    alcanzan. El flujo va horizontal porque es lo que hace el propio template en su
    lámina de flujo general (s02), y porque en 16:9 deja los bloques al doble de tamaño
    que la versión vertical.

    Rediseño 19-jul sobre el molde de las láminas de arquitectura del template (s11-s16),
    que es lo que pidió Ernesto. Cuatro cosas se copian de ahí y no estaban antes:
      · el diagrama OCUPA la lámina — sin subtítulo de prosa arriba ni barra de remate
        abajo; el template no las usa en sus láminas de arquitectura;
      · DOS tonos de bloque (oscuro el camino principal, claro el detalle interno), que
        es lo que da jerarquía; antes todo era del mismo oscuro y salía plano;
      · la forma del tensor como etiqueta suelta «[N × 512]» pegada al bloque, no dentro
        de una caja gris, con lo que los bloques pueden ser más grandes;
      · la EXPANSIÓN punteada: el bloque del flujo se abre abajo en un panel con su
        interior, igual que el template abre su «Transformer Block» en s12.
    El interior expandido es el mismo que la lámina siguiente detalla con números, así que
    acá va sin dimensiones: la función de esta lámina es ubicar, no cuantificar.
    """
    s = content(prs, "Dónde entra MAMMOTH en el pipeline")

    BW, GAP, BY, BH = 1.60, 0.32, 1.02, 0.82
    xs = [0.36 + i * (BW + GAP) for i in range(5)]

    etapas = [
        ("Lámina en parches", None),
        ("CONCH", "extractor de features"),
        ("MAMMOTH", "mezcla de expertos"),
        ("Atención CLAM", "backbone + cabeza"),
        ("Clasificador", "logits de la lámina"),
    ]
    for x, (txt, dim) in zip(xs, etapas):
        _proc(s, x, BY, BW, BH, txt, dim=dim, size=12)
    for i in range(4):
        _conn(s, xs[i] + BW, BY + BH / 2, xs[i + 1], BY + BH / 2)

    # Entra [N × 512] y sale [N × 512]: la evidencia visual de que es drop-in. Van como
    # etiqueta suelta bajo el bloque del que sale el tensor (molde del template).
    for x in (xs[1], xs[2], xs[3]):
        _dim(s, x, BY + BH + 0.04, BW, "[N × 512]", size=10)

    # --- expansión punteada del bloque MAMMOTH ---
    PT_, PH = 2.44, 1.92
    PL, PW = 0.36, 9.28
    _grupo(s, PL, PT_, PW, PH)
    _conn_dash(s, xs[2], BY + BH + 0.30, PL + 0.02, PT_)
    _conn_dash(s, xs[2] + BW, BY + BH + 0.30, PL + PW - 0.02, PT_)
    # La caja se declara ANCHA y baja: al rotar 270° el ancho pasa a ser el alto visible,
    # así que con 0.60 la palabra se partía en dos. Se recentra a mano sobre el borde.
    _rot_label(s, PL - 0.62, PT_ + PH / 2 - 0.16, 1.40, 0.32, "MAMMOTH")

    IBW, IGAP, IY, IH = 1.78, 0.36, 2.86, 0.78
    ixs = [0.90 + i * (IBW + IGAP) for i in range(4)]
    interior = [
        ("16 cabezas", "subespacios del parche"),
        ("ruteo", "similitud + softmax"),
        ("300 slots", "30 expertos × 10"),
        ("concat", "cabezas de vuelta"),
    ]
    for x, (txt, sub) in zip(ixs, interior):
        _proc_claro(s, x, IY, IBW, IH, txt, size=12, dim=sub)
    for i in range(3):
        _conn(s, ixs[i] + IBW, IY + IH / 2, ixs[i + 1], IY + IH / 2)
    add_textbox(s, PL + 0.30, IY + IH + 0.10, PW - 0.60, 0.34, [
        ("Cada slot se llena con el promedio ponderado de los parches que se le parecen.",
         10.5, False, ONCO_DARK, F_BODY, PP_ALIGN.CENTER)])

    add_textbox(s, 0.36, PT_ + PH + 0.14, 9.28, 0.40, [
        ("MAMMOTH reemplaza solo la primera capa lineal: entra [N × 512] y sale [N × 512]. "
         "Todo lo que viene después es CLAM sin cambios.",
         12.5, True, TEAL_TITLE, F_BODY, PP_ALIGN.CENTER)], anchor=MSO_ANCHOR.MIDDLE)
    return s


def dims_table(slide, l, t, w, rows, row_h=0.345, hdr=("paso", "forma del tensor")):
    """Tabla nativa de dimensiones (2 columnas: descripción del paso · forma)."""
    n = len(rows) + 1
    gtbl = slide.shapes.add_table(n, 2, Inches(l), Inches(t), Inches(w), Inches(row_h * n)).table
    gtbl.columns[0].width = Inches(w * 0.60)
    gtbl.columns[1].width = Inches(w * 0.40)
    gtbl.first_row = False; gtbl.horz_banding = False
    data = [hdr] + list(rows)
    for ri, (c0, c1) in enumerate(data):
        for ci, txt in enumerate((c0, c1)):
            cell = gtbl.cell(ri, ci)
            cell.margin_left = Inches(0.06); cell.margin_right = Inches(0.04)
            cell.margin_top = Inches(0.01); cell.margin_bottom = Inches(0.01)
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            if ri == 0:
                cell.fill.solid(); cell.fill.fore_color.rgb = TEAL_SQ
                col, bold, sz, fnt = WHITE, True, 9.5, F_BODY
            else:
                cell.fill.solid(); cell.fill.fore_color.rgb = TEAL_CARD2 if ri % 2 else TEAL_CARD
                col = INK if ci == 0 else TEAL_TITLE
                bold = ci == 1; sz = 9.5; fnt = F_BODY if ci == 0 else F_MONO
            tf = cell.text_frame; tf.word_wrap = True
            p = tf.paragraphs[0]; p.alignment = PP_ALIGN.LEFT
            r = p.add_run(); r.text = txt
            r.font.size = Pt(sz); r.font.bold = bold; r.font.name = fnt; r.font.color.rgb = col
    return gtbl


def simple_table(slide, l, t, w, headers, rows, col_fracs, row_h=0.32, fs=9.5):
    """Tabla nativa genérica (n columnas): header teal + filas con banding teal claro."""
    ncol = len(headers); nrow = len(rows) + 1
    tbl = slide.shapes.add_table(nrow, ncol, Inches(l), Inches(t), Inches(w),
                                 Inches(row_h * nrow)).table
    for ci, fr in enumerate(col_fracs):
        tbl.columns[ci].width = Inches(w * fr)
    tbl.first_row = False; tbl.horz_banding = False
    for ri, row in enumerate([headers] + list(rows)):
        for ci, txt in enumerate(row):
            cell = tbl.cell(ri, ci)
            cell.margin_left = Inches(0.06); cell.margin_right = Inches(0.04)
            cell.margin_top = Inches(0.015); cell.margin_bottom = Inches(0.015)
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            if ri == 0:
                cell.fill.solid(); cell.fill.fore_color.rgb = TEAL_SQ
                col, bold, fnt = WHITE, True, F_BODY
            else:
                cell.fill.solid(); cell.fill.fore_color.rgb = TEAL_CARD2 if ri % 2 else TEAL_CARD
                col, bold, fnt = INK, (ci == 0), F_BODY
            tf = cell.text_frame; tf.word_wrap = True
            p = tf.paragraphs[0]; p.alignment = PP_ALIGN.LEFT
            r = p.add_run(); r.text = txt
            r.font.size = Pt(fs); r.font.bold = bold; r.font.name = fnt; r.font.color.rgb = col
    return tbl


def panel(slide, l, t, w, h, title, tcol, lines, border, fill=TEAL_CARD2):
    """Panel con título + líneas de cuerpo (para 'dos tareas / demandas opuestas')."""
    sp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(l), Inches(t), Inches(w), Inches(h))
    sp.fill.solid(); sp.fill.fore_color.rgb = fill
    sp.line.color.rgb = border; sp.line.width = Pt(1.5); sp.shadow.inherit = False
    tb = slide.shapes.add_textbox(Inches(l + 0.18), Inches(t + 0.12), Inches(w - 0.36), Inches(h - 0.24))
    runs = [(title, 14.5, True, tcol, F_TITLE)] + [(ln, 12, False, INK, F_BODY) for ln in lines]
    _set_runs(tb.text_frame, runs)
    for p in tb.text_frame.paragraphs:
        p.space_after = Pt(3)


def nested_fields(slide, l, t, w, h):
    """Esquema nativo: dos campos concéntricos (contexto 512µm ⊃ fino 112µm) + leyenda."""
    _rect(slide, l, t, w, h, TEAL_CARD2, line=TEAL_SQ)
    side = h - 0.55
    cx = l + w * 0.30; cy = t + h / 2
    ctx = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(cx - side / 2), Inches(cy - side / 2),
                                 Inches(side), Inches(side))
    ctx.fill.background(); ctx.line.color.rgb = TEAL_SQ; ctx.line.width = Pt(2.5)
    ctx.shadow.inherit = False
    fside = side * 0.34                         # esquema (ratio real 112/512≈0.22, lo muestra la foto)
    fin = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(cx - fside / 2), Inches(cy - fside / 2),
                                 Inches(fside), Inches(fside))
    fin.fill.solid(); fin.fill.fore_color.rgb = ORA_ACC
    fin.line.color.rgb = ORA_T; fin.line.width = Pt(1.5); fin.shadow.inherit = False
    lx = l + w * 0.58
    add_textbox(slide, lx, t + 0.22, w * 0.42 - 0.12, h - 0.4, [
        ("Contexto  512 µm · ~5×", 12.5, True, TEAL_TITLE, F_TITLE),
        ("el conducto/lobulillo anfitrión", 10.5, False, GRIS_BODY, F_BODY),
        (" ", 8, False, INK, F_BODY),
        ("Fino  112 µm · ~20×", 12.5, True, ORA_T, F_TITLE),
        ("la calcificación + su forma", 10.5, False, GRIS_BODY, F_BODY),
    ], anchor=MSO_ANCHOR.MIDDLE)


# ============================================================================
# Arquetipos B4
# ============================================================================
def _set_solo_run(par, texto):
    """Deja el párrafo con UN solo run que lleva `texto`, conservando el formato del
    primero. El texto del template viene partido en varios runs ("OncoMets" + " - Spatial"),
    así que escribir solo runs[0] deja la cola del original pegada detrás."""
    runs = par.runs
    runs[0].text = texto
    for r in runs[1:]:
        r._r.getparent().remove(r._r)


def retitular_portada(prs):
    """Ajusta las 2 láminas de apertura HEREDADAS del template (s00 portada de marca,
    s01 título + fecha) sin redibujarlas: son vectoriales nativas y ya están branded.

    Sustituye a la portada JPG anterior. Se reescribe solo el texto, respetando el
    formato del run original (fuente, cuerpo, color) para no romper el diseño.
    Devuelve la lámina de título, que es la que lleva las notas de apertura."""
    # --- s00: dos defectos que vienen DEL TEMPLATE, no de este deck (se reproducen
    # abriendo Deep-LLM-V solo). El claim "Care in <code>" es un marcador sin reemplazar,
    # y el párrafo descriptivo arranca tan abajo (y=6.06 de 7.5) que se sale por el borde
    # inferior. Se corrigen acá porque es la primera lámina que se proyecta.
    portada = prs.slides[0]
    for sh in portada.shapes:
        if not sh.has_text_frame:
            continue
        txt = sh.text_frame.text.strip()
        if txt.startswith("Care in people"):
            # Queda solo el claim válido. La caja es autofit, así que al perder la 2ª
            # línea se encoge sola y libera el aire que necesita el párrafo de abajo.
            for par in list(sh.text_frame.paragraphs)[1:]:
                par._p.getparent().remove(par._p)
        elif txt.startswith("OncoMETS is an AI-powered platform"):
            # Se sube y se ensancha: con Barlow real entraba justo, pero cualquier
            # sustitución de fuente lo empuja fuera de la lámina. Con este alto libre
            # (hasta y=7.5) el texto tiene margen aunque PowerPoint lo redistribuya.
            sh.top = Inches(5.62)
            sh.left = Inches(0.28)
            sh.width = Inches(6.30)
    titulo = prs.slides[1]
    for sh in titulo.shapes:
        if not sh.has_text_frame:
            continue
        txt = sh.text_frame.text.strip()
        if txt == "OncoMets - Spatial":
            _set_solo_run(sh.text_frame.paragraphs[0], "OncoMets · MAMMOTH")
        elif txt == "14/11/2025":
            _set_solo_run(sh.text_frame.paragraphs[0], FECHA_REUNION)
    return titulo


def divider(prs, title, subtitle):
    s = new_slide(prs)
    s.background.fill.solid(); s.background.fill.fore_color.rgb = TEAL_DIV
    s.shapes.add_picture(LOGO, Inches(0.42), Inches(0.36), height=Inches(0.62))
    # La caja del título tiene que aguantar DOS líneas a 44 pt (~1.22"): con 1.1" el
    # título de dos líneas desbordaba y pisaba el subtítulo (pasaba en 3 de 4 portadillas).
    add_textbox(s, 0.8, 1.85, SW - 1.6, 1.45,
                [(title, 44, True, LAV_TITLE, F_TITLE, PP_ALIGN.CENTER)], anchor=MSO_ANCHOR.MIDDLE)
    add_textbox(s, 0.8, 3.42, SW - 1.6, 0.7,
                [(subtitle, 18, False, TEAL_SUB, F_BODY, PP_ALIGN.CENTER)], anchor=MSO_ANCHOR.TOP)
    return s


def content(prs, title, bar=True, size=26, style="onco"):
    """Lámina de contenido. style='onco' = cabecera técnica de Plantilla (default,
    es lo que son casi todas las nuestras); style='environ' = cabecera administrativa
    (barra gris + cuadro teal), reservada a agenda/recapitulación como en Plantilla."""
    s = new_slide(prs)
    if style == "environ":
        header(s, title, bar=bar, size=size)
    else:
        header_oncomets(s, title, size=size)
    return s


def _scale_el(el, s):
    """Escala geométrica (off/ext) y tipográfica (sz) de un shape XML por factor s."""
    for off in el.iter(qn("a:off")):
        off.set("x", str(int(int(off.get("x")) * s)))
        off.set("y", str(int(int(off.get("y")) * s)))
    for ext in el.iter(qn("a:ext")):
        ext.set("cx", str(int(int(ext.get("cx")) * s)))
        ext.set("cy", str(int(int(ext.get("cy")) * s)))
    for tag in ("a:rPr", "a:defRPr", "a:endParaRPr"):
        for rpr in el.iter(qn(tag)):
            sz = rpr.get("sz")
            if sz:
                rpr.set("sz", str(max(100, int(int(sz) * s))))


def _shift_el(el, dx, dy):
    """Desplaza un shape top-level. Solo el a:off del propio spPr/grpSpPr (el primero),
    porque los a:off internos de un grupo son coordenadas hijas y no deben moverse."""
    for off in el.iter(qn("a:off")):
        off.set("x", str(int(int(off.get("x")) + dx * 914400)))
        off.set("y", str(int(int(off.get("y")) + dy * 914400)))
        break


def copy_diagram_scaled(prs, src_path, idx, title=None, scale=0.75, bar=False,
                        style="environ", dx=0.0, dy=0.0):
    """Copia el diagrama standalone (13.333x7.5) escalado a 10x5.625 (x0.75, full-bleed).
    Conserva nativo/editable + remapea imágenes embebidas.

    style='onco' + (scale, dx, dy) achica y baja el diagrama para dejarle sitio a la
    cabecera técnica de Plantilla: a x0.75 el diagrama ocupa y=0.08..5.59, es decir la
    lámina entera, y una cabecera encima lo pisaría."""
    s = new_slide(prs)
    src = Presentation(src_path)
    src_slide = src.slides[idx]
    spTree = s.shapes._spTree
    for shp in src_slide.shapes:
        el = copy.deepcopy(shp._element)
        for cNvPr in el.iter(qn("p:cNvPr")):
            _uid[0] += 1; cNvPr.set("id", str(_uid[0]))
        for blip in el.iter(qn("a:blip")):
            r_embed = blip.get(qn("r:embed"))
            if r_embed:
                img_part = src_slide.part.related_part(r_embed)
                _, new_rid = s.part.get_or_add_image_part(io.BytesIO(img_part.blob))
                blip.set(qn("r:embed"), new_rid)
        _scale_el(el, scale)
        if dx or dy:
            _shift_el(el, dx, dy)
        spTree.append(el)
    if style == "onco":
        header_oncomets(s, title)
    else:
        header(s, title, bar=bar)  # logo mark (+ barra/título si se pide)
    return s


# ============================================================================
# Reflow: hacerle sitio a la cabecera LITERAL del template
# ============================================================================
# Antes la cabecera técnica iba compactada a la banda de 0.785 y el contenido arrancaba en
# y=0.86. La cabecera real del template termina en 1.171, o sea 0.31" más abajo. En vez de
# re-maquetar 17 láminas a mano, se baja el contenido en bloque y, si la lámina no tiene
# ese aire al pie, se la comprime verticalmente lo justo.
CONTENT_TOP_OLD = 0.86              # donde arrancaba el contenido de las técnicas
CONTENT_TOP_NEW = ONCO_BAND + 0.05  # 1.221 · un respiro bajo la línea teal
SAFE_BOTTOM = SH - 0.14             # margen inferior que no se debe invadir


def _scale_block(el, f):
    """Escala un shape COMPLETO por f: geometría, cuerpo tipográfico y métrica de tablas.

    Escalar solo la geometría no sirve para las tablas: en PowerPoint el alto real de una
    tabla lo manda el texto de sus filas, no el alto del shape, así que la tabla no
    encogía y se comía lo que tuviera debajo. Bajar también el `sz` y el alto de fila la
    encoge de verdad."""
    for off in el.iter(qn("a:off")):
        off.set("x", str(int(round(int(off.get("x")) * f))))
        off.set("y", str(int(round(int(off.get("y")) * f))))
    for ext in el.iter(qn("a:ext")):
        ext.set("cx", str(int(round(int(ext.get("cx")) * f))))
        ext.set("cy", str(int(round(int(ext.get("cy")) * f))))
    for tag in ("a:rPr", "a:defRPr", "a:endParaRPr"):
        for rpr in el.iter(qn(tag)):
            if rpr.get("sz"):
                rpr.set("sz", str(max(100, int(round(int(rpr.get("sz")) * f)))))
    for gc in el.iter(qn("a:gridCol")):
        if gc.get("w"):
            gc.set("w", str(int(round(int(gc.get("w")) * f))))
    for tr in el.iter(qn("a:tr")):
        if tr.get("h"):
            tr.set("h", str(int(round(int(tr.get("h")) * f))))


def reflow_onco(prs, skip=()):
    """Baja el contenido de las láminas con cabecera OncoMets para que no lo pise la banda.

    Por lámina: se toman los shapes que NO son la cabecera y se los desplaza dy. Si el
    bloque no entra, se lo escala entero (geometría + tipografía) por el factor justo y se
    lo reancla bajo la banda. El ajuste típico ronda el 8%: un cuerpo de 10.5pt queda en
    9.7pt, invisible a ojo y sin romper proporciones.
    """
    dy = CONTENT_TOP_NEW - CONTENT_TOP_OLD
    for slide in prs.slides:
        if id(slide._element) in skip:
            continue
        cuerpo = [sh for sh in slide.shapes if not (sh.name or "").startswith("ONCOHDR_")]
        if len(cuerpo) == len(list(slide.shapes)):
            continue                                  # sin cabecera OncoMets: no se toca
        tops = [Emu(sh.top).inches for sh in cuerpo]
        bots = [Emu(sh.top).inches + Emu(sh.height).inches for sh in cuerpo]
        if not tops:
            continue
        top0, bot0 = min(tops), max(bots)
        f = 1.0
        if bot0 + dy > SAFE_BOTTOM and bot0 > top0:
            f = max(0.80, (SAFE_BOTTOM - CONTENT_TOP_NEW) / (bot0 - top0))
        if f < 1.0:
            for sh in cuerpo:
                _scale_block(sh._element, f)
        # reanclar: el tope del bloque (ya escalado) va justo debajo de la banda
        desplaz = CONTENT_TOP_NEW - top0 * f
        for sh in cuerpo:
            sh.top = Inches(Emu(sh.top).inches + desplaz)


# ============================================================================
# Re-base al template de Sebastián: escala uniforme 10x5.625 → 13.333x7.5
# ============================================================================
def scale_deck_to_1610(prs, k=13.333 / 10.0, skip=()):
    """Escala TODO el deck (geometría + fuentes + tablas + líneas) por k y cambia el tamaño
    de slide a 13.333x7.5 (16:9, = el de los templates de Sebastián). Misma relación de
    aspecto que 10x5.625 → escala uniforme sin deformar. Los diagramas reusados (copiados
    a x0.75) quedan a 0.75·1.3333 = 1.0 = su tamaño nativo 13.333. Todas las shapes del
    deck son planas (top-level), así que escalar cada a:off/a:ext es seguro.

    `skip` = ids de elementos de láminas heredadas del template, que YA están a 13.333 y
    se romperían si se las volviera a escalar."""
    for slide in prs.slides:
        if id(slide._element) in skip:
            continue
        tree = slide.shapes._spTree
        for off in tree.iter(qn("a:off")):
            off.set("x", str(int(round(int(off.get("x")) * k))))
            off.set("y", str(int(round(int(off.get("y")) * k))))
        for ext in tree.iter(qn("a:ext")):
            ext.set("cx", str(int(round(int(ext.get("cx")) * k))))
            ext.set("cy", str(int(round(int(ext.get("cy")) * k))))
        for tag in ("a:rPr", "a:defRPr", "a:endParaRPr"):
            for rpr in tree.iter(qn(tag)):
                sz = rpr.get("sz")
                if sz:
                    rpr.set("sz", str(max(100, int(round(int(sz) * k)))))
        for gc in tree.iter(qn("a:gridCol")):
            if gc.get("w"):
                gc.set("w", str(int(round(int(gc.get("w")) * k))))
        for tr in tree.iter(qn("a:tr")):
            if tr.get("h"):
                tr.set("h", str(int(round(int(tr.get("h")) * k))))
        for ln in tree.iter(qn("a:ln")):
            if ln.get("w"):
                ln.set("w", str(int(round(int(ln.get("w")) * k))))
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)


# ============================================================================
# Build
# ============================================================================
def _leer_q1():
    """Números de Q1 desde el JSON que produce answer_q1_expertos_slots.py.

    Se leen del artefacto en vez de hardcodearse para que el deck no se desincronice
    del análisis. Si el JSON todavía no existe, la slide se construye con el hueco
    marcado en vez de con un número inventado.
    """
    import json
    if not os.path.exists(Q1_JSON):
        return {"exp": "pendiente", "slots": "pendiente",
                "pie": "pendiente: correr answer_q1_expertos_slots.py"}
    filas = json.load(open(Q1_JSON))
    exps = [r["expertos_efectivos"] for r in filas if r.get("expertos_efectivos")]
    slots = [r["slots_efectivos"] for r in filas if r.get("slots_efectivos")]
    media = lambda v: sum(v) / len(v)
    # El rango de slots va en el pie: la media sola sobrevende precisión (varía con el
    # tamaño de la lámina, no con la tarea). Los expertos no lo necesitan: 30.0 en todas.
    return {"exp": f"{media(exps):.1f}", "slots": f"{media(slots):.0f}",
            "pie": f"media sobre {len(filas)} láminas · slots de {min(slots):.0f} a "
                   f"{max(slots):.0f} según tamaño de lámina · efectivo = exp(entropía)"}


def build():
    # Base = template válido, con sus fuentes embebidas y sus 2 láminas de apertura
    # (portada de marca + lámina de título) conservadas NATIVAS a 13.333.
    prs, keep_ids = base_from_template()
    # El resto del deck se compone en 10x5.625 y se escala al final; las heredadas quedan
    # fuera de ese escalado (skip=keep_ids).
    prs.slide_width = Inches(SW); prs.slide_height = Inches(SH)

    # ---- 1. Apertura: se hereda del template, solo se ajusta el título y la fecha ----
    s = retitular_portada(prs)
    notes(s, "El objetivo de esta reunión no es demostrar que MAMMOTH mejore la métrica: esa "
             "pregunta ya quedó cerrada, y la respuesta fue que no. Lo que sigue abierto, y es de "
             "lo que trata esta presentación, es distinto: qué aprende MAMMOTH por dentro y en qué "
             "se fija cada experto. El recorrido tiene dos partes. Primero el mecanismo, esta vez "
             "explicado con calma y hasta el fondo. Después, la evidencia visual de qué morfología "
             "concentra cada experto, mirada sobre las propias slides de mama del proyecto.")

    # ---- 2. Recapitulación de objetivos (MISMO formato que B4 slide 2:
    #        título 32pt + lista numerada en un solo cuadro, 24pt Barlow bold gris) ----
    # Cabecera OncoMets como el resto: la Environ (barra gris + cuadro teal) es de
    # Plantilla.pptx y el template válido no la tiene en ninguna lámina.
    s = content(prs, "Recapitulación de objetivos", size=32)
    # Objetivos del eje (pedido de Benjamín, 29-jun): entender el mecanismo + interpretar.
    # Enunciados en infinitivo (sin 1ª persona), concisos, sin resultados.
    # (texto, estado): "done" = ticket verde (cerrado) · "prog" = pill "En progreso" (abierto)
    recap = [
        ("1. Dominar el mecanismo de MAMMOTH.", "done"),
        ("2. Precisar qué es una cabeza y el tensor de prototipos (30×16×10×16).", "done"),
        ("3. Distinguir MoE de PoE y situar el número de cabezas para mama.", "done"),
        ("4. Interpretar los expertos: qué región y qué morfología concentra cada uno.", "prog"),
        ("5. Comparar dónde mira cada modelo, entrenados sobre las mismas particiones.", "done"),
        ("6. Medir cuántos expertos y cuántos slots se usan de verdad.", "done"),
    ]
    # 6 filas: se mantiene el molde B4 (lista numerada + marcador de estado) bajando
    # la tipografía de 24 a 19pt para que entren las dos nuevas sin apretar el interlineado.
    row_tops = [1.10 + i * 0.68 for i in range(6)]
    row_h = 0.62
    for (it, st), rt in zip(recap, row_tops):
        add_textbox(s, 0.35, rt, 7.75, row_h, [(it, 19, True, GRIS_BODY, F_BODY)],
                    anchor=MSO_ANCHOR.MIDDLE)
        cy = rt + row_h / 2
        if st == "done":
            status_done(s, 8.98, cy)
        else:
            status_progress(s, 8.98, cy)
    # Guion en 6 puntos numerados, uno por objetivo de la lámina: sigue siendo prosa
    # hablada (convención vigente), pero partida para poder engancharla de un vistazo
    # mientras se presenta. Los "\n" son párrafos reales en el panel de notas.
    notes(s, "Antes de entrar en el mecanismo, fijemos qué nos propusimos entender: es el eje de "
             "la reunión.\n"
             "\n"
             "1. El primero era dominar el mecanismo de MAMMOTH a fondo.\n"
             "\n"
             "2. El segundo, precisar dos puntos que habían quedado sin respuesta: qué es "
             "exactamente una cabeza y cómo se lee el tensor de prototipos, el treinta por "
             "dieciséis por diez por dieciséis.\n"
             "\n"
             "3. El tercero, distinguir una mezcla de expertos de un producto de expertos, y "
             "ubicar cuántas cabezas convienen para mama. Esos tres quedaron cerrados.\n"
             "\n"
             "4. El cuarto abre la segunda parte y sigue en progreso: interpretar a los expertos, "
             "es decir, mirar sobre las propias slides qué región concentra cada uno y qué "
             "morfología hay ahí; falta el visto bueno de un patólogo.\n"
             "\n"
             "5. Los dos últimos son el trabajo nuevo desde la vez pasada. El quinto es comparar "
             "dónde mira cada modelo, entrenando los dos sobre exactamente las mismas particiones "
             "para que la comparación sea limpia.\n"
             "\n"
             "6. Y el sexto responde una pregunta concreta que había quedado abierta: cuántos "
             "expertos y cuántos slots se usan de verdad.")

    # ---- 3. Divisoria: MAMMOTH ----
    s = divider(prs, "MAMMOTH", "Mixture-of-Experts en la primera capa de CLAM (patch-embed)")
    notes(s, "El mecanismo primero. MAMMOTH es una intervención muy acotada: cambia una sola "
             "parte de CLAM, la primera capa, la que proyecta los parches al espacio interno del "
             "modelo. Todo el resto de CLAM queda intacto.")

    # ---- 4. Qué es y por qué (título = acrónimo completo; 3 tarjetas + Fig 1 + Fig 3) ----
    s = content(prs, "MAMMOTH: MAtrix-factorized Mixture Module of Transformation Heads", size=17)
    cards = [
        "Reemplaza la 1ª capa lineal de CLAM por una mezcla de expertos (MoE).",
        "Un router reparte cada parche entre expertos especializados por morfología.",
        "Objetivo: bajar la interferencia de gradientes entre parches distintos.",
    ]
    cy = 1.05
    for i, txt in enumerate(cards):
        add_card(s, 0.30, cy, 5.35, 0.95, i + 1, txt, size=13)
        cy += 1.10
    add_image_fit(s, FIG1_MAM, 5.95, 0.98, 3.85, 1.85, align="top")
    caption(s, 5.95, 2.86, 3.85,
            "Fig. 1: el espacio interno pasa de una nube continua a grupos por experto,\n"
            "y mejora a todos los agregadores MIL (Shao et al., ICLR 2026)", size=8)
    add_image_fit(s, FIG3_MAM, 5.95, 3.42, 3.85, 1.42, align="top")
    # El pie decía "cada color es el ruteo de un experto": mal en dos puntos. Los paneles
    # de la Fig. 3 están rotulados por PAR experto+slot ("Expert 16 Slot 4") y la barra de
    # color es "Patch-slot similarity". Decir "por experto" contradice además la lámina 17
    # (los 30 expertos salen uniformes; el margen está en los slots).
    caption(s, 5.95, 4.86, 3.85,
            "Fig. 3: cada mapa pinta el parecido de los parches contra un prototipo\n"
            "(experto + slot); abajo, el tejido que concentra cada uno", size=8)
    # Lámina ya vista la sesión pasada: el guion pasa rápido por el acrónimo y el porqué,
    # y gasta el tiempo en las dos cosas que quedaron sin cerrar, el ruteo y la Fig. 3.
    notes(s, "Esta lámina ya la vimos, así que paso rápido por lo de arriba y me detengo en dos "
             "cosas: el ruteo y la figura de abajo.\n"
             "\n"
             "1. En CLAM, una sola matriz proyecta todos los parches, y tiene que servir a la vez "
             "para epitelio, estroma y ductos, así que queda en un punto intermedio. MAMMOTH la "
             "reemplaza por treinta expertos.\n"
             "\n"
             "2. El ruteo. Cada experto tiene diez prototipos aprendidos, que el paper llama "
             "slots: trescientos en total. Todo el ruteo consiste en medir cuánto se parece cada "
             "parche a cada uno de esos trescientos prototipos, y eso arma una tabla de parches "
             "contra prototipos. Esos parecidos se convierten después en porcentajes con una "
             "softmax, y ahí vamos a volver con las dimensiones puestas. Lo que importa acá es "
             "que ningún parche va a un solo experto: participan los trescientos, unos con un "
             "porcentaje grande y otros con uno casi nulo.\n"
             "\n"
             "3. El motivo de fondo es bajar la interferencia: en una sola matriz, los gradientes "
             "de parches muy distintos se pisan entre sí.\n"
             "\n"
             "4. La figura de arriba es la evidencia del paper: el espacio interno pasa de una "
             "nube continua a grupos separados, y la mejora aparece en todos los agregadores.\n"
             "\n"
             "5. La figura de abajo es la que quiero que se lea bien. Cada panel es una lámina de "
             "pulmón pintada con el parecido de cada parche contra un solo prototipo, no contra un "
             "experto entero: por eso los títulos dicen experto veintiuno slot cinco, o experto "
             "dieciséis slot cuatro. Y ahí está el detalle que más dice: el experto dieciséis "
             "tiene un slot que se enciende en alvéolos y otro que se enciende en estroma, así que "
             "la unidad que captura una morfología es el slot y no el experto. Abajo del todo, la "
             "fila de parches muestra qué morfología concentra cada uno. Ese es exactamente el análisis que la "
             "segunda parte reproduce, pero sobre mama y con las láminas del proyecto.")

    # ---- 5. Diagrama: pipeline CLAM + punto de integración (NATIVO, Deep-LLM-V) ----
    # Antes se traía de B4 con copy_diagram_scaled(DIAG_MAM, 0, scale=0.63): venía en
    # Carlito a 6pt y con la paleta de B4. Ver pipeline_mammoth() para el detalle.
    s = pipeline_mammoth(prs)
    notes(s, "Sobre el pipeline completo de CLAM el mecanismo se ubica mejor. La slide entra "
             "troceada en parches; el encoder CONCH le asigna a cada parche un vector de "
             "quinientos doce; de ahí siguen la atención y el clasificador. El único bloque que "
             "cambia es el del medio: ahí MAMMOTH reemplaza la primera capa lineal. Fíjense en "
             "las tres formas de abajo, que son la misma: entra un vector de quinientos doce por "
             "parche y sale un vector de quinientos doce por parche. Por eso decimos que es "
             "drop-in. Todo lo que viene después es CLAM tal cual.")

    # ---- 6. Interior de MAMMOTH: matrices+dims Y código en paralelo (NATIVO) ----
    s = content(prs, "Por dentro de MAMMOTH: dimensiones y código")
    add_textbox(s, 0.30, 0.86, 4.7, 0.3, [("El tensor, paso a paso", 13, True, TEAL_TITLE, F_TITLE)])
    dims_table(s, 0.30, 1.20, 4.72, [
        ("z · features CONCH (parche)", "[N, 512]"),
        ("q = LN(Wq · z)", "[N, 256]"),
        ("split en 16 cabezas", "[N, 16, 16]"),
        ("S · prototipos (slot_embeds)", "[30,16,10,16]"),
        ("logits de ruteo ⟨q, S⟩", "[N,30,16,10]"),
        ("D = softmax sobre N (dispatch)", "[N,30,16,10]"),
        ("u · slots (media ponderada)", "[30,16,10,16]"),
        ("expertos LoRA + concat cabezas", "[300, 512]"),
        ("combine → salida drop-in", "[N, 512]"),
    ], row_h=0.30)
    add_textbox(s, 0.30, 4.24, 4.72, 0.44, [
        ("El 16 aparece dos veces: nº de cabezas y dimensión de cada prototipo (256/16), para "
         "compararlos con un producto interno. 30 expertos × 10 slots = 300.",
         8.5, False, GRIS_BODY, F_BODY)])
    add_textbox(s, 5.20, 0.86, 4.6, 0.3, [("El código real (mammoth.py)", 13, True, TEAL_TITLE, F_TITLE)])
    code_panel(s, 5.20, 1.20, 4.60, 3.00, [
        "q = norm(wq(z))              # [N,512]->[N,256]",
        'q = rearrange(q,             # -> [N,16,16]',
        '      "n (h d)->n h d", h=16)',
        'logits = einsum("n h d, e h s d',
        '                 -> n e h s", q, S)',
        "                             # S=[30,16,10,16]",
        "dispatch = softmax(logits, dim=N)",
        'u = einsum("n h d, n e h s',
        '            -> e h s d", q, dispatch)',
        "out = expert_heads(u)        # LoRA -> [300,512]",
        "# drop-in (keep_slots=False):",
        'out = einsum("h p d, n h p -> n h d",',
        "             out, combine)   # -> [N,512]",
    ])
    takeaway_bar(s, "Las cabezas son subespacios aprendidos (multi-head), no textura/forma/color: "
                    "la semántica de tejido vive en los slots.", t=4.80, size=12)
    notes(s, "El interior de MAMMOTH, con las dimensiones a la izquierda y el código real a la "
             "derecha, para seguir cada paso sin ambigüedad. La entrada es un parche, un vector de "
             "quinientos doce, la z. Primero se proyecta a un query de doscientos cincuenta y seis "
             "y se parte en dieciséis cabezas de dieciséis. En paralelo, el modelo guarda sus "
             "prototipos aprendidos en el tensor treinta por dieciséis por diez por dieciséis. "
             "Se lee de izquierda a "
             "derecha: treinta expertos, cada experto con dieciséis cabezas, cada cabeza con diez "
             "slots, y cada slot es un vector de dieciséis dimensiones. El dieciséis aparece dos "
             "veces por una razón concreta. El primero es el número de cabezas; el último es la "
             "dimensión de cada prototipo, que sale de dividir el query, doscientos cincuenta y "
             "seis, entre las dieciséis cabezas. Tienen que coincidir porque el ruteo compara cada "
             "prototipo con el query mediante un producto interno, y un producto interno exige que "
             "ambos vectores vivan en el mismo espacio de dieciséis. Treinta expertos por diez "
             "slots dan los trescientos slots que reaparecen más adelante. De ahí en adelante la "
             "tabla y el código van paso a paso: el ruteo compara query contra prototipos, una "
             "softmax reparte cada parche entre los slots, con esos pesos se arma el promedio "
             "ponderado que llena cada slot, cada experto lo transforma con una operación de bajo "
             "rango y se concatenan las cabezas para dar trescientos por quinientos doce. En la "
             "variante base, una segunda softmax recombina los trescientos slots y reconstruye los "
             "parches, así que la salida tiene la misma forma que tendría una capa lineal; por eso "
             "el reemplazo es directo. Y un punto que en su momento se respondió mal y hay "
             "que dejar afinado: las cabezas no son textura, forma ni color. Son subespacios "
             "aprendidos, igual que en atención multi-cabeza; la semántica de tejido no vive en "
             "las cabezas, vive en los slots.")

    # ---- 7. La arquitectura oficial del paper: figura GRANDE y SIN callouts encima que
    #        tapen las variables. Lleva cabecera OncoMets como el resto de las técnicas
    #        (antes iba sin marca alguna y se leía como lámina rota); la figura se achica
    #        de 9.42 a 8.96 de ancho para dejar la banda libre sin perder legibilidad.
    #        El trazo al pie y las notas se refieren a las VARIABLES DE LA PROPIA FIGURA. ----
    # SIN título (pedido de Ernesto, 19-jul): la figura es el contenido entero de la
    # lámina y el título le robaba alto. Se conservan logo y línea, así que la marca
    # sigue puesta — es lo que hace el propio template en s15/s16, donde el título es un
    # «LLM» mínimo o directamente no aporta. También se retiró la tira de bloques con las
    # dimensiones: era la única otra cosa que competía por el alto, y su contenido ya está
    # completo y legible en la lámina anterior (tabla «El tensor, paso a paso»).
    # Resultado: la figura pasa de 8.50 a 9.70 de ancho (+30% de área).
    s = content(prs, "")
    iw = 9.70; ih = iw / 2.489                      # aspect real de la Fig 2 (4375x1758)
    s.shapes.add_picture(FIG2_ARCH, Inches((SW - iw) / 2), Inches(0.86), Inches(iw), Inches(ih))
    add_textbox(s, 0.15, 0.86 + ih + 0.02, SW - 0.30, 0.30, [
        # Los índices van con las MISMAS letras de la figura (j, k). El pie decía «e» y «s»,
        # letras que no aparecen en el diagrama, así que no se podían seguir.
        ("La arquitectura de MAMMOTH, paso a paso  ·  N = parches de la lámina  ·  "
         "en la salida z, j = slot (S=10) y k = experto (E=30)",
         9.5, True, TEAL_TITLE, F_BODY, PP_ALIGN.CENTER)], anchor=MSO_ANCHOR.MIDDLE)
    notes(s, "Esta es la arquitectura completa del paper, y la vamos a recorrer de izquierda a "
             "derecha siguiendo sus variables, porque cada símbolo es un paso del pipeline. Se "
             "parte de la lámina, que se corta en parches, y cada parche pasa por el encoder y "
             "sale convertido en un vector de features, que en la figura es x_i. Hay uno por "
             "parche, así que para toda la slide tenemos el conjunto de x_i, de uno hasta N. En un "
             "MIL corriente cada x_i entraría por una única matriz lineal, la W del diagrama, que "
             "lo proyecta al espacio interno del agregador, y ese es el único lugar donde MAMMOTH "
             "interviene: reemplaza esa W. Todo lo de la izquierda, el encoder, y todo lo de la "
             "derecha, el agregador, quedan igual. Lo primero que hace MAMMOTH es partir esa "
             "proyección en cabezas. La salida de W se corta en trozos y a cada trozo la figura lo "
             "llama x̄, la porción del parche que ve cada cabeza; son subespacios aprendidos, como "
             "en atención multi-cabeza, no textura ni forma ni color. Cada cabeza corre su propia "
             "mezcla de expertos en paralelo, que son los bloques MoE. Dentro de cada experto "
             "están los slots, que son prototipos aprendidos. Cada parche se compara con cada "
             "prototipo con un producto interno, y eso es el ruteo: mide qué tan parecido es el "
             "parche a ese slot. Sobre todos los parches de la slide se aplica una softmax, que "
             "decide cuánto aporta cada parche a cada slot. Con esos pesos se arma un promedio "
             "ponderado de los parches, y ese promedio ponderado, que en la figura es el weighted "
             "average, es lo que llena cada slot. El orden completo es este: el parche entra "
             "como x_i, el ruteo calcula los pesos comparándolo con los prototipos, y el slot "
             "termina siendo la x ponderada, el promedio de todos los parches que se le parecen. "
             "Cada slot queda representando un fenotipo, un tipo de tejido. Cada slot pasa después "
             "por la transformación del experto, que está factorizada en bajo rango, la Φ por la "
             "W_low seguida de una no linealidad; esa factorización es lo que le permite tener "
             "treinta expertos sin gastar más parámetros que la matriz original. Al final se "
             "concatenan las cabezas, el cross-head concat, y sale la salida de MAMMOTH, que "
             "alimenta al agregador. En este proyecto el agregador es CLAM: su atención junta los "
             "parches en un vector de slide y el clasificador entrega el diagnóstico. Un detalle: "
             "la figura muestra como salida los propios slots, y en nuestra configuración base una "
             "segunda softmax los recombina para reconstruir los parches y devolver la misma forma "
             "que entró, así el reemplazo de la capa lineal es directo. Casi seguro va a salir por "
             "qué esto es una mezcla de expertos y no un producto de expertos, y la respuesta está "
             "en cómo se combinan. Acá los expertos se combinan sumando: las dos softmax que "
             "describí reparten pesos que suman uno, así que la salida es un promedio, una "
             "combinación convexa. Eso es una mezcla, se entrena directo y es estable. Un producto "
             "de expertos haría lo contrario: multiplicaría las distribuciones de los expertos, "
             "como si cada uno tuviera un veto, y para entrenarlo haría falta un normalizador que "
             "es intratable; además modela una distribución de probabilidad, no una transformación "
             "de features, que es lo que acá necesitamos. Todo el diseño es aditivo por "
             "construcción, así que la mezcla es la elección natural y un producto rompería la "
             "formulación y traería de vuelta justo la inestabilidad que el método busca evitar. "
             "Aclaro que el paper no menciona el producto de expertos: esta "
             "comparación es un razonamiento de arquitectura, no una cita del paper.")

    # ---- 7b. NUEVA (§2.4): cabezas × expertos × slots — responde la duda 16 vs 30 ----
    # Rediseño 19-jul: los dos paneles de prosa se cambiaron por un ESQUEMA ANIDADO, que es
    # lo que la lámina intenta decir — que los tres números no se multiplican en fila sino
    # que uno vive dentro del otro. Los bullets que quedan son de una línea cada uno, y la
    # tabla (que era lo más consultable de la lámina) crece de 10 a 12 pt con filas altas.
    s = content(prs, "Cabezas, expertos y slots: la relación 16 × 30 × 10", size=22)

    # --- izquierda: el anidamiento, dibujado ---
    add_textbox(s, 0.30, 0.95, 5.25, 0.30, [
        ("Un parche, tres ejes anidados", 13, True, TEAL_TITLE, F_TITLE, PP_ALIGN.CENTER)])
    _proc(s, 2.12, 1.30, 1.60, 0.50, "parche  x_i", size=12)
    _conn(s, 2.92, 1.80, 2.92, 2.04)

    _grupo(s, 0.30, 2.04, 5.25, 1.92)
    add_textbox(s, 0.30, 2.10, 5.25, 0.28, [
        ("Cabeza h  ·  ×16 en paralelo", 11, True, ONCO_DARK, F_BODY, PP_ALIGN.CENTER)])
    for x, txt in zip((0.42, 2.21, 4.00),
                      ("Experto 1", "Experto 2", "…   Experto 30")):
        _proc(s, x, 2.44, 1.50, 0.48, txt, size=11)
    # los 10 slots de un experto: prototipos, dibujados como celdas y no descritos
    for i in range(10):
        _rect(s, 0.68 + i * 0.46, 3.10, 0.36, 0.30, TEAL_CARD2, line=ONCO_CONN)
    add_textbox(s, 0.30, 3.46, 5.25, 0.30, [
        ("10 slots por experto  →  30 × 10 = 300 slots", 11, True, ONCO_DARK, F_BODY,
         PP_ALIGN.CENTER)])

    add_textbox(s, 0.30, 4.06, 5.25, 1.00, [
        ("Cabezas y expertos son ejes distintos.", 11.5, False, INK, F_BODY),
        ("No hay un experto por cabeza: hay 16 × 30 combinaciones.", 11.5, False, INK, F_BODY),
        ("Cada slot es un prototipo aprendido.", 11.5, False, INK, F_BODY),
    ])

    # --- derecha: la tabla, ahora el objeto consultable de la lámina ---
    simple_table(s, 5.75, 0.95, 3.95,
                 ["Eje", "Cuántos", "Qué es"],
                 [["Cabezas (H)", "16", "subespacios del parche"],
                  ["Expertos (E)", "30", "especialistas por morfología"],
                  ["Slots (S)", "10 c/u", "prototipos por experto"],
                  ["Slots totales", "300", "E · S"]],
                 col_fracs=[0.32, 0.24, 0.44], row_h=0.50, fs=12)
    add_textbox(s, 5.75, 3.62, 3.95, 0.30, [
        ("Las dos softmax del ruteo", 13, True, TEAL_TITLE, F_TITLE)])
    _proc_claro(s, 5.75, 3.96, 3.95, 0.50, "dispatch", size=11,
                dim="softmax sobre los N parches · llena cada slot")
    _proc_claro(s, 5.75, 4.54, 3.95, 0.50, "combine", size=11,
                dim="softmax sobre los 300 slots · rearma cada parche")
    notes(s, "Una duda que suele aparecer acá: si hay dieciséis cabezas, ¿son "
             "dieciséis expertos? No. Cabezas y expertos son ejes distintos. Las dieciséis cabezas "
             "son vistas paralelas sobre subespacios del parche, como en atención multi-cabeza. Los "
             "treinta expertos son especialistas por morfología. Y cada cabeza no manda a un solo "
             "experto: cada cabeza calcula, para cada parche, su parecido con los diez slots de cada "
             "uno de los treinta expertos, así que hay dieciséis por treinta combinaciones, no "
             "dieciséis expertos. Cada slot guarda un prototipo aprendido, un vector que resume un "
             "fenotipo, y se llena con el promedio ponderado de los parches que se le parecen. "
             "Treinta expertos por diez slots dan trescientos slots. Todo el ruteo se resuelve con "
             "dos softmax sobre los mismos parecidos: una reparte los parches hacia los slots y otra "
             "recombina los trescientos slots para reconstruir cada parche.")

    # ---- 8. keep_slots: la bifurcación, con math+código (NATIVO) ----
    # Rediseño 19-jul (pedido de Ernesto: «que fuese un diagrama»). La lámina afirma que
    # hay UNA bifurcación sobre un tronco común, y eso es topología: dibujada se entiende
    # sola, mientras que en dos paneles de código enfrentados había que reconstruirla
    # leyendo. Se conserva la línea de código DECISIVA de cada rama, en cuerpo chico bajo
    # su bloque — es la evidencia de que la bifurcación es literalmente una línea — pero
    # deja de ser el objeto principal de la lámina.
    s = content(prs, "La variante keep_slots: dónde cambia la salida")

    # --- tronco compartido ---
    TW, TG, TY, TH = 1.30, 0.28, 0.98, 0.62
    txs = [1.98 + i * (TW + TG) for i in range(4)]
    for x, txt in zip(txs, ("z · features", "ruteo", "expertos LoRA", "concat")):
        _proc(s, x, TY, TW, TH, txt, size=11)
    for i in range(3):
        _conn(s, txs[i] + TW, TY + TH / 2, txs[i + 1], TY + TH / 2)
    add_textbox(s, 0.30, TY + TH + 0.04, 9.4, 0.26, [
        ("tronco compartido  →  [300 × 512]", 11, True, ONCO_INK, F_BODY, PP_ALIGN.CENTER)])

    # --- la bifurcación: un stub al centro y dos bajadas ---
    _conn(s, 5.0, 1.92, 5.0, 2.06, arrow=False)
    _conn(s, 2.60, 2.06, 7.40, 2.06, arrow=False)
    _conn(s, 2.60, 2.06, 2.60, 2.32)
    _conn(s, 7.40, 2.06, 7.40, 2.32)

    ramas = [
        (0.45, "keep_slots = False  ·  drop-in (base)", TEAL_TITLE, True,
         "2ª softmax  ·  combine", "[N × 512]",
         "Reconstruye los N parches. CLAM agrega sobre los parches, igual que la capa "
         "lineal original.",
         'out = einsum("h p d, n h p -> n h d", out, combine)'),
        (5.25, "keep_slots = True  ·  variante nueva", ORA_T, False,
         "se omite la recombinación", "[300 × 512]",
         "Se queda con los 300 slots. CLAM agrega sobre 300 tokens, no sobre los parches.",
         "return out    # 300 slot-tokens"),
    ]
    for x, cab, ccol, hace, op, forma, nota, codigo in ramas:
        add_textbox(s, x, 2.34, 4.30, 0.28, [(cab, 12.5, True, ccol, F_TITLE)])
        # la rama que OPERA va en bloque oscuro; la que se saltea el paso, en el tono
        # claro: el contraste dice cuál hace algo sin necesidad de escribirlo.
        (_proc if hace else _proc_claro)(s, x, 2.66, 4.30, 0.56, op, size=12)
        _conn(s, x + 2.15, 3.22, x + 2.15, 3.46)
        _proc_claro(s, x, 3.46, 4.30, 0.50, forma, size=13)
        add_textbox(s, x, 4.00, 4.30, 0.50, [(nota, 11, False, GRIS_BODY, F_BODY)])
        add_textbox(s, x, 4.54, 4.30, 0.26, [(codigo, 8.5, False, ONCO_DARK, F_MONO)])

    add_textbox(s, 0.30, 4.86, 9.4, 0.26, [
        ("Misma cabeza CLAM y misma pérdida: cambia solo qué se agrega. La interpretabilidad "
         "se calcula sobre el tronco, antes de la bifurcación, así que vale para las dos ramas.",
         11.5, True, TEAL_TITLE, F_BODY, PP_ALIGN.CENTER)])
    notes(s, "Hay una variante, keep_slots, y lo relevante es dónde cambia. Hasta la parte de "
             "arriba todo es idéntico a lo anterior: proyección, ruteo, expertos, y quedan los "
             "trescientos slots concatenados. La bifurcación es una sola línea de código. En la "
             "versión base, keep_slots en falso, una segunda softmax recombina los trescientos "
             "slots y reconstruye los parches, así que CLAM sigue agregando sobre los parches y el "
             "reemplazo de la capa lineal es directo. En la versión nueva, keep_slots en "
             "verdadero, esa recombinación se omite y la salida son los trescientos slots; "
             "entonces CLAM agrega sobre trescientos tokens en vez de sobre los parches. Es una "
             "variante orientada a rendimiento, y para la reunión de hoy no cambia nada, porque la "
             "interpretabilidad se calcula antes de esta bifurcación y vale para las dos ramas. La "
             "cabeza de CLAM y la pérdida son las mismas; lo único que cambia es qué se agrega.")

    # ---- 9. Divisoria: Interpretabilidad (OBJ-A) ----
    s = divider(prs, "¿Qué mira cada experto?",
                "Interpretabilidad post-hoc sobre un checkpoint entrenado · 4 slides TCGA-BRCA")
    notes(s, "Con el mecanismo claro, empieza la segunda parte: mirar, sobre las propias slides "
             "del proyecto, qué región concentra cada experto y qué morfología hay ahí. Es un "
             "análisis post-hoc, en CPU, sobre un modelo ya entrenado; no reentrena nada.")

    # ---- 10. Ruteo espacial (dónde) + morfología top-k (qué) — FUSIÓN ----
    s = content(prs, "Qué mira cada experto: dónde y qué morfología")
    add_image_fit(s, HEATMAP_MONTAGE, 0.30, 0.98, 4.95, 3.35, align="top")
    caption(s, 0.30, 4.38, 4.95,
            "Ruteo espacial (WSI TCGA-E2-A14Q, tarea cdis, fold 0): los 30 expertos encienden zonas distintas",
            size=9)
    add_image_fit(s, TOPK_SUBSET, 5.42, 0.98, 2.55, 3.05, align="top")
    add_textbox(s, 8.05, 1.05, 1.80, 3.2, [
        ("Top-k a alta resolución:", 11.5, True, INK, F_BODY),
        ("cada experto concentra un patrón de tejido consistente", 11, True, TEAL_TITLE, F_BODY),
        ("los patrones (e8, e26, e3) se nombran por inspección visual; falta sign-off de patólogo",
         9.5, False, GRIS_BODY, F_BODY),
        ("Color = percentil por experto (estructura relativa, no magnitud de uso).",
         9.5, False, ORA_T, F_BODY),
    ], anchor=MSO_ANCHOR.TOP)
    caption(s, 5.42, 4.12, 2.55, "Morfología que cada experto rutea (Fig 3.2)", size=9)
    takeaway_bar(s, "Heatmap = dónde · top-k = qué morfología. Emergió sin supervisión de tejido "
                    "(el paper lo validó con patólogos).", t=4.82, size=12)
    notes(s, "Estas dos vistas cuentan una sola historia, así que las leo juntas. A la "
             "izquierda, cada uno de los treinta cuadros es la misma slide pintada según cuánto el "
             "router manda cada parche a un experto: rojo, mucho; azul, casi nada. Lo que se "
             "observa es que los treinta encienden zonas distintas de la lámina; es decir, la capa "
             "lineal única quedó reemplazada por especialistas que miran regiones diferentes del "
             "tejido. Una precisión sobre el color: está normalizado por percentil dentro "
             "de cada experto, así que sirve para ver qué regiones prefiere cada uno, no para "
             "decir cuál se usa más; medido aparte, el uso sale casi uniforme, sin expertos "
             "muertos ni acaparadores. Ahora bien, el mapa de calor dice dónde, no qué tejido hay "
             "en esas zonas. Eso lo cierra la derecha: para cada experto se toman los parches que "
             "más rutea y se recortan a alta resolución real. Ahí se ve directo: el experto ocho "
             "mira nidos de epitelio tumoral, el veintiséis mira estroma fibroso rosado, el tres "
             "mira ductos. Es la misma especialización que el paper validó con dos patólogos, y "
             "emerge sola: nadie le indicó al modelo qué es estroma. Los expertos "
             "mixtos que aparecen son esperables, porque el ruteo es suave y reparte cada parche "
             "entre varios.")

    # ---- 11. Morfología ≠ clase (cross-slide) + honestidad/cierre — FUSIÓN ----
    #   imágenes GRANDES lado a lado (grillas 4 slides × 5 parches; aspect 1.19) + texto compacto.
    s = content(prs, "El experto detecta tejido, no clase")
    xh = 2.92; xw = xh * (704 / 593)                 # aspect real de las cross-slide
    gap = 0.34
    x0 = (SW - (2 * xw + gap)) / 2
    s.shapes.add_picture(XSLIDE_E8, Inches(x0), Inches(0.86), Inches(xw), Inches(xh))
    s.shapes.add_picture(XSLIDE_E26, Inches(x0 + xw + gap), Inches(0.86), Inches(xw), Inches(xh))
    # El pie decía el sign-off y el párrafo lo repetía: dos líneas de más que hacían
    # que el pie pisara al párrafo. La honestidad completa vive en el párrafo.
    caption(s, 0.35, 0.86 + xh + 0.02, SW - 0.7,
            "el mismo experto (e8, e26) concentra el mismo patrón en las 4 slides TCGA-BRCA · tarea cdis",
            size=10.5, bold=True, col=INK)
    add_textbox(s, 0.35, 0.86 + xh + 0.30, SW - 0.7, 0.62, [
        ("Dos slides positivas y dos negativas.  "
         "«Negativo» = cdis sin microcalcificación, no «sin tumor» → el experto es un detector de "
         "TEJIDO, no de clase (la clase se decide después, en la atención + el clasificador).  "
         "Honestidad: 4 slides · 1 tarea · 1 fold · etiquetas provisionales (falta sign-off de "
         "patólogo).", 10.5, False, GRIS_BODY, F_BODY, PP_ALIGN.CENTER)], anchor=MSO_ANCHOR.TOP)
    takeaway_bar(s, "Que la especialización sea real y aun así no mueva la métrica es la evidencia "
                    "de que el cuello no está en la 1ª capa, sino en el dato.", t=4.82, size=12)
    notes(s, "Con esta prueba, la más fina de las tres, cierra la presentación. Como los prototipos "
             "son parámetros compartidos del modelo, el experto ocho es el mismo experto en todas "
             "las slides. Eso permite una prueba limpia: si su especialización es real, tiene que "
             "elegir la misma morfología en todas. Y así ocurre. En las dos grillas, cada fila es "
             "una de cuatro slides y cada columna un parche top del mismo experto: el ocho "
             "enciende epitelio en las cuatro, el veintiséis enciende estroma en las cuatro, "
             "incluidas las dos negativas. El matiz importa: negativo aquí significa cdis sin "
             "microcalcificación, no sin tumor, así que las negativas siguen siendo slides de mama "
             "con epitelio y estroma. Entonces el experto detecta tejido, no clase; que la slide "
             "sea positiva se decide después, en la atención y el clasificador. De ahí sale el "
             "punto para la discusión: que los expertos separen bien los tejidos y aun así MAMMOTH "
             "no le gane a CLAM en métrica es la evidencia de que el cuello de botella no está en "
             "la primera capa, sino en el dato. Los límites de la prueba: es una cala chica, "
             "cuatro slides, una tarea, un fold, y las etiquetas de tejido todavía son "
             "provisionales, a la espera del visto bueno de un patólogo. Son dos ejes distintos: "
             "que MAMMOTH no mejore la métrica está cerrado; qué aprende por dentro está abierto, "
             "y es lo que trae esta presentación. Mirar qué mira por dentro no reabre la pregunta "
             "del rendimiento; explica por qué dio lo que dio.")

    # ========================================================================
    # SECCIÓN COMPARACIÓN PAREADA (resultados nuevos: entrenamiento 3 tareas x 2
    # brazos x 5 folds + comparación de atención + conteo efectivo de expertos/slots)
    # ========================================================================

    # ---- 11b. Divisoria: comparación pareada ----
    s = divider(prs, "¿Dónde mira cada modelo?",
                "Entrenamiento pareado sobre tres tareas clínicas · comparación de la atención "
                "sobre las mismas láminas")
    notes(s, "Hasta acá vimos el mecanismo por dentro y qué concentra cada experto sobre una cala "
             "chica. Esta parte es el trabajo nuevo: entrenar los dos modelos sobre exactamente "
             "las mismas particiones de tres tareas clínicas, y después comparar, lámina por "
             "lámina, dónde pone la atención cada uno. La pregunta ya no es cuál mide mejor, es "
             "en qué se parecen y en qué se diferencian por dentro cuando ven el mismo tejido.")

    # ---- 11c. Diseño pareado + métricas (política de eval: balanceada Y AUC juntas) ----
    s = content(prs, "Comparación pareada: mismas particiones y datos", size=24)
    add_textbox(s, 0.30, 0.86, 9.4, 0.42, [
        ("Tres tareas · cinco particiones · dos brazos · treinta ejecuciones completas. El conjunto "
         "de prueba se verificó idéntico entre brazos por firma md5, partición por partición.",
         11.5, False, GRIS_BODY, F_BODY)])
    # Sin saltos de línea en las celdas: con wrap la tabla crece y pisa los paneles.
    simple_table(s, 0.30, 1.36, 9.4,
                 ["Tarea", "n", "Balance. CLAM", "Balance. MAM", "Δ pareada (folds+)",
                  "AUC CLAM", "AUC MAM", "Δ AUC"],
                 [["Tipo histológico (3 cl.)", "2027", "0.665 ± 0.056", "0.655 ± 0.047",
                   "−0.010 ± 0.017 (1/5)", "0.833 ± 0.043", "0.821 ± 0.056", "−0.012"],
                  ["Invasión linfovascular", "836", "0.657 ± 0.040", "0.634 ± 0.050",
                   "−0.023 ± 0.086 (2/5)", "0.720 ± 0.032", "0.684 ± 0.056", "−0.036"],
                  ["Ductal in situ presente", "862", "0.668 ± 0.098", "0.742 ± 0.099",
                   "+0.074 ± 0.033 (5/5)", "0.765 ± 0.111", "0.825 ± 0.086", "+0.060"]],
                 col_fracs=[0.175, 0.05, 0.125, 0.125, 0.165, 0.125, 0.125, 0.11],
                 row_h=0.42, fs=9)
    add_textbox(s, 0.30, 3.28, 4.55, 1.10, [
        ("Dos tareas confirman lo previsto", 12, True, TEAL_TITLE, F_TITLE),
        ("La diferencia queda dentro del ruido: la desviación iguala o supera a la media.",
         10.5, False, GRIS_BODY, F_BODY)])
    _rect(s, 5.10, 3.22, 4.60, 1.22, TEAL_CARD2, line=ORA_ACC)
    add_textbox(s, 5.26, 3.28, 4.30, 1.10, [
        ("Ductal in situ: a verificar, no a celebrar", 12, True, ORA_T, F_TITLE),
        ("Suben los dos recalls a la vez (0.477→0.569 y 0.860→0.915), así que no es haber "
         "corrido el umbral. Pero son 65 negativos en total.", 10.5, False, GRIS_BODY, F_BODY)])
    takeaway_bar(s, "Candidato a réplica con más semillas antes de contarlo como mejora.",
                 t=4.62, size=12)
    notes(s, "Empiezo por el diseño, porque es lo que hace que los números "
             "signifiquen algo. Los dos modelos corrieron sobre exactamente las mismas "
             "particiones, y eso no se asumió: se comprobó comparando la firma de los "
             "identificadores de lámina del conjunto de prueba, partición por partición. Así la "
             "diferencia entre brazos no arrastra el azar del sorteo. Treinta ejecuciones, todas "
             "completas. En tipo histológico y en invasión linfovascular el resultado es el que "
             "estaba anticipado por escrito antes de correr: la diferencia se queda dentro del "
             "ruido, la desviación entre particiones iguala o supera a la media, no hay señal. La "
             "tercera fila es distinta. En carcinoma ductal in "
             "situ MAMMOTH mide mejor en las cinco particiones, tanto en exactitud balanceada "
             "como en área bajo la curva. Y hay un detalle que importa: suben los dos recalls a la "
             "vez, el de la clase minoritaria y el de la mayoritaria. Cuando uno simplemente mueve "
             "el punto de corte, sube uno y baja el otro; acá no pasa eso, así que la mejora está "
             "en el ordenamiento, no en el umbral. Aun así freno antes de llamarlo mejora, por una "
             "razón de tamaño: hay sesenta y cinco negativos en total, unos trece por partición, "
             "y cada uno mueve casi ocho puntos de recall. Además es una formulación de la tarea "
             "que es nueva. Lo dejo como candidato a replicar con más semillas, no como un "
             "resultado cerrado.")

    # ---- 11d. Comparación de atención: el entregable central ----
    s = content(prs, "Mismo barrio, distintas casas", size=28)
    add_textbox(s, 0.30, 0.84, 9.4, 0.40, [
        ("Siete láminas, una por clase y tarea, todas del conjunto de prueba y bien clasificadas "
         "por los dos brazos: se compara el foco sin el ruido de un error.",
         11.5, False, GRIS_BODY, F_BODY)])
    add_image_fit(s, ATT_SBS, 0.30, 1.30, 5.75, 2.95, align="top")
    caption(s, 0.30, 4.30, 5.75, "atención CLAM | MAMMOTH | diferencia · misma lámina, mismo código",
            size=9.5, col=GRIS_TXT)
    _rect(s, 6.25, 1.26, 3.45, 3.02, TEAL_CARD2, line=TEAL_SQ)
    add_textbox(s, 6.42, 1.34, 3.15, 2.90, [
        ("Agregado sobre las siete", 12.5, True, TEAL_TITLE, F_TITLE),
        ("Correlación de rangos   0.805", 11, False, INK, F_BODY),
        ("Solapamiento 5% super.   0.172", 11, False, INK, F_BODY),
        ("Solapamiento 1% super.   0.073", 11, False, INK, F_BODY),
        ("Entropía CLAM   0.781", 11, False, INK, F_BODY),
        ("Entropía MAMMOTH   0.894", 11, False, ORA_T, F_BODY),
        ("Coinciden en el mapa grueso: ordenan el tejido igual, la región que importa es la misma.",
         10.5, False, GRIS_BODY, F_BODY),
        ("Difieren en los picos: los parches concretos que cada uno pone arriba son distintos.",
         10.5, False, GRIS_BODY, F_BODY),
        ("MAMMOTH reparte, CLAM concentra: más difusa en seis de siete.",
         10.5, False, GRIS_BODY, F_BODY)])
    takeaway_bar(s, "Que la mayor difusión aparezca donde MAMMOTH también mide mejor es "
                    "sugerente, pero con siete láminas es hipótesis, no explicación.",
                 t=4.62, size=11.5)
    notes(s, "El resultado central de la parte nueva sale de acá. Se eligieron siete láminas, una por "
             "cada clase de cada tarea, todas del conjunto de prueba, es decir nunca vistas en "
             "entrenamiento, y además bien clasificadas por los dos modelos. Esa segunda condición "
             "es deliberada: si uno de los dos se equivoca, la comparación de dónde mira queda "
             "contaminada por el error. La comparación además es directa porque MAMMOTH hereda "
             "el mismo método de atención del modelo base, así que los dos mapas salen del mismo "
             "código. A la izquierda está una de las siete: a un lado la atención de CLAM, "
             "al otro la de MAMMOTH, y al costado la diferencia. Los números de la derecha "
             "resumen las siete y cuentan tres cosas. La primera es que coinciden en el mapa "
             "grueso: la correlación de rangos es alta, los dos ordenan el tejido de forma "
             "parecida, la región que importa es la misma. La segunda es que difieren en los "
             "picos: si uno mira el cinco por ciento de parches más atendidos, el solapamiento es "
             "bajo, y si mira el uno por ciento, es casi nulo. Mismo barrio, distintas casas. La "
             "tercera es que la variante reparte la atención y CLAM la concentra: la "
             "entropía es mayor en seis de las siete láminas. Y hay algo que da ganas de conectar: "
             "la mayor difusión aparece justo en la tarea donde MAMMOTH también mide mejor. Es "
             "una hipótesis, nada más. Con siete láminas no alcanza para atribuir la "
             "diferencia de métrica a la forma de la atención.")

    # ---- 11e. Cuántos expertos y cuántos slots se usan de verdad ----
    # Rediseño 19-jul: el resultado de Q1 es una FRACCIÓN de presupuesto usado (30 de 30
    # contra 159 de 300), y una fracción se lee de un vistazo como barra. La tabla de dos
    # filas obligaba a hacer la división mentalmente; las barras ponen el contraste —
    # lleno contra medio lleno — que ES el hallazgo. La prosa de «Lectura» se disuelve en
    # la nota al pie de cada barra.
    s = content(prs, "¿Cuántos expertos y slots se usan de verdad?", size=25)
    q1 = _leer_q1()
    add_textbox(s, 0.30, 0.90, 9.4, 0.44, [
        ("Se mide sobre el peso de combinación, la segunda softmax sobre los 300 slots "
         "(30 expertos × 10). No es el conteo de parches más atendidos.",
         11.5, False, GRIS_BODY, F_BODY)])

    def _frac(v, total):
        try:
            return min(1.0, float(v) / total)
        except (TypeError, ValueError):
            return 0.0                      # Q1 sin correr: barra vacía, no un número inventado

    ratio_bar(s, 0.30, 1.46, 4.40, 0.42, _frac(q1["exp"], 30), "Expertos",
              q1["exp"], "30", "Se usan los treinta por igual: ninguno queda apagado.")
    ratio_bar(s, 5.30, 1.46, 4.40, 0.42, _frac(q1["slots"], 300), "Slots",
              q1["slots"], "300", "Cerca de la mitad del presupuesto aporta poco al peso final.")

    _grupo(s, 0.30, 3.26, 9.4, 0.86)
    add_textbox(s, 0.48, 3.32, 9.04, 0.74, [
        ("Número efectivo = exp(entropía): vale el total si el reparto fuera uniforme, y 1 si "
         "colapsara en un solo slot.", 11, True, ONCO_DARK, F_BODY),
        ("Se usa esta medida y no un conteo porque la softmax da peso positivo a todos.",
         10, False, ONCO_DARK, F_BODY)])
    caption(s, 0.30, 4.18, 9.4, q1["pie"], size=9.5, col=GRIS_TXT)
    takeaway_bar(s, "Si hay que ajustar capacidad, el parámetro a tocar son los slots, no los "
                    "expertos.", t=4.56, size=12)
    notes(s, "Quedaba una pregunta concreta de la vez pasada: cuántos expertos y cuántos slots usa "
             "realmente el modelo. Antes de responder tengo que precisar qué se mide, porque hay "
             "dos cosas que se confunden. El peso por slot es la segunda distribución softmax, la "
             "que combina las trescientas salidas, treinta expertos por diez slots cada uno. No es "
             "el conteo de parches que cada experto atiende, que es otra medida distinta. Y hay un "
             "detalle metodológico: como es una softmax, todos los slots reciben algo de peso, "
             "nunca cero. Entonces contar cuántos reciben peso no sirve, siempre daría "
             "trescientos. Por eso se usa el número efectivo, que es la exponencial de la "
             "entropía: da trescientos si el reparto fuera perfectamente parejo y uno si todo el "
             "peso cayera en un solo slot. Con esa medida, los expertos salen prácticamente "
             "uniformes, se usan los treinta, ninguno queda apagado; mientras que en los slots el "
             "número efectivo queda bastante por debajo del total. La conclusión práctica es que "
             "el número de expertos no está sobredimensionado, y que si en algún momento hay que "
             "recortar capacidad, el parámetro a tocar son los slots.")

    # ========================================================================
    # SECCIÓN MAGNIFICACIÓN MULTI-ESCALA (reunión Sebastián — decisión de escalas)
    # ========================================================================

    # ---- 12. Divisoria: Magnificación multi-escala ----
    # guion NO separable (U+2011): con el guion normal partía en "Magnificación multi-" / "escala"
    s = divider(prs, "Magnificación multi‑escala",
                "La única señal nueva tras cerrar la arquitectura: el contexto espacial · "
                "piloto microcalcificaciones")
    notes(s, "Cerrado el capítulo de la arquitectura, donde cuatro ejes distintos no movieron "
             "la métrica, queda una idea que todavía no probamos y que sí trae información "
             "nueva: la escala a la que miramos el tejido. De eso trata esta parte. Es un "
             "piloto sobre microcalcificaciones, que elegimos porque son pocas láminas y la "
             "extracción cabe en un fin de semana, pero la misma idea sirve para cualquier "
             "tarea que dependa del contexto.")

    # ---- 13. El problema es contexto + hallazgo físico (fusión de las dos slides) ----
    # Rediseño 19-jul: los dos paneles de bullets («detectar» vs «localizar») decían con
    # seis líneas algo que es una comparación de TAMAÑOS, así que pasa a un eje de escala
    # física. Ahí se ve de una que el parche de hoy cubre la calcificación y se queda corto
    # frente al conducto, que es el argumento entero de la lámina. El eje va en log porque
    # el rango útil abarca dos órdenes de magnitud (decenas de µm a milímetros).
    s = content(prs, "No es más zoom, es contexto: escalas por cohorte", size=21)
    add_textbox(s, 0.30, 0.90, 9.4, 0.34, [
        ("La etiqueta pregunta DÓNDE vive la microcalcificación, no si existe: es una pregunta "
         "de contexto, no de detalle celular.", 12.5, False, GRIS_BODY, F_BODY)])

    import math
    _LO, _HI = 10.0, 4000.0                       # µm, extremos del eje
    _f = lambda um: (math.log10(um) - math.log10(_LO)) / (math.log10(_HI) - math.log10(_LO))
    # Los dos tramos de arriba van en #3E6877 y #0E2841: teal contra azul marino. Con
    # #386271 (el color de conector) el contraste contra #3E6877 era nulo y el eje se leía
    # como un solo tramo largo, justo lo contrario de lo que la lámina quiere mostrar.
    scale_axis(s, 0.45, 1.86, 9.10, [
        (_f(50), _f(500), "la calcificación  50-500 µm", True, ONCO_DARK),
        (_f(500), _f(2000), "el conducto anfitrión  0.5-2 mm", True, ONCO_INK),
        (_f(60), _f(120), "el parche fino de hoy  ~100 µm", False, ONCO_CONN),
    ])
    add_textbox(s, 0.30, 2.36, 9.4, 0.30, [
        ("El parche de hoy resuelve la calcificación, pero no cubre la estructura que la aloja.",
         12, True, TEAL_TITLE, F_BODY, PP_ALIGN.CENTER)])

    add_textbox(s, 0.30, 2.74, 4.60, 0.30, [
        ("Y las cohortes están a distinta escala física", 12.5, True, TEAL_TITLE, F_TITLE)])
    simple_table(s, 0.30, 3.06, 4.60,
                 ["Cohorte", "Magnif.", "Parche 256 px"],
                 [["Pública (TCGA)", "~40×", "59 µm"],
                  ["Privada", "~20×", "119 µm"],
                  ["HistAI", "sin MPP", "excluida"]],
                 col_fracs=[0.44, 0.24, 0.32], row_h=0.32, fs=10.5)
    _grupo(s, 5.20, 2.74, 4.50, 1.60)
    add_textbox(s, 5.40, 2.84, 4.10, 1.42, [
        ("Medido en µm/px real,", 11.5, True, ONCO_DARK, F_BODY),
        ("no en la etiqueta del archivo.", 11.5, True, ONCO_DARK, F_BODY),
        (" ", 8, False, ONCO_DARK, F_BODY),   # separador; 8·1.3333 = 10.7 pt, sobre el mínimo
        ("Difieren 2×.", 12.5, True, ONCO_DARK, F_BODY),
        ("La pirámide se define en µm/px,", 11.5, False, ONCO_DARK, F_BODY),
        ("no en «level» del archivo.", 11.5, False, ONCO_DARK, F_BODY),
    ])
    takeaway_bar(s, "Hay que sumar una escala GRUESA que aporte el contexto del tejido, "
                    "y definirla en micras, no en «level».", t=4.46, size=12.5)
    notes(s, "Conviene empezar por algo que parece al revés. Uno pensaría que una "
             "microcalcificación, como es un objeto chico, pide más aumento, y es lo contrario. "
             "La etiqueta que queremos predecir no pregunta si hay una calcificación, sino en qué "
             "estructura vive: dentro de un carcinoma in situ, de un carcinoma invasor o de "
             "tejido no neoplásico. Eso es contexto, no detalle. Para responderlo el modelo tiene "
             "que ver la estructura anfitriona completa, el conducto, el nido invasor o el "
             "lobulillo, que miden entre medio milímetro y dos milímetros. Hoy extraemos un solo "
             "parche fino por región, que detecta la calcificación pero no alcanza a cubrir el "
             "conducto de alrededor. Por eso lo que falta no es más zoom, es sumar una escala "
             "gruesa que traiga ese contexto. A esto se agrega un hallazgo que apareció cuando "
             "medimos las láminas. Uno esperaría que un parche del mismo tamaño en píxeles "
             "cubriera lo mismo en las dos cohortes, y no. Con la resolución real de cada "
             "escáner, las láminas de la cohorte pública están cerca de cuarenta aumentos y las "
             "de la privada cerca de veinte, o sea el doble de resolución unas que otras, aunque "
             "la etiqueta del archivo diga otra cosa. Un mismo parche de doscientos cincuenta y "
             "seis píxeles cubre cincuenta y nueve micras en una y ciento diecinueve en la otra. "
             "De ahí salen dos consecuencias: la escala hay que definirla en micras por píxel y "
             "no en niveles del archivo, y el pipeline actual ya arrastra un sesgo, porque le "
             "entrega cada cohorte a un aumento distinto. Re-extraer a un campo físico común "
             "habilita la pirámide y de paso corrige ese sesgo. La tercera cohorte no tiene "
             "resolución confiable en su metadata, así que queda fuera.")

    # ---- 13b. NUEVA (§3): la matemática área ↔ magnificación ↔ tamaño de parche ----
    # Rediseño 19-jul: la lámina es una CUENTA, y una cuenta se muestra hecha. Las dos
    # cohortes pasan a ser dos cadenas de bloques con el mismo esqueleto (P × MPP = lado),
    # así el confound se ve en que la única celda distinta es el MPP y el resultado se va
    # al doble. El panel de cinco bullets se reduce a la cadena inversa que despeja P.
    s = content(prs, "La matemática: µm/px, área física y tamaño de parche", size=22)
    add_textbox(s, 0.30, 0.90, 9.4, 0.28, [
        ("La escala física la fija el MPP (micras por píxel), no la magnificación nominal ni el "
         "«level» del archivo.", 12, False, GRIS_BODY, F_BODY)])

    _grupo(s, 0.40, 1.24, 9.20, 0.50)
    add_textbox(s, 0.40, 1.24, 9.20, 0.50, [
        ("lado (µm)  =  P (px)  ×  MPP (µm/px)          área = (P × MPP)²",
         15, True, ONCO_DARK, F_BODY, PP_ALIGN.CENTER)], anchor=MSO_ANCHOR.MIDDLE)

    XL, WL = 0.45, 1.75
    XP, WP = 2.40, 1.15
    XM, WM = 4.15, 1.60
    XR, WR = 6.35, 1.55
    for y, coh, mpp, res, nota in (
            (1.88, "TCGA  (~40×)", "0.2325", "52 µm", "más resolución"),
            (2.54, "Privado  (~20×)", "0.465", "104 µm", "el doble de campo"),
    ):
        add_textbox(s, XL, y, WL, 0.52, [(coh, 12, True, TEAL_TITLE, F_TITLE)],
                    anchor=MSO_ANCHOR.MIDDLE)
        _proc_claro(s, XP, y, WP, 0.52, "224 px", size=12)
        _oper(s, (XP + WP + XM) / 2, y + 0.26, "×")
        _proc_claro(s, XM, y, WM, 0.52, mpp, size=12)
        _oper(s, (XM + WM + XR) / 2, y + 0.26, "=")
        _proc(s, XR, y, WR, 0.52, res, size=13)
        add_textbox(s, XR + WR + 0.15, y, 1.70, 0.52, [(nota, 10.5, False, GRIS_BODY, F_BODY)],
                    anchor=MSO_ANCHOR.MIDDLE)
    add_textbox(s, XL, 3.10, 9.15, 0.26, [
        ("Mismo P en píxeles, distinto campo físico: ése es el confound de escala entre cohortes.",
         11, False, ONCO_DARK, F_BODY)])

    _grupo(s, 0.40, 3.44, 9.20, 0.92)
    add_textbox(s, 0.60, 3.50, 2.60, 0.80, [
        ("Igualar el campo físico", 12.5, True, ONCO_DARK, F_BODY)], anchor=MSO_ANCHOR.MIDDLE)
    _proc_claro(s, 3.20, 3.62, 1.30, 0.52, "104 µm", size=12)
    _oper(s, 4.68, 3.88, "÷")
    _proc_claro(s, 4.86, 3.62, 1.30, 0.52, "0.2325", size=12)
    _oper(s, 6.34, 3.88, "=")
    _proc(s, 6.52, 3.62, 1.30, 0.52, "448 px", size=13)
    add_textbox(s, 7.95, 3.62, 1.50, 0.52, [
        ("224@×20 ≡ 448@×40", 10.5, True, ONCO_DARK, F_BODY)], anchor=MSO_ANCHOR.MIDDLE)

    takeaway_bar(s, "La pirámide se define en µm/px físicos, no en «level»; así el campo de visión "
                    "es comparable entre cohortes.", t=4.48, size=12.5)
    notes(s, "Acá respondo una pregunta que suele surgir: cómo se relacionan la "
             "magnificación, el área de tejido y el tamaño del parche en píxeles. La cantidad que "
             "manda es las micras por píxel, el eme pe pe. El lado físico de un parche es el número "
             "de píxeles por las micras por píxel, y el área es ese lado al cuadrado. Con la "
             "resolución real de cada escáner, un parche de doscientos veinticuatro píxeles cubre "
             "cincuenta y dos micras en la cohorte pública, que está a cuarenta aumentos, y ciento "
             "cuatro en la privada, que está a veinte. El mismo número de píxeles, distinto campo "
             "físico: ése es el sesgo de escala entre cohortes. Para igualarlo se agranda el parche "
             "en la cohorte de más resolución: para cubrir en la pública el mismo campo de ciento "
             "cuatro micras que doscientos veinticuatro píxeles dan a veinte aumentos, hacen falta "
             "cuatrocientos cuarenta y ocho píxeles. O sea, doscientos veinticuatro a veinte "
             "equivalen a cuatrocientos cuarenta y ocho a cuarenta: misma área física, distinta "
             "resolución. Por eso las escalas se definen en micras por píxel y no en niveles del "
             "archivo.")

    # ---- 14. Lo que estudiamos: patología + referencias ----
    # Rediseño 19-jul: el contraste entre los dos tipos de calcio pasa a dos tarjetas con
    # una barra de VISIBILIDAD, que es lo que decide si el modelo puede verlas — una tabla
    # de tres columnas hacía leer para llegar a eso. Las referencias, que ocupaban una
    # columna entera con catorce líneas apiladas, se reparten en tres grupos al pie: misma
    # información, un cuarto del peso visual.
    s = content(prs, "Lo que estudiamos: patología de la microcalcificación", size=24)
    add_textbox(s, 0.30, 0.90, 9.4, 0.28, [
        ("Dos tipos de calcio, y uno es invisible en la tinción de rutina.",
         12.5, False, GRIS_BODY, F_BODY)])

    for x, titulo, cuerpo, vis, visnota in (
            (0.30, "Tipo I  ·  oxalato",
             ["Casi invisible en H&E: solo aparece con luz polarizada.",
              "El modelo está ciego a cualquier escala → TECHO, sobre todo en "
              "tejido no neoplásico."], 0.06, "prácticamente nula"),
            (5.10, "Tipo II  ·  fosfato",
             ["Basófilo y visible, con los anillos laminados característicos.",
              "Acá sí juega la escala: es donde el campo grueso puede aportar."],
             0.90, "buena"),
    ):
        _grupo(s, x, 1.24, 4.60, 1.98)
        _proc(s, x + 0.15, 1.36, 4.30, 0.46, titulo, size=13)
        add_textbox(s, x + 0.22, 1.88, 4.16, 0.80,
                    [(ln, 11, False, ONCO_DARK, F_BODY) for ln in cuerpo])
        add_textbox(s, x + 0.22, 2.72, 2.10, 0.26,
                    [("visible en H&E", 9.5, True, ONCO_DARK, F_BODY)], anchor=MSO_ANCHOR.MIDDLE)
        _rect(s, x + 2.10, 2.78, 2.28, 0.16, TEAL_CARD2)
        _rect(s, x + 2.10, 2.78, 2.28 * vis, 0.16, ONCO_DARK)
        add_textbox(s, x + 2.10, 2.94, 2.28, 0.24,
                    [(visnota, 9.5, False, ONCO_DARK, F_BODY, PP_ALIGN.RIGHT)])

    add_textbox(s, 0.30, 3.32, 9.4, 0.30, [
        ("La calcificación (50-500 µm) entra en un parche fino; el conducto anfitrión "
         "(0.5-2 mm), no.", 12, True, TEAL_TITLE, F_BODY, PP_ALIGN.CENTER)])

    # referencias: tres grupos al pie, una línea por cita
    _rect(s, 0.30, 3.70, 9.4, 1.30, TEAL_CARD2, line=TEAL_SQ)
    for x, grupo, refs in (
            (0.48, "Clínica de la microcalcificación",
             ["Breast microcalcifications · Review 2022",
              "Calcification in breast histopath. · 2024",
              "Polyhedral = oxalato · Radiology 1993",
              "Size matters! · 2007",
              "Predictors of malignancy · BJC 2011"]),
            (3.68, "Multi-escala en patología",
             ["CPathAgent · NeurIPS 2025 (Ap. C.1.2)",
              "DSMIL 20×+5× · Li et al., CVPR 2021",
              "Deep Multi-Magnification Nets · 2021"]),
            (6.88, "Modelo base y etiqueta",
             ["CONCH · Lu et al., Nat Med 2024 (20×)",
              "CAP Invasive Breast, Nota D"]),
    ):
        add_textbox(s, x, 3.78, 3.00, 1.16,
                    [(grupo, 10.5, True, TEAL_TITLE, F_BODY)]
                    + [("· " + r, 9, False, INK, F_BODY) for r in refs])
    notes(s, "Antes de elegir números me apoyo en la patología, donde hay dos hechos que "
             "mandan. El primero es que no todas las calcificaciones se ven igual en la tinción "
             "de rutina. Las de oxalato de calcio son casi invisibles en campo claro, solo "
             "aparecen con luz polarizada; el modelo, que trabaja sobre la tinción normal, está "
             "ciego a ellas a cualquier escala, y eso pone un techo sobre todo en el tejido no "
             "neoplásico. Las de fosfato de calcio, en cambio, son basófilas, se ven bien, y ahí "
             "sí la escala puede ayudar. El segundo hecho es de tamaño: la calcificación en sí, "
             "de cincuenta a quinientas micras, entra en un parche fino; el conducto que la aloja, "
             "no. Estas dos ideas vienen de la literatura de patología mamaria que está citada en "
             "la lámina, junto con los trabajos de multi-escala en los que nos apoyamos.")

    # ---- 15. La decisión de escalas (LA slide para Sebastián) ----
    s = content(prs, "La decisión de escalas", size=28)
    # La pill vive DENTRO de la banda, a la derecha del título (que en esta lámina es corto
    # y no llega hasta ahí). Bajarla al contenido no es opción: ahí chocaba con los dos
    # crops de tejido. Se centra en la banda (0.33-1.07) y se la nombra ONCOHDR_ para que
    # reflow_onco la trate como cabecera y no la desplace con el cuerpo.
    _pill = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.60), Inches(0.46), Inches(3.10), Inches(0.48))
    _pill.name = "ONCOHDR_pill"
    _pill.fill.solid(); _pill.fill.fore_color.rgb = PROG_BG
    _pill.line.color.rgb = ONCO_CONN; _pill.line.width = Pt(1.25); _pill.shadow.inherit = False
    # Marcador de estado: mismo arquetipo gris que la pill "En progreso", y por la misma
    # razón el texto va NEGRO (el teal oscuro sobre #B7B7B7 no llega al contraste mínimo).
    _set_runs(_pill.text_frame, [("Escalas a definir", 11, True, BLACK, F_BODY, PP_ALIGN.CENTER)],
              anchor=MSO_ANCHOR.MIDDLE)
    # izquierda: esquema nativo de campos concéntricos + pipeline
    nested_fields(s, 0.30, 1.02, 4.55, 1.95)
    # tabla px por cohorte
    simple_table(s, 0.30, 3.08, 4.55,
                 ["Escala", "Campo", "px TCGA", "px privado"],
                 [["Fina", "112 µm", "482", "241"],
                  ["Contexto", "512 µm", "2202", "1101"]],
                 col_fracs=[0.28, 0.24, 0.24, 0.24], row_h=0.32, fs=10.5)
    add_textbox(s, 0.30, 4.10, 4.55, 0.34, [
        ("cada campo → resize 224 → CONCH → [512]", 9.5, False, GRIS_TXT, F_MONO)])
    # derecha: crop REAL a dos escalas
    add_image_fit(s, MS_FINE, 5.10, 1.02, 2.25, 2.25, align="top")
    add_image_fit(s, MS_CTX, 7.45, 1.02, 2.25, 2.25, align="top")
    caption(s, 5.10, 3.30, 2.25, "fino 112 µm · citología", size=9.5, col=ORA_T, bold=True)
    caption(s, 7.45, 3.30, 2.25, "contexto 512 µm · arquitectura", size=9.5, col=TEAL_TITLE, bold=True)
    add_textbox(s, 5.10, 3.66, 4.6, 0.62, [
        ("Región real de mama (TCGA-BRCA), mismo centro. La caja marca el campo fino dentro "
         "del contexto.", 10, False, GRIS_BODY, F_BODY, PP_ALIGN.CENTER)])
    takeaway_bar(s, "Fusión por promedio → un token [N,512] → CLAM_MB intacto, la comparación más "
                    "limpia.", t=4.50, size=12.5)
    notes(s, "Llego a la decisión concreta, y es donde pido guía. La propuesta son dos escalas. "
             "Una fina, de ciento doce micras de campo, cerca de veinte magnificaciones, que es "
             "justo donde el encoder fue entrenado: ahí se detecta la calcificación y su forma. Y "
             "una de contexto, de quinientas doce micras, cerca de cinco magnificaciones, que "
             "abarca el conducto o el lobulillo anfitrión, que es literalmente lo que la etiqueta "
             "pregunta. A la derecha está una región real de una lámina de mama: el mismo centro "
             "visto a las dos escalas. En el campo fino se ve la citología; en el de contexto "
             "aparece la arquitectura glandular, y la caja marca dónde cae el parche fino dentro "
             "del grueso. Como las cohortes están a distinta resolución, el tamaño en píxeles se "
             "calcula por lámina para dar en el mismo campo físico: en la pública el crop fino es "
             "de cuatrocientos ochenta y dos píxeles y el de contexto dos mil doscientos; en la "
             "privada, doscientos cuarenta y uno y mil ciento uno; todos se reescalan a doscientos "
             "veinticuatro antes del encoder. Los dos vectores se promedian en uno solo, así el "
             "agregador queda intacto y la comparación es limpia. Las dos escalas y la fusión por "
             "promedio quedan como la propuesta a discutir; los campos son ajustables.")

    # Re-base al template de Sebastián: escalar el deck terminado a 13.333x7.5.
    reflow_onco(prs, skip=keep_ids)
    scale_deck_to_1610(prs, skip=keep_ids)
    os.makedirs(OUT_DIR, exist_ok=True)
    prs.save(DST)
    print("Guardado:", DST, "·", len(prs.slides), "slides ·", "13.333x7.5")


if __name__ == "__main__":
    build()
