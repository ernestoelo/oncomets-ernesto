#!/usr/bin/env python
"""generate_b6_deck.py — deck de la reunión del VIERNES (B6): mammoth / interpretabilidad OBJ-A.

Formato = deck B4 (10x5.625 in, Barlow, paleta teal) — corrige el formato que Benjamín
rechazó del B5 (13.333, Carlito). Contenido = mecanismo mammoth (slides B5 re-formateadas)
+ interpretabilidad OBJ-A (material nuevo del run 30-jun). Eje = ENTENDIMIENTO (no rendimiento).

- Slides de OBJETIVOS/contenido → NATIVAS en Barlow a tamaños B4 (el fix de Benjamín).
- Diagramas de arquitectura → reusados de los .pptx standalone, ESCALADOS x0.75 a 10x5.625
  (mismo aspect ratio; conservan Carlito, convención aceptada para diagramas). Editables.
- Slides de mecanismo (interior/tensor, keep_slots) → NATIVAS con matrices+dims Y código real
  de mammoth.py en paralelo (pedido de Ernesto para la reunión).
- Figuras de interpretabilidad (heatmap/top-k/cross-slide) y del paper → como imagen (excepción).
- Notas del presentador → guion HABLADO en prosa (sin etiquetas de fase, sin nº de job/nombres),
  integrando las explicaciones del estudio (5 bloques de estudio_reunion_viernes.md).

Uso:
  PYTHONPATH=/media/administrador/Storage1/sdonoso/clam_testing2/.pylibs \
  /home/sdonoso/miniconda3/envs/clam_latest/bin/python \
    sprints/B6_sprint6/presentacion_viernes/generate_b6_deck.py
"""
import copy
import io
import os

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn
from PIL import Image

REPO = "/media/administrador/Storage1/sdonoso/clam_testing2/oncomets-ernesto"
PRES = os.path.join(REPO, "papers/presentations")
ASSETS_BRAND = os.path.join(PRES, "assets_branding")
PAPER_FIGS = os.path.join(ASSETS_BRAND, "paper_figs")
OBJA = os.path.join(REPO, "sprints/B5_sprint5/mammoth_entendimiento/interpretabilidad")
OUT_DIR = os.path.join(REPO, "sprints/B6_sprint6/presentacion_viernes")
DST = os.path.join(OUT_DIR, "CLAM_Reunion_Mammoth.pptx")

# --- diagramas standalone (13.333x7.5, sin header) que reusamos escalados ---
DIAG_MAM = os.path.join(PRES, "Diagrama_CLAM_mammoth.pptx")     # s0 pipeline · s1 interior/tensor
DIAG_FUSED = os.path.join(PRES, "Diagrama_mammoth_fused.pptx")  # s0 flujo oficial · s1 keep_slots
# --- figuras ---
FIG1_MAM = os.path.join(PAPER_FIGS, "mammoth_fig1_overview.png")  # t-SNE + barras (paper Fig 1)
FIG3_MAM = os.path.join(PAPER_FIGS, "mammoth_fig3_routing.png")   # ruteo→fenotipo (paper Fig 3)
FIG2_ARCH = os.path.join(PAPER_FIGS, "mammoth_fig2_arch.png")     # arquitectura oficial (Fig 2, 4375x1758)
A14Q = os.path.join(OBJA, "TCGA-E2-A14Q-01Z-00-DX1_cdis_f0")
HEATMAP_MONTAGE = os.path.join(A14Q, "heatmap_montage.png")
TOPK_SUBSET = os.path.join(A14Q, "topk_subset_6experts.png")
XSLIDE_E8 = os.path.join(OBJA, "cross_slide/expert_08_crossslide.png")
XSLIDE_E26 = os.path.join(OBJA, "cross_slide/expert_26_crossslide.png")

# ---- paleta B4 (valores reales del deck) ----
TEAL_TITLE = RGBColor(0x21, 0x75, 0x89)
TEAL_SQ    = RGBColor(0x31, 0x85, 0x9C)
BAR_GRIS   = RGBColor(0xF2, 0xF2, 0xF2)
TEAL_DIV   = RGBColor(0x2E, 0x7E, 0x8F)
LAV_TITLE  = RGBColor(0xCD, 0xD6, 0xF4)
TEAL_SUB   = RGBColor(0xB8, 0xD4, 0xD9)
TEAL_CARD  = RGBColor(0xDD, 0xEA, 0xEE)
TEAL_CARD2 = RGBColor(0xF3, 0xF8, 0xF9)
CODE_BG    = RGBColor(0x1E, 0x2A, 0x2E)   # panel de código (teal muy oscuro)
CODE_FG    = RGBColor(0xE6, 0xEE, 0xF0)   # texto de código
CODE_CMT   = RGBColor(0x8F, 0xB8, 0xC0)   # comentarios de código
GRIS_BODY  = RGBColor(0x59, 0x59, 0x59)
GRIS_TXT   = RGBColor(0x55, 0x55, 0x55)
ORA_T      = RGBColor(0xB4, 0x52, 0x1E)
ORA_ACC    = RGBColor(0xE2, 0x72, 0x3B)
VERDE      = RGBColor(0x1E, 0x84, 0x49)
ROJO       = RGBColor(0xC0, 0x39, 0x2B)
INK        = RGBColor(0x22, 0x22, 0x22)
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)

F_TITLE = "Barlow ExtraBold"
F_BODY  = "Barlow"
F_MONO  = "Consolas"

# --- geometría B4 (10 x 5.625), extraída del volcado real ---
SW, SH = 10.0, 5.625
HDR = 0.785
LOGO = os.path.join(ASSETS_BRAND, "logo_header.png")
PORTADA = os.path.join(ASSETS_BRAND, "portada_fullbleed.jpg")
CHECK_VERDE = os.path.join(ASSETS_BRAND, "check_verde.png")  # ticket "hecho" (reusado del deck B3)
PROG_BG   = RGBColor(0xFD, 0xEF, 0xE6)   # relleno pill "En progreso" (naranja muy claro)

# --- assets sección magnificación multi-escala (B6, reunión Sebastián) ---
MSCROP = os.path.join(ASSETS_BRAND, "multiscale_crop")
MS_FINE = os.path.join(MSCROP, "fine_112um.png")   # crop real TCGA-BRCA campo 112µm (~20×)
MS_CTX = os.path.join(MSCROP, "context_512um.png")  # crop real, mismo centro, campo 512µm (~5×)

_uid = [2000]


# ============================================================================
# Helpers base
# ============================================================================
def _blank(prs):
    for lay in prs.slide_layouts:
        if (lay.name or "").lower().strip() == "blank":
            return lay
    return prs.slide_layouts[6]


def _set_runs(tf, lines, anchor=MSO_ANCHOR.TOP):
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    for i, ln in enumerate(lines):
        txt, sz, bold, col = ln[0], ln[1], ln[2], ln[3]
        font = ln[4] if len(ln) > 4 else F_BODY
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = ln[5] if len(ln) > 5 else PP_ALIGN.LEFT
        p.space_after = Pt(4)
        r = p.add_run(); r.text = txt
        r.font.size = Pt(sz); r.font.bold = bold; r.font.name = font; r.font.color.rgb = col


def add_textbox(slide, l, t, w, h, lines, anchor=MSO_ANCHOR.TOP):
    tb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    _set_runs(tb.text_frame, lines, anchor=anchor)
    return tb


def notes(slide, text):
    """Guion HABLADO en prosa (sin etiquetas de fase, sin nº de job ni nombres)."""
    slide.notes_slide.notes_text_frame.text = text


def _rect(slide, l, t, w, h, color, line=None):
    sp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(l), Inches(t), Inches(w), Inches(h))
    sp.fill.solid(); sp.fill.fore_color.rgb = color
    if line is None:
        sp.line.fill.background()
    else:
        sp.line.color.rgb = line; sp.line.width = Pt(1)
    sp.shadow.inherit = False
    return sp


def logo_mark(slide, size=0.785, logo_h=0.645):
    _rect(slide, 0, 0, size, size, TEAL_SQ)
    slide.shapes.add_picture(LOGO, Inches(0.065), Inches(0.045), height=Inches(logo_h))


def header(slide, title, bar=True, size=26):
    """Header B4: barra gris + cuadrado teal + logo + título Barlow ExtraBold teal."""
    if bar:
        _rect(slide, 0.0, 0.0, SW, HDR, BAR_GRIS)
    logo_mark(slide)
    if title:
        tb = slide.shapes.add_textbox(Inches(0.99), Inches(0.10), Inches(8.9), Inches(0.62))
        _set_runs(tb.text_frame, [(title, size, True, TEAL_TITLE, F_TITLE)], anchor=MSO_ANCHOR.MIDDLE)


def add_image_fit(slide, path, l, t, w, h, align="center"):
    iw, ih = Image.open(path).size
    ar = iw / ih; box_ar = w / h
    if ar > box_ar:
        nw = w; nh = w / ar
    else:
        nh = h; nw = h * ar
    nl = l + (w - nw) / 2
    nt = t + (h - nh) / 2 if align != "top" else t
    slide.shapes.add_picture(path, Inches(nl), Inches(nt), Inches(nw), Inches(nh))


def add_card(slide, l, t, w, h, idx, text, size=14):
    """Tarjeta numerada estilo B4: óvalo naranja + texto Barlow."""
    fill = TEAL_CARD if idx % 2 == 0 else TEAL_CARD2
    sp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(l), Inches(t), Inches(w), Inches(h))
    sp.fill.solid(); sp.fill.fore_color.rgb = fill
    sp.line.color.rgb = TEAL_SQ; sp.line.width = Pt(1.25); sp.shadow.inherit = False
    cd = 0.44
    circ = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(l + 0.16), Inches(t + (h - cd) / 2),
                                  Inches(cd), Inches(cd))
    circ.fill.solid(); circ.fill.fore_color.rgb = ORA_ACC
    circ.line.fill.background(); circ.shadow.inherit = False
    _set_runs(circ.text_frame, [(str(idx), 15, True, WHITE, F_TITLE, PP_ALIGN.CENTER)],
              anchor=MSO_ANCHOR.MIDDLE)
    tb = slide.shapes.add_textbox(Inches(l + 0.74), Inches(t), Inches(w - 0.88), Inches(h))
    _set_runs(tb.text_frame, [(text, size, True, INK, F_BODY)], anchor=MSO_ANCHOR.MIDDLE)


def caption(slide, l, t, w, text, size=11, col=GRIS_TXT, align=PP_ALIGN.CENTER, bold=False):
    add_textbox(slide, l, t, w, 0.4, [(text, size, bold, col, F_BODY, align)])


def status_done(slide, cx, cy, size=0.42):
    """Ticket verde 'hecho/cerrado' (icono reusado del deck B3), centrado en (cx, cy)."""
    slide.shapes.add_picture(CHECK_VERDE, Inches(cx - size / 2), Inches(cy - size / 2),
                             Inches(size), Inches(size))


def status_progress(slide, cx, cy, w=1.5, h=0.44):
    """Pill 'En progreso' (naranja = 'lo abierto/nuevo' del deck B6), centrada en (cx, cy)."""
    sp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(cx - w / 2), Inches(cy - h / 2),
                                Inches(w), Inches(h))
    sp.fill.solid(); sp.fill.fore_color.rgb = PROG_BG
    sp.line.color.rgb = ORA_ACC; sp.line.width = Pt(1.25); sp.shadow.inherit = False
    _set_runs(sp.text_frame, [("En progreso", 11, True, ORA_T, F_BODY, PP_ALIGN.CENTER)],
              anchor=MSO_ANCHOR.MIDDLE)


def takeaway_bar(slide, text, t=4.85, col=TEAL_TITLE, size=14):
    """Barra de remate al pie (línea teal + frase centrada)."""
    _rect(slide, 0.35, t, SW - 0.7, 0.02, TEAL_SQ)
    add_textbox(slide, 0.35, t + 0.08, SW - 0.7, 0.62,
                [(text, size, True, col, F_BODY, PP_ALIGN.CENTER)], anchor=MSO_ANCHOR.MIDDLE)


def dim_pipeline(slide, blocks, t, h=0.62, bw=1.06, gap=0.25, arrow_sz=15):
    """Pipeline al pie: fila de bloques (variable + dimensión) conectados por flechas →.
    `blocks` = lista de (variable, dimensión). Centrado horizontalmente."""
    n = len(blocks)
    total = n * bw + (n - 1) * gap
    x = (SW - total) / 2
    for i, (var, dim) in enumerate(blocks):
        sp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(t), Inches(bw), Inches(h))
        sp.fill.solid(); sp.fill.fore_color.rgb = TEAL_CARD2 if i % 2 else TEAL_CARD
        sp.line.color.rgb = TEAL_SQ; sp.line.width = Pt(1.0); sp.shadow.inherit = False
        lines = [(var, 12, True, TEAL_TITLE, F_TITLE, PP_ALIGN.CENTER)]
        if dim:
            lines.append((dim, 9.5, False, INK, F_MONO, PP_ALIGN.CENTER))
        _set_runs(sp.text_frame, lines, anchor=MSO_ANCHOR.MIDDLE)
        for p in sp.text_frame.paragraphs:
            p.space_after = Pt(0)
        if i < n - 1:
            ar = slide.shapes.add_textbox(Inches(x + bw), Inches(t), Inches(gap), Inches(h))
            _set_runs(ar.text_frame, [("→", arrow_sz, True, TEAL_SQ, F_BODY, PP_ALIGN.CENTER)],
                      anchor=MSO_ANCHOR.MIDDLE)
        x += bw + gap


def code_panel(slide, l, t, w, h, lines, title=None):
    """Panel de código nativo: fondo oscuro + texto monoespaciado (Consolas).
    `lines` = lista de strings; los que empiezan con '#' se pintan como comentario."""
    _rect(slide, l, t, w, h, CODE_BG)
    tb = slide.shapes.add_textbox(Inches(l + 0.12), Inches(t + 0.08), Inches(w - 0.24), Inches(h - 0.16))
    tf = tb.text_frame; tf.word_wrap = True; tf.vertical_anchor = MSO_ANCHOR.TOP
    runs = []
    if title:
        runs.append((title, 10, True, CODE_CMT, F_MONO))
    for ln in lines:
        col = CODE_CMT if ln.strip().startswith("#") else CODE_FG
        runs.append((ln if ln else " ", 9.5, False, col, F_MONO))
    _set_runs(tf, runs)
    for p in tf.paragraphs:
        p.space_after = Pt(1)
    return tb


def dims_table(slide, l, t, w, rows, row_h=0.345, hdr=("paso", "forma del tensor")):
    """Tabla nativa de dimensiones (2 columnas: descripción del paso · forma)."""
    n = len(rows) + 1
    gtbl = slide.shapes.add_table(n, 2, Inches(l), Inches(t), Inches(w), Inches(row_h * n)).table
    gtbl.columns[0].width = Inches(w * 0.60)
    gtbl.columns[1].width = Inches(w * 0.40)
    gtbl.first_row = False; gtbl.horz_banding = False
    data = [hdr] + list(rows)
    for ri, (c0, c1) in enumerate(data):
        for ci, txt in enumerate((c0, c1)):
            cell = gtbl.cell(ri, ci)
            cell.margin_left = Inches(0.06); cell.margin_right = Inches(0.04)
            cell.margin_top = Inches(0.01); cell.margin_bottom = Inches(0.01)
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            if ri == 0:
                cell.fill.solid(); cell.fill.fore_color.rgb = TEAL_SQ
                col, bold, sz, fnt = WHITE, True, 9.5, F_BODY
            else:
                cell.fill.solid(); cell.fill.fore_color.rgb = TEAL_CARD2 if ri % 2 else TEAL_CARD
                col = INK if ci == 0 else TEAL_TITLE
                bold = ci == 1; sz = 9.5; fnt = F_BODY if ci == 0 else F_MONO
            tf = cell.text_frame; tf.word_wrap = True
            p = tf.paragraphs[0]; p.alignment = PP_ALIGN.LEFT
            r = p.add_run(); r.text = txt
            r.font.size = Pt(sz); r.font.bold = bold; r.font.name = fnt; r.font.color.rgb = col
    return gtbl


def simple_table(slide, l, t, w, headers, rows, col_fracs, row_h=0.32, fs=9.5):
    """Tabla nativa genérica (n columnas): header teal + filas con banding teal claro."""
    ncol = len(headers); nrow = len(rows) + 1
    tbl = slide.shapes.add_table(nrow, ncol, Inches(l), Inches(t), Inches(w),
                                 Inches(row_h * nrow)).table
    for ci, fr in enumerate(col_fracs):
        tbl.columns[ci].width = Inches(w * fr)
    tbl.first_row = False; tbl.horz_banding = False
    for ri, row in enumerate([headers] + list(rows)):
        for ci, txt in enumerate(row):
            cell = tbl.cell(ri, ci)
            cell.margin_left = Inches(0.06); cell.margin_right = Inches(0.04)
            cell.margin_top = Inches(0.015); cell.margin_bottom = Inches(0.015)
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            if ri == 0:
                cell.fill.solid(); cell.fill.fore_color.rgb = TEAL_SQ
                col, bold, fnt = WHITE, True, F_BODY
            else:
                cell.fill.solid(); cell.fill.fore_color.rgb = TEAL_CARD2 if ri % 2 else TEAL_CARD
                col, bold, fnt = INK, (ci == 0), F_BODY
            tf = cell.text_frame; tf.word_wrap = True
            p = tf.paragraphs[0]; p.alignment = PP_ALIGN.LEFT
            r = p.add_run(); r.text = txt
            r.font.size = Pt(fs); r.font.bold = bold; r.font.name = fnt; r.font.color.rgb = col
    return tbl


def panel(slide, l, t, w, h, title, tcol, lines, border, fill=TEAL_CARD2):
    """Panel con título + líneas de cuerpo (para 'dos tareas / demandas opuestas')."""
    sp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(l), Inches(t), Inches(w), Inches(h))
    sp.fill.solid(); sp.fill.fore_color.rgb = fill
    sp.line.color.rgb = border; sp.line.width = Pt(1.5); sp.shadow.inherit = False
    tb = slide.shapes.add_textbox(Inches(l + 0.18), Inches(t + 0.12), Inches(w - 0.36), Inches(h - 0.24))
    runs = [(title, 14.5, True, tcol, F_TITLE)] + [(ln, 12, False, INK, F_BODY) for ln in lines]
    _set_runs(tb.text_frame, runs)
    for p in tb.text_frame.paragraphs:
        p.space_after = Pt(3)


def nested_fields(slide, l, t, w, h):
    """Esquema nativo: dos campos concéntricos (contexto 512µm ⊃ fino 112µm) + leyenda."""
    _rect(slide, l, t, w, h, TEAL_CARD2, line=TEAL_SQ)
    side = h - 0.55
    cx = l + w * 0.30; cy = t + h / 2
    ctx = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(cx - side / 2), Inches(cy - side / 2),
                                 Inches(side), Inches(side))
    ctx.fill.background(); ctx.line.color.rgb = TEAL_SQ; ctx.line.width = Pt(2.5)
    ctx.shadow.inherit = False
    fside = side * 0.34                         # esquema (ratio real 112/512≈0.22, lo muestra la foto)
    fin = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(cx - fside / 2), Inches(cy - fside / 2),
                                 Inches(fside), Inches(fside))
    fin.fill.solid(); fin.fill.fore_color.rgb = ORA_ACC
    fin.line.color.rgb = ORA_T; fin.line.width = Pt(1.5); fin.shadow.inherit = False
    lx = l + w * 0.58
    add_textbox(slide, lx, t + 0.22, w * 0.42 - 0.12, h - 0.4, [
        ("Contexto  512 µm · ~5×", 12.5, True, TEAL_TITLE, F_TITLE),
        ("el conducto/lobulillo anfitrión", 10.5, False, GRIS_BODY, F_BODY),
        (" ", 8, False, INK, F_BODY),
        ("Fino  112 µm · ~20×", 12.5, True, ORA_T, F_TITLE),
        ("la calcificación + su forma", 10.5, False, GRIS_BODY, F_BODY),
    ], anchor=MSO_ANCHOR.MIDDLE)


# ============================================================================
# Arquetipos B4
# ============================================================================
def portada(prs):
    """Portada B4: imagen full-bleed limpia, sin overlay (el título va horneado en la imagen)."""
    s = prs.slides.add_slide(_blank(prs))
    s.shapes.add_picture(PORTADA, Inches(0), Inches(-0.02), Inches(SW), Inches(SH + 0.04))
    return s


def divider(prs, title, subtitle):
    s = prs.slides.add_slide(_blank(prs))
    s.background.fill.solid(); s.background.fill.fore_color.rgb = TEAL_DIV
    s.shapes.add_picture(LOGO, Inches(0.42), Inches(0.36), height=Inches(0.62))
    add_textbox(s, 0.8, 2.05, SW - 1.6, 1.1,
                [(title, 44, True, LAV_TITLE, F_TITLE, PP_ALIGN.CENTER)], anchor=MSO_ANCHOR.MIDDLE)
    add_textbox(s, 0.8, 3.25, SW - 1.6, 0.7,
                [(subtitle, 18, False, TEAL_SUB, F_BODY, PP_ALIGN.CENTER)], anchor=MSO_ANCHOR.TOP)
    return s


def content(prs, title, bar=True, size=26):
    s = prs.slides.add_slide(_blank(prs))
    header(s, title, bar=bar, size=size)
    return s


def _scale_el(el, s):
    """Escala geométrica (off/ext) y tipográfica (sz) de un shape XML por factor s."""
    for off in el.iter(qn("a:off")):
        off.set("x", str(int(int(off.get("x")) * s)))
        off.set("y", str(int(int(off.get("y")) * s)))
    for ext in el.iter(qn("a:ext")):
        ext.set("cx", str(int(int(ext.get("cx")) * s)))
        ext.set("cy", str(int(int(ext.get("cy")) * s)))
    for tag in ("a:rPr", "a:defRPr", "a:endParaRPr"):
        for rpr in el.iter(qn(tag)):
            sz = rpr.get("sz")
            if sz:
                rpr.set("sz", str(max(100, int(int(sz) * s))))


def copy_diagram_scaled(prs, src_path, idx, title=None, scale=0.75, bar=False):
    """Copia el diagrama standalone (13.333x7.5) escalado a 10x5.625 (x0.75, full-bleed).
    Conserva nativo/editable + remapea imágenes embebidas. Header opcional (logo only por defecto)."""
    s = prs.slides.add_slide(_blank(prs))
    src = Presentation(src_path)
    src_slide = src.slides[idx]
    spTree = s.shapes._spTree
    for shp in src_slide.shapes:
        el = copy.deepcopy(shp._element)
        for cNvPr in el.iter(qn("p:cNvPr")):
            _uid[0] += 1; cNvPr.set("id", str(_uid[0]))
        for blip in el.iter(qn("a:blip")):
            r_embed = blip.get(qn("r:embed"))
            if r_embed:
                img_part = src_slide.part.related_part(r_embed)
                _, new_rid = s.part.get_or_add_image_part(io.BytesIO(img_part.blob))
                blip.set(qn("r:embed"), new_rid)
        _scale_el(el, scale)
        spTree.append(el)
    header(s, title, bar=bar)  # logo mark (+ barra/título si se pide)
    return s


# ============================================================================
# Build
# ============================================================================
def build():
    prs = Presentation()
    prs.slide_width = Inches(SW); prs.slide_height = Inches(SH)

    # ---- 1. Portada (limpia, sin overlay) ----
    s = portada(prs)
    notes(s, "El objetivo de esta reunión no es demostrar que MAMMOTH mejore la métrica: esa "
             "pregunta ya quedó cerrada, y la respuesta fue que no. Lo que sigue abierto, y es de "
             "lo que trata esta presentación, es distinto: qué aprende MAMMOTH por dentro y en qué "
             "se fija cada experto. El recorrido tiene dos partes. Primero el mecanismo, esta vez "
             "explicado con calma y hasta el fondo. Después, la evidencia visual de qué morfología "
             "reclama cada experto, mirada sobre las propias slides de mama del proyecto.")

    # ---- 2. Recapitulación de objetivos (MISMO formato que B4 slide 2:
    #        título 32pt + lista numerada en un solo cuadro, 24pt Barlow bold gris) ----
    s = content(prs, "Recapitulación de objetivos", size=32)
    # Objetivos del eje (pedido de Benjamín, 29-jun): entender el mecanismo + interpretar.
    # Enunciados en infinitivo (sin 1ª persona), concisos, sin resultados.
    # (texto, estado): "done" = ticket verde (cerrado) · "prog" = pill "En progreso" (abierto)
    recap = [
        ("1. Dominar el mecanismo de MAMMOTH y explicarlo con una analogía simple.", "done"),
        ("2. Precisar qué es una cabeza y el tensor de prototipos (30×16×10×16).", "done"),
        ("3. Distinguir MoE de PoE y situar el número de cabezas para mama.", "done"),
        ("4. Interpretar los expertos: qué región y qué morfología reclama cada uno.", "prog"),
    ]
    row_tops = [1.16, 2.18, 3.20, 4.22]   # 4 filas, pitch 1.02; tipografía B4 (24pt Barlow)
    row_h = 0.90
    for (it, st), rt in zip(recap, row_tops):
        add_textbox(s, 0.35, rt, 7.75, row_h, [(it, 24, True, GRIS_BODY, F_BODY)],
                    anchor=MSO_ANCHOR.MIDDLE)
        cy = rt + row_h / 2
        if st == "done":
            status_done(s, 8.98, cy)
        else:
            status_progress(s, 8.98, cy)
    notes(s, "Antes de entrar en el mecanismo conviene fijar qué se propuso entender, porque es "
             "el eje de la reunión. El primer objetivo era dominar el mecanismo de MAMMOTH a fondo "
             "y poder explicarlo con una analogía simple. El segundo, precisar dos puntos que "
             "habían quedado sin respuesta: qué es exactamente una cabeza y cómo se lee el tensor "
             "de prototipos, el treinta por dieciséis por diez por dieciséis. El tercero, "
             "distinguir una mezcla de expertos de un producto de expertos, y ubicar cuántas "
             "cabezas convienen para mama. Esos tres quedaron cerrados. El cuarto abre la segunda "
             "parte y sigue en progreso: interpretar a los expertos, es decir, mirar sobre las "
             "propias slides qué región reclama cada uno y qué morfología hay ahí; falta el visto "
             "bueno de un patólogo.")

    # ---- 3. Divisoria: MAMMOTH ----
    s = divider(prs, "MAMMOTH", "Mixture-of-Experts en la primera capa de CLAM (patch-embed)")
    notes(s, "El mecanismo primero. MAMMOTH es una intervención muy acotada: cambia una sola "
             "parte de CLAM, la primera capa, la que proyecta los parches al espacio interno del "
             "modelo. Todo el resto de CLAM queda intacto.")

    # ---- 4. Qué es y por qué (título = acrónimo completo; 3 tarjetas + Fig 1 + Fig 3) ----
    s = content(prs, "MAMMOTH — MAtrix-factorized Mixture Module of Transformation Heads", size=17)
    cards = [
        "Reemplaza la 1ª capa lineal de CLAM por una mezcla de expertos (MoE).",
        "Un router reparte cada parche entre expertos especializados por morfología.",
        "Objetivo: bajar la interferencia de gradientes entre parches distintos.",
    ]
    cy = 1.05
    for i, txt in enumerate(cards):
        add_card(s, 0.30, cy, 5.35, 0.95, i + 1, txt, size=13)
        cy += 1.10
    add_image_fit(s, FIG1_MAM, 5.95, 0.98, 3.85, 1.85, align="top")
    caption(s, 5.95, 2.86, 3.85,
            "Fig. 1 — el espacio interno pasa de una nube continua a grupos por experto,\n"
            "y mejora a todos los agregadores MIL (Shao et al., ICLR 2026)", size=8)
    add_image_fit(s, FIG3_MAM, 5.95, 3.42, 3.85, 1.42, align="top")
    caption(s, 5.95, 4.86, 3.85,
            "Fig. 3 — cada color es el ruteo de un experto sobre la slide; abajo, la "
            "morfología que resume (tumor, estroma, alvéolos…)", size=8)
    notes(s, "El nombre es un acrónimo y cada parte nombra un paso del modelo: cabezas de "
             "transformación, una mezcla de expertos y una factorización de matrices de bajo "
             "rango. La idea de fondo es simple. En CLAM, una sola capa lineal proyecta todos los "
             "parches al espacio interno, y esa única matriz tiene que servir a la vez para "
             "epitelio, estroma y ductos, así que queda en un punto intermedio, mediocre para "
             "todos. El paper lo compara con un traductor obligado a traducir chino, árabe y ruso "
             "con la misma plantilla. MAMMOTH reemplaza esa matriz por treinta expertos "
             "especializados y un router que decide, para cada parche, cuánto mandarlo a cada "
             "experto. Las dos figuras de la derecha son la evidencia del paper. Arriba, el "
             "espacio interno pasa de una nube continua a grupos separados, uno por experto, y la "
             "mejora aparece en todos los agregadores MIL. Abajo, cada color es el ruteo de un "
             "experto sobre la lámina, con la morfología que resume justo debajo. Esa segunda "
             "figura es exactamente el análisis que la segunda parte reproduce sobre las slides de "
             "mama del proyecto.")

    # ---- 5. Diagrama: pipeline CLAM + punto de integración (reusado, escalado) ----
    s = copy_diagram_scaled(prs, DIAG_MAM, 0, title=None, bar=False)
    notes(s, "Sobre el pipeline completo de CLAM el mecanismo se ubica mejor. La slide entra "
             "troceada en parches; el encoder CONCH le asigna a cada parche un vector de "
             "quinientos doce; de ahí siguen la atención y el clasificador. El único bloque que "
             "cambia es el naranja: ahí MAMMOTH reemplaza la primera capa lineal. Todo lo que "
             "viene después es CLAM tal cual.")

    # ---- 6. Interior de MAMMOTH: matrices+dims Y código en paralelo (NATIVO) ----
    s = content(prs, "Por dentro de MAMMOTH: dimensiones y código")
    add_textbox(s, 0.30, 0.86, 4.7, 0.3, [("El tensor, paso a paso", 13, True, TEAL_TITLE, F_TITLE)])
    dims_table(s, 0.30, 1.20, 4.72, [
        ("z · features CONCH (parche)", "[N, 512]"),
        ("q = LN(Wq · z)", "[N, 256]"),
        ("split en 16 cabezas", "[N, 16, 16]"),
        ("S · prototipos (slot_embeds)", "[30,16,10,16]"),
        ("logits de ruteo ⟨q, S⟩", "[N,30,16,10]"),
        ("D = softmax sobre N (dispatch)", "[N,30,16,10]"),
        ("u · slots (media ponderada)", "[30,16,10,16]"),
        ("expertos LoRA + concat cabezas", "[300, 512]"),
        ("combine → salida drop-in", "[N, 512]"),
    ], row_h=0.30)
    add_textbox(s, 0.30, 4.24, 4.72, 0.44, [
        ("El 16 aparece dos veces: nº de cabezas y dimensión de cada prototipo (256/16), para "
         "compararlos con un producto interno. 30 expertos × 10 slots = 300.",
         8.5, False, GRIS_BODY, F_BODY)])
    add_textbox(s, 5.20, 0.86, 4.6, 0.3, [("El código real (mammoth.py)", 13, True, TEAL_TITLE, F_TITLE)])
    code_panel(s, 5.20, 1.20, 4.60, 3.00, [
        "q = norm(wq(z))              # [N,512]->[N,256]",
        'q = rearrange(q,             # -> [N,16,16]',
        '      "n (h d)->n h d", h=16)',
        'logits = einsum("n h d, e h s d',
        '                 -> n e h s", q, S)',
        "                             # S=[30,16,10,16]",
        "dispatch = softmax(logits, dim=N)",
        'u = einsum("n h d, n e h s',
        '            -> e h s d", q, dispatch)',
        "out = expert_heads(u)        # LoRA -> [300,512]",
        "# drop-in (keep_slots=False):",
        'out = einsum("h p d, n h p -> n h d",',
        "             out, combine)   # -> [N,512]",
    ])
    takeaway_bar(s, "Las cabezas son subespacios aprendidos (multi-head), no textura/forma/color: "
                    "la semántica de tejido vive en los slots.", t=4.80, size=12)
    notes(s, "El interior de MAMMOTH, con las dimensiones a la izquierda y el código real a la "
             "derecha, para seguir cada paso sin ambigüedad. La entrada es un parche, un vector de "
             "quinientos doce, la z. Primero se proyecta a un query de doscientos cincuenta y seis "
             "y se parte en dieciséis cabezas de dieciséis. En paralelo, el modelo guarda sus "
             "prototipos aprendidos en el tensor treinta por dieciséis por diez por dieciséis, y "
             "este es el número que conviene dejar claro de una vez. Se lee de izquierda a "
             "derecha: treinta expertos, cada experto con dieciséis cabezas, cada cabeza con diez "
             "slots, y cada slot es un vector de dieciséis dimensiones. El dieciséis aparece dos "
             "veces por una razón concreta. El primero es el número de cabezas; el último es la "
             "dimensión de cada prototipo, que sale de dividir el query, doscientos cincuenta y "
             "seis, entre las dieciséis cabezas. Tienen que coincidir porque el ruteo compara cada "
             "prototipo con el query mediante un producto interno, y un producto interno exige que "
             "ambos vectores vivan en el mismo espacio de dieciséis. Treinta expertos por diez "
             "slots dan los trescientos slots que reaparecen más adelante. De ahí en adelante la "
             "tabla y el código van paso a paso: el ruteo compara query contra prototipos, una "
             "softmax reparte cada parche entre los slots, con esos pesos se arma el promedio "
             "ponderado que llena cada slot, cada experto lo transforma con una operación de bajo "
             "rango y se concatenan las cabezas para dar trescientos por quinientos doce. En la "
             "variante base, una segunda softmax recombina los trescientos slots y reconstruye los "
             "parches, así que la salida tiene la misma forma que tendría una capa lineal; por eso "
             "el reemplazo es directo. Y un punto clave, que en su momento se respondió mal y hay "
             "que dejar afinado: las cabezas no son textura, forma ni color. Son subespacios "
             "aprendidos, igual que en atención multi-cabeza; la semántica de tejido no vive en "
             "las cabezas, vive en los slots.")

    # ---- 7. La arquitectura oficial del paper: figura GRANDE y limpia (sin logo, sin título,
    #        sin callouts encima que tapen las variables de la figura). El trazo al pie y las
    #        notas se refieren a las VARIABLES DE LA PROPIA FIGURA (W, x̄, s, Φ, W_low, ...). ----
    s = prs.slides.add_slide(_blank(prs))
    iw = 9.42; ih = iw / 2.489                      # aspect real de la Fig 2 (4375x1758)
    s.shapes.add_picture(FIG2_ARCH, Inches((SW - iw) / 2), Inches(0.26), Inches(iw), Inches(ih))
    caption(s, 0.35, 0.26 + ih + 0.05, SW - 0.7,
            "Dimensiones por bloque (variables de la figura)  ·  N = nº de parches de la slide",
            size=9.5, bold=True, col=TEAL_TITLE)
    # pipeline en bloques: variable de la figura + forma del tensor en cada paso
    dim_pipeline(s, [
        ("x_i", "[N,512]"),
        ("W → x̄", "[N,16,16]"),
        ("ruteo", "sim + softmax"),
        ("slots s", "300"),
        ("Φ·W_low", "[300,512]"),
        ("concat", "[N,512]"),
        ("CLAM", "logits"),
    ], t=0.26 + ih + 0.34)
    notes(s, "Esta es la arquitectura completa del paper, y conviene recorrerla de izquierda a "
             "derecha siguiendo sus variables, porque cada símbolo es un paso del pipeline. Se "
             "parte de la lámina, que se corta en parches, y cada parche pasa por el encoder y "
             "sale convertido en un vector de features, que en la figura es x_i. Hay uno por "
             "parche, así que para toda la slide tenemos el conjunto de x_i, de uno hasta N. En un "
             "MIL corriente cada x_i entraría por una única matriz lineal, la W del diagrama, que "
             "lo proyecta al espacio interno del agregador, y ese es el único lugar donde MAMMOTH "
             "interviene: reemplaza esa W. Todo lo de la izquierda, el encoder, y todo lo de la "
             "derecha, el agregador, quedan igual. Lo primero que hace MAMMOTH es partir esa "
             "proyección en cabezas. La salida de W se corta en trozos y a cada trozo la figura lo "
             "llama x̄, la porción del parche que ve cada cabeza; son subespacios aprendidos, como "
             "en atención multi-cabeza, no textura ni forma ni color. Cada cabeza corre su propia "
             "mezcla de expertos en paralelo, que son los bloques MoE. Dentro de cada experto "
             "están los slots, que son prototipos aprendidos, y acá viene la parte que quiero "
             "dejar clarísima, porque es el corazón del mecanismo. Cada parche se compara con cada "
             "prototipo con un producto interno, y eso es el ruteo: mide qué tan parecido es el "
             "parche a ese slot. Sobre todos los parches de la slide se aplica una softmax, que "
             "decide cuánto aporta cada parche a cada slot. Con esos pesos se arma un promedio "
             "ponderado de los parches, y ese promedio ponderado, que en la figura es el weighted "
             "average, es lo que llena cada slot. Dicho de una vez y en orden: el parche entra "
             "como x_i, el ruteo calcula los pesos comparándolo con los prototipos, y el slot "
             "termina siendo la x ponderada, el promedio de todos los parches que se le parecen. "
             "Cada slot queda representando un fenotipo, un tipo de tejido. Cada slot pasa después "
             "por la transformación del experto, que está factorizada en bajo rango, la Φ por la "
             "W_low seguida de una no linealidad; esa factorización es lo que le permite tener "
             "treinta expertos sin gastar más parámetros que la matriz original. Al final se "
             "concatenan las cabezas, el cross-head concat, y sale la salida de MAMMOTH, que "
             "alimenta al agregador. En este proyecto el agregador es CLAM: su atención junta los "
             "parches en un vector de slide y el clasificador entrega el diagnóstico. Un detalle: "
             "la figura muestra como salida los propios slots, y en nuestra configuración base una "
             "segunda softmax los recombina para reconstruir los parches y devolver la misma forma "
             "que entró, así el reemplazo de la capa lineal es directo. Casi seguro va a salir por "
             "qué esto es una mezcla de expertos y no un producto de expertos, y la respuesta está "
             "en cómo se combinan. Acá los expertos se combinan sumando: las dos softmax que "
             "describí reparten pesos que suman uno, así que la salida es un promedio, una "
             "combinación convexa. Eso es una mezcla, se entrena directo y es estable. Un producto "
             "de expertos haría lo contrario: multiplicaría las distribuciones de los expertos, "
             "como si cada uno tuviera un veto, y para entrenarlo haría falta un normalizador que "
             "es intratable; además modela una distribución de probabilidad, no una transformación "
             "de features, que es lo que acá necesitamos. Todo el diseño es aditivo por "
             "construcción, así que la mezcla es la elección natural y un producto rompería la "
             "formulación y traería de vuelta justo la inestabilidad que el método busca evitar. Y "
             "conviene decirlo con honestidad: el paper no menciona el producto de expertos, esta "
             "comparación es un razonamiento de arquitectura, no una cita del paper.")

    # ---- 8. keep_slots: la bifurcación, con math+código (NATIVO) ----
    s = content(prs, "La variante keep_slots: dónde cambia la salida")
    # tronco compartido
    trunk = _rect(s, 0.30, 0.98, 9.4, 0.66, TEAL_CARD, line=TEAL_SQ)
    add_textbox(s, 0.46, 1.00, 9.1, 0.62, [
        ("Tronco compartido — idéntico a la slide anterior:  z → q → ruteo → expertos  →  "
         "out = concat de expertos  [300, 512]", 12.5, True, INK, F_BODY)], anchor=MSO_ANCHOR.MIDDLE)
    # rama izquierda: keep_slots=False
    add_textbox(s, 0.30, 1.80, 4.6, 0.3, [("keep_slots = False  ·  drop-in (base)", 12.5, True, TEAL_TITLE, F_TITLE)])
    code_panel(s, 0.30, 2.14, 4.6, 1.5, [
        "# 2º softmax recombina los 300 slots",
        'out = einsum("h p d, n h p -> n h d",',
        "             out, combine)",
        'out = rearrange(out,"n h d->n (h d)")',
        "                          # -> [N, 512]",
    ])
    add_textbox(s, 0.30, 3.72, 4.6, 0.95, [
        ("Reconstruye los N parches (cardinalidad N → N). CLAM agrega sobre los N parches, "
         "igual que la capa lineal original.", 12, False, GRIS_BODY, F_BODY)])
    # rama derecha: keep_slots=True
    add_textbox(s, 5.10, 1.80, 4.6, 0.3, [("keep_slots = True  ·  variante nueva", 12.5, True, ORA_T, F_TITLE)])
    code_panel(s, 5.10, 2.14, 4.6, 1.5, [
        "# se salta la recombinación",
        "return out                # -> [300, 512]",
        "#     300 slot-tokens como salida",
        "#     (+ slot_dropout opcional)",
    ])
    add_textbox(s, 5.10, 3.72, 4.6, 0.95, [
        ("Se queda con los 300 slots (cuello de botella aprendido N → 300). CLAM agrega sobre "
         "los 300 slots, no sobre los parches.", 12, False, GRIS_BODY, F_BODY)])
    takeaway_bar(s, "Misma cabeza CLAM y misma pérdida — cambia solo qué se agrega. La "
                    "interpretabilidad se calcula antes de la bifurcación → vale para las dos.",
                 t=4.78, size=12)
    notes(s, "Hay una variante, keep_slots, y lo relevante es dónde cambia. Hasta la parte de "
             "arriba todo es idéntico a lo anterior: proyección, ruteo, expertos, y quedan los "
             "trescientos slots concatenados. La bifurcación es una sola línea de código. En la "
             "versión base, keep_slots en falso, una segunda softmax recombina los trescientos "
             "slots y reconstruye los parches, así que CLAM sigue agregando sobre los parches y el "
             "reemplazo de la capa lineal es directo. En la versión nueva, keep_slots en "
             "verdadero, esa recombinación se omite y la salida son los trescientos slots; "
             "entonces CLAM agrega sobre trescientos tokens en vez de sobre los parches. Es una "
             "variante orientada a rendimiento, y para la reunión de hoy no cambia nada, porque la "
             "interpretabilidad se calcula antes de esta bifurcación y vale para las dos ramas. La "
             "cabeza de CLAM y la pérdida son las mismas; lo único que cambia es qué se agrega.")

    # ---- 9. Divisoria: Interpretabilidad (OBJ-A) ----
    s = divider(prs, "¿Qué mira cada experto?",
                "Interpretabilidad post-hoc sobre un checkpoint entrenado · 4 slides TCGA-BRCA")
    notes(s, "Con el mecanismo claro, empieza la segunda parte: mirar, sobre las propias slides "
             "del proyecto, qué región reclama cada experto y qué morfología hay ahí. Es un "
             "análisis post-hoc, en CPU, sobre un modelo ya entrenado; no reentrena nada.")

    # ---- 10. Ruteo espacial (dónde) + morfología top-k (qué) — FUSIÓN ----
    s = content(prs, "Qué mira cada experto: dónde y qué morfología")
    add_image_fit(s, HEATMAP_MONTAGE, 0.30, 0.98, 4.95, 3.35, align="top")
    caption(s, 0.30, 4.38, 4.95, "Ruteo espacial: los 30 expertos encienden zonas distintas", size=9)
    add_image_fit(s, TOPK_SUBSET, 5.42, 0.98, 2.55, 3.05, align="top")
    add_textbox(s, 8.05, 1.05, 1.80, 3.2, [
        ("Top-k a alta resolución:", 11.5, True, INK, F_BODY),
        ("e8 → epitelio tumoral", 12, True, TEAL_TITLE, F_BODY),
        ("e26 → estroma fibroso", 12, True, TEAL_TITLE, F_BODY),
        ("e3 → epitelio ductal", 12, True, TEAL_TITLE, F_BODY),
        ("Color = percentil por experto (estructura relativa, no magnitud de uso).",
         9.5, False, ORA_T, F_BODY),
    ], anchor=MSO_ANCHOR.TOP)
    caption(s, 5.42, 4.12, 2.55, "Morfología que cada experto rutea (Fig 3.2)", size=9)
    takeaway_bar(s, "Heatmap = dónde · top-k = qué morfología. Emergió sin supervisión de tejido "
                    "(el paper lo validó con patólogos).", t=4.82, size=12)
    notes(s, "Estas dos vistas cuentan una sola historia y conviene leerlas juntas. A la "
             "izquierda, cada uno de los treinta cuadros es la misma slide pintada según cuánto el "
             "router manda cada parche a un experto: rojo, mucho; azul, casi nada. Lo que se "
             "observa es que los treinta encienden zonas distintas de la lámina; es decir, la capa "
             "lineal única quedó reemplazada por especialistas que miran regiones diferentes del "
             "tejido. Sobre el color, una precisión honesta: está normalizado por percentil dentro "
             "de cada experto, así que sirve para ver qué regiones prefiere cada uno, no para "
             "decir cuál se usa más; medido aparte, el uso sale casi uniforme, sin expertos "
             "muertos ni acaparadores. Ahora bien, el mapa de calor dice dónde, no qué tejido hay "
             "en esas zonas. Eso lo cierra la derecha: para cada experto se toman los parches que "
             "más rutea y se recortan a alta resolución real. Ahí se ve directo: el experto ocho "
             "mira nidos de epitelio tumoral, el veintiséis mira estroma fibroso rosado, el tres "
             "mira ductos. Es la misma especialización que el paper validó con dos patólogos, y lo "
             "notable es que emerge sola: nadie le indicó al modelo qué es estroma. Los expertos "
             "mixtos que aparecen son esperables, porque el ruteo es suave y reparte cada parche "
             "entre varios.")

    # ---- 11. Morfología ≠ clase (cross-slide) + honestidad/cierre — FUSIÓN ----
    #   imágenes GRANDES lado a lado (grillas 4 slides × 5 parches; aspect 1.19) + texto compacto.
    s = content(prs, "El experto detecta tejido, no clase")
    xh = 3.02; xw = xh * (704 / 593)                 # aspect real de las cross-slide
    gap = 0.34
    x0 = (SW - (2 * xw + gap)) / 2
    s.shapes.add_picture(XSLIDE_E8, Inches(x0), Inches(0.86), Inches(xw), Inches(xh))
    s.shapes.add_picture(XSLIDE_E26, Inches(x0 + xw + gap), Inches(0.86), Inches(xw), Inches(xh))
    caption(s, 0.35, 0.86 + xh + 0.02, SW - 0.7,
            "e8 → epitelio tumoral   ·   e26 → estroma   —   el mismo experto en las 4 slides "
            "(2 positivas, 2 negativas)", size=10.5, bold=True, col=INK)
    add_textbox(s, 0.35, 0.86 + xh + 0.30, SW - 0.7, 0.62, [
        ("«Negativo» = cdis sin microcalcificación, no «sin tumor» → el experto es un detector de "
         "TEJIDO, no de clase (la clase se decide después, en la atención + el clasificador).  "
         "Honestidad: 4 slides · 1 tarea · 1 fold · etiquetas provisionales (falta sign-off de "
         "patólogo).", 10.5, False, GRIS_BODY, F_BODY, PP_ALIGN.CENTER)], anchor=MSO_ANCHOR.TOP)
    takeaway_bar(s, "Que la especialización sea real y aun así no mueva la métrica es la evidencia "
                    "de que el cuello no está en la 1ª capa, sino en el dato.", t=4.82, size=12)
    notes(s, "Esta es la prueba más fina, y con ella cierra la presentación. Como los prototipos "
             "son parámetros compartidos del modelo, el experto ocho es el mismo experto en todas "
             "las slides. Eso permite una prueba limpia: si su especialización es real, tiene que "
             "elegir la misma morfología en todas. Y así ocurre. En las dos grillas, cada fila es "
             "una de cuatro slides y cada columna un parche top del mismo experto: el ocho "
             "enciende epitelio en las cuatro, el veintiséis enciende estroma en las cuatro, "
             "incluidas las dos negativas. El matiz importa: negativo aquí significa cdis sin "
             "microcalcificación, no sin tumor, así que las negativas siguen siendo slides de mama "
             "con epitelio y estroma. Entonces el experto detecta tejido, no clase; que la slide "
             "sea positiva se decide después, en la atención y el clasificador. De ahí sale el "
             "punto para la discusión: que los expertos separen bien los tejidos y aun así MAMMOTH "
             "no le gane a CLAM en métrica es la evidencia de que el cuello de botella no está en "
             "la primera capa, sino en el dato. Con la honestidad por delante: es una cala chica, "
             "cuatro slides, una tarea, un fold, y las etiquetas de tejido todavía son "
             "provisionales, a la espera del visto bueno de un patólogo. Son dos ejes distintos: "
             "que MAMMOTH no mejore la métrica está cerrado; qué aprende por dentro está abierto, "
             "y es lo que trae esta presentación. Mostrar qué mira no reabre el rendimiento, le "
             "pone un mecanismo.")

    # ========================================================================
    # SECCIÓN MAGNIFICACIÓN MULTI-ESCALA (reunión Sebastián — decisión de escalas)
    # ========================================================================

    # ---- 12. Divisoria: Magnificación multi-escala ----
    s = divider(prs, "Magnificación multi-escala",
                "La única señal nueva tras cerrar la arquitectura: el contexto espacial · "
                "piloto microcalcificaciones")
    notes(s, "Cerrado el capítulo de la arquitectura, donde cuatro ejes distintos no movieron "
             "la métrica, queda una idea que todavía no probamos y que sí trae información "
             "nueva: la escala a la que miramos el tejido. De eso trata esta parte. Es un "
             "piloto sobre microcalcificaciones, que elegimos porque son pocas láminas y la "
             "extracción cabe en un fin de semana, pero la misma idea sirve para cualquier "
             "tarea que dependa del contexto.")

    # ---- 13. El problema es contexto + hallazgo físico (fusión de las dos slides) ----
    s = content(prs, "No es más zoom, es contexto: cohortes a distinta escala", size=21)
    add_textbox(s, 0.30, 0.86, 9.4, 0.54, [
        ("La etiqueta pregunta DÓNDE vive la microcalcificación (¿CDIS?, ¿invasor?, "
         "¿no neoplásico?), no si existe. Es una pregunta de contexto, no de detalle celular.",
         13, False, GRIS_BODY, F_BODY)])
    # izquierda: las dos demandas (detectar vs localizar)
    panel(s, 0.30, 1.48, 4.55, 1.35, "Detectar la calcificación", TEAL_TITLE, [
        "Alta magnif. 20-40×, campo ~100 µm.",
        "Anillos laminados: pista de DCIS.",
        "Ya la tenemos: el parche fino la resuelve.",
    ], TEAL_SQ)
    panel(s, 0.30, 2.98, 4.55, 1.35, "Localizar: ¿en qué estructura?", ORA_T, [
        "Baja o media, 5-10×, campo 0.5-2 mm.",
        "El conducto, nido invasor o lobulillo.",
        "Es lo que hoy falta: el fino no lo cubre.",
    ], ORA_ACC, fill=PROG_BG)
    # derecha: el hallazgo físico (cohortes a distinta escala)
    add_textbox(s, 5.15, 1.30, 4.55, 0.66, [
        ("Y las cohortes están a distinta escala física", 13, True, TEAL_TITLE, F_TITLE),
        ("medido en µm/px real, no en la etiqueta del archivo:", 11.5, False, GRIS_BODY, F_BODY)])
    simple_table(s, 5.15, 2.06, 4.55,
                 ["Cohorte", "Magnif.", "Parche 256 px"],
                 [["Pública (TCGA)", "~40×", "59 µm"],
                  ["Privada", "~20×", "119 µm"],
                  ["HistAI", "sin MPP", "excluida"]],
                 col_fracs=[0.44, 0.24, 0.32], row_h=0.36, fs=10.5)
    add_textbox(s, 5.15, 3.60, 4.55, 0.64, [
        ("Difieren 2×: la pirámide se define en µm/px, no en «level» del archivo.",
         11.5, False, INK, F_BODY)])
    takeaway_bar(s, "Hay que sumar una escala GRUESA que aporte el contexto del tejido, "
                    "y definirla en micras, no en «level».", t=4.48, size=12.5)
    notes(s, "Conviene empezar por algo que parece al revés. Uno pensaría que una "
             "microcalcificación, como es un objeto chico, pide más aumento, y es lo contrario. "
             "La etiqueta que queremos predecir no pregunta si hay una calcificación, sino en qué "
             "estructura vive: dentro de un carcinoma in situ, de un carcinoma invasor o de "
             "tejido no neoplásico. Eso es contexto, no detalle. Para responderlo el modelo tiene "
             "que ver la estructura anfitriona completa, el conducto, el nido invasor o el "
             "lobulillo, que miden entre medio milímetro y dos milímetros. Hoy extraemos un solo "
             "parche fino por región, que detecta la calcificación pero no alcanza a cubrir el "
             "conducto de alrededor. Por eso lo que falta no es más zoom, es sumar una escala "
             "gruesa que traiga ese contexto. A esto se agrega un hallazgo que apareció cuando "
             "medimos las láminas. Uno esperaría que un parche del mismo tamaño en píxeles "
             "cubriera lo mismo en las dos cohortes, y no. Con la resolución real de cada "
             "escáner, las láminas de la cohorte pública están cerca de cuarenta aumentos y las "
             "de la privada cerca de veinte, o sea el doble de resolución unas que otras, aunque "
             "la etiqueta del archivo diga otra cosa. Un mismo parche de doscientos cincuenta y "
             "seis píxeles cubre cincuenta y nueve micras en una y ciento diecinueve en la otra. "
             "De ahí salen dos consecuencias: la escala hay que definirla en micras por píxel y "
             "no en niveles del archivo, y el pipeline actual ya arrastra un sesgo, porque le "
             "entrega cada cohorte a un aumento distinto. Re-extraer a un campo físico común "
             "habilita la pirámide y de paso corrige ese sesgo. La tercera cohorte no tiene "
             "resolución confiable en su metadata, así que queda fuera.")

    # ---- 14. Lo que estudiamos: patología + referencias ----
    s = content(prs, "Lo que estudiamos: patología de la microcalcificación", size=24)
    add_textbox(s, 0.30, 0.88, 5.55, 0.3, [("Dos tipos de calcio, y uno es invisible", 13, True,
                                            TEAL_TITLE, F_TITLE)])
    simple_table(s, 0.30, 1.22, 5.55,
                 ["Tipo", "En H&E de rutina", "Para el modelo"],
                 [["Tipo I: oxalato", "casi invisible\n(solo luz polarizada)", "ciego → TECHO\n(sobre todo no_neoplásico)"],
                  ["Tipo II: fosfato", "basófilo, visible\n(anillos laminados)", "visible → aquí\njuega la escala"]],
                 col_fracs=[0.26, 0.37, 0.37], row_h=0.72, fs=10.5)
    add_textbox(s, 0.30, 3.62, 5.55, 0.95, [
        ("Tamaño: la calcificación (50-500 µm) ENTRA en un parche fino;", 12, False, GRIS_BODY, F_BODY),
        ("el conducto anfitrión (0.5-2 mm) NO → por eso hace falta el campo grueso.",
         12, False, GRIS_BODY, F_BODY)])
    # panel de referencias
    _rect(s, 6.05, 0.86, 3.65, 4.06, TEAL_CARD2, line=TEAL_SQ)
    add_textbox(s, 6.22, 0.93, 3.35, 3.95, [
        ("Referencias", 13, True, TEAL_TITLE, F_TITLE),
        ("Clínica de la microcalcificación", 10.5, True, ORA_T, F_BODY),
        ("· Breast microcalcifications: past, present & future · Review 2022", 9.5, False, INK, F_BODY),
        ("· Calcification in breast histopathology · Diagn. Histopathol. 2024", 9.5, False, INK, F_BODY),
        ("· Polyhedral microcalc. = calcium oxalate · Radiology 1993", 9.5, False, INK, F_BODY),
        ("· Microcalcifications: size matters! · 2007", 9.5, False, INK, F_BODY),
        ("· Predictors of malignancy in microcalc. · Br J Cancer 2011", 9.5, False, INK, F_BODY),
        ("Multi-escala en patología", 10.5, True, ORA_T, F_BODY),
        ("· CPathAgent · NeurIPS 2025 (baseline MIL multi-escala, Ap. C.1.2)", 9.5, False, INK, F_BODY),
        ("· DSMIL 20×+5× · Li et al., CVPR 2021", 9.5, False, INK, F_BODY),
        ("· Deep Multi-Magnification Nets (mama/DCIS) · Ho et al., 2021", 9.5, False, INK, F_BODY),
        ("Modelo base", 10.5, True, ORA_T, F_BODY),
        ("· CONCH · Lu et al., Nat Med 2024 (nativo 20×)", 9.5, False, INK, F_BODY),
        ("· CAP Invasive Breast, Nota D · la etiqueta contextual", 9.5, False, INK, F_BODY),
    ])
    notes(s, "Antes de elegir números conviene apoyarse en la patología, y hay dos hechos que "
             "mandan. El primero es que no todas las calcificaciones se ven igual en la tinción "
             "de rutina. Las de oxalato de calcio son casi invisibles en campo claro, solo "
             "aparecen con luz polarizada; el modelo, que trabaja sobre la tinción normal, está "
             "ciego a ellas a cualquier escala, y eso pone un techo sobre todo en el tejido no "
             "neoplásico. Las de fosfato de calcio, en cambio, son basófilas, se ven bien, y ahí "
             "sí la escala puede ayudar. El segundo hecho es de tamaño: la calcificación en sí, "
             "de cincuenta a quinientas micras, entra en un parche fino; el conducto que la aloja, "
             "no. Estas dos ideas vienen de la literatura de patología mamaria que está citada en "
             "la lámina, junto con los trabajos de multi-escala en los que nos apoyamos.")

    # ---- 15. La decisión de escalas (LA slide para Sebastián) ----
    s = content(prs, "La decisión de escalas", size=28)
    _pill = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.60), Inches(0.19), Inches(3.10), Inches(0.48))
    _pill.fill.solid(); _pill.fill.fore_color.rgb = PROG_BG
    _pill.line.color.rgb = ORA_ACC; _pill.line.width = Pt(1.25); _pill.shadow.inherit = False
    _set_runs(_pill.text_frame, [("Decisión delegada: pido tu guía", 11, True, ORA_T, F_BODY, PP_ALIGN.CENTER)],
              anchor=MSO_ANCHOR.MIDDLE)
    # izquierda: esquema nativo de campos concéntricos + pipeline
    nested_fields(s, 0.30, 1.02, 4.55, 1.95)
    # tabla px por cohorte
    simple_table(s, 0.30, 3.08, 4.55,
                 ["Escala", "Campo", "px TCGA", "px privado"],
                 [["Fina", "112 µm", "482", "241"],
                  ["Contexto", "512 µm", "2202", "1101"]],
                 col_fracs=[0.28, 0.24, 0.24, 0.24], row_h=0.32, fs=10.5)
    add_textbox(s, 0.30, 4.10, 4.55, 0.34, [
        ("cada campo → resize 224 → CONCH → [512]", 9.5, False, GRIS_TXT, F_MONO)])
    # derecha: crop REAL a dos escalas
    add_image_fit(s, MS_FINE, 5.10, 1.02, 2.25, 2.25, align="top")
    add_image_fit(s, MS_CTX, 7.45, 1.02, 2.25, 2.25, align="top")
    caption(s, 5.10, 3.30, 2.25, "fino 112 µm · citología", size=9.5, col=ORA_T, bold=True)
    caption(s, 7.45, 3.30, 2.25, "contexto 512 µm · arquitectura", size=9.5, col=TEAL_TITLE, bold=True)
    add_textbox(s, 5.10, 3.66, 4.6, 0.62, [
        ("Región real de mama (TCGA-BRCA), mismo centro. La caja naranja = el campo fino dentro "
         "del contexto.", 10, False, GRIS_BODY, F_BODY, PP_ALIGN.CENTER)])
    takeaway_bar(s, "Fusión por promedio → un token [N,512] → CLAM_MB intacto (la comparación más "
                    "limpia). ¿Estas escalas y esta fusión?", t=4.50, size=12.5)
    notes(s, "Esta es la decisión concreta, y es donde pido guía. La propuesta son dos escalas. "
             "Una fina, de ciento doce micras de campo, cerca de veinte magnificaciones, que es "
             "justo donde el encoder fue entrenado: ahí se detecta la calcificación y su forma. Y "
             "una de contexto, de quinientas doce micras, cerca de cinco magnificaciones, que "
             "abarca el conducto o el lobulillo anfitrión, que es literalmente lo que la etiqueta "
             "pregunta. A la derecha está una región real de una lámina de mama: el mismo centro "
             "visto a las dos escalas. En el campo fino se ve la citología; en el de contexto "
             "aparece la arquitectura glandular, y la caja marca dónde cae el parche fino dentro "
             "del grueso. Como las cohortes están a distinta resolución, el tamaño en píxeles se "
             "calcula por lámina para dar en el mismo campo físico: en la pública el crop fino es "
             "de cuatrocientos ochenta y dos píxeles y el de contexto dos mil doscientos; en la "
             "privada, doscientos cuarenta y uno y mil ciento uno; todos se reescalan a doscientos "
             "veinticuatro antes del encoder. Los dos vectores se promedian en uno solo, así el "
             "agregador queda intacto y la comparación es limpia. Lo que quiero conversar es si "
             "estas dos escalas y esta fusión te parecen las adecuadas, o si conviene mover los "
             "campos.")

    os.makedirs(OUT_DIR, exist_ok=True)
    prs.save(DST)
    print("Guardado:", DST, "·", len(prs.slides), "slides")


if __name__ == "__main__":
    build()
