#!/usr/bin/env python
"""generate_b9_deck.py — el deck del período sobre la plantilla oficial nueva.

Registro: **executive deck**, no deck técnico de sprint. Cuenta un período con la gramática
de la plantilla: qué me propuse (OBJETIVOS), qué salió (cuatro láminas de contenido), qué
sigue (Tareas del próximo período).

**Va entero en español** por decisión de Ernesto del 27-ago, el nombre del proyecto incluido
(«Detección Nuclear», que también nombra el archivo). La única lámina que queda en inglés es
la portada, que es copy de la empresa y viene tal cual del molde. Los nombres de clase del
patólogo (`Mitosis`, `Tumour`, `AreaTubular`…) tampoco se traducen: son las etiquetas
literales de su geojson y traducirlas rompería el vínculo con el material.

El método es lo que lo separa de los seis decks anteriores: **se RELLENA en sitio**, no se
reconstruye. `base_from_template()` de `generate_b8_deck.py` abría el template y le BORRABA
las láminas para dibujar todo de cero; acá no, porque las tablas de s02 y s04 vienen con las
filas de cuerpo vacías pero **ya estilizadas** (relleno de cabecera, bordes por celda,
márgenes de 22860 EMU, `anchor="ctr"`, 16 pt bold blanco) y reproducir eso a mano es trabajo
perdido y una fuente de diferencias contra el molde. Detalle y verificación:
`docs/plantilla_oficial.md` §7.

Sigue en pie el motivo de siempre para abrir el `.pptx` en vez de usar `Presentation()`: la
plantilla **embebe Barlow** ([[deck-template-fuentes-embebidas]]) y un deck construido desde
cero la pierde. Y sigue en pie `forzar_barlow()`, porque el `fontScheme` del theme de ESTA
plantilla también es Arial.

Fuentes de los números, leídas por el script y no transcritas a mano:
  - `results/b9_cruce_94/por_lamina.csv` y `recall_por_tolerancia_agregado.csv`
  - `sprints/B9_sprint9/mitosis_12_laminas/cruce_94.md`
  - `sprints/B9_sprint9/hovernext_tareas/inventario_tareas.md`
  - `results/b9_nucleos/{regiones_epi_estroma.csv,regiones_nulo.npy,marcas_grado.csv}`, que las
    cuatro láminas de los ejes nucleares **recalculan** con las mismas primitivas que produjeron
    el número (`scripts/b9_epitelio_estroma.py` y `scripts/b9_pleomorfismo.py`), no reimplementan
    ([[hallazgo-necesita-forma-presentable]]). Los números que quedan dibujados se escriben a
    `ejes_nucleares/figuras/numeros_figuras.csv`, que sí se versiona.

El guion vive en `guion_b9.md` y se aplica desde ahí: ese archivo es la fuente y las notas
del `.pptx` son derivadas.

Uso:
    PYTHONPATH=/media/administrador/Storage1/sdonoso/clam_testing2/.pylibs \
      /home/sdonoso/miniconda3/envs/pruebas/bin/python generate_b9_deck.py

(`envs/pruebas` y no `clam_latest` desde el 28-ago: al incorporar los dos ejes nucleares el deck
pasa a depender de `zarr`, que `b9_pleomorfismo` importa al tope y `clam_latest` no tiene. Es el
único env con pptx, pandas, numpy, zarr y scipy en un solo proceso. Workaround B: binario
absoluto.)
"""
import copy
import csv
import json
import os
import re
import sys

import numpy as np
from lxml import etree
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE, MSO_SHAPE_TYPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.oxml.ns import qn
from pptx.util import Emu, Inches, Pt

# ---------------------------------------------------------------------------
# Constantes del período
# ---------------------------------------------------------------------------
FECHA_ARCHIVO = "20260907"          # fin del período = fecha de la reunión
AUTOR = "Ernesto Gamero"
PROYECTO = "Detección Nuclear"      # NO «Mitosis Detection»: ése es el rótulo con el que
                                    # otra persona del equipo presenta su detector, y es
                                    # literalmente el ejemplo que trae la plantilla. En
                                    # español desde el 27-ago, con el deck entero
PERIODO = "25/08/2026 - 07/09/2026"

RAIZ = os.path.abspath(os.path.join(os.path.dirname(__file__), "..", "..", ".."))
TPL = os.path.join(RAIZ, "papers", "presentations",
                   "[AAAAMMDD] [Nombre Apellido] [Image-to-text].pptx")
AQUI = os.path.dirname(os.path.abspath(__file__))
OUT = os.path.join(AQUI, "[%s] [%s] [%s].pptx" % (FECHA_ARCHIVO, AUTOR, PROYECTO))
GUION = os.path.join(AQUI, "guion_b9.md")
CSV_LAMINA = os.path.join(RAIZ, "results", "b9_cruce_94", "por_lamina.csv")
CSV_TOL = os.path.join(RAIZ, "results", "b9_cruce_94", "recall_por_tolerancia_agregado.csv")
PNG_ACIERTOS = os.path.join(AQUI, "assets", "mitosis_aciertos.png")
PNG_FALLADAS = os.path.join(AQUI, "assets", "mitosis_falladas.png")
LADO_RECORTE_UM = 59.5              # 128 px de nivel 0 a 0,465 µm/px (galeria_mitosis_12.py)

# Las cinco láminas nuevas del 07/09. Las tres imágenes de resultado y la figura del paper son
# la excepción declarada en CLAUDE.md §"Formato de entregables": una fotografía de tejido y una
# figura publicada no son dibujables con shapes. Todo lo que las acompaña es NATIVO.
PNG_HOVERNEXT = os.path.join(RAIZ, "papers", "presentations", "assets_branding",
                             "paper_figs", "hovernext_fig1_pipeline.png")
PNG_ATENCION = os.path.join(AQUI, "assets", "atencion_12_laminas.png")
PNG_REGIONES = os.path.join(AQUI, "assets", "regiones_epi.png")
PNG_NUCLEOS = os.path.join(AQUI, "assets", "nucleos_grado.png")
JSON_REGIONES = os.path.join(AQUI, "assets", "regiones_epi.json")
JSON_NUCLEOS = os.path.join(AQUI, "assets", "nucleos_grado.json")
DIR_ATN = os.path.join(RAIZ, "results", "b9_atencion_12")
CSV_ESCALERA = os.path.join(RAIZ, "results", "b9_escalera_area", "agregado.csv")

# Los dos ejes nucleares de CPU. `CSV_NUM` es el único artefacto versionado de las cuatro
# figuras (el `.pptx` es derivado y está gitignored), y **no se mueve de sitio**: es el path
# que citan `ejes_nucleares/resultados.md` §4 y el envoltorio de la hoja suelta.
EJES_DIR = os.path.join(RAIZ, "sprints", "B9_sprint9", "ejes_nucleares", "figuras")
CSV_REG = os.path.join(RAIZ, "results", "b9_nucleos", "regiones_epi_estroma.csv")
NPY_NUL = os.path.join(RAIZ, "results", "b9_nucleos", "regiones_nulo.npy")
CSV_MAR = os.path.join(RAIZ, "results", "b9_nucleos", "marcas_grado.csv")
CSV_NUM = os.path.join(EJES_DIR, "numeros_figuras.csv")

SW, SH = 13.3333, 7.5

# ---------------------------------------------------------------------------
# Paleta de la plantilla oficial (docs/plantilla_oficial.md §4)
# ---------------------------------------------------------------------------
TITULO = RGBColor(0x1A, 0x1A, 0x2E)
CUERPO = RGBColor(0x1B, 0x4F, 0x8C)
ACENTO = RGBColor(0x52, 0x93, 0xDE)
SEP = RGBColor(0x9A, 0xA3, 0xB4)
LINEA = RGBColor(0xE4, 0xE9, 0xEC)
BLANCO = RGBColor(0xFF, 0xFF, 0xFF)
# Los dos anillos de las láminas de recortes. Son los de `galeria_mitosis_12.py` y tienen que
# seguir siendo los mismos: la leyenda es nativa y los círculos están quemados en el PNG.
BLANCO_ANILLO = BLANCO
AMARILLO = RGBColor(0xFF, 0xC1, 0x07)

# La rampa del mapa de atención. Es **turbo** y no un degradé cualquiera: es el colormap con
# el que `build_overlay_rgba` pinta el mapa (`scripts/mammoth_interpretability.py:262`). Los
# 24 stops van quemados en vez de importar matplotlib, que le agregaría al deck una
# dependencia entera para dibujar una leyenda de dos pulgadas.
TURBO = (0x30123B, 0x3C3286, 0x4451BF, 0x476EE6, 0x458AFC, 0x38A5FB, 0x25C0E7, 0x18D7CA,
         0x20EAAC, 0x3FF68A, 0x69FD66, 0x92FF47, 0xB1F936, 0xCDEC34, 0xE5D938, 0xF6C33A,
         0xFEA732, 0xFC8725, 0xF46617, 0xE7490C, 0xD43305, 0xBC2002, 0x9E1001, 0x7A0403)


def _rgb(c):
    return RGBColor(c >> 16, (c >> 8) & 0xFF, c & 0xFF)

F = "Barlow"
PT_TITULO, PT_CEJILLA, PT_CUERPO = 40, 12, 16
PT_CELDA, PT_PIE = 12, 9

# Geometría verificada sobre el archivo (docs/plantilla_oficial.md §5)
G_CEJILLA = (0.580, 0.380, 9.00, 0.30)
G_TITULO_S03 = (0.575, 0.668, 12.44, 0.72)
G_CUERPO = (0.625, 1.639, 12.097)
T_PIE_S03 = 6.72                    # zona de pie de s03: nada la cruza
# Ojo: `Google Shape;287;p7` NO es una barra 1B4F8C, es un cuadro de texto de pie VACÍO
# (`<a:noFill/>` y `<a:ln><a:noFill/></a:ln>`, verificado en el XML). Lo único que se ve
# ahí es la línea E4E9EC de 7,079. Se respeta como límite igual, porque es el pie.

A = "{http://schemas.openxmlformats.org/drawingml/2006/main}"


# ===========================================================================
# Medición de texto real, con los TTF de Barlow bajo containment
# ===========================================================================
# Portado tal cual de generate_b8_deck.py:708. Es la única capa de QA que ve un desborde:
# un chequeo de bounding boxes da «todo limpio» sobre texto que se sale de su caja
# ([[deck-qa-puntos-ciegos-chequeo]]).
BARLOW_DIR = "/media/administrador/Storage1/sdonoso/clam_testing2/fonts/barlow"
_FCACHE = {}
_MEDIDA = 40


def num(x, dec=1):
    """Número con coma decimal. El deck está en español y el `%.1f` de Python escribe punto:
    sin esto la misma lámina decía «7,5 a 50 µm» en la tabla y «7.5» en el cuerpo."""
    return ("%.*f" % (dec, x)).replace(".", ",")


def _face(bold):
    from PIL import ImageFont
    key = bool(bold)
    if key not in _FCACHE:
        nombre = "Barlow-Bold.ttf" if bold else "Barlow-Regular.ttf"
        _FCACHE[key] = ImageFont.truetype(os.path.join(BARLOW_DIR, nombre), _MEDIDA)
    return _FCACHE[key]


def text_w(txt, size, bold=False):
    """Ancho del texto en pulgadas de la lámina (72 pt/in)."""
    f = _face(bold)
    limpio = "".join(c if f.getbbox(c) else "n" for c in txt)
    return f.getlength(limpio) / _MEDIDA * size / 72.0


def wrap_lines(txt, ancho, size, bold=False):
    """Cuántas líneas ocupa `txt` en `ancho` pulgadas (wrap por palabra)."""
    if not txt.strip():
        return 1
    n, actual = 1, ""
    for palabra in txt.split(" "):
        cand = (actual + " " + palabra).strip()
        if actual and text_w(cand, size, bold) > ancho:
            n += 1
            actual = palabra
        else:
            actual = cand
    return n


def wrap_lines_mixto(tramos, ancho, size):
    """Igual que `wrap_lines` pero con el peso de cada tramo, para párrafos MIXTOS.

    Los puntos del cuerpo son «arranque en bold + resto normal» y medirlos con un solo peso
    se equivoca en los dos sentidos: todo normal subestima (el bold es más ancho) y todo bold
    sobreestima. Con el inglés, más corto, la diferencia nunca llegaba a cambiar el número de
    líneas; en español sí, y el generador y el auditor daban alturas distintas para el mismo
    párrafo. Cada uno medía a su manera, y el desacuerdo era el aviso."""
    tokens = []
    for txt, bold in tramos:
        for w in txt.split(" "):
            if w:
                tokens.append((w, bool(bold)))
    if not tokens:
        return 1
    n, ancho_actual = 1, 0.0
    for w, bold in tokens:
        pieza = text_w(w, size, bold)
        sep = text_w(" ", size, bold) if ancho_actual else 0.0
        if ancho_actual and ancho_actual + sep + pieza > ancho:
            n += 1
            ancho_actual = pieza
        else:
            ancho_actual += sep + pieza
    return n


def _alto_bloque(lineas, ancho, size, bold=False, space_after=3):
    total = 0.0
    for ln in lineas:
        total += wrap_lines(ln, ancho, size, bold) * size * 1.22 / 72.0 + space_after / 72.0
    return total


# ===========================================================================
# Primitivas
# ===========================================================================
def _set_runs(tf, lines, anchor=MSO_ANCHOR.TOP, space_after=3):
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    for i, ln in enumerate(lines):
        txt, sz, bold, col = ln[0], ln[1], ln[2], ln[3]
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = ln[4] if len(ln) > 4 else PP_ALIGN.LEFT
        p.space_after = Pt(space_after)
        r = p.add_run()
        r.text = txt
        r.font.size = Pt(sz)
        r.font.bold = bold
        r.font.name = F
        r.font.color.rgb = col


def add_textbox(slide, l, t, w, h, lines, anchor=MSO_ANCHOR.TOP, space_after=3):
    tb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tb.text_frame.margin_left = Inches(0.02)
    tb.text_frame.margin_right = Inches(0.02)
    tb.text_frame.margin_top = Inches(0)
    tb.text_frame.margin_bottom = Inches(0)
    _set_runs(tb.text_frame, lines, anchor=anchor, space_after=space_after)
    return tb


def _rect(slide, l, t, w, h, color):
    sp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(l), Inches(t), Inches(w), Inches(h))
    sp.fill.solid()
    sp.fill.fore_color.rgb = color
    sp.line.fill.background()
    sp.shadow.inherit = False
    return sp


def pie_lineas(slide, l, t, w, lineas, size=PT_PIE, col=CUERPO):
    """Pie de figura o de tabla: mide exactamente lo que mide su texto y devuelve el alto."""
    h = _alto_bloque(lineas, w - 0.06, size, space_after=1.5)
    tb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tb.text_frame.margin_left = Inches(0.02)
    tb.text_frame.margin_right = Inches(0.02)
    tb.text_frame.margin_top = Inches(0)
    tb.text_frame.margin_bottom = Inches(0)
    _set_runs(tb.text_frame, [(ln, size, False, col) for ln in lineas], space_after=1.5)
    return h


def notes(slide, text):
    """Guion HABLADO en prosa. El texto viene de guion_b9.md, que es la fuente."""
    slide.notes_slide.notes_text_frame.text = text


# ===========================================================================
# Las tres maniobras que python-pptx no tiene (docs/plantilla_oficial.md §7.a)
# ===========================================================================
def clonar_s03(prs, src):
    """Maniobra 1: lámina nueva del layout DEFAULT + deepcopy del XML de las formas.

    De las siete formas de s03 **sólo la imagen tiene relación** (`r:embed` del `a:blip`).
    Como la figura de este deck es NATIVA, la imagen no se clona y no hay ni una rel que
    copiar. El `Text 2` vacío tampoco se clona: está en L=0,41 T=1,41, justo donde va el
    cuerpo."""
    s = prs.slides.add_slide(prs.slide_layouts[0])   # DEFAULT: sin shapes ni placeholders
    for ph in list(s.placeholders):
        ph._element.getparent().remove(ph._element)
    tree = s.shapes._spTree
    for sp in src.shapes:
        if sp.shape_type == MSO_SHAPE_TYPE.PICTURE:
            continue
        if sp.name == "Text 2":
            continue
        tree.append(copy.deepcopy(sp._element))
    return s


def borrar_slide(prs, slide):
    """Saca la lámina del `sldIdLst` y suelta su relación."""
    idx = list(prs.slides).index(slide)
    lst = prs.slides._sldIdLst
    sldId = list(lst)[idx]
    prs.part.drop_rel(sldId.rId)
    lst.remove(sldId)


def reordenar(prs, orden):
    """Maniobra 3: mueve los `sldId` dentro del `sldIdLst`. `orden` = lista de slides."""
    lst = prs.slides._sldIdLst
    actuales = list(prs.slides)
    elems = list(lst)
    porslide = {id(s._element): e for s, e in zip(actuales, elems)}
    for e in elems:
        lst.remove(e)
    for s in orden:
        lst.append(porslide[id(s._element)])


def _quitar_fila(gf, i):
    """Maniobra 2: saca el `<a:tr>` i-ésimo y le resta su alto al graphicFrame."""
    tbl = gf._element.graphic.graphicData.tbl
    trs = tbl.findall(A + "tr")
    tr = trs[i]
    h = int(tr.get("h"))
    tbl.remove(tr)
    gf.height = Emu(gf.height - h)


def _alto_fila(gf, i, alto_in):
    """Fija el alto de una fila y ajusta el del graphicFrame por la diferencia."""
    tbl = gf._element.graphic.graphicData.tbl
    tr = tbl.findall(A + "tr")[i]
    nuevo = int(round(alto_in * 914400))
    gf.height = Emu(gf.height - int(tr.get("h")) + nuevo)
    tr.set("h", str(nuevo))


# ===========================================================================
# Rellenado en sitio
# ===========================================================================
def _shape(slide, nombre):
    for sh in slide.shapes:
        if sh.name == nombre:
            return sh
    raise KeyError("no está la forma %r en la lámina" % nombre)


def _solo_run(par, texto):
    """Deja el párrafo con UN run, conservando el formato del primero.

    El texto del template viene partido en varios runs (el título de s02 llega en tres), así
    que escribir solo `runs[0]` dejaría la cola del original pegada detrás."""
    runs = par.runs
    runs[0].text = texto
    for r in runs[1:]:
        r._r.getparent().remove(r._r)


def set_titulo(slide, texto, nombre="Text 2", size=PT_TITULO):
    sh = _shape(slide, nombre)
    _solo_run(sh.text_frame.paragraphs[0], texto)
    for r in sh.text_frame.paragraphs[0].runs:
        r.font.size = Pt(size)
    ancho = Emu(sh.width).inches
    if text_w(texto, size, True) > ancho:
        raise SystemExit("título que no entra en una línea: %r (%.2f\" en %.2f\")"
                         % (texto, text_w(texto, size, True), ancho))
    return sh


def set_cejilla(slide, tema, nombre=None):
    """`ONCOMETS   ·   <tema>`, conservando los tres colores del molde.

    s02 y s04 traen CUATRO runs (el tema partido en dos) y s03 trae tres: se normaliza a
    tres y se reescriben los colores, que es lo único que hay que preservar."""
    sh = None
    for cand in (nombre, "Google Shape;409;p22", "Text 0"):
        if cand is None:
            continue
        try:
            sh = _shape(slide, cand)
            break
        except KeyError:
            continue
    if sh is None:
        raise KeyError("no está la cejilla")
    p = sh.text_frame.paragraphs[0]
    runs = p.runs
    for r in runs[3:]:
        r._r.getparent().remove(r._r)
    textos = ["ONCOMETS", "   ·   ", tema]
    colores = [CUERPO, SEP, ACENTO]
    for r, txt, col in zip(p.runs, textos, colores):
        r.text = txt
        r.font.color.rgb = col
        r.font.size = Pt(PT_CEJILLA)
        r.font.bold = True
        r.font.name = F
    return sh


def set_cuerpo(slide, puntos, nombre="CuadroTexto 6"):
    """Cuerpo de s03: pares (arranque_bold, resto), `algn="just"` conservado.

    El cuadro usa `spAutoFit` y python-pptx **no lo recalcula**: la altura guardada es la que
    dejó PowerPoint para el texto anterior. Se fija midiendo con los TTF de Barlow. Devuelve
    el borde inferior en pulgadas, para que quien llama apoye la figura debajo."""
    sh = _shape(slide, nombre)
    tf = sh.text_frame
    modelo = copy.deepcopy(tf.paragraphs[0]._p)
    for p in list(tf._txBody.findall(A + "p")):
        tf._txBody.remove(p)
    for bold, resto in puntos:
        nuevo = copy.deepcopy(modelo)
        tf._txBody.append(nuevo)
    for par, (bold, resto) in zip(tf.paragraphs, puntos):
        runs = par.runs
        for r in runs[2:]:
            r._r.getparent().remove(r._r)
        runs[0].text = bold
        runs[1].text = resto
        for r in par.runs:
            r.font.size = Pt(PT_CUERPO)
            r.font.name = F
            r.font.color.rgb = CUERPO

    l, t, w = G_CUERPO
    ancho_txt = w - 0.20                      # lIns/rIns por defecto del cuadro de texto
    alto = 0.05
    for bold, resto in puntos:
        alto += (wrap_lines_mixto([(bold, True), (resto, False)], ancho_txt, PT_CUERPO)
                 * PT_CUERPO * 1.22 / 72.0)
    alto += 0.05
    sh.left, sh.top, sh.width = Inches(l), Inches(t), Inches(w)
    sh.height = Inches(alto)
    return t + alto


def _util(cell, ancho_col):
    """Ancho utilizable de una celda: el de la columna menos sus márgenes REALES."""
    return ancho_col - Emu(cell.margin_left).inches - Emu(cell.margin_right).inches


def llenar_tabla(gf, filas, size=PT_CELDA, min_h=0.76):
    """Rellena las filas de cuerpo respetando el estilizado que ya traen las celdas.

    Quita las que sobren restándole el alto al `graphicFrame`, y le da a cada una el alto que
    su texto pide de verdad, medido con Barlow, en vez de confiar en que PowerPoint la crezca
    al abrir el archivo.

    `min_h` existe porque el molde trae CUATRO filas de cuerpo y acá van menos: con el alto
    justo del texto la tabla se encoge, queda pegada al título y deja media lámina vacía. Se
    calibra para conservar la huella del molde, 3,03" contando la cabecera de 0,745: con tres
    filas es 0,76 (el default) y con dos, 1,14. **Cambiar el número de filas obliga a pasarlo**:
    con dos filas y el default de tres la tabla de Tareas quedaba flotando arriba."""
    tbl = gf.table
    while len(tbl.rows) - 1 > len(filas):
        _quitar_fila(gf, len(tbl.rows) - 1)
    if len(tbl.rows) - 1 < len(filas):
        raise SystemExit("la tabla del molde tiene %d filas de cuerpo y se piden %d"
                         % (len(tbl.rows) - 1, len(filas)))
    anchos = [Emu(c.width).inches for c in tbl.columns]
    for ri, fila in enumerate(filas, start=1):
        nlin = 1
        for ci, txt in enumerate(fila):
            cell = tbl.cell(ri, ci)
            p = cell.text_frame.paragraphs[0]
            cell.text_frame.word_wrap = True
            r = p.add_run()
            r.text = txt
            r.font.size = Pt(size)
            r.font.name = F
            r.font.bold = False
            r.font.color.rgb = TITULO
            nlin = max(nlin, wrap_lines(txt, _util(cell, anchos[ci]), size))
        _alto_fila(gf, ri, max(min_h, nlin * size * 1.22 / 72.0 + 0.14))
    return Emu(gf.top).inches + Emu(gf.height).inches


# ===========================================================================
# Las dos figuras nativas
# ===========================================================================
def barras_marcas(slide, l, t, w, alto_max, datos):
    """Doce barras horizontales, una por lámina, ordenadas por marcas.

    Largo proporcional a las marcas; la parte reencontrada rellena, el resto en el gris de la
    plantilla. La 129741 va en el celeste de acento porque es la que aporta la mitad de los
    aciertos y la lámina lo dice. El `n of m` va a la derecha en **columna fija**: es lo único
    que se lee cuando la barra mide un píxel, y siete láminas tienen de 1 a 3 marcas."""
    w_lab, w_num, gap = 1.30, 1.30, 0.14
    x_bar = l + w_lab + gap
    w_bar = w - w_lab - w_num - 2 * gap
    x_num = l + w - w_num
    maxm = max(d["marcas"] for d in datos)

    # leyenda
    fs_leg = 8.5
    cx = x_bar
    for col, txt in ((CUERPO, "reencontradas a 30 µm"), (LINEA, "no reencontradas")):
        _rect(slide, cx, t + 0.045, 0.16, 0.105, col)
        add_textbox(slide, cx + 0.21, t, 1.9, 0.18,
                    [(txt, fs_leg, False, CUERPO)], anchor=MSO_ANCHOR.MIDDLE)
        cx += 0.21 + text_w(txt, fs_leg) + 0.30
    add_textbox(slide, x_num, t, w_num, 0.18,
                [("de sus marcas", fs_leg, True, CUERPO, PP_ALIGN.RIGHT)],
                anchor=MSO_ANCHOR.MIDDLE)

    y0 = t + 0.30
    paso = (alto_max - 0.30) / len(datos)
    h_bar = min(0.17, paso - 0.10)
    fs = 9.5
    for i, d in enumerate(datos):
        y = y0 + i * paso
        yc = y + (paso - h_bar) / 2.0
        destacada = d["slide_id"] == "129741"
        largo = w_bar * d["marcas"] / float(maxm)
        _rect(slide, x_bar, yc, largo, h_bar, LINEA)
        if d["tp"]:
            _rect(slide, x_bar, yc, largo * d["tp"] / float(d["marcas"]), h_bar,
                  ACENTO if destacada else CUERPO)
        add_textbox(slide, l, y, w_lab, paso,
                    [(d["slide_id"], fs, destacada, CUERPO, PP_ALIGN.RIGHT)],
                    anchor=MSO_ANCHOR.MIDDLE)
        add_textbox(slide, x_num, y, w_num, paso,
                    [("%d de %d" % (d["tp"], d["marcas"]), fs, destacada, CUERPO,
                      PP_ALIGN.RIGHT)], anchor=MSO_ANCHOR.MIDDLE)
    return t + alto_max


def tabla_ejes(slide, l, t, w, headers, filas, fracs, fs=9.5, row_h=0.30):
    """Tabla nativa con la paleta de la plantilla: cabecera 1B4F8C con texto blanco y banding
    en E4E9EC. Es `simple_table` de generate_b8_deck.py:660 con la paleta nueva."""
    ncol, nrow = len(headers), len(filas) + 1
    gf = slide.shapes.add_table(nrow, ncol, Inches(l), Inches(t), Inches(w),
                                Inches(row_h * nrow))
    tbl = gf.table
    for ci, fr in enumerate(fracs):
        tbl.columns[ci].width = Inches(w * fr)
    tbl.first_row = False
    tbl.horz_banding = False
    for ri, fila in enumerate([headers] + list(filas)):
        for ci, txt in enumerate(fila):
            cell = tbl.cell(ri, ci)
            cell.margin_left = Inches(0.07)
            cell.margin_right = Inches(0.05)
            cell.margin_top = Inches(0.01)
            cell.margin_bottom = Inches(0.01)
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            cell.fill.solid()
            if ri == 0:
                cell.fill.fore_color.rgb = CUERPO
                col, bold = BLANCO, True
            else:
                cell.fill.fore_color.rgb = LINEA if ri % 2 else BLANCO
                col, bold = TITULO, (ci == 0)
            tf = cell.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.alignment = PP_ALIGN.LEFT
            r = p.add_run()
            r.text = txt
            r.font.size = Pt(fs)
            r.font.bold = bold
            r.font.name = F
            r.font.color.rgb = col
    for ri in range(nrow):
        tbl.rows[ri].height = Inches(row_h)
    gf.height = Inches(row_h * nrow)
    return t + row_h * nrow


# ===========================================================================
# Primitivas de dibujo que la plantilla no traía
# ===========================================================================
# Las usan las cuatro láminas de los ejes nucleares. Vivían en la hoja suelta
# `ejes_nucleares/figuras/generate_figuras_ejes34.py`, que importaba de acá; con las láminas
# adentro del deck la dependencia se invirtió, porque al revés habría ciclo y correr el deck
# como `__main__` cargaría una segunda copia de este módulo.
def eje_x(slide, l, t, w, vmin, vmax, ticks, dec=2, fs=8.5, col=SEP):
    """Eje horizontal: línea de base y etiquetas debajo. Devuelve el borde inferior."""
    _rect(slide, l, t, w, 0.012, col)
    for v in ticks:
        x = l + w * (v - vmin) / float(vmax - vmin)
        _rect(slide, x - 0.006, t, 0.012, 0.07, col)
        txt = num(v, dec)
        add_textbox(slide, x - 0.45, t + 0.08, 0.90, 0.17,
                    [(txt, fs, False, CUERPO, PP_ALIGN.CENTER)], anchor=MSO_ANCHOR.TOP)
    return t + 0.08 + 0.17


def eje_y(slide, l, t, h, vmin, vmax, ticks, dec=0, fs=8.5, col=SEP):
    """Eje vertical con el 0 abajo. `l` es la x de la línea; las etiquetas van a su izquierda."""
    _rect(slide, l, t, 0.012, h, col)
    for v in ticks:
        y = t + h * (1.0 - (v - vmin) / float(vmax - vmin))
        _rect(slide, l - 0.06, y - 0.006, 0.07, 0.012, col)
        add_textbox(slide, l - 0.62, y - 0.09, 0.50, 0.18,
                    [(num(v, dec), fs, False, CUERPO, PP_ALIGN.RIGHT)],
                    anchor=MSO_ANCHOR.MIDDLE)


def punto(slide, cx, cy, d, color, hueco=False):
    """Marcador circular. Hueco = anillo, para las láminas con `alineada: false`."""
    ov = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(cx - d / 2.0), Inches(cy - d / 2.0),
                                Inches(d), Inches(d))
    if hueco:
        ov.fill.background()
        ov.line.color.rgb = color
        ov.line.width = Pt(1.75)
    else:
        ov.fill.solid()
        ov.fill.fore_color.rgb = color
        ov.line.fill.background()
    ov.shadow.inherit = False
    return ov


def histograma(slide, l, t, w, h, vals, vmin, vmax, nbins, color):
    """Histograma nativo: una barra por bin, apoyadas en la base. Devuelve el conteo máximo."""
    bordes = np.linspace(vmin, vmax, nbins + 1)
    cuenta, _ = np.histogram(np.asarray(vals, float), bins=bordes)
    top = max(int(cuenta.max()), 1)
    wb = w / float(nbins)
    for i, c in enumerate(cuenta):
        if not c:
            continue
        hb = h * c / float(top)
        _rect(slide, l + i * wb + 0.006, t + h - hb, max(wb - 0.012, 0.012), hb, color)
    return top


def marcador_vertical(slide, x, t, h, texto, color, fs=9.5, lado="der"):
    """Línea vertical rotulada: el valor observado sobre un histograma."""
    _rect(slide, x - 0.012, t, 0.024, h, color)
    aw = text_w(texto, fs, True) + 0.12
    xl = x + 0.09 if lado == "der" else x - 0.09 - aw
    al = PP_ALIGN.LEFT if lado == "der" else PP_ALIGN.RIGHT
    add_textbox(slide, xl, t - 0.02, aw, 0.20, [(texto, fs, True, color, al)],
                anchor=MSO_ANCHOR.MIDDLE)


def leyenda_puntos(slide, l, t, items, fs=9, d=0.13):
    """Leyenda de marcadores: (color, hueco, texto)."""
    x = l
    for color, hueco, txt in items:
        punto(slide, x + d / 2.0, t + 0.09, d, color, hueco)
        add_textbox(slide, x + d + 0.08, t, text_w(txt, fs) + 0.10, 0.18,
                    [(txt, fs, False, CUERPO)], anchor=MSO_ANCHOR.MIDDLE)
        x += d + 0.08 + text_w(txt, fs) + 0.30
    return t + 0.18


def caja_figura(fin_cuerpo, pie, w):
    """Reparte lo que queda entre el cuerpo y la zona de pie. Devuelve (top, alto, y_pie)."""
    h_pie = _alto_bloque(pie, w - 0.06, PT_PIE, space_after=1.5)
    top = fin_cuerpo + 0.16
    alto = T_PIE_S03 - 0.12 - h_pie - 0.12 - top
    return top, alto, top + alto + 0.12


# ===========================================================================
# Datos: se LEEN, no se transcriben
# ===========================================================================
def leer_datos():
    with open(CSV_LAMINA) as fh:
        filas = list(csv.DictReader(fh))
    datos = [{"slide_id": r["slide_id"], "marcas": int(r["marcas"]),
              "tp": int(r["tp_30um"]), "det": int(r["detecciones"])} for r in filas]
    datos.sort(key=lambda d: (-d["marcas"], d["slide_id"]))
    with open(CSV_TOL) as fh:
        tol = list(csv.DictReader(fh))
    d = {"n_laminas": len(datos),
         "marcas": sum(x["marcas"] for x in datos),
         "tp": sum(x["tp"] for x in datos),
         "det": sum(x["det"] for x in datos)}
    fila30 = [r for r in tol if float(r["tolerancia_um"]) == 30.0][0]
    d["recall"] = float(fila30["recall"])
    assert int(fila30["tp"]) == d["tp"] and int(fila30["n"]) == d["marcas"], \
        "el agregado y el por-lámina no coinciden"
    # el tramo plano de la escalera, que es lo que la lámina afirma
    plano = [float(r["tolerancia_um"]) for r in tol if int(r["tp"]) == d["tp"]]
    d["plano"] = (min(plano), max(plano))
    ref = [x for x in datos if x["slide_id"] == "129741"][0]
    d["ref_tp"], d["ref_marcas"] = ref["tp"], ref["marcas"]
    d["resto_tp"] = d["tp"] - ref["tp"]
    d["resto_marcas"] = d["marcas"] - ref["marcas"]
    d["resto_recall"] = d["resto_tp"] / float(d["resto_marcas"])
    return datos, d


# El prefijo y NO el nombre completo: el `json_out` trae la tarea con el sufijo
# `_pth_balance` y los dos brazos de checkpoint con `_combined_5fold`. Un filtro por igualdad
# contra un nombre fijo devuelve CERO filas en dos de los tres archivos, y sin error
# (`atencion_12_laminas/csv_audit.md` §1, trampa 5).
TAREA_ATN = "grado_histologico_mitotic_rate"
BRAZOS_ATN = (
    ("json_out", "auc_por_lamina.csv", "meta.json",
     "ensemble de los cinco folds, rama predicha"),
    ("ckpt_todos", "auc_por_lamina_ckpt_todos.csv", "meta_ckpt_todos.json",
     "los cinco folds, rama verdadera"),
    ("ckpt_limpios", "auc_por_lamina_ckpt_limpios.csv", "meta_ckpt_limpios.json",
     "sólo los folds limpios de cada lámina"),
)


def leer_atencion():
    """Los tres brazos sobre las NUEVE láminas del primario, con su media y su dispersión.

    Aborta si algún brazo no da nueve láminas, si los tres no cubren exactamente las mismas, o
    si el ordenamiento `json_out >= todos >= limpios` no se sostiene: ése es el ordenamiento
    que el pre-registro fijó ANTES de correr, y si se rompiera sería bug de la tabla de
    membresía y no resultado."""
    brazos = []
    for clave, arch, meta_arch, etiqueta in BRAZOS_ATN:
        with open(os.path.join(DIR_ATN, arch)) as fh:
            filas = [r for r in csv.DictReader(fh)
                     if r["universo"] == "lamina" and r["primario"] == "True"
                     and r["tarea"].startswith(TAREA_ATN)]
        if len(filas) != 9:
            raise SystemExit("%s: %d láminas del primario, se esperaban 9 (¿filtro por "
                             "igualdad contra el nombre de la tarea?)" % (arch, len(filas)))
        aucs = [float(r["auc"]) for r in filas]
        med = sum(aucs) / len(aucs)
        sd = (sum((a - med) ** 2 for a in aucs) / (len(aucs) - 1)) ** 0.5
        with open(os.path.join(DIR_ATN, meta_arch)) as fh:
            meta = json.load(fh)
        brazos.append({"clave": clave, "rotulo": etiqueta, "media": med, "sd": sd,
                       "n": len(filas), "meta": meta, "familia": filas[0]["tarea"],
                       "slides": tuple(sorted(r["slide"] for r in filas)),
                       "sobre_azar": sum(1 for a in aucs if a > 0.5)})
    if len({b["slides"] for b in brazos}) != 1:
        raise SystemExit("los tres brazos no cubren las mismas nueve láminas")
    if not brazos[0]["media"] >= brazos[1]["media"] >= brazos[2]["media"]:
        raise SystemExit("el ordenamiento pre-registrado no se sostiene: %s"
                         % [num(b["media"], 3) for b in brazos])
    c = brazos[0]["meta"]["conteos"]
    for b in brazos:
        if b["meta"]["conteos"] != c:
            raise SystemExit("los tres meta*.json no declaran los mismos conteos")
    return {"brazos": brazos, "conteos": c}


# Los cinco peldaños, de más área a menos. El orden es el de la escalera, y el chequeo de
# monotonía lo da por sentado.
PELDANOS = (("lamina_entera", "lámina entera"), ("30.0", "30 mm² por lámina"),
            ("10.0", "10 mm² por lámina"), ("3.0", "3 mm² por lámina"),
            ("1.0", "1 mm² por lámina"))


def leer_escalera():
    """Los cinco peldaños de las dos cabezas y del azar, más los dos controles.

    `clam_combinado` NO entra en la lámina: es exploratorio y así está declarado en
    `atencion_12_laminas/resultados.md` §5. Aborta si las acreditadas no bajan al bajar el
    presupuesto, o si alguno de los conteos duros se movió."""
    with open(CSV_ESCALERA) as fh:
        filas = list(csv.DictReader(fh))

    def fila(brazo, pres):
        r = [x for x in filas if x["brazo"] == brazo and x["presupuesto_mm2"] == pres]
        if len(r) != 1:
            raise SystemExit("agregado.csv: %d filas para %s / %s" % (len(r), brazo, pres))
        return r[0]

    peld = []
    for pres, etiqueta in PELDANOS:
        m, g, a = fila("clam_mitosis", pres), fila("clam_gate", pres), fila("azar", pres)
        peld.append({"etiqueta": etiqueta, "area": float(m["area_real_mm2"]),
                     "k": int(float(m["k_parches"])),
                     "mitosis": int(float(m["acreditadas_dentro"])),
                     "gate": int(float(g["acreditadas_dentro"])),
                     "azar": float(a["acreditadas_dentro"]),
                     "destacado": pres == "3.0"})
    mit = [x["mitosis"] for x in peld]
    if mit != sorted(mit, reverse=True):
        raise SystemExit("las acreditadas no son monótonas al bajar el presupuesto: %s" % mit)
    sf, te = fila("sin_filtro", "lamina_entera"), fila("teselado", "lamina_entera")
    d = {"peldanos": peld,
         "marcas": int(float(fila("azar", "lamina_entera")["marcas_dentro"])),
         "acreditadas": int(float(te["acreditadas_dentro"])),
         "laminas": int(te["n_laminas"]),
         "parches": int(float(te["k_parches"])),
         "area_total": float(te["area_real_mm2"]),
         "sin_filtro": (int(float(sf["detecciones_dentro"])),
                        int(float(sf["acreditadas_dentro"]))),
         "teselado": (int(float(te["detecciones_dentro"])),
                      int(float(te["acreditadas_dentro"])))}
    duros = ((d["marcas"], 94), (d["acreditadas"], 26), (d["laminas"], 12),
             (d["parches"], 49832), (d["sin_filtro"], (732, 26)), (d["teselado"], (707, 26)))
    for visto, esperado in duros:
        if visto != esperado:
            raise SystemExit("conteo duro movido: %r, se esperaba %r" % (visto, esperado))
    if abs(d["area_total"] - 706.1) > 0.1 or abs(peld[0]["area"] - d["area_total"]) > 1e-6:
        raise SystemExit("el área total se movió: %.3f mm²" % d["area_total"])
    return d


def leer_guion():
    """Parsea `guion_b9.md` por sus marcadores `## [sNN]`. El .md es la fuente."""
    txt = open(GUION, encoding="utf-8").read()
    bloques, orden = {}, []
    for m in re.finditer(r"^## \[([^\]]+)\][^\n]*\n(.*?)(?=^## \[|\Z)",
                         txt, re.M | re.S):
        clave = m.group(1).strip()
        bloques[clave] = m.group(2).strip()
        orden.append(clave)
    return bloques, orden


def datos_eje4():
    """Eje 4. Recalcula el AUC con `rank_auc` del propio script del eje, no con una copia.

    Los dos módulos de los ejes se importan **acá adentro** y no al tope: `b9_pleomorfismo`
    trae `zarr`, y con el import a nivel de módulo bastaría abrir este archivo desde un env
    sin zarr para que ni siquiera cargue. El deck lo necesita para CORRER; importarlo no."""
    import pandas as pd
    sys.path.insert(0, os.path.join(RAIZ, "scripts"))
    import b9_epitelio_estroma as E
    df = pd.read_csv(CSV_REG)
    nul = np.load(NPY_NUL)
    pos = (df["grupo"] == "epitelio").to_numpy()
    obs = E.rank_auc(df["f_epi"].to_numpy(), pos)
    aucs = np.array([E.rank_auc(nul[:, i], pos) for i in range(nul.shape[1])])
    aucs = aucs[~np.isnan(aucs)]
    p = (1 + int((aucs >= obs).sum())) / (1.0 + len(aucs))
    sub = df[df["alineada"]]
    clases = []
    for cl, g in df.groupby("clase"):
        v = g["f_epi"].dropna()
        if not len(v):
            continue
        clases.append(dict(clase=cl, grupo=g["grupo"].iloc[0], n=len(g), n_val=len(v),
                           med=float(v.median()), p25=float(v.quantile(.25)),
                           p75=float(v.quantile(.75))))
    # epitelio primero, cada grupo por mediana descendente: el orden de la tabla de §1.a
    clases.sort(key=lambda d: (d["grupo"] != "epitelio", -d["med"]))
    return dict(df=df, clases=clases, obs=obs, nulo=aucs, p=p,
                n_reg=len(df), n_epi=int(pos.sum()), n_est=int((~pos).sum()),
                sin_nucleos=int(df["f_epi"].isna().sum()),
                auc_alin=E.rank_auc(sub["f_epi"].to_numpy(),
                                    (sub["grupo"] == "epitelio").to_numpy()),
                n_alin=len(sub))


def datos_eje3():
    """Eje 3. `spearman` y `permutacion_exacta` son las del script del eje: el `p` que queda
    dibujado sale de la MISMA enumeración que produjo el de `resultados.md`.

    `ORDEN` viaja en el diccionario de vuelta para que `lamina_f3` no tenga que importar el
    módulo (y con él zarr) sólo para saber el orden de los tres grados."""
    import pandas as pd
    sys.path.insert(0, os.path.join(RAIZ, "scripts"))
    import b9_pleomorfismo as P
    m = pd.read_csv(CSV_MAR)
    COL = "pct_area_um2"                      # el PRIMARIO: percentil intra-lámina, no el área
    # Las dos láminas cuyo offset quedó con `alineada: false`. Se dibujan con el punto
    # HUECO: una de ellas es el peldaño del medio entero, y eso tiene que verse en la
    # figura y no en el pie.
    NO_ALIN = P.NO_ALINEADAS

    def bloque(df, etiqueta):
        pl = (df.groupby(["slide", "grado", "grado_ord"], as_index=False)
                .agg(n=(COL, "size"), med=(COL, "median")))
        pl["alineada"] = ~pl["slide"].isin(NO_ALIN)
        rho, _ = P.spearman(pl["grado_ord"], pl["med"])
        perm = P.permutacion_exacta(pl["grado_ord"], pl["med"], con_rhos=True)
        grados = []
        for g in P.ORDEN:
            s = df[df["grado"] == g][COL].dropna()
            grados.append(dict(grado=g, n_marcas=len(s),
                               n_lam=int(pl[pl["grado"] == g].shape[0]),
                               med=float(s.median()), p25=float(s.quantile(.25)),
                               p75=float(s.quantile(.75))))
        pares = []
        for g1, g2 in (("bajo", "moderado"), ("moderado", "alto")):
            sub = df[df["grado"].isin([g1, g2])]
            auc, npos, nneg = P.rank_auc(sub[COL], (sub["grado"] == g2).to_numpy())
            pares.append(dict(par="%s > %s" % (g2, g1), auc=auc, n_neg=nneg, n_pos=npos))
        return dict(etiqueta=etiqueta, pl=pl, rho=rho, obs=perm[0], p_bi=perm[1],
                    p_uni=perm[2], k=perm[3], rhos=perm[4], grados=grados, pares=pares,
                    n_marcas=len(df), n_lam=len(pl))

    r = m[m["clase_hovernext"] == "epithelial-cell"]
    return dict(restringida=bloque(r, "restringida a epitelio"),
                completa=bloque(m, "completa"), orden=list(P.ORDEN))


# ===========================================================================
# Tipografía y auditoría
# ===========================================================================
def forzar_barlow(prs, fuente=F):
    """Deja Barlow como única tipografía del archivo. Portado de generate_b8_deck.py:1185.

    **No se puede saltear**: el `fontScheme` del theme de esta plantilla también es el de
    Office (Arial) y las láminas heredadas traen Calibri en `endParaRPr`/`buFont`, que
    gobierna lo que se escriba encima."""
    tags = tuple(A + t for t in ("latin", "ea", "cs", "sym", "buFont"))

    def normaliza(root):
        n = 0
        for el in root.iter():
            if el.tag in tags and el.get("typeface") != fuente:
                el.set("typeface", fuente)
                n += 1
        return n

    cola = [prs.part] + [s.part for s in prs.slides]
    cola += [s.notes_slide.part for s in prs.slides if s.has_notes_slide]
    vistos, partes = set(), []
    while cola:
        p = cola.pop()
        if id(p) in vistos:
            continue
        vistos.add(id(p))
        partes.append(p)
        for rel in p.rels.values():
            if not rel.is_external and rel.reltype.endswith(
                    ("slideLayout", "slideMaster", "notesMaster", "theme")):
                cola.append(rel.target_part)
    n = 0
    for p in partes:
        el = getattr(p, "_element", None)
        if el is not None:
            n += normaliza(el)
            continue
        if str(p.partname).startswith("/ppt/theme/"):
            raiz = etree.fromstring(p.blob)
            tocados = 0
            for cual in ("majorFont", "minorFont"):
                for fs in raiz.iter(A + cual):
                    for lat in fs.findall(A + "latin"):
                        if lat.get("typeface") != fuente:
                            lat.set("typeface", fuente)
                            tocados += 1
            if tocados:
                p._blob = etree.tostring(raiz, xml_declaration=True,
                                         encoding="UTF-8", standalone=True)
                n += tocados
    print("  tipografía: %d referencias forzadas a %s" % (n, fuente))


def auditar(prs, saltar_idx=(1,)):
    """Texto que no entra en su caja, formas fuera del lienzo, cuerpos por debajo del mínimo
    del template (7 pt) y objetos que cruzan la zona de pie de s03.

    `saltar_idx` = láminas heredadas que no escribimos (la portada). No reemplaza mirar las
    láminas ([[deck-qa-puntos-ciegos-chequeo]])."""
    problemas = []
    for idx, slide in enumerate(prs.slides, start=1):
        if idx in saltar_idx:
            continue
        pie = None
        for sh in slide.shapes:
            if sh.name == "Google Shape;287;p7":
                pie = Emu(sh.top).inches
        for sh in slide.shapes:
            try:
                l, t = Emu(sh.left).inches, Emu(sh.top).inches
                w, h = Emu(sh.width).inches, Emu(sh.height).inches
            except TypeError:
                continue
            if not getattr(sh, "rotation", 0):
                if l < -0.02 or t < -0.02 or l + w > SW + 0.02 or t + h > SH + 0.02:
                    problemas.append("s%02d  fuera del lienzo: %s (%.2f, %.2f, %.2f x %.2f)"
                                     % (idx, sh.name, l, t, w, h))
            if sh.has_table:
                tbl = sh.table
                anchos = [Emu(c.width).inches for c in tbl.columns]
                acum = 0.0
                for ri, row in enumerate(tbl.rows):
                    hr = Emu(row.height).inches
                    nl = 1
                    for ci, cell in enumerate(row.cells):
                        util = _util(cell, anchos[ci])
                        for p in cell.text_frame.paragraphs:
                            if not p.runs:
                                continue
                            sz = max((r.font.size.pt for r in p.runs if r.font.size),
                                     default=12)
                            bold = any(r.font.bold for r in p.runs)
                            txt = "".join(r.text for r in p.runs)
                            nl = max(nl, wrap_lines(txt, util, sz, bold))
                            # una celda de UNA línea que llega al borde no cuenta como
                            # desborde y aun así se ve pegada al filete: 0,06" de guarda
                            ancho = text_w(txt, sz, bold)
                            if util >= ancho > util - 0.06:
                                problemas.append(
                                    "s%02d  celda al límite de su columna (%.2f de %.2f in): «%s»"
                                    % (idx, ancho, util, txt[:44]))
                            if sz < 7.0:
                                problemas.append("s%02d  celda por debajo del mínimo: %.1f pt"
                                                 % (idx, sz))
                    need = nl * sz * 1.22 / 72.0
                    if need > hr + 0.02:
                        problemas.append("s%02d  fila %d de tabla corta: pide %.2f\", tiene %.2f\""
                                         % (idx, ri, need, hr))
                    acum += max(hr, need)
                if pie is not None and t + acum > pie + 0.02:
                    problemas.append("s%02d  la tabla cruza la zona de pie (%.2f > %.2f)"
                                     % (idx, t + acum, pie))
                continue
            if not sh.has_text_frame:
                if pie is not None and sh.name not in ("Google Shape;287;p7",
                                                       "Google Shape;286;p7") \
                        and t + h > pie + 0.02:
                    problemas.append("s%02d  %s cruza la zona de pie (%.2f > %.2f)"
                                     % (idx, sh.name, t + h, pie))
                continue
            # Los márgenes REALES del cuadro, no un 0,20" supuesto: los cuadros propios
            # usan 0,02" por lado y suponer el default del template daba falso positivo.
            tf = sh.text_frame
            ins = (Emu(tf.margin_left).inches + Emu(tf.margin_right).inches)
            alto, chico = 0.0, None
            for p in tf.paragraphs:
                if not p.runs:
                    continue
                sz = max((r.font.size.pt for r in p.runs if r.font.size), default=12)
                if sz >= 6 and (chico is None or sz < chico):
                    chico = sz
                tramos = [(r.text, bool(r.font.bold)) for r in p.runs]
                alto += (wrap_lines_mixto(tramos, max(w - ins, 0.2), sz)
                         * sz * 1.22 / 72.0)
            if alto > h + 0.06:
                problemas.append("s%02d  texto que no entra: sobra %.2f\" en «%s…»"
                                 % (idx, alto - h, sh.text_frame.text[:44].replace("\n", " ")))
            if chico is not None and chico < 7.0:
                problemas.append("s%02d  cuerpo por debajo del mínimo: %.1f pt" % (idx, chico))
            if pie is not None and t + alto > pie + 0.02:
                problemas.append("s%02d  texto que cruza la zona de pie (%.2f > %.2f)"
                                 % (idx, t + alto, pie))
    if problemas:
        print("  AUDITORÍA: %d avisos" % len(problemas))
        for p in problemas:
            print("   ·", p)
    else:
        print("  AUDITORÍA: sin avisos")
    return problemas


def barrer_rayas(prs, saltar_idx=(1,)):
    """Cero «—» y «–» en lo que escribimos, y nada de «palanca»
    ([[deck-estilo-sin-rayas-ni-palanca]]).

    **Excluye la portada**: su titular trae un «—» y es copy de la empresa, no contenido
    nuestro, así que sin la exclusión el barrido reporta un falso positivo en cada corrida."""
    malos = []
    for idx, slide in enumerate(prs.slides, start=1):
        if idx in saltar_idx:
            continue
        textos = []
        for sh in slide.shapes:
            if sh.has_text_frame:
                textos.append(sh.text_frame.text)
            if sh.has_table:
                for row in sh.table.rows:
                    for c in row.cells:
                        textos.append(c.text)
        if slide.has_notes_slide:
            textos.append(slide.notes_slide.notes_text_frame.text)
        for txt in textos:
            for tok in ("—", "–", "palanca"):
                if tok in txt:
                    malos.append("s%02d  %r en «%s…»" % (idx, tok, txt[:50]))
    if malos:
        print("  ESTILO: %d avisos" % len(malos))
        for m in malos:
            print("   ·", m)
    else:
        print("  ESTILO: sin rayas ni «palanca» (portada excluida)")
    return malos


# ===========================================================================
# Las láminas
# ===========================================================================
def set_encabezado(gf, textos):
    """Traduce la fila de cabecera de una tabla del molde conservando su estilizado.

    El molde trae `Objective | Deliverable | Date | Status` a 16 pt bold blanco sobre 1B4F8C.
    Se reescribe el texto y nada más: relleno, bordes y márgenes son los que ya venían."""
    tbl = gf.table
    for ci, txt in enumerate(textos):
        _solo_run(tbl.cell(0, ci).text_frame.paragraphs[0], txt)
    return gf


def leyenda_circulos(slide, l, t, items, fs=9.5, lado=0.20):
    """Leyenda NATIVA de las dos láminas de recortes: el anillo sobre una muestra de tejido.

    Un anillo blanco sobre el fondo blanco de la lámina es invisible, así que cada muestra va
    dentro de un cuadradito del rosa del tejido. Es la única forma de que la leyenda sea
    nativa y editable en vez de quedar quemada dentro del PNG."""
    TEJIDO = RGBColor(0xD6, 0xB2, 0xC8)
    x = l
    for color, txt in items:
        _rect(slide, x, t, lado, lado, TEJIDO)
        ov = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x + 0.035), Inches(t + 0.035),
                                    Inches(lado - 0.07), Inches(lado - 0.07))
        ov.fill.background()
        ov.line.color.rgb = color
        ov.line.width = Pt(1.5)
        ov.shadow.inherit = False
        add_textbox(slide, x + lado + 0.09, t - 0.02, text_w(txt, fs) + 0.10, lado + 0.04,
                    [(txt, fs, False, CUERPO)], anchor=MSO_ANCHOR.MIDDLE)
        x += lado + 0.09 + text_w(txt, fs) + 0.34
    return t + lado


def poner_figura(slide, png, l, t, w_max, h_max):
    """Inserta el PNG escalado para entrar entero en la caja, centrado en el ancho.

    La figura es una **fotografía de un resultado** y por eso va como imagen y no como
    shapes: es la excepción declarada en CLAUDE.md §"Formato de entregables". Lo que la
    acompaña (título, cuerpo, leyenda y pie) sí es nativo."""
    from PIL import Image
    with Image.open(png) as im:
        wi, hi = im.size
    esc = min(w_max / float(wi), h_max / float(hi))
    w, h = wi * esc, hi * esc
    slide.shapes.add_picture(png, Inches(l + (w_max - w) / 2.0), Inches(t),
                             Inches(w), Inches(h))
    return t + h


# ---------------------------------------------------------------------------
def lamina_objetivos(s, guion):
    set_cejilla(s, PROYECTO)
    set_titulo(s, "OBJETIVOS " + PERIODO)
    gf = _shape(s, "Google Shape;196;p29")
    set_encabezado(gf, ["Objetivo", "Entregable", "Fecha", "Estado"])
    llenar_tabla(gf, [
        ("Escalar la detección de mitosis a 12 láminas",
         "Cruzada contra las 94 marcas del patólogo sobre las 12 láminas: 26 reencontradas "
         "a 30 µm, emparejamiento uno a uno, plano de 7,5 a 50 µm",
         "25/08", "Cerrado"),
        ("Evaluar HoVer-NeXt para necrosis",
         "GO con los pesos PanNuke, los únicos con clase de célula muerta. Cinco láminas, "
         "7 a 9 h de GPU. Sin lanzar: el nodo está tomado",
         "08/09", "Bloqueado"),
        ("Evaluar métricas nuevas para mitosis",
         "Ocho ejes puntuados, y los dos de procesador ya medidos: el control positivo separa "
         "epitelio de estroma con AUC 0,906 sobre 209 regiones, y el grado nuclear ordena con "
         "ρ +0,809 sobre 10 láminas",
         "08/09", "Cerrado"),
    ])
    notes(s, guion)


def lamina_mitosis(s, datos, agg, guion):
    set_cejilla(s, PROYECTO)
    set_titulo(s, "Mitosis: %d de %d marcas del patólogo" % (agg["tp"], agg["marcas"]),
               nombre="Text 1")
    fin = set_cuerpo(s, [
        ("Doce láminas anotadas, %d marcas: " % agg["marcas"],
         "%d reencontradas emparejando cada marca con la detección más cercana, si esa "
         "distancia no pasa de 30 µm. El número es plano de %s a %d µm de distancia."
         % (agg["tp"], num(agg["plano"][0]), int(agg["plano"][1]))),
        ("Una sola lámina aporta la mitad de los aciertos: ",
         "la 129741 aporta %d de los %d aciertos del total. Sin ella, las otras once "
         "reencuentran %d de sus %d marcas, o sea %s %%."
         % (agg["ref_tp"], agg["tp"], agg["resto_tp"],
                                      agg["resto_marcas"],
                                      num(100 * agg["resto_recall"]))),
    ])
    pie = ["Los 30 µm son una distancia entre la marca y el centroide de la detección, no el "
           "tamaño de nada.",
           "Las doce láminas están a 0,465 µm/px, así que 30 µm son 65 píxeles de nivel 0.",
           "El detector corrió sobre la lámina completa: nada se filtró antes de contar.",
           "El denominador es lo que el patólogo marcó, no las mitosis que hay en la lámina. "
           "Contra positivos parciales no se calcula precisión ni F1."]
    l, _, w = G_CUERPO
    h_pie = _alto_bloque(pie, w - 0.06, PT_PIE, space_after=1.5)
    top = fin + 0.16
    alto_fig = T_PIE_S03 - 0.12 - h_pie - 0.10 - top
    barras_marcas(s, l, top, w, alto_fig, datos)
    pie_lineas(s, l, top + alto_fig + 0.10, w, pie)
    notes(s, guion)


def lamina_aciertos(s, agg, guion):
    """Los 26 recortes acreditados: la marca del patólogo y la detección que la acredita."""
    set_cejilla(s, PROYECTO)
    set_titulo(s, "Las %d marcas que el detector reencuentra" % agg["tp"], nombre="Text 1")
    fin = set_cuerpo(s, [
        ("Las %d marcas acreditadas, agrupadas por lámina: " % agg["tp"],
         "cinco de las doce ponen alguna, y la 129741 pone la mitad."),
    ])
    l, _, w = G_CUERPO
    pie = ["Cada recorte mide %s µm de lado, a la magnificación nativa de la lámina. "
           "El emparejamiento es uno a uno: ninguna detección acredita dos marcas."
           % num(LADO_RECORTE_UM)]
    h_pie = _alto_bloque(pie, w - 0.06, PT_PIE, space_after=1.5)
    top = leyenda_circulos(s, l, fin + 0.10,
                           [(BLANCO_ANILLO, "marca del patólogo"),
                            (AMARILLO, "detección de HoVer-NeXt")]) + 0.12
    alto_fig = T_PIE_S03 - 0.12 - h_pie - 0.10 - top
    fin_fig = poner_figura(s, PNG_ACIERTOS, l, top, w, alto_fig)
    pie_lineas(s, l, fin_fig + 0.12, w, pie)
    notes(s, guion)


def lamina_falladas(s, agg, guion):
    """Las 68 que se escapan: es la mitad del cruce que muestra en qué se equivoca."""
    n = agg["marcas"] - agg["tp"]
    set_cejilla(s, PROYECTO)
    set_titulo(s, "Las %d marcas que se escapan" % n, nombre="Text 1")
    fin = set_cuerpo(s, [
        ("Las %d marcas de las doce que quedaron sin detección: " % n,
         "no hubo ninguna mitosis detectada dentro de la tolerancia."),
    ])
    l, _, w = G_CUERPO
    pie = ["Que una marca no se reencuentre no la vuelve un error del patólogo, y una "
           "detección sin marca encima no es un falso positivo: no hay con qué distinguirla "
           "de una mitosis real sin marcar."]
    h_pie = _alto_bloque(pie, w - 0.06, PT_PIE, space_after=1.5)
    top = leyenda_circulos(s, l, fin + 0.10,
                           [(BLANCO_ANILLO, "marca del patólogo")]) + 0.12
    alto_fig = T_PIE_S03 - 0.12 - h_pie - 0.10 - top
    fin_fig = poner_figura(s, PNG_FALLADAS, l, top, w, alto_fig)
    pie_lineas(s, l, fin_fig + 0.12, w, pie)
    notes(s, guion)


def lamina_f1(s, d4, guion):
    """Eje 4: la fracción epitelial por clase. Es el chequeo que no depende del estadístico."""
    set_cejilla(s, PROYECTO)
    set_titulo(s, "El control positivo separa epitelio de estroma", nombre="Text 1")
    fin = set_cuerpo(s, [
        ("%d regiones del patólogo en las doce láminas: " % d4["n_reg"],
         "la fracción epitelial cuenta núcleos de HoVer-NeXt dentro del polígono, y las dos "
         "clases de estroma dan cero exacto en la mediana."),
    ])
    l, _, w = G_CUERPO
    pie = ["Unidad: una REGIÓN por fila. La barra es el rango intercuartil y la marca oscura "
           "la mediana; %d regiones quedaron sin ningún núcleo adentro y salen del cálculo. "
           "«CDIS_papilar» se comporta como estroma estando en el grupo epitelio, y son 6 "
           "marcas de una sola lámina: se deja anotado, no se interpreta."
           % d4["sin_nucleos"]]
    top, alto, y_pie = caja_figura(fin, pie, w)

    w_lab, w_num, gap = 2.05, 1.15, 0.16
    x0 = l + w_lab + gap
    w_bar = w - w_lab - w_num - 2 * gap
    filas = d4["clases"]
    # 0,52 y no 0,25: la línea con sus ticks mide 0,25 y debajo va el rótulo del eje. Sin
    # reservarlo, el rótulo se monta sobre el pie.
    h_eje = 0.52
    paso = (alto - h_eje - 0.30) / len(filas)
    h_b = min(0.20, paso - 0.10)

    y_leg = top
    cx = x0
    for col, txt in ((CUERPO, "grupo epitelio"), (SEP, "grupo estroma")):
        _rect(s, cx, y_leg + 0.045, 0.16, 0.11, col)
        # el ancho es el MEDIDO, no un 1,7" holgado: una caja de más se cruza con la del
        # item siguiente y ensucia cualquier chequeo de intersecciones con falsos positivos
        add_textbox(s, cx + 0.21, y_leg, text_w(txt, 8.5) + 0.10, 0.19,
                    [(txt, 8.5, False, CUERPO)], anchor=MSO_ANCHOR.MIDDLE)
        cx += 0.21 + text_w(txt, 8.5) + 0.32
    add_textbox(s, l + w - w_num, y_leg, w_num, 0.19,
                [("regiones", 8.5, True, CUERPO, PP_ALIGN.RIGHT)], anchor=MSO_ANCHOR.MIDDLE)

    y0 = top + 0.30
    for i, f in enumerate(filas):
        y = y0 + i * paso
        yc = y + (paso - h_b) / 2.0
        col = CUERPO if f["grupo"] == "epitelio" else SEP
        _rect(s, x0, yc, w_bar, h_b, LINEA)
        xa = x0 + w_bar * f["p25"]
        xb = x0 + w_bar * f["p75"]
        _rect(s, xa, yc, max(xb - xa, 0.02), h_b, col)
        _rect(s, x0 + w_bar * f["med"] - 0.022, yc - 0.035, 0.044, h_b + 0.07, TITULO)
        add_textbox(s, l, y, w_lab, paso, [(f["clase"], 9.5, True, CUERPO, PP_ALIGN.RIGHT)],
                    anchor=MSO_ANCHOR.MIDDLE)
        add_textbox(s, l + w - w_num, y, w_num, paso,
                    [("%d" % f["n"], 9.5, False, CUERPO, PP_ALIGN.RIGHT)],
                    anchor=MSO_ANCHOR.MIDDLE)
    y_eje = y0 + len(filas) * paso + 0.06
    fin_eje = eje_x(s, x0, y_eje, w_bar, 0.0, 1.0, [0.0, 0.25, 0.5, 0.75, 1.0], dec=2)
    # Todo eje dibujado lleva rótulo. Es el chequeo barato que ataca el defecto real: un eje
    # sin nombre pasa las cuatro capas de QA y aun así no se entiende
    # ([[deck-qa-puntos-ciegos-chequeo]]).
    add_textbox(s, x0, fin_eje + 0.02, w_bar, 0.20,
                [("fracción de núcleos epiteliales dentro de la región "
                  "(0 = ninguno · 1 = todos)", 8.5, False, CUERPO, PP_ALIGN.CENTER)],
                anchor=MSO_ANCHOR.MIDDLE)
    pie_lineas(s, l, y_pie, w, pie)
    notes(s, guion)


def lamina_f2(s, d4, guion):
    """Eje 4: el estadístico contra su propio nulo. El número necesita su figura."""
    set_cejilla(s, PROYECTO)
    set_titulo(s, "Ninguna traslación del nulo llega al observado", nombre="Text 1")
    fin = set_cuerpo(s, [
        ("AUC de rango %s, contra un nulo por traslación rígida: " % num(d4["obs"], 3),
         "las %d traslaciones se centran en %s y su percentil 97,5 queda en %s, más de veinte "
         "puntos por debajo." % (len(d4["nulo"]), num(d4["nulo"].mean(), 3),
                                 num(float(np.percentile(d4["nulo"], 97.5)), 3))),
    ])
    l, _, w = G_CUERPO
    pie = ["Unidad: una REGIÓN por observación. El nulo traslada la máscara de cada región "
           "sobre el tejido de su lámina y nunca permuta etiquetas: los polígonos son "
           "contiguos.",
           "El p es el PISO alcanzable: con %d traslaciones el mínimo es 1 entre %d, así que "
           "el resultado se lee «por debajo de 1 en %d» y no como un valor exacto. Y el nulo "
           "se centra en %s y no en 0,5 porque las regiones de epitelio son más grandes que "
           "las de estroma, así que al trasladarlas no muestrean el mismo fondo."
           % (len(d4["nulo"]), len(d4["nulo"]) + 1, len(d4["nulo"]) + 1,
              num(d4["nulo"].mean(), 3))]
    top, alto, y_pie = caja_figura(fin, pie, w)

    l_h, w_h = l + 0.70, w - 1.10
    h_h = alto - 0.55
    histograma(s, l_h, top + 0.22, w_h, h_h - 0.22, d4["nulo"], 0.0, 1.0, 50, LINEA)
    for v, col, txt, lado in (
            (float(d4["nulo"].mean()), SEP, "nulo %s" % num(d4["nulo"].mean(), 3), "izq"),
            (float(np.percentile(d4["nulo"], 97.5)), SEP, "p97,5 %s"
             % num(float(np.percentile(d4["nulo"], 97.5)), 3), "der"),
            (float(d4["obs"]), ACENTO, "observado %s" % num(d4["obs"], 3), "izq")):
        marcador_vertical(s, l_h + w_h * v, top + 0.22, h_h - 0.22, txt, col, lado=lado)
    eje_x(s, l_h, top + h_h + 0.06, w_h, 0.0, 1.0, [0.0, 0.25, 0.5, 0.75, 1.0], dec=2)
    add_textbox(s, l_h, top + h_h + 0.32, w_h, 0.20,
                [("AUC de rango (epitelio > estroma)", 8.5, False, CUERPO, PP_ALIGN.CENTER)],
                anchor=MSO_ANCHOR.TOP)
    pie_lineas(s, l, y_pie, w, pie)
    notes(s, guion)


def lamina_f3(s, d3, guion):
    """Eje 3: el ordenamiento, con el `n` por lámina VISIBLE. Es lo que exige el handoff §6.3."""
    R = d3["restringida"]
    set_cejilla(s, PROYECTO)
    set_titulo(s, "Los tres grados ordenan sobre diez láminas", nombre="Text 1")
    fin = set_cuerpo(s, [
        ("Percentil del área dentro de la población epitelial de la propia lámina: ",
         "%s · %s · %s de mediana, con ρ = %s entre las %d láminas."
         % (num(R["grados"][0]["med"]), num(R["grados"][1]["med"]),
            num(R["grados"][2]["med"]), num(R["rho"], 3), R["n_lam"])),
    ])
    l, _, w = G_CUERPO
    pie = ["Unidad: una LÁMINA por punto, en la mediana de sus marcas. El grado está "
           "confundido con la lámina sin un solo cruce, así que comparar grados es comparar "
           "láminas y el n honesto son láminas, no las %d marcas." % R["n_marcas"],
           "Los AUC por lámina de esta población no se citan: valen 1,000 con n = 1 de un "
           "lado. Y los percentiles son altos en los tres grados, «bajo» incluido: el "
           "patólogo marca núcleos grandes para su lámina en cualquier grado, y lo que "
           "separa es cuánto."]
    top, alto, y_pie = caja_figura(fin, pie, w)

    # -- izquierda: el strip vertical, un punto por lámina
    w_izq = 5.30
    x_eje = l + 0.72
    h_pl = alto - 0.62
    t_pl = top + 0.24
    eje_y(s, x_eje, t_pl, h_pl, 0, 100, [0, 25, 50, 75, 100])
    add_textbox(s, l - 0.10, t_pl - 0.26, 3.2, 0.20,
                [("percentil intra-lámina", 8.5, False, CUERPO)], anchor=MSO_ANCHOR.TOP)
    w_pl = w_izq - (x_eje - l) - 0.20
    d = 0.155
    for gi, g in enumerate(d3["orden"]):
        sub = R["pl"][R["pl"]["grado"] == g].sort_values("med")
        cx = x_eje + w_pl * (gi + 0.5) / 3.0
        # mediana del grado: tick horizontal, y el conector que muestra el ordenamiento
        mg = R["grados"][gi]["med"]
        ym = t_pl + h_pl * (1.0 - mg / 100.0)
        _rect(s, cx - 0.34, ym - 0.011, 0.68, 0.022, ACENTO)
        n = len(sub)
        for j, (_, row) in enumerate(sub.iterrows()):
            jitter = 0.0 if n == 1 else (j / float(n - 1) - 0.5) * min(0.62, 0.14 * n)
            y = t_pl + h_pl * (1.0 - row["med"] / 100.0)
            punto(s, cx + jitter, y, d, CUERPO, hueco=not row["alineada"])
        add_textbox(s, cx - 0.75, t_pl + h_pl + 0.10, 1.50, 0.19,
                    [(g, 9.5, True, CUERPO, PP_ALIGN.CENTER)], anchor=MSO_ANCHOR.TOP)
        add_textbox(s, cx - 0.75, t_pl + h_pl + 0.29, 1.50, 0.19,
                    [("%d lámina%s · %d marcas" % (n, "" if n == 1 else "s",
                                                   R["grados"][gi]["n_marcas"]),
                      8.0, False, CUERPO, PP_ALIGN.CENTER)], anchor=MSO_ANCHOR.TOP)
    # La leyenda va DENTRO del panel y ABAJO. Arriba no: los puntos de `alto` viven en el
    # percentil 95 a 100 por construcción, así que el borde superior es la única banda que
    # está ocupada siempre, y ahí la leyenda quedaba a 0,08" de la nube, alineada con ella
    # y leyéndose como dos puntos más. No se cruzaban, así que ni `auditar` ni un chequeo de
    # intersecciones lo veían: es la proximidad, no el solape ([[deck-qa-puntos-ciegos-chequeo]]).
    # El hueco de abajo no se supone, se mide: se apoya bajo el punto más bajo que haya.
    y_leg = t_pl + h_pl * (1.0 - min(R["pl"]["med"]) / 100.0) + 0.30
    leyenda_puntos(s, x_eje + 0.30, min(y_leg, t_pl + h_pl - 0.24),
                   [(CUERPO, False, "offset alineado"),
                    (CUERPO, True, "alineada: false")], fs=8.5)

    # -- derecha: las diez láminas, que es de dónde sale cada punto
    x_t = l + w_izq + 0.30
    w_t = w - w_izq - 0.30
    filas = []
    for g in d3["orden"]:
        sub = R["pl"][R["pl"]["grado"] == g].sort_values("med")
        for _, row in sub.iterrows():
            filas.append([g, row["slide"], "%d" % row["n"], num(row["med"]),
                          "sí" if row["alineada"] else "NO"])
    tabla_ejes(s, x_t, top, w_t,
               ["grado", "lámina", "marcas", "percentil", "alineada"], filas,
               [0.21, 0.26, 0.17, 0.21, 0.15], fs=9.0, row_h=0.265)
    pie_lineas(s, l, y_pie, w, pie)
    notes(s, guion)


def lamina_f4(s, d3, guion):
    """Eje 3: los dos nulos exactos, que es donde se ve qué población despega y cuál no."""
    R, C = d3["restringida"], d3["completa"]
    set_cejilla(s, PROYECTO)
    set_titulo(s, "El nulo exacto, y la población que no despega", nombre="Text 1")
    fin = set_cuerpo(s, [
        ("Las dos poblaciones se declararon antes y se reportan juntas: ",
         "la restringida a epitelio da p = %s sobre %d láminas y la completa queda en %s "
         "sobre %d, y el pre-registro dice que manda el nivel lámina."
         % (num(R["p_bi"], 4), R["n_lam"], num(C["p_bi"], 4), C["n_lam"])),
    ])
    l, _, w = G_CUERPO
    n_igual = int((np.abs(R["rhos"]) >= abs(R["obs"]) - 1e-12).sum())
    pie = ["Unidad: una ASIGNACIÓN del grado a las láminas por observación. El nulo permuta a "
           "nivel lámina, donde las unidades son intercambiables bajo la nula; a nivel marca "
           "no lo sería, porque las marcas de una lámina no son independientes.",
           "En la restringida el p también es el PISO: la asignación observada es la única de "
           "las %d que ordena perfecto, así que %s de las %d igualan o superan su ρ en valor "
           "absoluto y %s es el mínimo que este reparto puede dar. La lectura honesta es "
           "«ordena, con p = %s sobre %d láminas en la población limpia y %s sobre %d en la "
           "completa», no «da significativo»."
           % (R["k"], n_igual, R["k"], num(R["p_bi"], 4), num(R["p_bi"], 4), R["n_lam"],
              num(C["p_bi"], 4), C["n_lam"])]
    top, alto, y_pie = caja_figura(fin, pie, w)

    l_h, w_h = l + 0.70, w - 1.30
    h_par = (alto - 0.30) / 2.0
    for i, B in enumerate((R, C)):
        t_h = top + i * h_par
        # 0,60 y no 0,46: con 0,46 las etiquetas del eje del panel de arriba entraban 0,05"
        # en la caja del rótulo del de abajo. Es alto de histograma que se cede a propósito.
        h_b = h_par - 0.60
        histograma(s, l_h, t_h + 0.22, w_h, h_b, B["rhos"], -1.0, 1.0, 41, LINEA)
        x_obs = l_h + w_h * (B["obs"] + 1.0) / 2.0
        marcador_vertical(s, x_obs, t_h + 0.22, h_b,
                          "ρ = %s   p = %s" % (num(B["obs"], 3), num(B["p_bi"], 4)),
                          ACENTO, lado="izq")
        add_textbox(s, l_h, t_h, w_h, 0.20,
                    [("%s · %d láminas · %d asignaciones"
                      % (B["etiqueta"], B["n_lam"], B["k"]), 9.5, True, CUERPO)],
                    anchor=MSO_ANCHOR.MIDDLE)
        eje_x(s, l_h, t_h + 0.22 + h_b + 0.04, w_h, -1.0, 1.0, [-1.0, -0.5, 0.0, 0.5, 1.0], dec=1)
    # Debajo de los ticks del eje de abajo, no encima: a `- 0.10` pisaba el «0,0».
    add_textbox(s, l_h, top + 2 * h_par + 0.06, w_h, 0.20,
                [("ρ de Spearman entre el grado y el percentil, por lámina",
                  8.5, False, CUERPO, PP_ALIGN.CENTER)], anchor=MSO_ANCHOR.TOP)
    pie_lineas(s, l, y_pie, w, pie)
    notes(s, guion)


def lamina_tareas(s, guion):
    set_cejilla(s, PROYECTO)
    set_titulo(s, "Tareas del próximo período")
    gf = _shape(s, "Google Shape;196;p29")
    set_encabezado(gf, ["Objetivo", "Entregable", "Fecha"])
    llenar_tabla(gf, [
        ("Correr el eje de necrosis",
         "Densidad de detecciones de célula muerta dentro contra fuera de los polígonos del "
         "patólogo, con el nulo por traslación de la máscara. Antes hay que unificar el "
         "vocabulario",
         "08/09"),
        ("Punto caliente mitótico",
         "Conteo por mm², mapa de densidad y la ventana de área fija que lo maximiza: la "
         "primera versión de una zona para proponerle al patólogo",
         "15/09"),
    ], min_h=1.14)
    notes(s, guion)


def escribir_numeros(d4, d3):
    """Todo número que quede dibujado, en un CSV versionado: la figura tiene que ser auditable
    sin abrir el .pptx, que es derivado y está gitignored."""
    filas = []
    for c in d4["clases"]:
        filas.append(["eje4", "f_epi por clase", c["clase"], "region", c["n"],
                      num(c["med"], 3), num(c["p25"], 3), num(c["p75"], 3)])
    filas.append(["eje4", "AUC de rango", "epitelio > estroma", "region", d4["n_reg"],
                  num(d4["obs"], 3), "", ""])
    filas.append(["eje4", "nulo por traslacion", "media / p97,5", "region", len(d4["nulo"]),
                  num(float(d4["nulo"].mean()), 3),
                  num(float(np.percentile(d4["nulo"], 97.5)), 3), ""])
    filas.append(["eje4", "p", "piso 1/(n+1)", "region", len(d4["nulo"]),
                  num(d4["p"], 4), "", ""])
    filas.append(["eje4", "AUC solo alineadas", "epitelio > estroma", "region", d4["n_alin"],
                  num(d4["auc_alin"], 3), "", ""])
    for k, B in (("restringida", d3["restringida"]), ("completa", d3["completa"])):
        for g in B["grados"]:
            filas.append(["eje3", "percentil intra-lamina (%s)" % k, g["grado"], "marca",
                          g["n_marcas"], num(g["med"]), num(g["p25"]), num(g["p75"])])
        for _, row in B["pl"].iterrows():
            filas.append(["eje3", "mediana por lamina (%s)" % k,
                          "%s / %s" % (row["grado"], row["slide"]), "lamina", int(row["n"]),
                          num(row["med"]), "alineada" if row["alineada"] else "NO alineada",
                          ""])
        filas.append(["eje3", "rho por lamina (%s)" % k, "grado vs percentil", "lamina",
                      B["n_lam"], num(B["rho"], 3), "", ""])
        filas.append(["eje3", "nulo exacto (%s)" % k, "p bilateral / unilateral",
                      "asignacion", B["k"], num(B["p_bi"], 4), num(B["p_uni"], 4), ""])
        for pr in B["pares"]:
            filas.append(["eje3", "AUC por marca (%s)" % k, pr["par"], "marca",
                          pr["n_neg"] + pr["n_pos"], num(pr["auc"], 3), "", ""])
    with open(CSV_NUM, "w", newline="") as fh:
        wr = csv.writer(fh)
        wr.writerow(["eje", "medida", "grupo", "unidad", "n", "valor", "p25_o_aux",
                     "p75_o_aux"])
        wr.writerows(filas)
    return len(filas)


# ===========================================================================
def main():
    datos, agg = leer_datos()
    d4 = datos_eje4()
    d3 = datos_eje3()
    R = d3["restringida"]
    guion, _ = leer_guion()
    print("Deck del período · %s · %s" % (PROYECTO, PERIODO))
    print("  mitosis: %d láminas, %d marcas, %d detecciones, %d TP (%.1f %%)"
          % (agg["n_laminas"], agg["marcas"], agg["det"], agg["tp"], 100 * agg["recall"]))
    print("  eje 4: %d regiones (%d epi / %d estroma), AUC %.3f, p %.4f"
          % (d4["n_reg"], d4["n_epi"], d4["n_est"], d4["obs"], d4["p"]))
    print("  eje 3: %d marcas en %d láminas, rho %.3f, p exacto %.4f sobre %d asignaciones"
          % (R["n_marcas"], R["n_lam"], R["rho"], R["p_bi"], R["k"]))

    prs = Presentation(TPL)
    s01, s02, s03, s04 = list(prs.slides)

    sA = clonar_s03(prs, s03)       # el número
    sB = clonar_s03(prs, s03)       # los 26 recortes acreditados
    sC = clonar_s03(prs, s03)       # las 68 que se escapan
    sE = clonar_s03(prs, s03)       # F1  el control positivo por clase
    sG = clonar_s03(prs, s03)       # F3  los tres grados sobre diez láminas

    lamina_objetivos(s02, guion["s02"])
    lamina_mitosis(sA, datos, agg, guion["s03a"])
    lamina_aciertos(sB, agg, guion["s03b"])
    lamina_falladas(sC, agg, guion["s03c"])
    lamina_f1(sE, d4, guion["s03d"])
    lamina_f3(sG, d3, guion["s03f"])
    lamina_tareas(s04, guion["s04"])
    notes(s01, guion["s01"])

    borrar_slide(prs, s03)                      # trae el ejemplo de otra persona
    # Tres láminas salen por pedido de Ernesto tras la reunión del 1-sep (D1): los dos
    # histogramas de nulo (F2 y F4) y la tabla de los ocho ejes. `lamina_f2` y `lamina_f4`
    # NO se borran del archivo: las importa `ejes_nucleares/figuras/generate_figuras_ejes34.py`,
    # que es el envoltorio que re-renderiza las cuatro figuras de los ejes para QA visual.
    # `lamina_ejes` sí se borró (nadie más la usa); su tabla de los ocho ejes queda como
    # única fuente en `hovernext_tareas/inventario_tareas.md` §4.
    reordenar(prs, [s01, s02, sA, sB, sC, sE, sG, s04])

    forzar_barlow(prs)
    problemas = auditar(prs)
    problemas += barrer_rayas(prs)

    n = escribir_numeros(d4, d3)
    prs.save(OUT)
    print("  %d números dibujados -> %s" % (n, os.path.basename(CSV_NUM)))
    print("  escrito: %s" % os.path.basename(OUT))
    return 1 if problemas else 0


if __name__ == "__main__":
    sys.exit(main())
