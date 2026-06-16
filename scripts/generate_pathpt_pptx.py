#!/usr/bin/env python
"""generate_pathpt_pptx.py — diagrama de ARQUITECTURA de PathPT (.pptx EDITABLE).

Rev. 6 (feedback Ernesto 15-jun): RAIL DE ÁLGEBRA DE DIMENSIONES a la derecha de la
cascada — la transformación de forma en CADA flecha/transición, estilo CLAM
(`[N×512]·[512×C]=[N×C]`), que es lo que faltaba. Rama de texto a la IZQUIERDA (como en
Fig 1b del paper, donde converge en el matching). Cada bloque lleva su ECUACIÓN del paper
(ec. 4/5/6/8). Sin título interno (lo pone el header del deck), sin bullets, sin solapes.

Ecuaciones (paper He et al. 2025 — Methods §4.3-4.8, Fig 1b; verbatim del PDF):
  vᵢ = Φᵥ(xᵢ) ∈ ℝ^d                                            (4.1/4.2, extractor visión, congelado)
  v̄₁..v̄_M = Ψ(Φᵥ(x₁)..Φᵥ(x_M))                                (ec. 4, módulo espacial: conv 3/5/7 ⊕ transformer)
  c̄ = [T]₁…[T]ₖ [CLASS] ; [T]ᵢ ∈ ℝ^(d×1) ; K=32               (ec. 5, prompt-tuning CoOp)
  p(y=j|xᵢ) = softmaxⱼ( ⟨Φₜ(c̄ⱼ), v̄ᵢ⟩ / τ )                    (ec. 6, clasificación por parche)
  ŷ = tumor-ratio(P) ;  mᵢ = argmaxⱼ Pᵢⱼ                       (ec. 8 BACC eval-protocol + grounding)
Notación del paper: M=nº tiles, d=dim, N=nº clases. En el deck usamos N=nº parches,
512=d (CONCH), C=nº clases (coherente con Diagrama_CLAM.pptx que usa N=parches).

Color: azul = pipeline (como CLAM); gris = CONGELADO (CONCH); naranjo = ENTRENABLE (θ).

Uso: PYTHONPATH=/media/administrador/Storage1/sdonoso/clam_testing2/.pylibs \
     /home/sdonoso/miniconda3/envs/clam_latest/bin/python scripts/generate_pathpt_pptx.py
"""
import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn

REPO = "/media/administrador/Storage1/sdonoso/clam_testing2/oncomets-ernesto"
DST = os.path.join(REPO, "papers/presentations/Diagrama_PathPT.pptx")

BLUE_F = RGBColor(0x6E, 0x9B, 0xC5)
BLUE_E = RGBColor(0x3C, 0x6A, 0x95)
NAVY   = RGBColor(0x14, 0x2A, 0x42)
FRZ_F  = RGBColor(0xE3, 0xE7, 0xEB)
FRZ_E  = RGBColor(0x9A, 0xA6, 0xB2)
FRZ_T  = RGBColor(0x4A, 0x54, 0x60)
ORA_F  = RGBColor(0xFB, 0xE2, 0xC8)
ORA_E  = RGBColor(0xE2, 0x72, 0x3B)
ORA_T  = RGBColor(0xB4, 0x52, 0x1E)
DIMF   = RGBColor(0xEC, 0xF1, 0xF5)
DIME   = RGBColor(0xB7, 0xC6, 0xD4)
DIM_T  = RGBColor(0x1F, 0x4E, 0x5F)
INK    = RGBColor(0x22, 0x22, 0x22)
GREY_T = RGBColor(0x55, 0x55, 0x55)
CARLITO = "Carlito"


def set_text(shape, lines, align=PP_ALIGN.CENTER):
    tf = shape.text_frame
    tf.word_wrap = True; tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_left = tf.margin_right = Pt(3); tf.margin_top = tf.margin_bottom = Pt(1)
    for i, (txt, sz, bold, col) in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        r = p.add_run(); r.text = txt
        r.font.size = Pt(sz); r.font.bold = bold; r.font.name = CARLITO; r.font.color.rgb = col


def box(slide, l, t, w, h, lines, fill, edge, lw=1.5):
    sp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(l), Inches(t), Inches(w), Inches(h))
    sp.fill.solid(); sp.fill.fore_color.rgb = fill
    sp.line.color.rgb = edge; sp.line.width = Pt(lw); sp.shadow.inherit = False
    set_text(sp, lines)
    return sp


def callout(slide, l, t, w, h, lines):
    return box(slide, l, t, w, h, lines, DIMF, DIME, lw=1)


def arrow(slide, l, t, w, h, direction="down"):
    shp = {"down": MSO_SHAPE.DOWN_ARROW, "right": MSO_SHAPE.RIGHT_ARROW,
           "left": MSO_SHAPE.LEFT_ARROW}[direction]
    a = slide.shapes.add_shape(shp, Inches(l), Inches(t), Inches(w), Inches(h))
    a.fill.solid(); a.fill.fore_color.rgb = BLUE_E
    a.line.fill.background(); a.shadow.inherit = False
    return a


def connector(slide, x1, y1, x2, y2):
    c = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(x1), Inches(y1), Inches(x2), Inches(y2))
    c.line.color.rgb = DIME; c.line.width = Pt(1.25); c.shadow.inherit = False
    return c


def edge(slide, x1, y1, x2, y2, color=BLUE_E, w=2.25):
    """Conector recto con punta de flecha al final (aristas del árbol/cascada)."""
    c = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(x1), Inches(y1), Inches(x2), Inches(y2))
    c.line.color.rgb = color; c.line.width = Pt(w); c.shadow.inherit = False
    ln = c.line._get_or_add_ln()
    ln.append(ln.makeelement(qn('a:tailEnd'), {'type': 'triangle', 'w': 'med', 'len': 'med'}))
    return c


def textbox(slide, l, t, w, h, lines, align=PP_ALIGN.LEFT):
    tb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    set_text(tb, lines, align=align)
    return tb


# ============================================================================
# Slide 1 — CASCADA / 3 RAMAS (rev 14): igual que rev 13 (modular + callouts + paneles)
# pero TEXTO MÁS GRANDE y bloques más grandes (visión fusiona Φᵥ+W_proj → 4 bloques,
# simétrico con texto), y el centro vacío se llena con el BACKBONE CONCH compartido
# (Φᵥ y Φₜ = los dos encoders del mismo CONCH congelado), conectado a ambas ramas.
# Fuentes de código: spatial.py / prompt.py / pathpt.py (pin 0ab7f1b).
# ============================================================================

PVIS = RGBColor(0xEC, 0xF2, 0xF8)   # panel visión (azul claro)
PTXT = RGBColor(0xFB, 0xF3, 0xEA)   # panel texto (naranjo claro)
PMAT = RGBColor(0xEE, 0xF1, 0xF3)   # panel matching (gris claro)


def build_slide1(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    hb, pit = 0.66, 0.88

    def mod(x, t, w, label, kind):
        f, e, tc = {"b": (BLUE_F, BLUE_E, NAVY), "f": (FRZ_F, FRZ_E, FRZ_T),
                    "o": (ORA_F, ORA_E, ORA_T)}[kind]
        box(s, x, t, w, hb, [(label, 13, True, tc)], f, e, lw=2)

    def panel(x, t, w, h, fill, title, tcol):
        box(s, x, t, w, h, [], fill, DIME, lw=1)
        textbox(s, x, t + 0.07, w, 0.34, [(title, 14, True, tcol)], align=PP_ALIGN.CENTER)

    # ===================== RAMA VISIÓN (panel izq) =====================
    panel(0.10, 0.92, 5.05, 4.36, PVIS, "RAMA VISIÓN · θᵥ", BLUE_E)
    xV, wV, cV = 2.66, 2.32, 3.82
    xLc, wLc = 0.24, 2.28
    yv = [1.55 + i * pit for i in range(4)]
    mod(xV, yv[0], wV, "N parches + coords", "b")
    mod(xV, yv[1], wV, "Φᵥ · CONCH + W_proj", "f")
    mod(xV, yv[2], wV, "θᵥ · conv 3/5/7 (local)", "o")
    mod(xV, yv[3], wV, "θᵥ · NyströmAttn (global)", "o")
    for i in range(3):
        edge(s, cV, yv[i] + hb, cV, yv[i + 1])
    for i, lines in [
        (0, [("xₖ ∈ ℝ^(256×256×3)", 10.5, True, INK), ("+ coords (x,y)", 9.5, False, GREY_T)]),
        (1, [("vᵢ = hᵢ · W_proj ∈ ℝ⁵¹²", 10.5, True, INK), ("V ∈ ℝ^(N×512)", 9.5, True, DIM_T)]),
        (2, [("oₖ = ReLU(LN(Convₖ V))", 10, True, INK), ("k=3,5,7 · o₁+o₂+o₃+V", 9, False, GREY_T)]),
        (3, [("V̄ = LN(NyströmAttn(·))", 10, True, INK), ("8h · 256 landmk · 6 pinv", 9, False, GREY_T)]),
    ]:
        yc = yv[i] + hb / 2
        callout(s, xLc, yv[i] + 0.02, wLc, 0.62, lines)
        connector(s, xLc + wLc, yc, xV, yc)

    # ===================== RAMA TEXTO (panel der) =====================
    panel(8.18, 0.92, 5.04, 4.36, PTXT, "RAMA TEXTO · θₜ", ORA_T)
    xT, wT, cT = 8.38, 2.32, 9.54
    xRc, wRc = 10.88, 2.30
    yt = [1.55 + i * pit for i in range(4)]
    mod(xT, yt[0], wT, "θₜ · ctx CoOp", "o")
    mod(xT, yt[1], wT, "ensamblar prompt", "o")
    mod(xT, yt[2], wT, "Φₜ · text-transformer", "f")
    mod(xT, yt[3], wT, "proyección · W_text", "f")
    for i in range(3):
        edge(s, cT, yt[i] + hb, cT, yt[i + 1])
    for i, lines in [
        (0, [("ctx ∈ ℝ^(C×32×768)", 10.5, True, INK), ('init "a histopath. image of"', 9, False, GREY_T)]),
        (1, [("[SOS | ctx×32 | CLS+EOS]", 10, True, INK), ("→ ℝ^(C×127×768)", 9.5, True, DIM_T)]),
        (2, [("+pos +CLS +attn causal", 10, True, INK), ("ln_final · pooled ∈ ℝ⁷⁶⁸", 9, False, GREY_T)]),
        (3, [("tⱼ = pooled · W_text ∈ ℝ⁵¹²", 10, True, INK), ("T ∈ ℝ^(C×512)", 9.5, True, DIM_T)]),
    ]:
        yc = yt[i] + hb / 2
        callout(s, xRc, yt[i] + 0.02, wRc, 0.62, lines)
        connector(s, xT + wT, yc, xRc, yc)

    # ===================== BACKBONE CONCH (centro-arriba, llena el medio) =====================
    box(s, 5.45, 1.55, 2.42, 1.34,
        [("CONCH", 16, True, NAVY), ("modelo visión–lenguaje", 10.5, False, GREY_T),
         ("Φᵥ + Φₜ · congelados", 11, True, FRZ_T)], FRZ_F, FRZ_E, lw=2)
    connector(s, 4.98, yv[1] + hb / 2, 5.45, 2.10)   # CONCH → Φᵥ (visión)
    connector(s, 7.87, 2.30, 8.38, yt[2] + hb / 2)   # CONCH → Φₜ (texto)

    # ===================== RAMA MATCHING (panel centro-abajo) =====================
    panel(4.96, 5.00, 3.40, 2.42, PMAT, "RAMA MATCHING", NAVY)
    xM, wM, cM = 5.14, 3.04, 6.66
    ym = [5.46 + i * 0.62 for i in range(3)]
    box(s, xM, ym[0], wM, 0.54,
        [("MATCHING POR PARCHE · ec. 6", 11.5, True, NAVY),
         ("L2-norm → logits = V̄·Tᵀ ∈ ℝ^(N×C)", 10, True, DIM_T)], BLUE_F, BLUE_E, lw=2)
    box(s, xM, ym[1], wM, 0.54,
        [("P = softmax(logits · 10)", 12, True, NAVY), ("τ = 0.1", 9.5, False, NAVY)], BLUE_F, BLUE_E, lw=2)
    box(s, xM, ym[2], wM, 0.54,
        [("P ∈ ℝ^(N×C) — una clase por parche", 11, True, NAVY)], DIMF, DIME, lw=1.5)
    edge(s, cM, ym[0] + 0.54, cM, ym[1])
    edge(s, cM, ym[1] + 0.54, cM, ym[2])

    # convergencia VISIÓN/TEXTO → MATCHING
    edge(s, cV, yv[3] + hb, cM - 0.95, ym[0])
    edge(s, cT, yt[3] + hb, cM + 0.95, ym[0])
    textbox(s, 3.55, 4.86, 1.50, 0.26, [("V̄ : N×512", 10.5, True, DIM_T)], align=PP_ALIGN.CENTER)
    textbox(s, 8.30, 4.86, 1.50, 0.26, [("T : C×512", 10.5, True, DIM_T)], align=PP_ALIGN.CENTER)

    # ===================== leyenda (esquina inf-izq) =====================
    box(s, 0.34, 5.62, 0.30, 0.26, [], FRZ_F, FRZ_E, lw=1)
    textbox(s, 0.72, 5.56, 2.6, 0.38, [("congelado (CONCH)", 11.5, False, GREY_T)], align=PP_ALIGN.LEFT)
    box(s, 0.34, 6.16, 0.30, 0.26, [], ORA_F, ORA_E, lw=1)
    textbox(s, 0.72, 6.10, 2.8, 0.38, [("entrenable (θᵥ, θₜ)", 11.5, False, GREY_T)], align=PP_ALIGN.LEFT)
    textbox(s, 0.34, 6.66, 4.4, 0.60,
            [("N = parches · C = clases · 512 = dim contrastivo", 10, False, GREY_T),
             ("768 = dim token (ctx) · τ = 0.1 (escala = 10)", 10, False, GREY_T)],
            align=PP_ALIGN.LEFT)

    return s


# ============================================================================
# Slide 2 — los 3 componentes entrenables (zoom, referencia)
# ============================================================================

def _panel(s, x, header, steps):
    W = 4.0
    box(s, x, 1.20, W, 0.72, [(header, 14, True, ORA_T)], ORA_F, ORA_E, lw=2.25)
    y = 2.30
    for i, st in enumerate(steps):
        h = 0.82
        box(s, x + 0.25, y, W - 0.5, h, st, DIMF, DIME, lw=1)
        y += h
        if i < len(steps) - 1:
            arrow(s, x + W / 2 - 0.17, y - 0.02, 0.34, 0.30)
            y += 0.30


def build_slide2(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _panel(s, 0.45, "θᵥ · MÓDULO ESPACIAL", [
        [("parches en grilla 2D (x, y)", 12, True, INK)],
        [("conv 3×3 / 5×5 / 7×7", 12, True, INK), ("vecindario local", 9.5, False, GREY_T)],
        [("transformer (global)", 12, True, INK)],
        [("ṽᵢ refinado · N×512", 12, True, DIM_T)],
    ])
    _panel(s, 4.65, "θₜ · PROMPT-TUNING", [
        [("[T]₁ [T]₂ ⋯ [T]ₖ [CLASE]", 12, True, INK)],
        [("Φₜ · CONCH texto (congelado)", 11.5, True, FRZ_T)],
        [("tⱼ vector de clase · C×512", 12, True, DIM_T)],
        [("aprende la frase, no la escribe", 11, True, ORA_T)],
    ])
    _panel(s, 8.85, "PSEUDO-LABELS (tile)", [
        [("coseno( parche , frase )", 12, True, INK)],
        [("clase tentativa del parche", 12, True, INK)],
        [("retener si concuerda con slide", 11.5, True, INK)],
        [("→ CE balanceada por clase", 12, True, DIM_T)],
    ])
    return s


def main():
    prs = Presentation()
    prs.slide_width = Inches(13.333); prs.slide_height = Inches(7.5)
    build_slide1(prs); build_slide2(prs)
    prs.save(DST)
    chk = Presentation(DST)
    has = any(x.has_text_frame and "MATCHING" in x.text for x in chk.slides[0].shapes)
    print(f"OK  {os.path.relpath(DST, REPO)}  slides={len(chk.slides)}  "
          f"s1_shapes={len(chk.slides[0].shapes)}  matching={has}")


if __name__ == "__main__":
    main()
