#!/usr/bin/env python
"""Extrae el guion del presentador de un .pptx a UN archivo markdown.

Por qué existe: las notas viven repartidas en una llamada `notes()` por lámina dentro del
generador, y leídas ahí, de a una, **los defectos entre láminas son invisibles** — racimos de
ritmo, una promesa que la lámina siguiente desmiente, y sobre todo el mismo contenido contado
dos veces. Extraerlas a un archivo corrido no es comodidad, es el método
([[humanizer-es-skill]] §3). En la pasada del 5-ago-2026 el hallazgo mayor fue justamente uno
de esos: dos láminas del B8 contaban el mismo montaje con veinte segundos de diferencia.

Se venía reescribiendo a mano en cada sesión de QA de deck (B7, B8). Acá queda una vez.

Uso:
    PYTHONPATH=/media/administrador/Storage1/sdonoso/clam_testing2/.pylibs \
    /home/sdonoso/miniconda3/envs/clam_latest/bin/python \
      scripts/extraer_guion_deck.py <deck.pptx> [salida.md]

Sin `salida.md` escribe al lado del .pptx como `<nombre>_guion.md`.

Lee el .pptx, no el generador: así refleja lo que efectivamente quedó en el deck.
"""
import sys
import os

sys.path.insert(0, "/media/administrador/Storage1/sdonoso/clam_testing2/.pylibs")

from pptx import Presentation


def titulo_de(slide):
    """Primer texto no trivial de la lámina, para poder ubicarla sin abrir el deck."""
    for sh in slide.shapes:
        if sh.has_text_frame and sh.text_frame.text.strip():
            t = sh.text_frame.text.strip().split("\n")[0]
            if len(t) > 2:
                return t[:80]
    return "(sin texto)"


def notas_de(slide):
    if not slide.has_notes_slide:
        return ""
    return slide.notes_slide.notes_text_frame.text.strip()


def main():
    if len(sys.argv) < 2:
        sys.exit(__doc__)
    deck = sys.argv[1]
    salida = sys.argv[2] if len(sys.argv) > 2 else os.path.splitext(deck)[0] + "_guion.md"

    prs = Presentation(deck)
    bloques, total = [], 0

    for i, s in enumerate(prs.slides, 1):
        n = notas_de(s)
        pal = len(n.split())
        total += pal
        bloques.append(
            "## Lámina {} · {}\n\n_({} palabras)_\n\n{}\n".format(
                i, titulo_de(s), pal, n if n else "**(SIN NOTAS)**"))

    cabecera = (
        "# Guion del presentador · {}\n\n"
        "{} láminas, {} palabras (~{:.0f} min hablados a 130 pal/min).\n\n"
        "Extraído del `.pptx` para QA de prosa. **La fuente canónica son las llamadas "
        "`notes(s, ...)` del generador**: editar acá no cambia el deck.\n\n"
        "Al leerlo de corrido, buscar primero **arcos** (qué se cuenta dos veces, qué promesa "
        "quedó sin dueño) y recién después tells de prosa.\n\n---\n\n"
    ).format(os.path.basename(deck), len(prs.slides), total, total / 130.0)

    with open(salida, "w", encoding="utf-8") as fh:
        fh.write(cabecera + "\n---\n\n".join(bloques))

    print("{} láminas, {} palabras -> {}".format(len(prs.slides), total, salida))
    for i, s in enumerate(prs.slides, 1):
        print("  {:2d}  {:4d} pal  {}".format(i, len(notas_de(s).split()), titulo_de(s)[:60]))


if __name__ == "__main__":
    main()
