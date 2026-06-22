#!/usr/bin/env python
"""generate_mammoth_fused_slide.py — slide DEDICADA a la figura OFICIAL de mammoth
con callouts de entrada/salida dibujados ENCIMA de la imagen (overlay sobre los bloques).

Evolución de la "Opción A" (figura arriba / cinta de dims abajo). Ernesto pidió ahora:
  - figura oficial `mammoth_architecture.png` a (casi) SLIDE COMPLETA (es panorámica, aspecto 2.49),
  - anotaciones de entrada→salida de cada paso dibujadas JUSTO ENCIMA de la imagen, posicionadas
    sobre los bloques del pipeline (encoder → W → cabezas → MoE → cross-head concat → MIL/CLAM).

La figura es la EXCEPCIÓN bendecida del deck para figuras de paper → va como imagen embebida.
Los callouts son shapes NATIVOS editables (familia teal = forma del tensor; naranjo = mammoth).

Dims code-accurate (verificadas contra MAMMOTH/src/mammoth/mammoth.py + models_mammoth/, keep_slots=False):
  z[N,512] → q=LN(Wq·z)[N,256]=16cab×16 → S[30,16,10,16], D=softmax↓N → 300 slots →
  expertos LoRA (A:16→8 compartido, Be:8→32 por experto, rank 8) → o[300,512] →
  C=softmax↓300 → h[N,512] (keep_slots=False preserva los N) → CLAM (MIL).

DEBUG_GRID=True dibuja una grilla de fracciones sobre la imagen para medir posiciones de bloques.

Requiere python-pptx (en /media/administrador/Storage1/sdonoso/clam_testing2/.pylibs).
Uso:
  PYTHONPATH=/media/administrador/Storage1/sdonoso/clam_testing2/.pylibs \
    /home/sdonoso/miniconda3/envs/clam_latest/bin/python scripts/generate_mammoth_fused_slide.py
"""
import os
import sys

sys.path.insert(0, os.path.dirname(os.path.abspath(__file__)))  # para importar helpers

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.enum.text import PP_ALIGN

from generate_clam_mammoth_pptx import (
    add_box, add_arrow, add_label, add_connector, set_text,
    TEAL_F, ORA_F, ORA_E, BLU_F, BLU_E, INK, GREY_T, DIM_F, DIM_E,
)

REPO = "/media/administrador/Storage1/sdonoso/clam_testing2/oncomets-ernesto"
SRC = os.path.join(REPO, "papers/presentations/Diagrama_CLAM.pptx")  # plantilla (tamaño + master)
DST = os.path.join(REPO, "papers/presentations/Diagrama_mammoth_fused.pptx")
FIG = "/media/administrador/Storage1/sdonoso/clam_testing2/MAMMOTH/assets/mammoth_architecture.png"

TEAL_E = RGBColor(0x2C, 0x7A, 0x8C)
TITLE_TEAL = RGBColor(0x21, 0x75, 0x89)   # = TEAL_TITLE del deck (armoniza títulos al acoplar)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
GRID_C = RGBColor(0xCC, 0x33, 0x33)

SLIDE_W, SLIDE_H = 13.3333, 7.5
FIG_ASPECT = 4375 / 1758  # = 2.489

# ---- placement de la figura (casi slide completa) ----
FIG_MARGIN = 0.12
FIG_W = SLIDE_W - 2 * FIG_MARGIN          # 13.093
FIG_H = FIG_W / FIG_ASPECT                # 5.260
FIG_X = FIG_MARGIN
# FIG_Y deja el título arriba Y despeja el logo (0.82) del header del deck al acoplar
# (copy_diagram con bar=False igual estampa el logo top-left).
FIG_Y = 0.88

DEBUG_GRID = False


def fx2in(fx):
    """fracción horizontal [0,1] de la figura → pulgadas en la slide."""
    return FIG_X + fx * FIG_W


def fy2in(fy):
    """fracción vertical [0,1] de la figura → pulgadas en la slide."""
    return FIG_Y + fy * FIG_H


def _blank(prs):
    for lay in prs.slide_layouts:
        if lay.name and lay.name.lower() == "blank":
            return lay
    return prs.slide_layouts[-1]


def draw_grid(s):
    """Grilla de fracciones sobre la figura (solo DEBUG): verticales cada 0.05, horizontales cada 0.10."""
    # verticales
    fx = 0.0
    while fx <= 1.0001:
        c = s.shapes.add_connector(MSO_CONNECTOR.STRAIGHT,
                                   Inches(fx2in(fx)), Inches(FIG_Y),
                                   Inches(fx2in(fx)), Inches(FIG_Y + FIG_H))
        c.line.color.rgb = GRID_C
        c.line.width = Pt(0.75 if abs((fx * 20) % 2) < 0.01 else 0.4)
        c.shadow.inherit = False
        add_label(s, fx2in(fx) - 0.12, FIG_Y - 0.20, 0.24, 0.18,
                  [(f"{int(round(fx*100))}", 7, True, GRID_C)], align=PP_ALIGN.CENTER)
        fx += 0.05
    # horizontales
    fy = 0.0
    while fy <= 1.0001:
        c = s.shapes.add_connector(MSO_CONNECTOR.STRAIGHT,
                                   Inches(FIG_X), Inches(fy2in(fy)),
                                   Inches(FIG_X + FIG_W), Inches(fy2in(fy)))
        c.line.color.rgb = GRID_C
        c.line.width = Pt(0.4)
        c.shadow.inherit = False
        add_label(s, FIG_X - 0.30, fy2in(fy) - 0.09, 0.26, 0.18,
                  [(f"{int(round(fy*100))}", 7, True, GRID_C)], align=PP_ALIGN.CENTER)
        fy += 0.10


def add_callout(s, cx_fx, top_fy, lines, anchor_fx, anchor_fy, edge,
                cw=1.46, ch=0.54, fill=WHITE, lw=1.5):
    """Chip de entrada→salida en overlay sobre la figura, centrado en cx_fx y con su tope en
    top_fy [fracciones de la figura], con una flecha fina al bloque que anota en
    (anchor_fx, anchor_fy). Funciona arriba o abajo del bloque: el conector nace del centro
    del chip y queda tapado por el propio chip hasta su borde → la flecha emerge hacia el
    bloque. cw/ch en pulgadas."""
    left = fx2in(cx_fx) - cw / 2.0
    top = fy2in(top_fy)
    cy_in = top + ch / 2.0
    add_connector(s, fx2in(cx_fx), cy_in, fx2in(anchor_fx), fy2in(anchor_fy),
                  color=edge, w=1.25, arrow=True)
    box = add_box(s, left, top, cw, ch, lines, fill, edge, lw=lw)  # tapa el origen del conector
    return box


def build(prs):
    s = prs.slides.add_slide(_blank(prs))

    # ---- título (compacto, deja aire a la fila de callouts) ----
    add_label(s, 0.0, 0.05, SLIDE_W, 0.36,
              [("MAMMOTH — flujo de datos sobre la arquitectura oficial", 16, True, TITLE_TEAL)],
              align=PP_ALIGN.CENTER)

    # ---- figura OFICIAL del paper (imagen), casi slide completa ----
    s.shapes.add_picture(FIG, Inches(FIG_X), Inches(FIG_Y), Inches(FIG_W), Inches(FIG_H))

    if DEBUG_GRID:
        draw_grid(s)
        return s

    # ---- callouts overlay (mixto): fila superior para los pasos lineales del pipeline +
    #      pegados donde hay espacio claro (MIL·CLAM y el zoom del experto).
    #      Posiciones en fracciones de la figura, calibradas con la grilla DEBUG. ----
    TOP = 0.006
    # fila superior: encoder → W → ruteo → expertos → concat
    add_callout(s, 0.135, TOP,
                [("ENCODER (CONCH)", 9.5, True, TEAL_E),
                 ("parches → z [N, 512]", 9.5, True, INK)],
                anchor_fx=0.165, anchor_fy=0.16, edge=TEAL_E)
    add_callout(s, 0.255, TOP,
                [("PROYECCIÓN  W_q", 9.5, True, ORA_E),
                 ("z → q [N, 256] = 16×16", 9, True, INK)],
                anchor_fx=0.205, anchor_fy=0.14, edge=ORA_E)
    add_callout(s, 0.375, TOP,
                [("RUTEO POR SLOTS", 9.5, True, ORA_E),
                 ("D = softmax_N⟨q,S⟩ → 300", 8.5, True, INK)],
                anchor_fx=0.400, anchor_fy=0.18, edge=ORA_E)
    add_callout(s, 0.495, TOP,
                [("EXPERTOS  (LoRA r=8)", 9.5, True, BLU_E),
                 ("300 slots → o [300, 512]", 9, True, INK)],
                anchor_fx=0.460, anchor_fy=0.16, edge=BLU_E)
    add_callout(s, 0.615, TOP,
                [("CROSS-HEAD CONCAT", 9, True, ORA_E),
                 ("o → h [N, 512]", 9.5, True, INK)],
                anchor_fx=0.515, anchor_fy=0.13, edge=ORA_E)
    add_callout(s, 0.735, TOP,
                [("MIL · CLAM", 9.5, True, INK),
                 ("h → logits clínicos", 9.5, True, TEAL_E)],
                anchor_fx=0.620, anchor_fy=0.205, edge=TEAL_E)
    # pegado (mixto): zoom del experto LoRA en el hueco inferior-izq, apunta a la derecha
    #                 al detalle MoE (Expert 1: slots → media pond. → Φ → W → no-linealidad)
    add_callout(s, 0.250, 0.595,
                [("DENTRO DEL EXPERTO  (LoRA)", 9.5, True, BLU_E),
                 ("slots → media pond. → Φ·W → no-lin.", 8.5, True, INK),
                 ("A: 16→8 (comp.) · B_e: 8→32 · r = 8", 8.5, True, GREY_T)],
                anchor_fx=0.315, anchor_fy=0.700, edge=BLU_E, cw=2.28, ch=0.78)

    # ---- nota de fuente + leyenda (excepción de paper = imagen; sin nombres propios) ----
    add_label(s, 0.0, FIG_Y + FIG_H + 0.04, SLIDE_W, 0.24,
              [("figura oficial del paper  ·  callouts = forma del tensor entrada → salida por paso  "
                "(verificado contra el código · keep-slots = False · N = nº de parches de la slide)",
                10, False, GREY_T)], align=PP_ALIGN.CENTER)
    add_label(s, 0.0, FIG_Y + FIG_H + 0.28, SLIDE_W, 0.22,
              [("teal = forma del tensor   ·   naranjo = MAMMOTH (entrenable)   ·   azul = expertos",
                9.5, False, GREY_T)], align=PP_ALIGN.CENTER)
    return s


def build_variant(prs):
    """Slide 2 (companion): la variante keep_slots. Tronco compartido (recap con la misma
    notación/colores de la fused) → bifurcación en los 2 modos de salida. NATIVO, sin imagen.
    Mecanismo code-accurate (clam_mammoth.py + prereg §1); SIN resultados, SIN nombres/jobs."""
    s = prs.slides.add_slide(_blank(prs))
    GREY_E = RGBColor(0x9A, 0xA6, 0xB2)
    PANEL_BASE_F = RGBColor(0xF1, 0xF4, 0xF6)   # gris claro = modo base (ya cerrado)
    PANEL_NEW_F = RGBColor(0xFD, 0xF1, 0xE5)    # naranjo claro = variante nueva (resalta)
    SOFT = RGBColor(0x6B, 0x74, 0x80)
    FORK_F = RGBColor(0xFF, 0xF3, 0xD6)

    # ---- título ----
    add_label(s, 0.0, 0.12, SLIDE_W, 0.46,
              [("MAMMOTH — la variante keep\\_slots: dónde se bifurca la salida", 18, True, TITLE_TEAL)],
              align=PP_ALIGN.CENTER)

    # ---- tronco compartido (recap compacto, = slide anterior) ----
    add_label(s, 0.0, 0.66, SLIDE_W, 0.26,
              [("TRONCO COMPARTIDO  ·  idéntico a la slide anterior  (encoder → W → ruteo → expertos)",
                11.5, True, SOFT)], align=PP_ALIGN.CENTER)
    bw, bh, gap, n = 2.62, 0.92, 0.30, 4
    x0 = (SLIDE_W - (n * bw + (n - 1) * gap)) / 2.0
    y_tr = 1.00
    trunk = [
        (TEAL_F, TEAL_E, [("PARCHES (CONCH)", 12, True, TEAL_E), ("z : [N, 512]", 13, True, INK)]),
        (ORA_F, ORA_E, [("PROYECCIÓN  W_q", 12, True, ORA_E), ("q : [N, 256] = 16×16", 12, True, INK)]),
        (ORA_F, ORA_E, [("RUTEO POR SLOTS", 12, True, ORA_E), ("D = softmax_N⟨q,S⟩ → 300", 10.5, True, INK)]),
        (BLU_F, BLU_E, [("EXPERTOS (LoRA r=8)", 12, True, BLU_E), ("o : [300, 512]", 13, True, INK)]),
    ]
    for i, (f, e, lines) in enumerate(trunk):
        x = x0 + i * (bw + gap)
        add_box(s, x, y_tr, bw, bh, lines, f, e, lw=1.75)
        if i < n - 1:
            add_arrow(s, x + bw + 0.005, y_tr + bh / 2 - 0.13, gap - 0.01, 0.26)
    last_cx = x0 + (n - 1) * (bw + gap) + bw / 2.0   # centro de EXPERTOS

    # ---- nodo de bifurcación (centrado, bajo el tronco) ----
    fork_w, fork_h = 2.60, 0.64
    fork_x = (SLIDE_W - fork_w) / 2.0
    fork_y = 2.18
    fork_cx = fork_x + fork_w / 2.0
    add_connector(s, last_cx, y_tr + bh, fork_cx, fork_y, color=GREY_E, w=2.0)
    add_box(s, fork_x, fork_y, fork_w, fork_h,
            [("keep\\_slots = ?", 13, True, INK),
             ("¿recombinar a parches o quedarse con los slots?", 8.5, False, SOFT)],
            FORK_F, ORA_E, lw=2)

    # ---- dos paneles ----
    pw, ph, py = 5.78, 2.74, 3.28
    lx = 0.60
    rx = SLIDE_W - 0.60 - pw
    add_box(s, lx, py, pw, ph, [("", 1, False, INK)], PANEL_BASE_F, GREY_E, lw=1.5)
    add_box(s, rx, py, pw, ph, [("", 1, False, INK)], PANEL_NEW_F, ORA_E, lw=2.75)
    # fan del fork a cada panel
    add_connector(s, fork_cx, fork_y + fork_h, lx + pw / 2.0, py, color=GREY_E, w=2.0)
    add_connector(s, fork_cx, fork_y + fork_h, rx + pw / 2.0, py, color=ORA_E, w=2.5)

    def fill_panel(px, header, boxA_lines, boxB_lines, caption, hcol, boxA_fill, boxA_edge):
        sw, sx = pw - 0.66, px + 0.33
        add_label(s, px + 0.18, py + 0.09, pw - 0.36, 0.30, [header], align=PP_ALIGN.CENTER)
        add_box(s, sx, py + 0.46, sw, 0.84, boxA_lines, boxA_fill, boxA_edge, lw=1.75)
        add_arrow(s, px + pw / 2.0 - 0.16, py + 1.34, 0.32, 0.18, down=True)
        add_box(s, sx, py + 1.56, sw, 0.52, boxB_lines, TEAL_F, TEAL_E, lw=1.5)
        add_label(s, px + 0.18, py + 2.14, pw - 0.36, 0.52, [caption], align=PP_ALIGN.CENTER)

    # panel izquierdo: keep_slots = FALSE (base, ya cerrado)
    fill_panel(
        lx,
        ("keep\\_slots = FALSE  —  mammoth normal (base)", 12.5, True, GREY_T),
        [("COMBINACIÓN   C = softmax↓300", 11.5, True, TEAL_E),
         ("recombina los 300 slots →", 10.5, False, INK),
         ("h : [N, 512]", 12.5, True, INK)],
        [("CLAM agrega sobre los N PARCHES", 12, True, INK)],
        ("drop-in de la capa lineal  ·  cardinalidad N → N  ·  misma atención e instance-loss que el baseline",
         9.5, False, SOFT),
        GREY_T, TEAL_F, TEAL_E)

    # panel derecho: keep_slots = TRUE (variante nueva, en entrenamiento)
    fill_panel(
        rx,
        ("keep\\_slots = TRUE  —  variante NUEVA (en entrenamiento)", 12.5, True, ORA_E),
        [("SE SALTA la recombinación →", 11.5, True, ORA_E),
         ("300 SLOT-TOKENS : [300, 512]", 12.5, True, INK),
         ("+ slot\\_dropout : enmascara el ruteo en train (no-op en eval)", 9, False, ORA_E)],
        [("CLAM agrega sobre los 300 SLOTS", 12, True, INK)],
        ("cuello de botella aprendido N → 300  ·  atención en 2 etapas (slots → CLAM)",
         9.5, False, SOFT),
        ORA_E, ORA_F, ORA_E)

    # ---- nota inferior (mecanismo + por qué; SIN resultados) ----
    add_label(s, 0.40, 6.18, SLIDE_W - 0.80, 0.62,
              [("Misma cabeza CLAM y misma pérdida — cambia SOLO qué agrega:  N parches (False)  vs  300 slot-tokens (True).",
                11.5, True, INK),
               ("Por qué la variante: un slot puede concentrar su masa en pocos parches raros → capacidad dedicada a la "
                "clase minoritaria (objetivo: recuperar recall).  Se evalúa pareado vs el modo base.",
                10, False, GREY_T)], align=PP_ALIGN.CENTER)
    return s


def main():
    prs = Presentation(SRC)
    build(prs)            # slide 1 = figura oficial fusionada (overlay)
    build_variant(prs)    # slide 2 = variante keep_slots (nativa)
    # eliminar la slide 0 heredada de la plantilla → quedan [fused, variante]
    xml_slides = prs.slides._sldIdLst
    xml_slides.remove(list(xml_slides)[0])
    prs.save(DST)

    chk = Presentation(DST)
    info = []
    for sl in chk.slides:
        npic = sum(1 for sh in sl.shapes if sh.shape_type == 13)
        nsp = sum(1 for sh in sl.shapes if sh.has_text_frame)
        info.append(f"(pics={npic}, text={nsp})")
    print(f"OK  {os.path.relpath(DST, REPO)}  slides={len(chk.slides)}  " + "  ".join(info)
          + f"  DEBUG_GRID={DEBUG_GRID}")


if __name__ == "__main__":
    main()
