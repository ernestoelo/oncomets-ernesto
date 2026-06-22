#!/usr/bin/env python
"""generate_pathpt_fused_slide.py — slide DEDICADA a la figura OFICIAL de PathPT (panel b)
con callouts de entrada/salida dibujados ENCIMA de la imagen (overlay sobre los bloques).

Mismo tratamiento que la fused de mammoth (scripts/generate_mammoth_fused_slide.py): la figura
oficial del paper va como IMAGEN embebida (excepción bendecida del deck para figuras de paper)
y las anotaciones de entrada→salida de cada paso van como shapes NATIVOS editables encima.

Figura: pathpt_fig1_b.png = panel b (PathPT) del paper (He et al., Nat. Commun. 2025).
Arquitectura code-accurate (funcionamiento_pathpt.md §2-3, models_pathpt pin 0ab7f1b):
  VISIÓN:  parches → Φ_v (CONCH, congelado) → v[512] → θ_v (módulo espacial: grilla 2D +
           conv 3/5/7 + transformer, ENTRENABLE) → v̄[512] refinado por parche (ec. 4)
  TEXTO:   prompt aprendible c̄ = [T]₁..[T]_K [CLASE] (θ_t, K=32, ENTRENABLE) → Φ_t (CONCH
           texto, congelado) → t[512] por clase (ec. 5)
  MATCHING: p(y=j|x) = softmax(⟨Φ_t(c̄_j), v̄⟩/τ) → clase POR PARCHE (ec. 6) →
           Tile Classification + grounding del tumor (lo que CLAM no da)
Convención de color (= diagrama nativo del deck): gris = CONGELADO (CONCH Φ_v/Φ_t);
naranjo = ENTRENABLE (θ_v, θ_t); teal = forma del tensor / salida.

DEBUG_GRID=True dibuja una grilla de fracciones para medir posiciones de bloques.

Requiere python-pptx (en /media/administrador/Storage1/sdonoso/clam_testing2/.pylibs).
Uso:
  PYTHONPATH=/media/administrador/Storage1/sdonoso/clam_testing2/.pylibs \
    /home/sdonoso/miniconda3/envs/clam_latest/bin/python scripts/generate_pathpt_fused_slide.py
"""
import os
import sys

sys.path.insert(0, os.path.dirname(os.path.abspath(__file__)))  # para importar helpers

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_CONNECTOR
from pptx.enum.text import PP_ALIGN

from generate_clam_mammoth_pptx import (
    add_box, add_connector, add_label,
    ORA_F, ORA_E, INK, GREY_T,
)

REPO = "/media/administrador/Storage1/sdonoso/clam_testing2/oncomets-ernesto"
SRC = os.path.join(REPO, "papers/presentations/Diagrama_CLAM.pptx")  # plantilla (tamaño + master)
DST = os.path.join(REPO, "papers/presentations/Diagrama_pathpt_fused.pptx")
FIG = os.path.join(REPO, "papers/presentations/assets_branding/paper_figs/pathpt_fig1_b.png")

TITLE_TEAL = RGBColor(0x21, 0x75, 0x89)
TEAL_E = RGBColor(0x2C, 0x7A, 0x8C)
FROZEN_E = RGBColor(0x6B, 0x74, 0x80)     # gris = congelado (CONCH)
FROZEN_F = RGBColor(0xF1, 0xF4, 0xF6)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
GRID_C = RGBColor(0xCC, 0x33, 0x33)

SLIDE_W, SLIDE_H = 13.3333, 7.5
# fig1_b.png (935x550) trae la arquitectura (panel b) arriba + barras de resultados (c/d)
# abajo → recortamos el 32% inferior y nos quedamos SOLO con la arquitectura.
CROP_BOTTOM = 0.19   # la arquitectura (incl. Tile Classification + grounding) llega a ~y0.81;
                     # los resultados c/d empiezan después → recortar solo eso
FIG_ASPECT = 935 / (550 * (1 - CROP_BOTTOM))   # ~2.10

# ---- placement de la figura (centrada, con franjas arriba/abajo para los callouts) ----
# La figura tiene 2 ramas (visión arriba / texto abajo) → callouts FUERA de la imagen:
# franja superior para la rama visión (apuntan ↓), franja inferior para la rama texto (apuntan ↑).
FIG_W = 10.00
FIG_H = FIG_W / FIG_ASPECT                  # = 4.00
FIG_X = (SLIDE_W - FIG_W) / 2.0             # 1.667
FIG_Y = 1.36                                # franja superior arriba; logo del deck despejado

DEBUG_GRID = False


def fx2in(fx):
    return FIG_X + fx * FIG_W


def fy2in(fy):
    return FIG_Y + fy * FIG_H


def _blank(prs):
    for lay in prs.slide_layouts:
        if lay.name and lay.name.lower() == "blank":
            return lay
    return prs.slide_layouts[-1]


def draw_grid(s):
    fx = 0.0
    while fx <= 1.0001:
        c = s.shapes.add_connector(MSO_CONNECTOR.STRAIGHT,
                                   Inches(fx2in(fx)), Inches(FIG_Y),
                                   Inches(fx2in(fx)), Inches(FIG_Y + FIG_H))
        c.line.color.rgb = GRID_C
        c.line.width = Pt(0.75 if abs((fx * 20) % 2) < 0.01 else 0.4)
        c.shadow.inherit = False
        add_label(s, fx2in(fx) - 0.12, FIG_Y - 0.20, 0.24, 0.18,
                  [(f"{int(round(fx*100))}", 7, True, GRID_C)], align=PP_ALIGN.CENTER)
        fx += 0.05
    fy = 0.0
    while fy <= 1.0001:
        c = s.shapes.add_connector(MSO_CONNECTOR.STRAIGHT,
                                   Inches(FIG_X), Inches(fy2in(fy)),
                                   Inches(FIG_X + FIG_W), Inches(fy2in(fy)))
        c.line.color.rgb = GRID_C
        c.line.width = Pt(0.4)
        c.shadow.inherit = False
        add_label(s, FIG_X - 0.30, fy2in(fy) - 0.09, 0.26, 0.18,
                  [(f"{int(round(fy*100))}", 7, True, GRID_C)], align=PP_ALIGN.CENTER)
        fy += 0.10


def add_callout(s, cx_in, top_in, lines, anchor_fx, anchor_fy, edge,
                cw=1.95, ch=0.56, fill=WHITE, lw=1.5):
    """Chip de entrada→salida colocado FUERA de la figura, centrado en cx_in con su tope en
    top_in [pulgadas en la slide], con flecha fina al bloque en (anchor_fx, anchor_fy)
    [fracciones de la figura]. El conector nace del centro del chip y queda tapado por el
    chip → la flecha emerge hacia el bloque (funciona arriba o abajo de la figura)."""
    left = cx_in - cw / 2.0
    cy_in = top_in + ch / 2.0
    add_connector(s, cx_in, cy_in, fx2in(anchor_fx), fy2in(anchor_fy),
                  color=edge, w=1.25, arrow=True)
    add_box(s, left, top_in, cw, ch, lines, fill, edge, lw=lw)


def build(prs):
    s = prs.slides.add_slide(_blank(prs))

    add_label(s, 0.0, 0.05, SLIDE_W, 0.40,
              [("PathPT — flujo de datos sobre la arquitectura oficial", 18, True, TITLE_TEAL)],
              align=PP_ALIGN.CENTER)

    pic = s.shapes.add_picture(FIG, Inches(FIG_X), Inches(FIG_Y), Inches(FIG_W), Inches(FIG_H))
    pic.crop_bottom = CROP_BOTTOM   # quita las barras de resultados (c/d), deja la arquitectura

    if DEBUG_GRID:
        draw_grid(s)
        return s

    # ---- callouts FUERA de la figura, ordenados por rama ----
    # FRANJA SUPERIOR (rama VISIÓN + matching): chips arriba de la figura, apuntan ↓
    top_y = 0.56
    add_callout(s, fx2in(0.135), top_y,
                [("VISIÓN · Φ_v  (congelado)", 9.5, True, FROZEN_E),
                 ("parches → v [512] por parche", 9, True, INK)],
                anchor_fx=0.135, anchor_fy=0.30, edge=FROZEN_E, cw=2.05)
    add_callout(s, fx2in(0.355), top_y,
                [("MÓDULO ESPACIAL  θ_v  (entrenable)", 9, True, ORA_E),
                 ("refina cada parche (grilla 2D + transformer)", 8.5, True, INK)],
                anchor_fx=0.345, anchor_fy=0.28, edge=ORA_E, cw=2.55)
    add_callout(s, fx2in(0.605), top_y,
                [("MATCHING  (coseno)", 9.5, True, TEAL_E),
                 ("softmax ⟨t, v⟩ / τ → clase por parche", 8, True, INK)],
                anchor_fx=0.520, anchor_fy=0.27, edge=TEAL_E, cw=2.25)
    # FRANJA INFERIOR (rama TEXTO + salida): chips debajo de la figura, apuntan ↑
    bot_y = FIG_Y + FIG_H + 0.10            # ~5.46
    add_callout(s, fx2in(0.175), bot_y,
                [("RAMA TEXTO · θ_t  (entrenable)", 9, True, ORA_E),
                 ("prompt [T]_1 .. [T]_(32) [CLASE]", 9, True, INK),
                 ("→ Φ_t (CONCH, congelado) → t [512]", 8.5, True, FROZEN_E)],
                anchor_fx=0.235, anchor_fy=0.80, edge=ORA_E, cw=2.70, ch=0.74)
    add_callout(s, fx2in(0.620), bot_y,
                [("SALIDA POR PARCHE", 9.5, True, TEAL_E),
                 ("Tile Classification + grounding", 9, True, INK),
                 ("localiza el tumor (CLAM no lo da)", 8.5, False, GREY_T)],
                anchor_fx=0.600, anchor_fy=0.82, edge=TEAL_E, cw=2.45, ch=0.74)

    # ---- nota de fuente + leyenda (excepción de paper = imagen; sin nombres propios) ----
    add_label(s, 0.0, bot_y + 0.86, SLIDE_W, 0.22,
              [("figura oficial del paper (panel b)  ·  CONCH congelado: solo se entrenan θ_v y θ_t   "
                "·   gris = congelado (Φ_v / Φ_t)  ·  naranjo = entrenable (θ_v, θ_t)  ·  teal = salida",
                9.5, False, GREY_T)], align=PP_ALIGN.CENTER)
    return s


def build_merged_left(prs):
    """Versión 'columna izquierda' para la slide UNIDA del deck (12+14): la figura oficial
    anotada ocupa el ~68% izquierdo (deja la derecha para las 3 tarjetas de piezas que agrega
    el deck). SIN título propio (lo pone el header del deck) y callouts compactos para el ancho
    menor. Misma estructura por ramas que la standalone: visión arriba, texto abajo."""
    s = prs.slides.add_slide(_blank(prs))
    FX, FW = 0.30, 8.70
    FH = FW / FIG_ASPECT                # ~4.14 (figura COMPLETA, aspect ~2.10)
    FY = 1.46

    def lx(fx):
        return FX + fx * FW

    def ly(fy):
        return FY + fy * FH

    def callout(cx_in, top_in, lines, afx, afy, edge, cw, ch=0.56):
        add_connector(s, cx_in, top_in + ch / 2.0, lx(afx), ly(afy), color=edge, w=1.25, arrow=True)
        add_box(s, cx_in - cw / 2.0, top_in, cw, ch, lines, WHITE, edge, lw=1.5)

    pic = s.shapes.add_picture(FIG, Inches(FX), Inches(FY), Inches(FW), Inches(FH))
    pic.crop_bottom = CROP_BOTTOM      # solo recorta las barras de resultados c/d

    # franja superior (rama visión + matching) → apuntan ↓ (chips corridos para no encimarse)
    ty = 0.90
    callout(1.50, ty, [("Φ_v · VISIÓN (congelado)", 9, True, FROZEN_E),
                       ("parches → v [512]", 8.5, True, INK)], 0.15, 0.30, FROZEN_E, cw=1.85)
    callout(3.55, ty, [("θ_v · ESPACIAL (entrenable)", 8.5, True, ORA_E),
                       ("refina cada parche", 8.5, True, INK)], 0.36, 0.30, ORA_E, cw=1.95)
    callout(5.70, ty, [("MATCHING (coseno)", 9, True, TEAL_E),
                       ("clase por parche", 8.5, True, INK)], 0.50, 0.32, TEAL_E, cw=1.80)
    # franja inferior (rama texto + salida) → apuntan ↑
    by = FY + FH + 0.08                # ~5.68
    callout(2.55, by, [("θ_t · TEXTO (entrenable)", 8.5, True, ORA_E),
                       ("prompt → Φ_t (congelado)", 8.5, True, FROZEN_E),
                       ("→ t [512] por clase", 8.5, True, INK)], 0.33, 0.70, ORA_E, cw=2.35, ch=0.70)
    callout(6.10, by, [("SALIDA POR PARCHE", 9, True, TEAL_E),
                       ("tile + grounding del tumor", 8.5, True, INK)], 0.66, 0.74, TEAL_E, cw=2.15)
    # leyenda slim (bajo las franjas, columna izquierda)
    add_label(s, FX, by + 0.78, FW, 0.20,
              [("gris = congelado (CONCH)   ·   naranjo = entrenable (θ_v, θ_t)   ·   teal = salida",
                9, False, GREY_T)], align=PP_ALIGN.CENTER)
    return s


def main():
    prs = Presentation(SRC)
    build(prs)                 # slide 0 = fused standalone (full width)
    build_merged_left(prs)     # slide 1 = versión columna-izquierda para la slide unida del deck
    xml_slides = prs.slides._sldIdLst
    xml_slides.remove(list(xml_slides)[0])
    prs.save(DST)
    chk = Presentation(DST)
    info = []
    for sl in chk.slides:
        npic = sum(1 for sh in sl.shapes if sh.shape_type == 13)
        nsp = sum(1 for sh in sl.shapes if sh.has_text_frame)
        info.append(f"(pics={npic}, text={nsp})")
    print(f"OK  {os.path.relpath(DST, REPO)}  slides={len(chk.slides)}  " + "  ".join(info)
          + f"  DEBUG_GRID={DEBUG_GRID}")


if __name__ == "__main__":
    main()
