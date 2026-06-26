#!/usr/bin/env python
"""generate_b5_deck.py — ensambla CLAM_Sprint_B5.pptx (deck completo, branding Environ).

Rev. 4 (feedback Ernesto, 3ª ronda): TABLAS NATIVAS (no imágenes) y escalables en S7/
microcalc/cierre; diagrama PathPT matemático estilo Diagrama_CLAM; header consistente
(gris + logo + título Barlow) también en los diagramas; S2 sin "lecciones de B4";
S8 con gráfico formal (bal_acc + AUC); S10 (idea) fusionada con los 3 componentes;
títulos de resultados minimalistas/profesionales.

Deck 13.333x7.5 in. Speaker notes por slide (sin nº de job, sin nombres).

Uso: PYTHONPATH=/media/administrador/Storage1/sdonoso/clam_testing2/.pylibs \
     /home/sdonoso/miniconda3/envs/clam_latest/bin/python scripts/generate_b5_deck.py
"""
import copy
import io
import os

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn
from PIL import Image

REPO = "/media/administrador/Storage1/sdonoso/clam_testing2/oncomets-ernesto"
PRES = os.path.join(REPO, "papers/presentations")
ASSETS_BRAND = os.path.join(PRES, "assets_branding")
PAPER_FIGS = os.path.join(ASSETS_BRAND, "paper_figs")
# Fig.3 del paper de MAMMOTH (Shao et al., ICLR 2026), panels A+B sin caption:
# ruteo de parches a slots → especialización de expertos por fenotipo. Recortada de la pág. 8.
FIG3_MAM = os.path.join(PAPER_FIGS, "mammoth_fig3_routing.png")
# Fig.1 del mismo paper (pág. 2), panels A+B sin caption: A = t-SNE del espacio interno
# (lineal=continuo vs MAMMOTH=clusters por experto) · B = Δ rendimiento vs el lineal en 8 MIL.
FIG1_MAM = os.path.join(PAPER_FIGS, "mammoth_fig1_overview.png")
DST = os.path.join(PRES, "CLAM_Sprint_B5.pptx")

A_MAM = os.path.join(REPO, "sprints/B5_sprint5/objetivo_2_mammoth_patron_invasion/figuras/slide_assets")
A_PPT = os.path.join(REPO, "sprints/B5_sprint5/pathpt/figuras/slide_assets")
DIAG_MAM = os.path.join(PRES, "Diagrama_CLAM_mammoth.pptx")
DIAG_PPT = os.path.join(PRES, "Diagrama_PathPT.pptx")
# slide 0 = figura oficial fusionada (overlay de dims) · slide 1 = variante keep_slots.
# Generado por scripts/generate_mammoth_fused_slide.py (correr ANTES del deck).
DIAG_FUSED = os.path.join(PRES, "Diagrama_mammoth_fused.pptx")
# figura oficial de PathPT (panel b) con callouts overlay del forward.
# Generado por scripts/generate_pathpt_fused_slide.py (correr ANTES del deck).
DIAG_PPT_FUSED = os.path.join(PRES, "Diagrama_pathpt_fused.pptx")

# ---- paleta ----
TEAL_TITLE = RGBColor(0x21, 0x75, 0x89)
TEAL_SQ    = RGBColor(0x31, 0x85, 0x9C)
BAR_GRIS   = RGBColor(0xF2, 0xF2, 0xF2)
TEAL_DIV   = RGBColor(0x2E, 0x7E, 0x8F)
TEAL_CARD  = RGBColor(0xDD, 0xEA, 0xEE)
TEAL_CARD2 = RGBColor(0xF3, 0xF8, 0xF9)
LAV_TITLE  = RGBColor(0xCD, 0xD6, 0xF4)
TEAL_SUB   = RGBColor(0xB8, 0xD4, 0xD9)
ROW_ALT    = RGBColor(0xF2, 0xF4, 0xF5)
GRIS_BODY  = RGBColor(0x59, 0x59, 0x59)
VERDE      = RGBColor(0x1E, 0x84, 0x49)
ROJO       = RGBColor(0xC0, 0x39, 0x2B)
AMBAR      = RGBColor(0xB9, 0x77, 0x0E)
GRIS_TXT   = RGBColor(0x55, 0x55, 0x55)
ORA_T      = RGBColor(0xB4, 0x52, 0x1E)
INK        = RGBColor(0x22, 0x22, 0x22)
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)

F_TITLE = "Barlow ExtraBold"
F_BODY  = "Barlow"

SW, SH = 13.333, 7.5
HDR = 0.82
LOGO = os.path.join(ASSETS_BRAND, "logo_header.png")
PORTADA = os.path.join(ASSETS_BRAND, "portada_fullbleed.jpg")
_uid = [1000]


# ============================================================================
# Helpers base
# ============================================================================

def _blank(prs):
    for lay in prs.slide_layouts:
        if (lay.name or "").lower().strip() == "blank":
            return lay
    return prs.slide_layouts[6]


def _set_runs(tf, lines, anchor=MSO_ANCHOR.TOP):
    tf.word_wrap = True; tf.vertical_anchor = anchor
    for i, ln in enumerate(lines):
        txt, sz, bold, col = ln[0], ln[1], ln[2], ln[3]
        font = ln[4] if len(ln) > 4 else F_BODY
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = ln[5] if len(ln) > 5 else PP_ALIGN.LEFT
        p.space_after = Pt(6)
        r = p.add_run(); r.text = txt
        r.font.size = Pt(sz); r.font.bold = bold; r.font.name = font; r.font.color.rgb = col


def add_textbox(slide, l, t, w, h, lines, anchor=MSO_ANCHOR.TOP):
    tb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    _set_runs(tb.text_frame, lines, anchor=anchor)
    return tb


def add_notes(slide, text):
    slide.notes_slide.notes_text_frame.text = text


def set_notes(slide, proposito, sections, nsz=13):
    """Notas del presentador como GUION DIDÁCTICO por fases (texto BLANCO sobre panel oscuro).
       proposito : frase que sigue a 'PROPÓSITO — ' (objetivo de la slide en una línea).
       sections  : lista ORDENADA de (LABEL, body) — el guion para exponer enseñando:
           body str            -> párrafo (frase literal, lista para leer/decir).
           body list[str]      -> ítems con marcador '→  ' (pasos / puntos).
           body list[(anc,txt)]-> ítem ANCLADO: '→ <elemento en pantalla> — <qué decir>'
                                  (el ancla va en NEGRITA → recorrido escaneable mientras señalás).
       Etiquetas habituales: ABRIR · RECORRIDO · EXPLICAR · ANALOGÍA · EJEMPLO · PUNTO CLAVE · TRANSICIÓN.
       RECORRIDO = guía elemento-por-elemento de la slide (cada caja/columna/panel nombrada,
       en orden visual); en los diagramas, el paso a paso del flujo. Sin nº de job, sin nombres."""
    tf = slide.notes_slide.notes_text_frame
    tf.clear(); tf.word_wrap = True

    def _para(first=False):
        return tf.paragraphs[0] if first else tf.add_paragraph()

    def _run(p, text, bold=False):
        r = p.add_run(); r.text = text
        r.font.bold = bold; r.font.size = Pt(nsz); r.font.name = F_BODY; r.font.color.rgb = WHITE

    # PROPÓSITO (etiqueta + frase)
    p = _para(first=True); p.space_after = Pt(8)
    _run(p, "PROPÓSITO — ", bold=True); _run(p, proposito)

    # Fases del guion: etiqueta en negrita + cuerpo (párrafo literal o ítems con flecha).
    for label, body in sections:
        ph = _para(); ph.space_before = Pt(7); ph.space_after = Pt(2)
        _run(ph, label, bold=True)
        is_list = isinstance(body, (list, tuple))
        for item in (body if is_list else [body]):
            pb = _para(); pb.space_after = Pt(2)
            if is_list and isinstance(item, (list, tuple)):
                # ítem ANCLADO: "→ <elemento en pantalla> — <qué decir>" (ancla en negrita)
                _run(pb, "→  "); _run(pb, item[0], bold=True); _run(pb, " — " + item[1])
            else:
                _run(pb, ("→  " + item) if is_list else item)


def _rect(slide, l, t, w, h, color):
    sp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(l), Inches(t), Inches(w), Inches(h))
    sp.fill.solid(); sp.fill.fore_color.rgb = color
    sp.line.fill.background(); sp.shadow.inherit = False
    return sp


def logo_mark(slide):
    _rect(slide, 0, 0, HDR, HDR, TEAL_SQ)
    slide.shapes.add_picture(LOGO, Inches(0.09), Inches(0.09), height=Inches(0.64))


def header(slide, title, bar=True):
    """Header consistente: (barra gris) + cuadrado teal + logo + título Barlow teal.
    title=None → solo logo (para diagramas a sangre que no deben taparse con el título)."""
    if bar:
        _rect(slide, 0, 0, SW, HDR, BAR_GRIS)
    logo_mark(slide)
    if title:
        tb = slide.shapes.add_textbox(Inches(1.05), Inches(0.04), Inches(12.0), Inches(0.74))
        _set_runs(tb.text_frame, [(title, 28, True, TEAL_TITLE, F_TITLE)], anchor=MSO_ANCHOR.MIDDLE)


def add_bullets(slide, l, t, w, h, items, size=26, sub_size=22, anchor=MSO_ANCHOR.TOP):
    tb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tb.text_frame; tf.word_wrap = True; tf.vertical_anchor = anchor
    for i, it in enumerate(items):
        txt, lvl = (it if isinstance(it, tuple) else (it, 0))
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(int(size * 0.7) if lvl == 0 else int(sub_size * 0.45))
        p.level = lvl
        b = "•  " if lvl == 0 else "–  "
        r = p.add_run(); r.text = ("     " * lvl) + b + txt
        r.font.size = Pt(size if lvl == 0 else sub_size); r.font.bold = (lvl == 0)
        r.font.name = F_BODY; r.font.color.rgb = INK if lvl == 0 else GRIS_BODY
    return tb


def add_card(slide, l, t, w, h, text, idx=0, size=21):
    fill = TEAL_CARD if idx % 2 == 0 else TEAL_CARD2
    sp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(l), Inches(t), Inches(w), Inches(h))
    sp.fill.solid(); sp.fill.fore_color.rgb = fill
    sp.line.color.rgb = TEAL_SQ; sp.line.width = Pt(1.5); sp.shadow.inherit = False
    cd = 0.62
    circ = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(l + 0.22), Inches(t + (h - cd) / 2),
                                  Inches(cd), Inches(cd))
    circ.fill.solid(); circ.fill.fore_color.rgb = TEAL_SQ
    circ.line.fill.background(); circ.shadow.inherit = False
    _set_runs(circ.text_frame, [(str(idx + 1), 20, True, WHITE, F_TITLE, PP_ALIGN.CENTER)],
              anchor=MSO_ANCHOR.MIDDLE)
    tb = slide.shapes.add_textbox(Inches(l + 1.05), Inches(t), Inches(w - 1.25), Inches(h))
    _set_runs(tb.text_frame, [(text, size, True, INK, F_BODY)], anchor=MSO_ANCHOR.MIDDLE)


def add_vcard(slide, l, t, w, h, title, body, tcol=ORA_T, tsize=17, bsize=14):
    """Tarjeta vertical: título color + cuerpo gris (para los 3 componentes)."""
    sp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(l), Inches(t), Inches(w), Inches(h))
    sp.fill.solid(); sp.fill.fore_color.rgb = TEAL_CARD2
    sp.line.color.rgb = TEAL_SQ; sp.line.width = Pt(1.5); sp.shadow.inherit = False
    _set_runs(sp.text_frame, [(title, tsize, True, tcol, F_BODY, PP_ALIGN.CENTER),
                              (body, bsize, False, GRIS_BODY, F_BODY, PP_ALIGN.CENTER)],
              anchor=MSO_ANCHOR.MIDDLE)


def add_image_fit(slide, path, l, t, w, h, align="center"):
    iw, ih = Image.open(path).size
    ar = iw / ih; box_ar = w / h
    if ar > box_ar:
        nw = w; nh = w / ar
    else:
        nh = h; nw = h * ar
    nl = l + (w - nw) / 2
    nt = t + (h - nh) / 2 if align != "top" else t
    slide.shapes.add_picture(path, Inches(nl), Inches(nt), Inches(nw), Inches(nh))


def _style_cell(cell, text, size, bold, color, align=PP_ALIGN.CENTER):
    cell.vertical_anchor = MSO_ANCHOR.MIDDLE
    cell.margin_left = cell.margin_right = Inches(0.06)
    cell.margin_top = cell.margin_bottom = Inches(0.02)
    tf = cell.text_frame; tf.word_wrap = True
    parts = str(text).split("\n")
    p = tf.paragraphs[0]; p.alignment = align
    r = p.add_run(); r.text = parts[0]
    r.font.size = Pt(size); r.font.bold = bold; r.font.name = F_BODY; r.font.color.rgb = color
    for extra in parts[1:]:
        p2 = tf.add_paragraph(); p2.alignment = align
        r2 = p2.add_run(); r2.text = extra
        r2.font.size = Pt(size); r2.font.bold = bold; r2.font.name = F_BODY; r2.font.color.rgb = color


def add_table(slide, headers, rows, l, t, w, h, col_fracs=None, cell_colors=None,
              fontsize=14, header_fontsize=14):
    """Tabla NATIVA pptx (escalable/editable) con estilo Environ."""
    nrows, ncols = len(rows) + 1, len(headers)
    gf = slide.shapes.add_table(nrows, ncols, Inches(l), Inches(t), Inches(w), Inches(h))
    table = gf.table
    table.first_row = False; table.horz_banding = False
    # quitar el tableStyle del tema (para que manden mis fills)
    tbl = table._tbl
    tblPr = tbl.find(qn("a:tblPr"))
    if tblPr is not None:
        sid = tblPr.find(qn("a:tableStyleId"))
        if sid is not None:
            tblPr.remove(sid)
    if col_fracs:
        tot = sum(col_fracs)
        for j, fr in enumerate(col_fracs):
            table.columns[j].width = Inches(w * fr / tot)
    for j, htxt in enumerate(headers):
        c = table.cell(0, j); c.fill.solid(); c.fill.fore_color.rgb = TEAL_SQ
        _style_cell(c, htxt, header_fontsize, True, WHITE)
    for i, row in enumerate(rows):
        for j, val in enumerate(row):
            c = table.cell(i + 1, j)
            bg = WHITE if i % 2 == 0 else ROW_ALT
            fg = INK; bold = (j == 0)
            if cell_colors and (i, j) in cell_colors:
                cc = cell_colors[(i, j)]
                fg = cc.get("fg", fg); bold = cc.get("bold", bold); bg = cc.get("bg", bg)
            c.fill.solid(); c.fill.fore_color.rgb = bg
            _style_cell(c, val, fontsize, bold, fg)
    rh = h / nrows
    for r in table.rows:
        r.height = Inches(rh)
    return gf


def _blue_shade(frac):
    frac = max(0.0, min(1.0, frac))
    return RGBColor(int(255 - (255 - 0x31) * frac),
                    int(255 - (255 - 0x85) * frac),
                    int(255 - (255 - 0x9C) * frac))


def add_confusion(slide, conf, classes, l, t, w, h, caption=None, vfont=15):
    """Matriz de confusión como TABLA NATIVA con celdas tipo heatmap (no imagen)."""
    n = len(classes)
    mx = max(max(r) for r in conf) or 1
    gf = slide.shapes.add_table(n + 1, n + 1, Inches(l), Inches(t), Inches(w), Inches(h))
    table = gf.table; table.first_row = False; table.horz_banding = False
    tblPr = table._tbl.find(qn("a:tblPr"))
    if tblPr is not None:
        sid = tblPr.find(qn("a:tableStyleId"))
        if sid is not None:
            tblPr.remove(sid)
    lblw = w * 0.20
    for j in range(n + 1):
        table.columns[j].width = Inches(lblw if j == 0 else (w - lblw) / n)
    c = table.cell(0, 0); c.fill.solid(); c.fill.fore_color.rgb = WHITE; _style_cell(c, "", 10, False, INK)
    for j in range(n):
        c = table.cell(0, j + 1); c.fill.solid(); c.fill.fore_color.rgb = TEAL_SQ
        _style_cell(c, f"pred\n{classes[j]}", 12, True, WHITE)
    for i in range(n):
        c = table.cell(i + 1, 0); c.fill.solid(); c.fill.fore_color.rgb = TEAL_SQ
        _style_cell(c, f"true\n{classes[i]}", 12, True, WHITE)
        for j in range(n):
            v = conf[i][j]; frac = v / mx
            cc = table.cell(i + 1, j + 1); cc.fill.solid(); cc.fill.fore_color.rgb = _blue_shade(frac)
            _style_cell(cc, str(v), vfont, True, WHITE if frac > 0.55 else INK)
    rh = h / (n + 1)
    for r in table.rows:
        r.height = Inches(rh)
    if caption:
        add_textbox(slide, l - 0.3, t + h + 0.06, w + 0.6, 0.5,
                    [(caption, 12, False, GRIS_TXT, F_BODY, PP_ALIGN.CENTER)])
    return gf


def add_chart_grouped(slide, cats, series, colors, l, t, w, h, ymax=1.0, hline=None):
    """Barras agrupadas como GRÁFICO NATIVO pptx (editable, no imagen)."""
    from pptx.chart.data import CategoryChartData
    from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION
    cd = CategoryChartData(); cd.categories = cats
    for name, vals in series:
        cd.add_series(name, vals)
    gf = slide.shapes.add_chart(XL_CHART_TYPE.COLUMN_CLUSTERED,
                                Inches(l), Inches(t), Inches(w), Inches(h), cd)
    chart = gf.chart; chart.has_title = False
    chart.has_legend = True; chart.legend.position = XL_LEGEND_POSITION.BOTTOM
    chart.legend.include_in_layout = False; chart.legend.font.size = Pt(13)
    plot = chart.plots[0]; plot.has_data_labels = True; plot.gap_width = 80
    dl = plot.data_labels; dl.number_format = "0.000"; dl.number_format_is_linked = False
    dl.font.size = Pt(12); dl.font.bold = True
    for srs, col in zip(chart.series, colors):
        srs.format.fill.solid(); srs.format.fill.fore_color.rgb = col
    va = chart.value_axis; va.minimum_scale = 0; va.maximum_scale = ymax
    va.tick_labels.font.size = Pt(11); va.has_major_gridlines = True
    chart.category_axis.tick_labels.font.size = Pt(13)
    return gf


def _node(slide, l, t, w, h, lines, fill, edge, lw=1.5):
    sp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(l), Inches(t), Inches(w), Inches(h))
    sp.fill.solid(); sp.fill.fore_color.rgb = fill
    sp.line.color.rgb = edge; sp.line.width = Pt(lw); sp.shadow.inherit = False
    _set_runs(sp.text_frame, lines, anchor=MSO_ANCHOR.MIDDLE)
    return sp


def _arrow(slide, l, t, w, h, direction="down"):
    shp = {"down": MSO_SHAPE.DOWN_ARROW, "right": MSO_SHAPE.RIGHT_ARROW}[direction]
    a = slide.shapes.add_shape(shp, Inches(l), Inches(t), Inches(w), Inches(h))
    a.fill.solid(); a.fill.fore_color.rgb = RGBColor(0x8A, 0x8A, 0x8A)
    a.line.fill.background(); a.shadow.inherit = False
    return a


def draw_mammoth_concept(slide, ox, oy):
    """Concepto antes/después (1 capa lineal → MoE) como BLOQUES nativos (no imagen)."""
    TEALF = RGBColor(0xD7, 0xE7, 0xEB); TEALE = TEAL_SQ
    ORAF = RGBColor(0xFC, 0xE5, 0xCD); ORAE = RGBColor(0xE2, 0x72, 0x3B); ORATT = ORA_T
    BLUF = RGBColor(0xEA, 0xF2, 0xFB); BLUE2 = RGBColor(0x5B, 0x8F, 0xB9)
    # --- CLAM (antes), columna izquierda ---
    add_textbox(slide, ox, oy, 2.4, 0.4, [("CLAM (antes)", 16, True, TEAL_TITLE, F_TITLE, PP_ALIGN.CENTER)])
    _node(slide, ox + 0.2, oy + 0.5, 2.0, 0.75, [("z ∈ ℝ⁵¹²", 13, True, INK)], TEALF, TEALE)
    _arrow(slide, ox + 1.05, oy + 1.28, 0.3, 0.32)
    _node(slide, ox, oy + 1.66, 2.4, 0.9, [("1 capa lineal", 14, True, ORATT),
                                           ("H = ReLU(W·z)", 11, False, GRIS_BODY)], ORAF, ORAE, lw=2)
    _arrow(slide, ox + 1.05, oy + 2.6, 0.3, 0.32)
    _node(slide, ox + 0.2, oy + 2.98, 2.0, 0.75, [("resto de CLAM", 13, True, INK)], TEALF, TEALE)
    # divisor
    ln = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(ox + 2.75), Inches(oy + 0.3),
                                Pt(1.4), Inches(3.4))
    ln.fill.solid(); ln.fill.fore_color.rgb = RGBColor(0xCC, 0xCC, 0xCC); ln.line.fill.background()
    ln.shadow.inherit = False
    # --- MAMMOTH (después), columna derecha ---
    rx = ox + 3.1
    add_textbox(slide, rx, oy, 2.9, 0.4, [("MAMMOTH (después)", 16, True, ORATT, F_TITLE, PP_ALIGN.CENTER)])
    _node(slide, rx + 0.9, oy + 0.5, 1.7, 0.62, [("z ∈ ℝ⁵¹²", 12, True, INK)], TEALF, TEALE)
    _arrow(slide, rx + 1.65, oy + 1.14, 0.26, 0.28)
    _node(slide, rx + 0.85, oy + 1.45, 1.8, 0.6, [("ROUTER g(z)", 12, True, ORATT)], ORAF, ORAE, lw=2)
    for i, lab in enumerate(["E₁", "E₂", "E₃"]):
        bx = rx + 0.1 + i * 1.15
        _node(slide, bx, oy + 2.3, 0.95, 0.6, [(lab, 13, True, INK)], BLUF, BLUE2)
        _arrow(slide, rx + 1.7, oy + 2.07, 0.22, 0.22) if i == 1 else None
    _node(slide, rx + 0.7, oy + 3.2, 2.1, 0.7, [("Σ ponderada", 12, True, INK),
                                                ("h ∈ ℝ⁵¹²", 11, False, GRIS_BODY)], TEALF, TEALE)
    add_textbox(slide, rx, oy + 3.95, 2.9, 0.5,
                [("expertos especializados por parche", 11, False, GRIS_TXT, F_BODY, PP_ALIGN.CENTER)])


# ---- datos necrosis (paired, para tabla nativa) ----
NEC_HEAD = ["Fold", "CLAM\nAUC / bal", "PathPT\nAUC / bal", "Δ AUC", "Δ bal_acc"]
NEC = [
    ("0", "0.798 / 0.653", "0.798 / 0.663", "+0.000", "+0.010"),
    ("1", "0.782 / 0.607", "0.798 / 0.698", "+0.016", "+0.091"),
    ("2", "0.598 / 0.656", "0.588 / 0.641", "−0.010", "−0.016"),
    ("3", "0.754 / 0.610", "0.599 / 0.502", "−0.155", "−0.108"),
    ("4", "0.703 / 0.641", "0.522 / 0.563", "−0.182", "−0.078"),
    ("media", "0.727 / 0.633", "0.661 / 0.613", "−0.066 ± 0.094", "−0.020 ± 0.078"),
]


def divider(prs, title, subtitle):
    s = prs.slides.add_slide(_blank(prs))
    s.background.fill.solid(); s.background.fill.fore_color.rgb = TEAL_DIV
    s.shapes.add_picture(LOGO, Inches(0.45), Inches(0.40), height=Inches(0.72))
    add_textbox(s, 1.0, 2.7, 11.3, 1.4,
                [(title, 50, True, LAV_TITLE, F_TITLE, PP_ALIGN.CENTER)], anchor=MSO_ANCHOR.MIDDLE)
    add_textbox(s, 1.0, 4.15, 11.3, 0.9,
                [(subtitle, 22, False, TEAL_SUB, F_BODY, PP_ALIGN.CENTER)], anchor=MSO_ANCHOR.TOP)
    return s


def content(prs, title):
    s = prs.slides.add_slide(_blank(prs)); header(s, title); return s


def copy_diagram(prs, src_path, idx, title, notes=None, bar=True):
    s = prs.slides.add_slide(_blank(prs))
    src = Presentation(src_path)
    src_slide = src.slides[idx]
    spTree = s.shapes._spTree
    for shp in src_slide.shapes:
        el = copy.deepcopy(shp._element)
        for cNvPr in el.iter(qn("p:cNvPr")):
            _uid[0] += 1; cNvPr.set("id", str(_uid[0]))
        # imágenes embebidas: el deepcopy NO arrastra la relación r:embed → copiar el
        # image part al slide destino y remapear el rId (si no, la figura desaparece).
        for blip in el.iter(qn("a:blip")):
            r_embed = blip.get(qn("r:embed"))
            if r_embed:
                img_part = src_slide.part.related_part(r_embed)
                _, new_rid = s.part.get_or_add_image_part(io.BytesIO(img_part.blob))
                blip.set(qn("r:embed"), new_rid)
        spTree.append(el)
    header(s, title, bar=bar)
    if notes:
        add_notes(s, notes)
    return s


# ============================================================================
# Datos de tablas
# ============================================================================

# Columnas UNIFICADAS (idénticas en las 3 tablas comparativas). El baseline CLAM es el MISMO
# (mismos splits k=5) → carcinoma/CDIS/tejido comparten bal_acc/AUC de CLAM en las 3 tablas.
# Tupla por fila: (tarea, dataset, bal_arrow, Δbal, auc_arrow, Δauc, color_de_los_Δ).
def _uni_head(base, arm):
    return ["Tarea", "Dataset\n(n: sí / no)", f"bal_acc\n{base} → {arm}",
            "Δ bal_acc\n(pareado)", f"AUC\n{base} → {arm}", "Δ AUC\n(pareado)"]

# Tabla 1 — Mammoth drop-in (keep_slots=False) vs CLAM [8 tareas]
MAM8_HEAD = _uni_head("CLAM", "+Mam")
MAM8 = [
    ("Microcalc · carcinoma inv.",    "n=328\n68 / 260",   "0.639 → 0.585", "−0.054 ± 0.125", "0.732 → 0.722", "−0.010 ± 0.065", GRIS_TXT),
    ("Microcalc · CDIS",              "n=328\n118 / 210",  "0.595 → 0.509", "−0.086 ± 0.113", "0.652 → 0.618", "−0.035 ± 0.104", ROJO),
    ("Microcalc · tejido no neopl.",  "n=328\n192 / 136",  "0.577 → 0.626", "+0.049 ± 0.077", "0.646 → 0.678", "+0.032 ± 0.084", VERDE),
    ("Patrón · cribiforme",           "n=513\n252 / 261",  "0.650 → 0.694", "+0.044 ± 0.048", "0.710 → 0.732", "+0.022 ± 0.042", VERDE),
    ("Patrón · sólido",               "n=513\n388 / 125",  "0.647 → 0.632", "−0.014 ± 0.064", "0.700 → 0.679", "−0.022 ± 0.055", GRIS_TXT),
    ("Patrón · micropapilar",         "n=513\n34 / 479",   "0.617 → 0.561", "−0.056 ‡",       "0.707 → 0.710", "+0.003 ‡",       GRIS_TXT),
    ("Patrón · papilar",              "n=513\n32 / 481",   "0.531 → 0.506", "−0.025 ‡",       "0.583 → 0.599", "+0.016 ‡",       GRIS_TXT),
    ("Invasión linfovascular (3 cl.)","n=2814\nno_id ~70%","0.622 → 0.575", "−0.047 ± 0.064", "0.828 → 0.818", "−0.011 ± 0.005", ROJO),
]
# Tabla 2 — Mammoth keep_slots=True (cuello de botella aprendido N→300) vs CLAM [4 tareas]
MAM_KST_HEAD = _uni_head("CLAM", "keep_slots")
MAM_KST = [
    ("Invasión linfovascular (3 cl.)","n=2814\nno_id ~70%","0.622 → 0.591", "−0.031 ± 0.047", "0.828 → 0.825", "−0.003 ± 0.015", GRIS_TXT),
    ("Microcalc · carcinoma inv.",    "n=328\n68 / 260",   "0.639 → 0.620", "−0.020 ± 0.121", "0.732 → 0.738", "+0.005 ± 0.087", GRIS_TXT),
    ("Microcalc · CDIS",              "n=328\n118 / 210",  "0.595 → 0.572", "−0.023 ± 0.120", "0.652 → 0.652", "+0.000 ± 0.108", AMBAR),
    ("Microcalc · tejido no neopl.",  "n=328\n192 / 136",  "0.577 → 0.623", "+0.046 ± 0.106", "0.646 → 0.647", "+0.002 ± 0.122", GRIS_TXT),
]
# Tabla 3 — CLAM + loss rebalanceada (focal / class_balanced) vs CLAM+CE [3 binarias × 2 losses].
# OJO (regla 5): la pérdida se aplica a CLAM_MB INTACTO, NO a mammoth. Baseline = CLAM+CE (el MISMO).
LOSS3_HEAD = ["Tarea  ·  pérdida", "Dataset\n(n: sí / no)", "bal_acc\nCLAM+CE → loss",
              "Δ bal_acc\n(pareado)", "AUC\nCLAM+CE → loss", "Δ AUC\n(pareado)"]
LOSS3 = [
    ("Carcinoma inv.  ·  focal",          "n=328\n68 / 260",  "0.639 → 0.597", "−0.042 ± 0.081", "0.732 → 0.697", "−0.036 ± 0.048", ROJO),
    ("Carcinoma inv.  ·  class_balanced", "n=328\n68 / 260",  "0.639 → 0.648", "+0.009 ± 0.074", "0.732 → 0.758", "+0.026 ± 0.099", GRIS_TXT),
    ("CDIS  ·  focal",                    "n=328\n118 / 210", "0.595 → 0.531", "−0.064 ± 0.093", "0.652 → 0.591", "−0.062 ± 0.074", ROJO),
    ("CDIS  ·  class_balanced",           "n=328\n118 / 210", "0.595 → 0.635", "+0.039 ± 0.091", "0.652 → 0.651", "−0.001 ± 0.045", GRIS_TXT),
    ("Tejido no neopl.  ·  focal",        "n=328\n192 / 136", "0.577 → 0.590", "+0.013 ± 0.052", "0.646 → 0.636", "−0.010 ± 0.034", GRIS_TXT),
    ("Tejido no neopl.  ·  class_balanced","n=328\n192 / 136","0.577 → 0.584", "+0.007 ± 0.032", "0.646 → 0.652", "+0.006 ± 0.009", GRIS_TXT),
]
MICRO_HEAD = ["Tarea binaria (microcalc)", "n  (sí / no)", "AUC zero-shot\n(mejor prompt)", "vs azar (0.5)", "Veredicto"]
MICRO = [
    ("en carcinoma invasivo", "68 / 260", "0.629", "leve", "NO-GO", AMBAR),
    ("en CDIS", "118 / 210", "0.533", "≈ azar", "NO-GO", ROJO),
    ("en tejido no neoplásico", "192 / 136", "0.444", "bajo azar", "NO-GO", ROJO),
]
EJES_HEAD = ["Eje atacado", "Modelo / mecanismo", "Tareas (pareadas, MC-CV)", "Resultado"]
EJES = [
    ("Agregador\n(cómo se fusionan los parches)", "DSMIL\n(dual-stream)", "microcalc binarias + fusionado", "0 palancas\n(CDIS regresión leve)"),
    ("Patch-embed\n(1ª capa de proyección)", "Mammoth\n(mixture-of-experts)", "12 tareas (microcalc, patrón, invasión)", "0 palancas\n(8 drop-in + 4 keep_slots)"),
    ("Lenguaje + supervisión tile", "PathPT-CONCH\n(prompt-tuning + θ espacial)", "necrosis, mitótica, microcalc", "0 palancas\n(grounding zero-shot débil)"),
]


# ============================================================================
# Construcción
# ============================================================================

def build():
    prs = Presentation(); prs.slide_width = Inches(SW); prs.slide_height = Inches(SH)

    # 1 · Portada
    s = prs.slides.add_slide(_blank(prs))
    s.shapes.add_picture(PORTADA, 0, 0, width=Inches(SW), height=Inches(SH))
    add_textbox(s, 0.8, 0.45, 11.7, 1.6,
                [("Sprint B5 — Integración de MAMMOTH y PathPT", 34, True, WHITE, F_TITLE, PP_ALIGN.CENTER),
                 ("OncoMets · CLAM sobre features CONCH · junio 2026", 17, False, WHITE, F_BODY, PP_ALIGN.CENTER)])
    set_notes(s,
        "abrir la presentación y encuadrar todo el sprint en una sola pregunta.",
        [("ABRIR  (decir tal cual)",
          "\"Esto es el cierre del sprint B5. Todo el trimestre giró en torno a una sola pregunta, "
          "y quiero que se la lleven desde el primer minuto.\""),
         ("RECORRIDO  (lo que se ve en la portada)",
          [("El título, 'Integración de MAMMOTH y PathPT'",
            "los dos modelos que integramos y medimos este sprint; el deck es, en el fondo, esas dos "
            "integraciones contadas en orden."),
           ("El subtítulo, 'CLAM sobre features CONCH'",
            "el baseline contra el que se compara TODO es CLAM, y todos corren sobre las mismas "
            "features CONCH: misma cancha para los tres.")]),
         ("EXPLICAR",
          ["La pregunta de fondo: ¿alguna modificación de la ARQUITECTURA de CLAM mueve la aguja en "
           "nuestras tareas clínicas, o el límite está en otro lado?",
           "Para responderla integramos y medimos los dos modelos con un mismo protocolo de "
           "evaluación, común y honesto."]),
         ("PUNTO CLAVE",
          "El hilo conductor del deck es uno solo: separar el 'efecto del modelo' del 'efecto del dato'."),
         ("TRANSICIÓN",
          "\"Antes de cualquier número, dejemos clara la vara con la que vamos a medir.\"")])

    # 2 · Objetivos
    s = content(prs, "Objetivos del sprint")
    add_bullets(s, 0.9, 1.05, 11.7, 6.0, [
        ("MAMMOTH: cómo funciona e integra en CLAM — diagrama + resultados.", 0),
        ("PathPT: integrar y explicar el modelo (visión + lenguaje) — diagrama + resultados.", 0),
        ("Marco de evaluación común:", 0),
        ("Comparación PAREADA con validación cruzada k=5 (mismos splits).", 1),
        ("Métrica honesta: balanced accuracy + AUC + matriz de confusión.", 1),
    ], size=30, sub_size=25, anchor=MSO_ANCHOR.MIDDLE)
    set_notes(s,
        "fijar la vara de medición ANTES de mostrar cualquier resultado, siguiendo los bullets.",
        [("ABRIR  (decir tal cual)",
          "\"Antes de los resultados, cómo los vamos a juzgar. Esto es lo que hace creíble todo lo "
          "que viene después.\""),
         ("RECORRIDO  (seguí los bullets en pantalla, de arriba hacia abajo)",
          [("Bullet 1, 'MAMMOTH'",
            "primer eje: cómo funciona y cómo se integra en CLAM; lo veremos con diagrama y luego con "
            "resultados."),
           ("Bullet 2, 'PathPT'",
            "segundo eje, de otra familia: visión + lenguaje; también con diagrama y resultados."),
           ("Bullet 3, 'Marco de evaluación común'",
            "lo que hace creíble todo lo demás; tiene dos sub-puntos debajo."),
           ("Sub-bullet, 'Comparación PAREADA, k=5'",
            "cada modelo se mide contra CLAM sobre exactamente los mismos splits, fold a fold; así el "
            "Δ por fold cancela la suerte del sorteo. Y k=5, no un único split optimista."),
           ("Sub-bullet, 'Métrica honesta'",
            "siempre balanced accuracy junto al AUC y la matriz de confusión; nunca una métrica sola.")]),
         ("EJEMPLO  (por qué balanced accuracy)",
          "Con clases desbalanceadas, un modelo que dice siempre 'lo más común' puede tener accuracy "
          "alto y ser inútil; la balanced accuracy lo desenmascara."),
         ("PUNTO CLAVE",
          "Comparación pareada + k=5 + balanced_acc/AUC/confusión = no nos engañamos solos."),
         ("TRANSICIÓN",
          "\"Con la vara fijada, abrimos el primer eje: MAMMOTH.\"")])

    # ===== MAMMOTH =====
    s = divider(prs, "MAMMOTH", "Mixture-of-Experts en la primera capa de CLAM (patch-embed)")
    set_notes(s,
        "abrir el primer bloque y anticipar exactamente qué eje ataca MAMMOTH.",
        [("ABRIR  (decir tal cual)",
          "\"Primer eje: MAMMOTH. Una intervención quirúrgica en la PRIMERA capa de CLAM, la que "
          "proyecta los parches.\""),
         ("RECORRIDO  (el subtítulo de esta portadilla)",
          [("'Mixture-of-Experts en la primera capa de CLAM (patch-embed)'",
            "ahí está todo el bloque en una línea: QUÉ es —una mezcla de expertos— y DÓNDE va —la "
            "primera capa, el patch-embed, no el clasificador final.")]),
         ("EXPLICAR",
          ["La pregunta del bloque: si reemplazamos esa capa por una mezcla de expertos, ¿bajamos la "
           "interferencia de gradientes y mejora la predicción?",
           "El recorrido va a ser: qué es → dónde entra → qué hace por dentro → resultados."]),
         ("PUNTO CLAVE",
          "Eje atacado = la capa de proyección de los parches (el patch-embed)."),
         ("TRANSICIÓN",
          "\"Empecemos por la intuición, sin una sola fórmula.\"")])

    s = content(prs, "MAMMOTH — qué es y por qué")
    cards = [
        "Reemplaza la 1ª capa lineal de CLAM por una mezcla de expertos (MoE).",
        "Un router envía cada parche a expertos especializados.",
        "Objetivo: reducir la interferencia de gradientes entre parches.",
        "El resto de CLAM no cambia → comparación limpia.",
    ]
    cy = 1.15
    for i, c in enumerate(cards):
        add_card(s, 0.5, cy, 6.2, 1.25, c, idx=i, size=21); cy += 1.42
    # Derecha: DOS figuras del paper apiladas, SIN leyenda (se explican de palabra).
    # Fig.1 arriba (qué hace + que mejora a todo MIL) · Fig.3 abajo (especialización por fenotipo).
    add_textbox(s, 6.95, 1.04, 6.18, 0.4,
                [("MAMMOTH estructura el espacio y mejora a todo MIL", 13, True, TEAL_TITLE, F_TITLE, PP_ALIGN.CENTER)])
    add_image_fit(s, FIG1_MAM, 7.58, 1.45, 4.79, 3.30, align="top")   # ~4.79 × 3.30 (aspecto 1.45)
    add_textbox(s, 6.95, 4.80, 6.18, 0.36,
                [("Cada experto se especializa en un fenotipo", 13, True, TEAL_TITLE, F_TITLE, PP_ALIGN.CENTER)])
    add_image_fit(s, FIG3_MAM, 7.00, 5.18, 6.00, 1.85, align="center")  # ~5.03 × 1.85 (aspecto 2.72)
    add_textbox(s, 6.95, 7.06, 6.18, 0.3,
                [("Figuras 1 y 3 — Shao et al., ICLR 2026", 9, False, GRIS_TXT, F_BODY, PP_ALIGN.CENTER)])
    set_notes(s,
        "dar la intuición de MAMMOTH —qué es, por qué, y qué evidencia hay de que funciona— "
        "con las 4 tarjetas y las dos figuras del paper, sin una sola fórmula en pantalla.",
        [("ABRIR  (decir tal cual)",
          "\"En CLAM, una sola capa lineal proyecta TODOS los parches al espacio interno del modelo. "
          "Esa única capa es el cuello que MAMMOTH ataca — y las dos figuras de la derecha son la "
          "prueba de que especializarla tiene sentido.\""),
         ("EXPLICAR  (qué es, en concreto, 'esa capa lineal')",
          ["Esa 'capa lineal' es, literalmente, una matriz W que multiplica al vector del parche: si z "
           "son los 512 números que CONCH le asigna a un parche, la capa calcula z' = W·z. Nada más: "
           "re-mezcla esos 512 números en otra combinación.",
           "El problema según el paper: hay UNA sola W para TODOS los parches, sin importar qué "
           "muestran. Y en una lámina los parches son heterogéneos — uno es tumor invasivo, otro "
           "estroma, otro grasa, otro necrosis: distintos FENOTIPOS (su 'tipo visual' de tejido).",
           "MAMMOTH cambia esa W única por VARIAS transformaciones especializadas, y a cada parche le "
           "asigna la que corresponde a su fenotipo. Eso es 'mixture of experts'. El CÓMO —router, "
           "slots, expertos de bajo rango— lo abrimos en las próximas dos slides; acá solo la intuición.",
           "Por qué una sola capa limita: al ser compartida, esa W debe servir simultáneamente a "
           "fenotipos opuestos, de modo que su ajuste óptimo es un compromiso promedio, subóptimo para "
           "cada tipo de tejido. Asignar una transformación especializada a cada fenotipo levanta esa "
           "restricción; ese desacople de los gradientes en conflicto es lo que el paper denomina "
           "reducir la interferencia de gradientes entre instancias."]),
         ("RECORRIDO  (Figura 1, arriba a la derecha — qué hace y que funciona)",
          [("Panel A, los dos mapas de puntos t-SNE",
            "es el espacio interno de embeddings proyectado a 2D. Arriba, con la capa lineal original, "
            "una nube continua sin estructura; abajo, con MAMMOTH, el espacio se ordena en grupos "
            "nítidos, cada color un experto. Es la contraparte cuantitativa de la idea de fenotipos."),
           ("Panel B, el gráfico de rendimiento",
            "el eje horizontal es el rendimiento promedio de ocho métodos MIL distintos. Para cada uno, "
            "el punto rojo es con MAMMOTH y el negro sin él: MAMMOTH desplaza el rendimiento a la "
            "derecha en los ocho, sin excepción, y con el mismo presupuesto de parámetros."),
           ("La lectura de la figura",
            "MAMMOTH es un módulo plug-and-play que mejora a cualquier agregador MIL; el titular del "
            "paper es que, equipado con MAMMOTH, incluso un pooling simple supera al mejor agregador "
            "que usa la capa lineal estándar.")]),
         ("RECORRIDO  (Figura 3, abajo a la derecha — por qué: especialización por fenotipo)",
          [("El título, 'Cada experto se especializa en un fenotipo'",
            "este resultado proviene del paper, no es un diseño nuestro: es la evidencia de que la "
            "especialización efectivamente ocurre."),
           ("Panel A, las dos láminas con mapa de calor",
            "cada lámina de pulmón está coloreada según a qué experto-slot se rutea cada parche; se "
            "observan regiones nítidas: distintas zonas del tejido activan distintos expertos."),
           ("Panel B, los parches inferiores",
            "son los parches que más activan cada slot, agrupados: uno reúne tumor, otro alvéolos, otro "
            "estroma, otro linfocitos, otro glóbulos rojos; dos patólogos certificados validaron que "
            "cada grupo corresponde a un fenotipo morfológico coherente."),
           ("La lectura de la figura",
            "la especialización por fenotipo es emergente: surge sola durante el entrenamiento, sin que "
            "nadie etiquete los tejidos — la evidencia de que separar la capa por fenotipo captura "
            "estructura morfológica real.")]),
         ("PUNTO CLAVE",
          "Cambio quirúrgico en una sola capa → comparación limpia contra CLAM; las figuras muestran "
          "que esa capa, especializada, estructura el espacio en fenotipos (Fig. 1A y Fig. 3) y que "
          "eso mejora a todo agregador MIL (Fig. 1B)."),
         ("TRANSICIÓN",
          "\"Veamos en el pipeline DÓNDE entra exactamente esa capa.\"")])

    s = copy_diagram(prs, DIAG_MAM, 0, None, bar=False)   # sin título: tapaba el diagrama
    set_notes(s,
        "ubicar visualmente DÓNDE entra MAMMOTH siguiendo el pipeline de CLAM bloque por bloque.",
        [("ABRIR  (decir tal cual)",
          "\"Este es nuestro diagrama de CLAM completo. El bloque naranja es el ÚNICO punto que cambia.\""),
         ("RECORRIDO  (seguí el pipeline de CLAM, en orden de flujo)",
          [("'INPUT PATCH'",
            "la slide entra como N parches independientes."),
           ("'FEATURE EXTRACTOR (CONCH)'",
            "cada parche se vuelve un vector de 512 — las features congeladas de CONCH."),
           ("'FULLY CONNECTED' (el bloque naranja)",
            "ACÁ y SOLO acá entra MAMMOTH: reemplaza esta 1ª capa lineal por la mezcla de expertos. El "
            "callout 'INTEGRACIÓN DE MAMMOTH' lo señala."),
           ("'ATTENTION HEAD / BACKBONE'",
            "CLAM aprende cuánto pesa cada parche para cada clase — la atención multi-rama."),
           ("'ATTENTION POOLING'",
            "junta los N parches en un único vector de slide, ponderado por la atención."),
           ("'TOP-B / BOTTOM-B' + 'INSTANCE CLASSIFIER'",
            "la rama de instancia: toma los parches más y menos atendidos y los clasifica (instance loss)."),
           ("'CAPA LINEAL por clase' + 'SOFTMAX SLIDE-LEVEL'",
            "el clasificador final. OJO con la confusión típica: esta capa lineal del FINAL no es la "
            "que cambia MAMMOTH; MAMMOTH toca la de ENTRADA.")]),
         ("PUNTO CLAVE",
          "Todo lo de aguas abajo del bloque naranja es CLAM sin tocar → es un drop-in: misma forma de "
          "entrada y de salida que la capa que reemplaza."),
         ("TRANSICIÓN",
          "\"Ahora abramos esa caja naranja y veamos qué hace por dentro.\"")])

    s = copy_diagram(prs, DIAG_MAM, 1, "MAMMOTH — qué hace (zoom de la 1ª capa)")
    set_notes(s,
        "abrir la 1ª capa de CLAM y seguir un parche, bloque por bloque del árbol, hasta la salida.",
        [("ABRIR",
          "\"Esta es la caja naranja del diagrama anterior, abierta por dentro. Cada parche entra como "
          "512 números de CONCH; sigamos UNO por todo el árbol.\""),
         ("RECORRIDO  (señalá cada bloque en pantalla, de arriba hacia abajo)",
          [("Arriba a la izquierda, 'MAMMOTH = MoE bajo rango'",
            "el recordatorio de qué reemplaza: la 1ª capa lineal de CLAM, como drop-in 512→512."),
           ("Tope del tronco, 'PARCHES DE LA SLIDE (CONCH)'",
            "z : [N × 512]. Son los N parches de la slide, cada uno un vector de 512 números."),
           ("Bloque naranja 'PROYECCIÓN A QUERY'",
            "W_q reescribe cada parche de 512 a 256: el query q. Ese 256 son 16 cabezas de 16."),
           ("Panel derecho '16 CABEZAS = 16 criterios'",
            "qué es una cabeza: el MISMO query mirado por 16 criterios en paralelo —textura, forma, "
            "densidad…— que al final se concatenan de vuelta en 512."),
           ("Bloque naranja 'RUTEO POR SLOTS (eq. 3)'",
            "comparo el query contra 300 claves aprendidas (los slots). La softmax D normaliza sobre "
            "los N parches: para cada slot, qué parches junta."),
           ("Panel derecho-abajo 'ZOOM: RUTEO POR SLOTS (eq. 3)'",
            "el mismo bloque RUTEO, abierto en sus 3 términos —① prototipo, ② producto interno, "
            "③ softmax + promedio ponderado—. Lo explico en detalle en el bloque EXPLICAR de abajo."),
           ("Etiqueta gris '300 slots = 30 expertos × 10'",
            "esos 300 slots se reparten en 30 expertos, 10 cada uno."),
           ("Las 3 cajas azules 'EXPERTO 1 / 2 / 30'",
            "cada experto transforma sus 10 slots. El callout de la izquierda dice qué es por dentro: "
            "un LoRA — baja de 16 a 8 (compartido) y sube de 8 a 32 (propio de cada experto)."),
           ("Bloque teal 'COMBINACIÓN'",
            "la 2ª softmax, C, normaliza sobre los 300 slots: para cada parche, a qué expertos "
            "escucha. Recompone la salida h : [N × 512]."),
           ("Banner arriba a la derecha, 'MISMO PRESUPUESTO'",
            "todo este aparato pesa casi lo mismo que la capa lineal 512×512 que reemplaza.")]),
         ("EXPLICAR  ·  cómo se arma 1 slot, los 3 términos de la eq. 3 (panel derecho-abajo)",
          [("Término 1 — PROTOTIPO s_j",
            "cada slot tiene un vector-clave APRENDIDO, su «afiche de búsqueda»: codifica qué morfología "
            "busca ese slot. Arranca al azar y el entrenamiento lo afina; hay 300 claves (30 expertos × "
            "10 slots), una por slot y por cabeza."),
           ("Término 2 — PRODUCTO INTERNO ⟨q, s_j⟩",
            "mido cuánto se parece el query del parche a esa clave: multiplico componente a componente y "
            "sumo. Más alto = más parecido. Da UN número por cada par (parche, slot)."),
           ("Término 3 — SOFTMAX + PROMEDIO PONDERADO (la eq. 3 completa)",
            "la softmax sobre los N parches convierte esos puntajes en pesos D que suman 1 (D_i = "
            "softmax_N ⟨q, s_j⟩); y el slot u_j es el promedio ponderado de los parches con esos pesos "
            "(u_j = Σ_i D_i · q_i). Es decir: cada slot resume los N parches, dándole más peso a los que "
            "se parecen a su prototipo.")]),
         ("EJEMPLO NUMÉRICO  (término 3, solo si lo piden)",
          ["3 parches con query q1=[1,0], q2=[0,1], q3=[1,1] y un prototipo de slot s=[1,0].",
           "Productos internos: 1, 0, 1 → exp: 2.72, 1.0, 2.72 → softmax D = [0.42, 0.155, 0.42].",
           "Slot u = 0.42·[1,0] + 0.155·[0,1] + 0.42·[1,1] = [0.84, 0.575]: prioriza los parches "
           "alineados con la clave; el que no matchea (q2) casi no aporta.",
           "Por qué es «soft»: ningún peso es 0 → TODOS los parches contribuyen a cada slot. Eso es lo "
           "que estabiliza el entrenamiento, a diferencia del ruteo duro (cada parche a un solo experto)."]),
         ("OJO  ·  las dos softmax",
          "D (ruteo) y C (combinación) usan la MISMA matriz de puntajes ⟨q, S⟩, pero normalizan por ejes "
          "distintos: D sobre los N parches (arma los slots), C sobre los 300 slots (recompone los parches)."),
         ("PUNTO CLAVE",
          "Entra [N, 512] y sale [N, 512] —reemplazo transparente—, y con rango 8 automático pesa casi "
          "lo mismo que la capa lineal: si no mejora, NO es por falta de capacidad."),
         ("TRANSICIÓN",
          "\"Esto mismo, ahora sobre la figura oficial del paper, con las dimensiones reales.\"")])

    # 6b · La MISMA mecánica del zoom, ahora sobre la figura OFICIAL del paper (+ dims overlay)
    s = copy_diagram(prs, DIAG_FUSED, 0, None, bar=False)   # full-bleed (lleva su propio título); logo only
    set_notes(s,
        "recorrer la figura OFICIAL del paper cartel por cartel, con las dimensiones reales del tensor.",
        [("ABRIR  (decir tal cual)",
          "\"Esta es la figura oficial de MAMMOTH, la de los autores. Es exactamente lo que abrimos en "
          "el zoom, pero de punta a punta; le pusimos encima un cartel por paso.\""),
         ("RECORRIDO  (los carteles sobre la figura, de izquierda a derecha)",
          [("Cartel 1, 'ENCODER (CONCH)'",
            "los parches entran y salen como z [N, 512] — el punto de partida, igual que en el zoom."),
           ("Cartel 2, 'PROYECCIÓN W_q'",
            "z → query q [N, 256]; ese 256 son 16 cabezas de 16."),
           ("Cartel 3, 'RUTEO POR SLOTS'",
            "la softmax sobre los N parches arma los 300 slots."),
           ("Cartel 4, 'EXPERTOS (LoRA r=8)'",
            "los 300 slots pasan por los expertos de bajo rango → o [300, 512]."),
           ("Cartel 5, 'CROSS-HEAD CONCAT'",
            "se recombinan de vuelta a los parches → h [N, 512]."),
           ("Cartel 6, 'MIL · CLAM'",
            "h entra a CLAM, que produce los logits clínicos. De acá en adelante es CLAM intacto."),
           ("Cartel de abajo, 'DENTRO DEL EXPERTO (LoRA)'",
            "el zoom del experto: A baja de 16 a 8 (compartido), B_e sube de 8 a 32 (por experto), rango 8.")]),
         ("PUNTO CLAVE",
          "Es la MISMA primera capa del zoom, ahora en el lenguaje visual de los autores; cada cartel "
          "es la forma real del tensor, verificada contra el código."),
         ("TRANSICIÓN",
          "\"Y de esta misma figura sale una variante que también medimos: keep_slots.\"")])

    # 6c · La variante keep_slots que también medimos (mismo tronco, bifurca la salida)
    s = copy_diagram(prs, DIAG_FUSED, 1, None, bar=False)   # nativo (lleva su propio título); logo only
    set_notes(s,
        "recorrer la bifurcación de keep_slots de arriba abajo: tronco compartido → fork → los 2 paneles.",
        [("ABRIR  (decir tal cual)",
          "\"El tronco —encoder, query, ruteo, expertos— es idéntico a la slide anterior. Cambia un "
          "único interruptor: keep_slots.\""),
         ("RECORRIDO  (de arriba hacia abajo)",
          [("La fila de arriba, 'TRONCO COMPARTIDO'",
            "encoder → W_q → ruteo → expertos: idéntico a la slide anterior. Hasta acá nada nuevo."),
           ("El nodo del medio, 'keep_slots = ?'",
            "el punto de bifurcación: ¿recombinar a parches o quedarse con los slots? De acá salen dos "
            "caminos."),
           ("Panel izquierdo (gris), 'keep_slots = FALSE'",
            "la base que ya medimos: recombina los 300 slots → h [N,512]; CLAM agrega sobre los N "
            "parches. Es el drop-in transparente."),
           ("Panel derecho (naranja), 'keep_slots = TRUE'",
            "la variante nueva: SE SALTA la recombinación, se queda con los 300 slot-tokens; CLAM "
            "agrega sobre los 300 slots. Un cuello de botella aprendido, más slot_dropout."),
           ("La nota de abajo",
            "misma cabeza CLAM y misma pérdida: cambia SOLO qué se agrega; el objetivo es dar "
            "capacidad dedicada a la clase minoritaria que la base tiende a colapsar.")]),
         ("PUNTO CLAVE",
          "Un solo interruptor decide la salida; ya medimos las dos posiciones en pareado k=5."),
         ("TRANSICIÓN",
          "\"Empecemos por los resultados de la variante base, el drop-in, en ocho tareas.\"")])

    # 7 · Resultados 8 tareas drop-in (TABLA NATIVA)
    s = content(prs, "MAMMOTH drop-in — resultados (8 tareas, k=5)")
    rows = [list(d[:6]) for d in MAM8]
    cc = {}
    for i, d in enumerate(MAM8):
        cc[(i, 3)] = {"fg": d[6], "bold": True}; cc[(i, 5)] = {"fg": d[6], "bold": True}
    add_table(s, MAM8_HEAD, rows, 0.25, 1.0, 12.85, 5.55,
              col_fracs=[0.215, 0.135, 0.165, 0.165, 0.155, 0.165],
              cell_colors=cc, fontsize=13, header_fontsize=12)
    add_textbox(s, 0.3, 6.62, 12.75, 0.75,
                [("0 palancas consistentes: en las 6 celdas Δ el ± (varianza) ≥ |Δ|. El lean+ leve "
                  "solo asoma en las 2 tareas más balanceadas (tejido, cribiforme) → lo gobierna el "
                  "balance de clase (col. Dataset), no la arquitectura.", 13, False, GRIS_TXT),
                 ("‡ régimen ciego (micropapilar/papilar, 3 positivos por test): Δ pooled de 5 folds, "
                  "sin ± por fold.", 11, False, GRIS_TXT)])
    set_notes(s,
        "leer la tabla columna por columna para mostrar que el drop-in de MAMMOTH no es palanca en 8 tareas.",
        [("ABRIR  (decir tal cual)",
          "\"Ocho tareas, todas pareadas, mismo baseline CLAM. Lean las dos columnas Δ y van a ver el "
          "patrón de una.\""),
         ("RECORRIDO  (cómo leer la tabla, columna por columna)",
          [("Columna 1, 'Tarea'",
            "las 8 tareas pareadas: 3 de microcalcificaciones, 4 de patrón arquitectónico, 1 de invasión."),
           ("Columna 2, 'Dataset (n: sí / no)'",
            "el tamaño y, sobre todo, cuántos POSITIVOS hay: es la clave para leer el resto, ya van a "
            "ver por qué."),
           ("Columnas 'bal_acc' y 'AUC' (CLAM → +Mam)",
            "la flecha muestra el antes y el después: CLAM, y CLAM con Mammoth. Casi no se mueven."),
           ("Las dos columnas 'Δ' (pareado, coloreadas)",
            "el efecto real fold a fold, en balanced accuracy y en AUC. Verde = leve mejora, rojo = "
            "leve regresión, gris = nulo. Fíjense que el ± (la varianza) es casi siempre ≥ que el "
            "número: no hay señal por encima del ruido. Solo hay verde en las 2 filas más balanceadas: "
            "tejido y cribiforme.")]),
         ("PUNTO CLAVE",
          "Lo que predice el signo del efecto es el BALANCE de la clase (columna Dataset), no la "
          "arquitectura: donde hay pocos positivos, Mammoth no rescata nada."),
         ("TRANSICIÓN",
          "\"Vamos al caso con más datos y la medición más limpia: invasión linfovascular.\"")])

    # 8 · Invasión (gráfico nativo + confusión nativa)
    s = content(prs, "MAMMOTH — invasión linfovascular")
    add_chart_grouped(s, ["balanced_acc", "macro-OVR AUC"],
                      [("CLAM", (0.622, 0.828)), ("+ Mammoth", (0.575, 0.818))],
                      [RGBColor(0x2C, 0x7A, 0x8C), RGBColor(0xE2, 0x72, 0x3B)],
                      0.5, 1.15, 6.0, 4.95, ymax=1.0)
    add_confusion(s, [[115, 64, 62], [97, 803, 85], [47, 56, 79]], ["aus", "no_id", "pres"],
                  7.7, 1.45, 4.9, 3.55,
                  caption="+Mammoth (5 folds): recall  aus 0.48 · no_id 0.82↑ · pres 0.43↓")
    add_textbox(s, 0.5, 6.25, 12.4, 0.95,
                [("n = 2814 (el caso más sano): el drop-in agrava el colapso a la mayoritaria; AUC y "
                  "bal_acc bajan levemente. La variante keep_slots ataca justo este modo de falla → "
                  "siguiente slide.", 15, False, GRIS_TXT)])
    set_notes(s,
        "leer el caso con más datos elemento por elemento, y mostrar el modo de falla del drop-in: "
        "el colapso a la clase mayoritaria.",
        [("ABRIR",
          "\"Esta es la tarea con más datos —2814 slides— y la evaluación más sana: cada clase tiene "
          "casos suficientes en cada fold. Si MAMMOTH iba a brillar, era acá.\""),
         ("RECORRIDO  (señalá cada elemento de la slide)",
          [("Gráfico de la izquierda, dos grupos de barras",
            "balanced accuracy a la izquierda, AUC macro a la derecha. En cada grupo, la barra teal es "
            "CLAM y la naranja es +Mammoth (lo dice la leyenda de abajo)."),
           ("Grupo balanced_acc",
            "CLAM 0.622 contra Mammoth 0.575: la barra naranja queda por DEBAJO. Mammoth empeora."),
           ("Grupo macro-OVR AUC",
            "0.828 contra 0.818: también baja, leve, pero en la misma dirección."),
           ("Matriz de la derecha, la diagonal",
            "filas = clase verdadera, columnas = clase predicha; la diagonal son los aciertos."),
           ("La fila del medio, 'no_id'",
            "está cargadísima (803 en la diagonal): el modelo manda casi todo a la mayoritaria "
            "'no identificado'."),
           ("La fila de abajo, 'pres'",
            "se vacía: a la clase 'presente' apenas le acierta (79). Y es la clase que importa "
            "clínicamente."),
           ("El caption bajo la matriz (recalls)",
            "lo resume: recall no_id 0.82 SUBE (flecha arriba), pres 0.43 BAJA (flecha abajo). Eso es "
            "el colapso a la mayoritaria.")]),
         ("PUNTO CLAVE",
          "Más poder estadístico no rescató a MAMMOTH: lo afinó. Dejó ver un signo negativo consistente "
          "en los 5 folds que en las tareas chicas quedaba tapado por el ruido."),
         ("TRANSICIÓN",
          "\"Ese colapso a la mayoritaria es justo lo que la variante keep_slots fue diseñada para "
          "revertir. Vamos a verlo.\"")])

    # 8b · keep_slots=True — resultados (4 tareas, TABLA NATIVA) → cierre del hilo MAMMOTH
    s = content(prs, "MAMMOTH keep_slots — resultados (4 tareas, k=5)")
    rows = [list(d[:6]) for d in MAM_KST]
    cc = {}
    for i, d in enumerate(MAM_KST):
        cc[(i, 3)] = {"fg": d[6], "bold": True}; cc[(i, 5)] = {"fg": d[6], "bold": True}
    add_table(s, MAM_KST_HEAD, rows, 0.25, 1.0, 12.85, 3.15,
              col_fracs=[0.215, 0.135, 0.175, 0.16, 0.16, 0.155],
              cell_colors=cc, fontsize=14, header_fontsize=12)
    add_textbox(s, 0.3, 4.45, 12.75, 2.7,
                [("Misma estructura y mismo baseline CLAM que la tabla drop-in → comparables fila a fila. "
                  "keep_slots NO supera a CLAM en ninguna (las 6 Δ en cero/negativo dentro del ruido).",
                  14, False, GRIS_TXT),
                 ("Matiz mecanístico (no visible en la tabla): el cuello de botella de slots revierte "
                  "PARCIALMENTE el colapso a la mayoritaria del drop-in — ↑ recall de la clase rara en "
                  "3/4 tareas (invasión 0.43→0.52 · carcinoma 0.29→0.31 · CDIS * 0.28→0.44). "
                  "slot_dropout: descartado (net-negativo en las 4).", 13, False, GRIS_TXT),
                 ("→ Hilo MAMMOTH completo: 12 tareas (8 drop-in + 4 keep_slots), 0 palancas.",
                  16, True, INK),
                 ("* CDIS: keep_slots sube el recall de la clase rara por encima de CLAM, pero a costa "
                  "de la mayoritaria → la bal_acc total no mejora.", 11, False, GRIS_TXT)])
    set_notes(s,
        "leer la tabla columna por columna: misma estructura que la drop-in; la mejora mecanística parcial "
        "NO alcanza para superar a CLAM; cierra el hilo MAMMOTH en 12 tareas.",
        [("ABRIR  (decir tal cual)",
          "\"Esta es la variante hecha a medida para atacar ese colapso. Misma tabla que la drop-in, "
          "mismo baseline CLAM, para que se comparen de una.\""),
         ("RECORRIDO  (la tabla, columna por columna)",
          [("Columnas 'bal_acc' y 'AUC' (CLAM → keep_slots)",
            "el antes y el después contra el MISMO CLAM de la tabla anterior; casi no se mueven."),
           ("Las dos columnas 'Δ' (pareado, la pregunta que importa)",
            "¿supera a CLAM? En las 4 filas, no — todas en cero o negativo dentro del ruido, igual que "
            "el drop-in."),
           ("El matiz que NO está en la tabla (lo digo yo, está en el pie)",
            "la buena noticia parcial: keep_slots recupera el recall de la clase rara que el drop-in "
            "colapsaba, en 3 de 4 tareas. Pero lo paga con la mayoritaria."),
           ("El caso CDIS, marcado con *",
            "el ilustrativo: keep_slots sube el recall de la clase rara incluso por encima de CLAM, "
            "pero a costa de la mayoritaria → la balanced accuracy total igual no mejora."),
           ("El texto de cierre abajo",
            "la conclusión del bloque: hilo MAMMOTH completo en 12 tareas (8 drop-in + 4 keep_slots), "
            "0 palancas; slot_dropout descartado.")]),
         ("PUNTO CLAVE",
          "Con esto el hilo MAMMOTH cierra en 12 tareas y cero palancas → el cuello sigue siendo el "
          "dato, no la arquitectura."),
         ("TRANSICIÓN",
          "\"Cambiamos de eje por completo: ahora visión + lenguaje, con PathPT.\"")])

    # ===== PATHPT =====
    s = divider(prs, "PathPT", "Visión + lenguaje: clasificación parche-a-parche sobre CONCH congelado")
    set_notes(s,
        "abrir el segundo bloque y marcar que cambiamos de eje por completo.",
        [("ABRIR  (decir tal cual)",
          "\"Segundo eje, y es uno totalmente distinto: PathPT no toca el agregador ni la capa de "
          "proyección.\""),
         ("RECORRIDO  (el subtítulo de esta portadilla)",
          [("'Visión + lenguaje: clasificación parche-a-parche sobre CONCH congelado'",
            "tres ideas en una línea: suma LENGUAJE (no solo visión), clasifica PARCHE a PARCHE (no la "
            "slide entera), y lo hace sobre CONCH CONGELADO (no lo re-entrena).")]),
         ("EXPLICAR",
          ["Suma una palanca nueva: lenguaje y supervisión a nivel de parche, sobre el mismo CONCH "
           "congelado.",
           "Mientras CLAM clasifica la slide entera, PathPT clasifica parche a parche comparando "
           "contra texto, y además localiza el hallazgo."]),
         ("PUNTO CLAVE",
          "Eje atacado = visión + lenguaje + supervisión por parche, no el flujo de CLAM."),
         ("TRANSICIÓN",
          "\"Veamos la idea completa con la figura del paper.\"")])

    # 10 · Figura 1 COMPLETA del paper (a=MIL · b=PathPT · c=tareas · d=benchmarks) + caption de las 3 piezas
    s = content(prs, "PathPT — la idea y el alcance (figura del paper)")
    add_image_fit(s, os.path.join(PAPER_FIGS, "pathpt_fig1_full.png"), 0.5, 0.95, 12.33, 5.55, align="center")
    add_textbox(s, 0.5, 6.58, 12.33, 0.78,
                [("Arriba:  a) MIL clásico   vs   b) PathPT (clasifica parche-a-parche; visión + texto sobre "
                  "CONCH congelado).   Abajo:  c) abanico de tareas   ·   d) benchmarks del paper.",
                  13, False, GRIS_TXT),
                 ("3 piezas livianas entrenables:   θᵥ contexto espacial   ·   θₜ prompt-tuning   ·   "
                  "pseudo-labels (tile).   Todo lo demás (CONCH) queda congelado.", 13, True, ORA_T)],
                anchor=MSO_ANCHOR.MIDDLE)
    set_notes(s,
        "recorrer la figura OFICIAL del paper panel por panel (a, b, c, d) + el pie de las 3 piezas entrenables.",
        [("ABRIR  (decir tal cual)",
          "\"Esta es la figura del paper, completa. Tiene cuatro paneles; los recorremos en orden.\""),
         ("RECORRIDO  (los 4 paneles de la figura + el pie)",
          [("Panel a) arriba-izquierda, 'MIL clásico'",
            "el modo tradicional: comprime toda la slide a un vector y predice UNA etiqueta de slide."),
           ("Panel b) arriba-derecha, 'PathPT'",
            "el modo nuevo: clasifica parche a parche comparándolos contra TEXTO, y además localiza el "
            "hallazgo en la slide."),
           ("Panel c) abajo-izquierda, 'abanico de tareas'",
            "el alcance: la variedad de tareas de patología que el método cubre."),
           ("Panel d) abajo-derecha, 'benchmarks'",
            "dónde compite contra los MIL clásicos en el paper."),
           ("El pie naranja, 'θᵥ · θₜ · pseudo-labels'",
            "lo único entrenable: θᵥ (contexto espacial), θₜ (prompt-tuning) y los pseudo-labels por "
            "tile. Todo lo demás —CONCH— queda congelado.")]),
         ("PUNTO CLAVE",
          "PathPT = CONCH congelado + tres piezas chicas entrenables; es barato de entrenar."),
         ("TRANSICIÓN",
          "\"Abramos el forward completo, con el mismo estilo de diagrama que usamos para CLAM.\"")])

    # 11 · Diagrama arquitectura (matemático)
    s = copy_diagram(prs, DIAG_PPT, 0, "PathPT — arquitectura (forward)")
    set_notes(s,
        "recorrer las 3 ramas del forward —visión, texto, matching— fiel al código, con el estilo del de CLAM.",
        [("ABRIR  (decir tal cual)",
          "\"Tres ramas que convergen: visión a la izquierda, texto a la derecha, y abajo donde se "
          "encuentran. Sigámoslas.\""),
         ("RECORRIDO  (las 3 ramas, en orden: visión, texto, y dónde se encuentran)",
          [("Rama izquierda, 'RAMA VISIÓN · θᵥ'",
            "los N parches con sus coordenadas pasan por CONCH (Φᵥ) y luego por lo entrenable: "
            "convoluciones 3/5/7 (contexto LOCAL) y una atención Nyström (contexto GLOBAL). Sale V̄, "
            "un vector por parche."),
           ("Rama derecha, 'RAMA TEXTO · θₜ'",
            "el prompt aprendible (ctx CoOp) se ensambla, pasa por el transformer de texto de CONCH "
            "(Φₜ) y se proyecta. Sale T, un vector por CLASE."),
           ("Nodo del centro, 'CONCH'",
            "el backbone visión-lenguaje, Φᵥ + Φₜ, CONGELADO: las dos ramas comparten el mismo modelo "
            "base; solo se entrenan las piezas naranjas a los lados."),
           ("Abajo, 'RAMA MATCHING'",
            "las dos ramas se encuentran: por cada parche se compara su vector de visión V̄ contra "
            "cada vector de texto T (coseno), y un softmax da P — una clase por parche.")]),
         ("EJEMPLO  (el concepto clave)",
          "Coseno = parecido de dirección entre dos vectores: 1 si apuntan igual, 0 si no tienen "
          "relación. Acá: cuánto se parece un parche a la descripción en texto de cada clase."),
         ("PUNTO CLAVE",
          "Dos dimensiones a no confundir (las marca la leyenda): 512 es la dimensión contrastiva —el "
          "matching—; 768 es la del token donde viven los contextos del prompt."),
         ("TRANSICIÓN",
          "\"Con el modelo claro, vamos a los tres resultados. Primero, necrosis.\"")])

    # 12 · Necrosis (tabla nativa + confusión nativa)
    s = content(prs, "PathPT — Necrosis")
    rows = [list(r) for r in NEC]
    cc = {}
    for i, r in enumerate(NEC):
        for col in (3, 4):
            v = r[col]
            if v.startswith("+") and v != "+0.000":
                cc[(i, col)] = {"fg": VERDE, "bold": True}
            elif v.startswith("−"):
                cc[(i, col)] = {"fg": ROJO, "bold": True}
    add_table(s, NEC_HEAD, rows, 0.5, 1.2, 7.6, 4.4,
              col_fracs=[0.11, 0.27, 0.27, 0.18, 0.17], cell_colors=cc, fontsize=15, header_fontsize=13)
    add_confusion(s, [[19, 21], [39, 120]], ["aus", "pres"], 8.75, 1.55, 4.0, 3.0,
                  caption="PathPT (5 folds): recall  aus 0.48 · pres 0.76")
    add_textbox(s, 0.5, 6.1, 12.4, 1.0,
                [("H_alt: PathPT no aporta. El Δ pareado cruza 0 (bal_acc); apenas supera el teacher "
                  "zero-shot ~0.62, mientras CLAM llega a 0.727.", 15, False, GRIS_TXT)])
    set_notes(s,
        "recorrer la tabla por fold y la matriz: primer resultado de PathPT, no aporta sobre CLAM.",
        [("ABRIR  (decir tal cual)",
          "\"Necrosis, comparación pareada contra CLAM sobre los mismos splits. El veredicto es H_alt: "
          "PathPT no aporta.\""),
         ("RECORRIDO  (la tabla a la izquierda, la matriz a la derecha)",
          [("Columnas 'CLAM' y 'PathPT' (AUC / bal)",
            "cada fila es un fold; se compara modelo contra modelo sobre el MISMO split."),
           ("Las dos columnas 'Δ' (coloreadas)",
            "el efecto por fold: verde a favor, rojo en contra. Está repartido — no hay un signo "
            "dominante."),
           ("La fila 'media' (abajo de la tabla)",
            "el resumen: Δ bal_acc −0.020 ± 0.078 cruza el cero; Δ AUC −0.066, levemente negativo."),
           ("La matriz de la derecha (aus / pres)",
            "recall: a 'presente' le acierta 0.76, pero a 'ausente' solo 0.48; no es un clasificador "
            "redondo."),
           ("El texto de abajo",
            "el dato revelador: PathPT apenas se despega de su 'teacher' zero-shot (~0.62), mientras "
            "CLAM llega a 0.727. El entrenamiento casi no agrega.")]),
         ("PUNTO CLAVE",
          "Mismo split que CLAM (pareado) → la comparación es limpia, y PathPT no mueve la aguja."),
         ("TRANSICIÓN",
          "\"El segundo caso es más instructivo, porque colapsa: tasa mitótica.\"")])

    # 13 · Mitótica (dos confusiones nativas)
    s = content(prs, "PathPT — Tasa mitótica")
    add_textbox(s, 0.6, 1.0, 5.6, 0.4, [("CLAM — usa las 3 clases", 16, True, TEAL_TITLE, F_BODY, PP_ALIGN.CENTER)])
    add_confusion(s, [[231, 53, 36], [77, 27, 38], [26, 28, 72]], ["s1", "s2", "s3"],
                  1.1, 1.55, 4.7, 3.5, caption="bal_acc 0.494 · recall [0.72, 0.19, 0.57]")
    add_textbox(s, 7.1, 1.0, 5.6, 0.4, [("PathPT — colapsa a score_1", 16, True, ORA_T, F_BODY, PP_ALIGN.CENTER)])
    add_confusion(s, [[320, 0, 0], [142, 0, 0], [126, 0, 0]], ["s1", "s2", "s3"],
                  7.6, 1.55, 4.7, 3.5, caption="bal_acc 0.333 (trivial) · 0 predicciones de s2 / s3")
    add_textbox(s, 0.5, 6.1, 12.4, 1.0,
                [("No es bug: CLAM no colapsa sobre los mismos splits. Es la formulación ordinal "
                  "(clase 0 = score_1 basal). El AUC (0.66) sobrevive → calibración del punto de operación.",
                  15, False, GRIS_TXT)])
    set_notes(s,
        "comparar las dos matrices lado a lado: un colapso instructivo — y de dónde sale la próxima palanca real.",
        [("ABRIR  (decir tal cual)",
          "\"Comparen las dos matrices. A la izquierda CLAM usa las tres clases; a la derecha PathPT "
          "predice SIEMPRE la mayoritaria.\""),
         ("RECORRIDO  (las dos matrices, izquierda y derecha)",
          [("Matriz izquierda, 'CLAM — usa las 3 clases'",
            "la diagonal tiene valores en las tres filas: CLAM reparte predicciones entre score_1, 2 y "
            "3. bal_acc 0.494."),
           ("Matriz derecha, 'PathPT — colapsa a score_1'",
            "solo la PRIMERA columna tiene números: PathPT predice SIEMPRE la mayoritaria, cero "
            "predicciones de score_2 o score_3. bal_acc 0.333, el trivial exacto."),
           ("El texto de abajo",
            "no es un bug: CLAM, sobre los mismos splits, no colapsa. Es la formulación ordinal —clase "
            "0 = score basal— la que domina el pseudo-etiquetado. Y OJO: el AUC sobrevive (~0.66).")]),
         ("PUNTO CLAVE",
          "Que el AUC sobreviva mientras el argmax colapsa señala la próxima palanca real: CALIBRAR el "
          "punto de operación, no cambiar de modelo."),
         ("TRANSICIÓN",
          "\"Tercer caso, y un ejemplo de cómo descartar barato: microcalcificaciones.\"")])

    # 14 · Microcalc (TABLA NATIVA)
    s = content(prs, "PathPT — Microcalcificaciones (go/no-go)")
    rows = [[d[0], d[1], d[2], d[3], d[4]] for d in MICRO]
    cc = {}
    for i, d in enumerate(MICRO):
        cc[(i, 2)] = {"fg": d[5], "bold": True}; cc[(i, 4)] = {"fg": d[5], "bold": True}
    add_table(s, MICRO_HEAD, rows, 1.1, 1.4, 11.1, 3.8,
              col_fracs=[0.30, 0.16, 0.22, 0.16, 0.16], cell_colors=cc,
              fontsize=18, header_fontsize=16)
    add_textbox(s, 0.9, 5.6, 11.5, 1.4,
                [("Chequeo zero-shot en CPU (~min, sin GPU): CONCH no groundea microcalcificaciones.",
                  18, True, INK),
                 ("Iterar prompts con más morfología EMPEORÓ. NO-GO → ahorró ~18–24 h de GPU.",
                  16, False, GRIS_BODY)])
    set_notes(s,
        "leer la tabla columna por columna: un descarte barato y bien hecho — el patrón 'CPU antes que GPU'.",
        [("ABRIR  (decir tal cual)",
          "\"Antes de gastar GPU, hicimos un chequeo zero-shot en CPU, en minutos. Y nos ahorró el viaje.\""),
         ("RECORRIDO  (la tabla, de columna a columna)",
          [("Columna 'AUC zero-shot (mejor prompt)'",
            "el corazón del chequeo: qué tan bien separa CONCH cada microcalcificación SIN entrenar. "
            "0.44 a 0.63 según la tarea."),
           ("Columna 'vs azar (0.5)'",
            "la referencia: una de las tres incluso queda POR DEBAJO del azar."),
           ("Columna 'Veredicto' (coloreada)",
            "las 3 en NO-GO: CONCH no 'groundea' microcalcificaciones."),
           ("El texto de abajo",
            "el detalle contraintuitivo: afinar los prompts con MÁS morfología lo EMPEORÓ —mismo "
            "patrón que en necrosis— y el patrón operativo: este chequeo en CPU, en minutos, ahorró "
            "18–24 h de GPU.")]),
         ("PUNTO CLAVE",
          "Patrón a institucionalizar: un zero-shot barato ANTES del entrenamiento caro; si CONCH no "
          "'groundea' el hallazgo, no vale la GPU."),
         ("TRANSICIÓN",
          "\"Con los tres ejes ya medidos, juntemos todo en una sola lectura.\"")])

    # ===== CIERRE =====
    s = divider(prs, "Cierre y próximos pasos", "Tres ejes de arquitectura → dónde sí hay palanca")
    set_notes(s,
        "abrir la síntesis y los próximos pasos.",
        [("ABRIR  (decir tal cual)",
          "\"Hasta acá medimos tres ejes de arquitectura distintos. Toca juntarlos en una sola "
          "conclusión.\""),
         ("RECORRIDO  (el subtítulo de esta portadilla)",
          [("'Tres ejes de arquitectura → dónde sí hay palanca'",
            "la flecha es el arco del cierre: primero el veredicto de los tres ejes, después hacia "
            "dónde SÍ conviene mover el esfuerzo.")]),
         ("EXPLICAR",
          "La idea: poner los tres veredictos lado a lado y, a partir de eso, decir dónde SÍ conviene "
          "invertir esfuerzo."),
         ("PUNTO CLAVE",
          "La conclusión converge — por eso el último mensaje del deck son los próximos pasos."),
         ("TRANSICIÓN",
          "\"La tesis del sprint cabe en una sola tabla.\"")])

    # 16 · Cierre 3 ejes (TABLA NATIVA)
    s = content(prs, "Cierre — 3 ejes, 0 palancas → el cuello es el dato")
    rows = [list(r) for r in EJES]
    cc = {(i, 3): {"fg": GRIS_TXT, "bold": True} for i in range(len(EJES))}
    add_table(s, EJES_HEAD, rows, 0.5, 1.15, 12.3, 4.7,
              col_fracs=[0.24, 0.22, 0.30, 0.24], cell_colors=cc,
              fontsize=16, header_fontsize=15)
    add_textbox(s, 0.6, 6.15, 12.1, 1.0,
                [("Tres familias independientes de cambio arquitectónico, todas pareadas y con eval "
                  "honesta → ninguna mueve la aguja.", 16, True, INK),
                 ("El cuello convergente es el DATO: desbalance, pocos positivos, contexto espacial.",
                  15, False, GRIS_BODY)])
    set_notes(s,
        "recorrer las 3 filas de la tabla, una por eje — la tesis del sprint: 3 ejes independientes, 0 palancas.",
        [("ABRIR  (decir tal cual)",
          "\"Tres familias de cambio arquitectónico, completamente independientes entre sí. Una fila "
          "cada una.\""),
         ("RECORRIDO  (las 3 filas de la tabla, una por eje)",
          [("Fila 1, 'Agregador' → DSMIL",
            "cómo se fusionan los parches; probado con DSMIL → 0 palancas (en CDIS hasta una leve "
            "regresión)."),
           ("Fila 2, 'Patch-embed' → MAMMOTH",
            "la 1ª capa de proyección; 12 tareas (8 drop-in + 4 keep_slots) → 0 palancas."),
           ("Fila 3, 'Lenguaje + supervisión tile' → PathPT",
            "visión + texto por parche; necrosis, mitótica, microcalc → 0 palancas (grounding "
            "zero-shot débil)."),
           ("La columna 'Resultado' + el texto de cierre",
            "tres familias independientes, las tres con el mismo veredicto: ninguna mueve la aguja. "
            "Esa triangulación ES el resultado, no el accidente de una sola prueba.")]),
         ("PUNTO CLAVE",
          "El cuello de botella no es el modelo, es el DATO: desbalance, pocos positivos, contexto "
          "espacial. Deja de tener sentido buscar la palanca dentro de la arquitectura."),
         ("TRANSICIÓN",
          "\"Y antes de los próximos pasos, una coda honesta: probamos un 4º cambio, esta vez NO en la "
          "arquitectura sino en la pérdida.\"")])

    # 16b · CLAM + loss rebalanceada (4º intento: la PÉRDIDA, no la arquitectura) — TABLA NATIVA
    s = content(prs, "CLAM + loss rebalanceada — tocar la pérdida, no la arquitectura")
    rows = [list(d[:6]) for d in LOSS3]
    cc = {}
    for i, d in enumerate(LOSS3):
        cc[(i, 3)] = {"fg": d[6], "bold": True}; cc[(i, 5)] = {"fg": d[6], "bold": True}
    add_table(s, LOSS3_HEAD, rows, 0.25, 1.0, 12.85, 4.05,
              col_fracs=[0.25, 0.12, 0.165, 0.16, 0.155, 0.15],
              cell_colors=cc, fontsize=13, header_fontsize=12)
    add_textbox(s, 0.3, 5.35, 12.75, 1.85,
                [("Mismo CLAM, misma evaluación: solo cambia la PÉRDIDA (CLAM_MB intacto, no es Mammoth). "
                  "Las 6 Δ dentro del ruido (± ≥ |Δ|) → tampoco es palanca.", 14, False, GRIS_TXT),
                 ("class_balanced SÍ sube el recall de la minoritaria (carcinoma sí 0.37→0.71, CDIS "
                  "0.38→0.62) pero HUNDE la mayoritaria en igual medida (carcinoma no 0.91→0.59) → "
                  "balanced accuracy neta sin cambio = re-balanceo del punto de operación, no señal nueva "
                  "(el AUC no se mueve). focal ni siquiera rescata.", 13, False, GRIS_TXT),
                 ("→ Es la versión EN-ENTRENAMIENTO de la calibración → por eso el próximo paso real es "
                  "calibrar el punto de operación POST-HOC.", 14, True, INK)])
    set_notes(s,
        "mostrar que el 4º cambio (re-ponderar la pérdida, no la arquitectura) tampoco es palanca, y que su "
        "modo de fallar (re-balanceo) apunta directo al primer próximo paso: la calibración post-hoc.",
        [("PROPÓSITO",
          "cerrar el último ángulo barato y convertir su 'no-resultado' en la motivación del próximo paso."),
         ("ABRIR  (decir tal cual)",
          "\"Hasta acá movimos la arquitectura. Acá probamos algo distinto: dejar CLAM igual y cambiar "
          "solo la función de pérdida, para que no colapse a la clase mayoritaria.\""),
         ("EXPLICAR",
          ["Dos pérdidas, ambas sobre el MISMO CLAM y el MISMO baseline: focal (baja el peso de los "
           "ejemplos fáciles) y class_balanced (sube el peso de la clase minoritaria).",
           "Importante (y es una corrección): esto se aplica a CLAM intacto, NO a Mammoth — por eso la "
           "slide dice 'CLAM + loss', no 'Mammoth + loss'."]),
         ("RECORRIDO  (la tabla, columna por columna)",
          [("Columna 1, 'Tarea · pérdida'",
            "las 3 binarias de microcalcificaciones, cada una con sus dos variantes de pérdida: focal y "
            "class_balanced."),
           ("Las dos columnas 'Δ' (pareado, coloreadas)",
            "el efecto contra CLAM+CE. Las 6 celdas con el ± ≥ que el número: ruido, no palanca."),
           ("La fila class_balanced de carcinoma",
            "el ejemplo del mecanismo: sube el recall de la minoritaria de 0.37 a 0.71 —¡el triple!— "
            "pero la mayoritaria se hunde de 0.91 a 0.59. La balanced accuracy neta queda igual.")]),
         ("ANALOGÍA",
          "es mover la línea de decisión a mano: ganás de un lado exactamente lo que perdés del otro. "
          "El AUC —que mide el ordenamiento, no el umbral— no se mueve."),
         ("PUNTO CLAVE",
          "Re-pesar la pérdida re-balancea el punto de operación, no agrega señal. Es la versión "
          "en-entrenamiento de la calibración → confirma que la palanca barata real es calibrar "
          "POST-HOC, sin re-entrenar."),
         ("TRANSICIÓN",
          "\"Con eso, los próximos pasos —y el primero sale justo de acá.\"")])

    # 17 · Próximos pasos
    s = content(prs, "Próximos pasos — dónde sí hay palanca")
    add_bullets(s, 0.9, 1.05, 11.7, 6.1, [
        ("Calibración del punto de operación (umbral por clase): salió del colapso mitótico "
         "(el AUC sobrevive, el argmax colapsa). Barato: re-umbralizar sin re-entrenar.", 0),
        ("Recuperación de casos similares (CBIR sobre features CONCH): valor clínico, sin re-entrenar.", 0),
        ("Magnificación / nº de parches (contexto espacial): atacar el cuello REAL = el dato.", 0),
        ("Más positivos y mejor balance en las tareas hambrientas de datos.", 0),
    ], size=24, anchor=MSO_ANCHOR.MIDDLE)
    set_notes(s,
        "recorrer los 4 bullets en orden de prioridad: convertir el 'no-resultado' en una agenda con palancas reales.",
        [("ABRIR  (decir tal cual)",
          "\"Que la arquitectura no sea palanca no nos deja sin jugadas: nos las empuja hacia cosas "
          "ORTOGONALES al modelo.\""),
         ("RECORRIDO  (los 4 bullets, en orden de prioridad)",
          [("Bullet 1, 'Calibración del punto de operación'",
            "la más barata, y salió sola del colapso mitótico: el AUC ya está, el argmax no; se "
            "re-umbraliza por clase SIN re-entrenar."),
           ("Bullet 2, 'Recuperación de casos similares (CBIR)'",
            "sobre las features CONCH que ya tenemos: valor clínico inmediato, tampoco re-entrena."),
           ("Bullet 3, 'Magnificación / nº de parches'",
            "acá atacamos el cuello REAL que encontramos: el contexto espacial, el dato."),
           ("Bullet 4, 'Más positivos y mejor balance'",
            "la otra cara del cuello: conseguir más positivos en las tareas hambrientas de datos.")]),
         ("PUNTO CLAVE",
          "El sprint cambia la pregunta: no '¿qué modelo?', sino '¿qué dato y qué punto de operación?'."),
         ("TRANSICIÓN",
          "\"Con eso cierro. Quedo para preguntas.\"")])

    prs.save(DST)
    return prs


def main():
    prs = build()
    print(f"OK  {os.path.relpath(DST, REPO)}  slides={len(prs.slides)}")


if __name__ == "__main__":
    main()
