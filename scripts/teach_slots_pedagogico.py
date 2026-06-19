#!/usr/bin/env python
"""teach_slots_pedagogico.py — secuencia pedagógica del zoom MoE (deck-ready).

Dos slides EDITABLES (formas nativas, estilo Diagrama_CLAM — Carlito/Unicode, rasteriza
limpio, sin OMML), pensadas para INSERTAR antes de la slide del MoE en la presentación
oficial (sin números de job ni nombres — convención del deck):

  SLIDE 1 — ANATOMÍA: construye el tensor de slots `S : [30, 16, 10, 16]` eje por eje
            (slot → 10 slots → 30 expertos → 16 cabezas) + caja-crux: los DOS 16.
  SLIDE 2 — MECANISMO: qué HACE una cabeza — el matching `⟨q, S⟩` (producto interno por
            lente = un puntaje) que alimenta el ruteo. Zoom a 1 lente + mini-ejemplo.

Dims CODE-ACCURATE (clam_testing2/MAMMOTH/src/mammoth/mammoth.py, pin fe36d4e):
slot_embeds = [experts=30, heads=16, slots=10, head_dim=16]; 256 = 16 cab × 16
(head_dim = slot_dim // num_heads = 256 // 16); Wq: 512→256; producto interno por cabeza
sobre rebanadas de 16 → logits [N,30,16,10]; softmax → ruteo.

Reusa los helpers del generador del deck (paleta + add_box/label/connector/arrow).

Uso:
  PYTHONPATH=/media/administrador/Storage1/sdonoso/clam_testing2/.pylibs \
    /home/sdonoso/miniconda3/envs/clam_latest/bin/python scripts/teach_slots_pedagogico.py
"""
import os
import sys

sys.path.insert(0, os.path.dirname(os.path.abspath(__file__)))

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN

import generate_clam_mammoth_pptx as g  # helpers + paleta del deck

REPO = "/media/administrador/Storage1/sdonoso/clam_testing2/oncomets-ernesto"
DST = os.path.join(REPO, "papers/presentations/zoom_slots_pedagogico.pptx")

SLOT_F = g.ORA_F                         # ficha-slot / query: familia naranja (entrenable)
SLOT_E = g.ORA_E
GHOST_F = RGBColor(0xFB, 0xEF, 0xE3)     # copias "lente" detrás, más tenue
SLOT_F2 = RGBColor(0xFA, 0xD9, 0xBC)     # slot un poco más saturado (contraste vs query)
BLU_F, BLU_E = g.BLU_F, g.BLU_E
TEAL_F = g.TEAL_F
TEAL_E = RGBColor(0x2C, 0x7A, 0x8C)
INK, GREY_T = g.INK, g.GREY_T
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
PANEL_BORDER = RGBColor(0xC8, 0xC8, 0xC8)
BANNER_F = RGBColor(0xFD, 0xF1, 0xE5)


# ----------------------------------------------------------------------------- helpers
def new_slide(prs):
    blank = next((lay for lay in prs.slide_layouts
                  if lay.name and "blank" in lay.name.lower()), prs.slide_layouts[-1])
    return prs.slides.add_slide(blank)


def title(s, tt, sub):
    g.add_label(s, 0.4, 0.20, 12.5, 0.5, [(tt, 22, True, SLOT_E)], align=PP_ALIGN.CENTER)
    g.add_label(s, 0.4, 0.72, 12.5, 0.30, [(sub, 12.5, False, GREY_T)], align=PP_ALIGN.CENTER)
    rule = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(1.06),
                              Inches(12.33), Inches(0.030))
    rule.fill.solid(); rule.fill.fore_color.rgb = TEAL_E
    rule.line.fill.background(); rule.shadow.inherit = False


def badge(s, x, y, n, color=SLOT_E, d=0.36):
    o = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x), Inches(y), Inches(d), Inches(d))
    o.fill.solid(); o.fill.fore_color.rgb = color
    o.line.fill.background(); o.shadow.inherit = False
    g.set_text(o, [(str(n), 15, True, WHITE)])
    return o


def card(s, x, y, w, h, accent=SLOT_E, rule_y=0.62):
    box = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y),
                             Inches(w), Inches(h))
    box.fill.solid(); box.fill.fore_color.rgb = WHITE
    box.line.color.rgb = accent; box.line.width = Pt(1.75); box.shadow.inherit = False
    rule = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x + 0.16), Inches(y + rule_y),
                              Inches(w - 0.32), Inches(0.022))
    rule.fill.solid(); rule.fill.fore_color.rgb = accent
    rule.line.fill.background(); rule.shadow.inherit = False
    return box


def cells_grid(s, l, t, ncols, nrows, fill, edge, cw=0.17, ch=0.15, gap=0.035, lw=0.75):
    for r in range(nrows):
        for c in range(ncols):
            sq = s.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                    Inches(l + c * (cw + gap)), Inches(t + r * (ch + gap)),
                                    Inches(cw), Inches(ch))
            sq.fill.solid(); sq.fill.fore_color.rgb = fill
            sq.line.color.rgb = edge; sq.line.width = Pt(lw); sq.shadow.inherit = False
    return l + ncols * (cw + gap) - gap, t + nrows * (ch + gap) - gap


def lens_icon(s, x, y, d, color=SLOT_E):
    """Mini lupa = anillo (óvalo sin relleno) + mango diagonal → metáfora de 'lente'."""
    ring = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x), Inches(y), Inches(d), Inches(d))
    ring.fill.background(); ring.line.color.rgb = color; ring.line.width = Pt(2)
    ring.shadow.inherit = False
    g.add_connector(s, x + d * 0.80, y + d * 0.80, x + d * 1.20, y + d * 1.20,
                    color=color, w=2.5, arrow=False)


# ------------------------------------------------------------------- SLIDE 1: anatomía
def slide_anatomia(prs):
    s = new_slide(prs)
    title(s, "ANATOMÍA DE LOS SLOTS    S : [30, 16, 10, 16]",
          "se arma eje por eje — cada paso agrega UNA dimensión (son las MISMAS fichas, solo mejor organizadas)")

    pw, gapp, pl0, pt, ph = 2.85, 0.27, 0.45, 1.28, 3.78
    xs = [pl0 + i * (pw + gapp) for i in range(4)]

    def panel(x, n, header, tag, caption, tag_note, tag_size=19):
        card(s, x, pt, pw, ph)
        badge(s, x + 0.12, pt + 0.16, n)
        g.add_label(s, x + 0.50, pt + 0.18, pw - 0.60, 0.34, [(header, 13.5, True, INK)],
                    align=PP_ALIGN.LEFT)
        g.add_label(s, x + 0.08, pt + ph - 1.04, pw - 0.16, 0.42, [(tag, tag_size, True, SLOT_E)],
                    align=PP_ALIGN.CENTER)
        g.add_label(s, x + 0.08, pt + ph - 0.62, pw - 0.16, 0.26, [(tag_note, 10.5, False, GREY_T)],
                    align=PP_ALIGN.CENTER)
        g.add_label(s, x + 0.08, pt + ph - 0.36, pw - 0.16, 0.30, [(caption, 11, False, GREY_T)],
                    align=PP_ALIGN.CENTER)

    draw_t = pt + 1.18                                   # tope común del área de dibujo

    # ① 1 slot
    panel(xs[0], 1, "1 SLOT", "[ 16 ]", "una ficha de referencia", "16 números = la 'ficha'")
    r, _ = cells_grid(s, xs[0] + 0.62, draw_t + 0.30, 5, 1, SLOT_F, SLOT_E)
    g.add_label(s, r + 0.06, draw_t + 0.25, 0.55, 0.30, [("…", 18, True, SLOT_E)], align=PP_ALIGN.LEFT)

    # ② 10 slots = 1 experto
    panel(xs[1], 2, "10 SLOTS  (1 experto)", "[ 10 , 16 ]", "apilo las 10 fichas", "10 fichas × 16")
    cells_grid(s, xs[1] + 0.62, draw_t - 0.05, 5, 4, SLOT_F, SLOT_E)
    g.add_label(s, xs[1] + 0.62, draw_t - 0.05 + 4 * 0.185 + 0.01, 1.2, 0.26,
                [("⋮  (×10)", 12, True, SLOT_E)], align=PP_ALIGN.LEFT)

    # ③ 30 expertos
    panel(xs[2], 3, "30 EXPERTOS", "[ 30 , 10 , 16 ]", "cada experto: sus 10 slots", "30 platos de 10×16")
    for off in [(0.0, 0.0), (0.16, 0.16), (0.32, 0.32)]:
        cells_grid(s, xs[2] + 0.50 + off[0], draw_t - 0.08 + off[1], 4, 3, SLOT_F2, SLOT_E,
                   cw=0.13, ch=0.12, gap=0.03, lw=0.6)
    g.add_label(s, xs[2] + 1.55, draw_t + 0.62, 1.2, 0.34, [("× 30", 15, True, BLU_E)], align=PP_ALIGN.LEFT)

    # ④ 16 cabezas = lentes en paralelo
    panel(xs[3], 4, "16 CABEZAS  (lentes)", "[30, 16, 10, 16]", "16 lentes en paralelo",
          "← este es el S real", tag_size=16)
    for k, off in enumerate([(0.30, 0.30), (0.15, 0.15), (0.0, 0.0)]):
        fill = GHOST_F if k < 2 else SLOT_F2
        cells_grid(s, xs[3] + 0.48 + off[0], draw_t - 0.08 + off[1], 4, 3, fill, SLOT_E,
                   cw=0.13, ch=0.12, gap=0.03, lw=0.6)
    lens_icon(s, xs[3] + 1.95, draw_t - 0.02, 0.30)
    g.add_label(s, xs[3] + 1.55, draw_t + 0.62, 1.2, 0.34, [("× 16", 15, True, SLOT_E)], align=PP_ALIGN.LEFT)

    # flechas entre paneles
    for i in range(3):
        g.add_arrow(s, xs[i] + pw + 0.005, draw_t + 0.30, 0.255, 0.34)

    # ---- caja CRUX: los dos 16 ----
    cy = pt + ph + 0.26
    crux = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.45), Inches(cy),
                              Inches(12.43), Inches(1.62))
    crux.fill.solid(); crux.fill.fore_color.rgb = BANNER_F
    crux.line.color.rgb = SLOT_E; crux.line.width = Pt(2.25); crux.shadow.inherit = False
    g.add_label(s, 0.6, cy + 0.08, 12.1, 0.34,
                [("OJO  ·  hay DOS  16  en  [30, 16, 10, 16]  y NO significan lo mismo", 15, True, SLOT_E)],
                align=PP_ALIGN.CENTER)
    g.add_label(s, 0.8, cy + 0.52, 5.7, 0.62,
                [("16  (2º eje) = CABEZAS", 13.5, True, INK),
                 ("cuántas 'lentes' miran el parche en paralelo", 11.5, False, GREY_T)], align=PP_ALIGN.CENTER)
    sep = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(6.665), Inches(cy + 0.50), Inches(0.012), Inches(0.66))
    sep.fill.solid(); sep.fill.fore_color.rgb = RGBColor(0xD8, 0xC4, 0xB0)
    sep.line.fill.background(); sep.shadow.inherit = False
    g.add_label(s, 6.85, cy + 0.52, 5.7, 0.62,
                [("16  (4º eje) = DIM / CABEZA", 13.5, True, INK),
                 ("el largo de la ficha de cada slot, por lente", 11.5, False, GREY_T)], align=PP_ALIGN.CENTER)
    g.add_label(s, 0.6, cy + 1.18, 12.1, 0.32,
                [("coinciden en 16 sólo por la config:   256 = 16 cabezas × 16   (256 = dimensión de la query)",
                  11.5, False, GREY_T)], align=PP_ALIGN.CENTER)

    # leyenda compacta (esquina inferior izq)
    g.add_label(s, 0.5, 7.16, 12.3, 0.26,
                [("[ . ] = forma del tensor    ·    naranja = slots aprendidos (fijos tras entrenar)    ·    lupa = 'lente' (cabeza)",
                  10.5, False, GREY_T)], align=PP_ALIGN.CENTER)
    return s


# ------------------------------------------------------------------ SLIDE 2: mecanismo
def slide_matching(prs):
    s = new_slide(prs)
    title(s, "QUÉ HACE UNA CABEZA:  el matching  ⟨q, S⟩",
          "cada 'lente' compara la query del parche con un slot — un producto interno = un puntaje de parecido")

    # ---- pipeline superior (3 chips + flechas) ----
    def chip(x, w, lines, fill, edge):
        return g.add_box(s, x, 1.28, w, 0.92, lines, fill, edge, lw=1.75)

    chip(0.55, 2.75, [("1 PARCHE  (CONCH)", 13, True, INK), ("z : [512]", 12, False, GREY_T)], TEAL_F, TEAL_E)
    chip(3.85, 3.30, [("QUERY del parche", 13, True, SLOT_E),
                      ("q = LN(W_q · z) : [256]", 12, False, INK),
                      ("256 = 16 lentes × 16", 10.5, False, GREY_T)], SLOT_F, SLOT_E)
    chip(7.70, 5.10, [("COMPARAR vs los 300 slots, por cada lente", 13, True, SLOT_E),
                      ("logits = ⟨q, S⟩ : [N, 30, 16, 10]   →   softmax = RUTEO", 11.5, False, INK)],
         BANNER_F, SLOT_E)
    g.add_arrow(s, 3.34, 1.62, 0.42, 0.30)
    g.add_arrow(s, 7.20, 1.62, 0.42, 0.30)

    # ---- ZOOM central: dentro de 1 lente ----
    zx, zy, zw, zh = 0.85, 2.55, 7.55, 3.30
    zoom = card(s, zx, zy, zw, zh, accent=SLOT_E)
    lens_icon(s, zx + 0.22, zy + 0.26, 0.34)
    g.add_label(s, zx + 0.70, zy + 0.22, zw - 0.9, 0.40,
                [("DENTRO DE 1 LENTE  ·  producto interno = un solo número", 14.5, True, SLOT_E)],
                align=PP_ALIGN.LEFT)

    # q_k : trozo de la query (1 lente) — fila de 16
    g.add_label(s, zx + 0.30, zy + 0.86, 2.0, 0.26, [("q  (1 lente)", 12.5, True, INK)], align=PP_ALIGN.LEFT)
    cells_grid(s, zx + 2.05, zy + 0.84, 8, 1, SLOT_F, SLOT_E, cw=0.20, ch=0.20)
    g.add_label(s, zx + 2.05 + 8 * 0.235 + 0.05, zy + 0.80, 1.6, 0.28,
                [("…  [16]", 13, True, SLOT_E)], align=PP_ALIGN.LEFT)

    # s : un slot — fila de 16
    g.add_label(s, zx + 0.30, zy + 1.36, 2.0, 0.26, [("s  (1 slot)", 12.5, True, INK)], align=PP_ALIGN.LEFT)
    cells_grid(s, zx + 2.05, zy + 1.34, 8, 1, SLOT_F2, SLOT_E, cw=0.20, ch=0.20)
    g.add_label(s, zx + 2.05 + 8 * 0.235 + 0.05, zy + 1.30, 1.6, 0.28,
                [("…  [16]", 13, True, SLOT_E)], align=PP_ALIGN.LEFT)

    # operación
    g.add_label(s, zx + 0.30, zy + 1.92, zw - 0.6, 0.30,
                [("⟨q, s⟩  =  Σ  q[i] · s[i]   →   1 número  (puntaje de parecido)", 14, True, INK)],
                align=PP_ALIGN.LEFT)

    # mini-ejemplo numérico chiquito (3 elementos)
    ex = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(zx + 0.30), Inches(zy + 2.32),
                            Inches(zw - 0.60), Inches(0.82))
    ex.fill.solid(); ex.fill.fore_color.rgb = RGBColor(0xF3, 0xF7, 0xF4)
    ex.line.color.rgb = RGBColor(0xBF, 0xD3, 0xC8); ex.line.width = Pt(1); ex.shadow.inherit = False
    g.add_label(s, zx + 0.45, zy + 2.38, zw - 0.9, 0.34,
                [("ejemplo (3 números):   q = [0.9, 0.1, 0.3]    s = [1.0, 0.0, 0.2]", 12.5, False, INK)],
                align=PP_ALIGN.LEFT)
    g.add_label(s, zx + 0.45, zy + 2.72, zw - 0.9, 0.34,
                [("⟨q, s⟩ = 0.9·1.0 + 0.1·0.0 + 0.3·0.2 = 0.96   → se parece MUCHO a este slot",
                  12.5, True, TEAL_E)], align=PP_ALIGN.LEFT)

    # ---- columna derecha: de 1 puntaje a todo el ruteo ----
    rx, rw = 8.75, 4.10
    rcard = card(s, rx, zy, rw, zh, accent=TEAL_E)
    g.add_label(s, rx + 0.18, zy + 0.20, rw - 0.36, 0.30, [("…Y SE REPITE", 13.5, True, TEAL_E)],
                align=PP_ALIGN.CENTER)
    g.add_label(s, rx + 0.20, zy + 0.66, rw - 0.40, 0.60,
                [("× 16 lentes  (todas las cabezas)", 12.5, True, INK),
                 ("× 300 slots  (30 expertos × 10)", 12.5, True, INK)], align=PP_ALIGN.LEFT)
    # barra de 300 puntajes
    g.add_label(s, rx + 0.20, zy + 1.42, rw - 0.40, 0.24, [("→ vector de puntajes [300]", 11.5, False, GREY_T)],
                align=PP_ALIGN.LEFT)
    cells_grid(s, rx + 0.22, zy + 1.70, 14, 1, SLOT_F2, SLOT_E, cw=0.205, ch=0.20, gap=0.0, lw=0.5)
    g.add_arrow(s, rx + rw / 2 - 0.13, zy + 2.02, 0.26, 0.26, down=True)
    g.add_label(s, rx + 0.20, zy + 2.36, rw - 0.40, 0.28,
                [("softmax  →  % que suman 1", 12.5, True, SLOT_E)], align=PP_ALIGN.CENTER)
    g.add_box(s, rx + 0.30, zy + 2.70, rw - 0.60, 0.44,
              [("\"a qué slots se parece el parche\"", 11.5, False, INK)], BANNER_F, SLOT_E, lw=1.5)

    # ---- nota inferior: conexión con el bloque MoE ----
    g.add_label(s, 0.5, 6.08, 12.3, 0.34,
                [("estos puntajes son los logits  ⟨q, S⟩  que entran al RUTEO del bloque MoE  (las dos softmax: reparto D ↓N  y  mezcla C ↓300)",
                  12, False, GREY_T)], align=PP_ALIGN.CENTER)
    g.add_label(s, 0.5, 7.16, 12.3, 0.26,
                [("naranja = query / slots (aprendidos)    ·    teal = puntajes / ruteo    ·    lupa = 1 lente (cabeza)",
                  10.5, False, GREY_T)], align=PP_ALIGN.CENTER)
    return s


def main():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    slide_anatomia(prs)
    slide_matching(prs)
    prs.save(DST)
    chk = Presentation(DST)
    print(f"OK  {os.path.relpath(DST, REPO)}  slides={len(chk.slides)}  "
          f"s1_shapes={len(chk.slides[0].shapes)}  s2_shapes={len(chk.slides[1].shapes)}")


if __name__ == "__main__":
    main()
