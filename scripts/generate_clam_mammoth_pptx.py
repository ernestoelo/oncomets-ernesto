#!/usr/bin/env python
"""generate_clam_mammoth_pptx.py — copia editable de Diagrama_CLAM.pptx con mammoth.

Entrega .pptx EDITABLE (no PNG), mismo molde que la referencia Diagrama_CLAM.pptx:
  - Slide 1: el diagrama CLAM original, con el bloque de la 1ª capa lineal
    (FULLY CONNECTED) resaltado y re-rotulado como MAMMOTH (MoE) — muestra DÓNDE va.
  - Slide 2 (nueva, lienzo limpio): zoom de QUÉ HACE mammoth (router + expertos de
    bajo rango → suma ponderada). Sin proceso de entrenamiento, sin nombres (convención
    de la presentación). Editable a mano en PowerPoint.

NOTA de arquitectura (codebase models_mammoth/): mammoth reemplaza la **1ª capa lineal**
de CLAM (el FULLY CONNECTED tras el extractor), NO el bloque final rotulado "CAPA LINEAL"
(que es el clasificador por clase). Por eso se modifica el FULLY CONNECTED.

Requiere python-pptx (instalado en /media/administrador/Storage1/sdonoso/clam_testing2/.pylibs).
Uso: PYTHONPATH=/media/administrador/Storage1/sdonoso/clam_testing2/.pylibs \
     /home/sdonoso/miniconda3/envs/clam_latest/bin/python scripts/generate_clam_mammoth_pptx.py
"""
import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

REPO = "/media/administrador/Storage1/sdonoso/clam_testing2/oncomets-ernesto"
SRC = os.path.join(REPO, "papers/presentations/Diagrama_CLAM.pptx")
DST = os.path.join(REPO, "papers/presentations/Diagrama_CLAM_mammoth.pptx")

TEAL_F = RGBColor(0xB8, 0xD4, 0xD9)   # teal de la referencia
ORA_F  = RGBColor(0xFC, 0xE5, 0xCD)   # highlight naranjo (mammoth = lo que cambia)
ORA_E  = RGBColor(0xE2, 0x72, 0x3B)
BLU_F  = RGBColor(0xEA, 0xF2, 0xFB)
BLU_E  = RGBColor(0x5B, 0x8F, 0xB9)
GREY_E = RGBColor(0x9A, 0xA6, 0xB2)
INK    = RGBColor(0x22, 0x22, 0x22)
GREY_T = RGBColor(0x55, 0x55, 0x55)
DIM_F  = RGBColor(0xE7, 0xF1, 0xF2)   # fill teal claro = callout de dimensiones (forma del tensor)
DIM_E  = RGBColor(0x2C, 0x7A, 0x8C)   # borde teal (= "features in/out" de la leyenda)
CARLITO = "Carlito"


def _add_runs(p, text, size, bold, color):
    """Renderiza `text` con mini-markup de sub/superíndices REALES (baseline OOXML,
    no caracteres Unicode → consistente para cualquier letra y limpio en PowerPoint
    y LibreOffice):  `_x` / `_(xx)` = subíndice ;  `^x` / `^(xx)` = superíndice.
    (No usar `_`/`^` literales en los strings fuente; ej. 'auto_rank' → 'automático'.)"""
    def emit(s, base=None):
        if not s:
            return
        r = p.add_run(); r.text = s
        r.font.name = CARLITO; r.font.bold = bold; r.font.color.rgb = color
        r.font.size = Pt(size * 0.74 if base is not None else size)
        if base is not None:
            r._r.get_or_add_rPr().set("baseline", str(base))
    i, n, buf = 0, len(text), ""
    while i < n:
        c = text[i]
        if c in "_^" and i + 1 < n:
            base = -25000 if c == "_" else 30000
            nxt = text[i + 1]
            if nxt == "(":
                j = text.find(")", i + 2)
                if j != -1:
                    emit(buf); buf = ""
                    emit(text[i + 2:j], base)
                    i = j + 1; continue
                buf += c; i += 1; continue
            emit(buf); buf = ""
            emit(nxt, base)
            i += 2; continue
        buf += c; i += 1
    emit(buf)


def set_text(shape, lines, align=PP_ALIGN.CENTER):
    """lines = [(text, size, bold, color), ...] -> párrafos (con markup _/^)."""
    tf = shape.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_left = tf.margin_right = Pt(2)
    tf.margin_top = tf.margin_bottom = Pt(1)
    for i, (txt, sz, bold, col) in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        _add_runs(p, txt, sz, bold, col)


def add_box(slide, l, t, w, h, lines, fill, edge, lw=1.5, shape=MSO_SHAPE.ROUNDED_RECTANGLE):
    sp = slide.shapes.add_shape(shape, Inches(l), Inches(t), Inches(w), Inches(h))
    sp.fill.solid(); sp.fill.fore_color.rgb = fill
    sp.line.color.rgb = edge; sp.line.width = Pt(lw)
    sp.shadow.inherit = False
    set_text(sp, lines)
    return sp


def add_arrow(slide, l, t, w, h, down=False):
    shp = MSO_SHAPE.DOWN_ARROW if down else MSO_SHAPE.RIGHT_ARROW
    a = slide.shapes.add_shape(shp, Inches(l), Inches(t), Inches(w), Inches(h))
    a.fill.solid(); a.fill.fore_color.rgb = RGBColor(0x8A, 0x8A, 0x8A)
    a.line.fill.background(); a.shadow.inherit = False
    return a


def edit_slide1(slide):
    """Resalta el bloque FULLY CONNECTED, lo re-rotula MAMMOTH (con dims) y agrega un
    callout matemático de la integración, consistente con el zoom de la slide 6."""
    # process box detrás del FC (geometría ~L4.85 T1.45)
    for sh in slide.shapes:
        try:
            if (sh.shape_type == 1 and 4.7 < Emu(sh.left).inches < 5.05
                    and 1.40 < Emu(sh.top).inches < 1.55 and 2.8 < Emu(sh.width).inches < 3.3):
                sh.fill.solid(); sh.fill.fore_color.rgb = ORA_F
                sh.line.color.rgb = ORA_E; sh.line.width = Pt(2.5)
        except Exception:
            pass
    # text box "FULLY CONNECTED" -> mammoth (rótulo limpio y legible; el detalle va en el callout)
    for sh in slide.shapes:
        if sh.has_text_frame and "FULLY CONNECTED" in sh.text:
            sh.text_frame.clear()
            set_text(sh, [
                ("MAMMOTH  (MoE)", 10.5, True, ORA_E),
                ("bajo rango · reemplaza la 1ª capa lineal", 7, False, GREY_T),
            ])
            break
    # callout de integración en la franja del título (que se eliminó en el deck), conectado al bloque
    co = add_box(slide, 0.95, 0.12, 3.58, 0.76,
                 [("INTEGRACIÓN DE MAMMOTH", 10.5, True, ORA_E),
                  ("la 1ª capa lineal (FC) → mezcla de expertos", 9, False, INK),
                  ("ℝ⁵¹²→ℝ⁵¹² · 30 expertos · rank 8 · #params ≈ FC", 8.5, False, GREY_T)],
                 RGBColor(0xFD, 0xF1, 0xE5), ORA_E, lw=2)
    co.adjustments[0] = 0.06
    add_connector(slide, 2.95, 0.88, 4.82, 1.50, color=ORA_E, w=1.75)


def add_connector(slide, x1, y1, x2, y2, color=RGBColor(0x8A, 0x8A, 0x8A), w=1.75, arrow=True):
    """Conector recto. arrow=True → punta de flecha (flujo); False → línea de callout."""
    from pptx.enum.shapes import MSO_CONNECTOR
    from pptx.oxml.ns import qn
    cxn = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT,
                                     Inches(x1), Inches(y1), Inches(x2), Inches(y2))
    cxn.line.color.rgb = color; cxn.line.width = Pt(w)
    cxn.shadow.inherit = False
    if arrow:
        ln = cxn.line._get_or_add_ln()
        tail = ln.makeelement(qn("a:tailEnd"), {"type": "triangle", "w": "med", "len": "med"})
        ln.append(tail)
    return cxn


def add_label(slide, l, t, w, h, lines, align=PP_ALIGN.LEFT):
    tb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    set_text(tb, lines, align=align)
    return tb


def add_heads_panel(slide, l, t, w, h):
    """Panel 'lentes paralelas' que explica las CABEZAS (multi-head): el MISMO query
    mirado por varios criterios en paralelo (textura/forma/densidad ... ×16) que luego
    se concatenan. Misma familia naranja del query (vs azul de los expertos)."""
    PANEL_F = RGBColor(0xFD, 0xF1, 0xE5)
    BLU_F2  = RGBColor(0xEA, 0xF2, 0xFB)
    GREYF   = RGBColor(0x7A, 0x86, 0x92)
    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                 Inches(l), Inches(t), Inches(w), Inches(h))
    box.fill.solid(); box.fill.fore_color.rgb = PANEL_F
    box.line.color.rgb = ORA_E; box.line.width = Pt(2); box.shadow.inherit = False
    add_label(slide, l + 0.12, t + 0.05, w - 0.24, 0.34,
              [("16 CABEZAS = 16 criterios", 13, True, ORA_E)], align=PP_ALIGN.CENTER)
    y0, gap = t + 0.46, 0.30
    for i, cr in enumerate(["textura", "forma", "densidad"]):
        yy = y0 + i * gap
        add_label(slide, l + 0.18, yy, 1.10, 0.24, [(cr, 11.5, False, INK)], align=PP_ALIGN.LEFT)
        bar = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                     Inches(l + 1.30), Inches(yy + 0.02),
                                     Inches(w - 1.52), Inches(0.18))
        bar.fill.solid(); bar.fill.fore_color.rgb = BLU_F2
        bar.line.color.rgb = BLU_E; bar.line.width = Pt(1); bar.shadow.inherit = False
    add_label(slide, l + 0.12, y0 + 3 * gap - 0.02, w - 0.24, 0.26,
              [("· · ·   (×16 lentes)", 11.5, True, GREYF)], align=PP_ALIGN.CENTER)
    add_label(slide, l + 0.10, t + h - 0.46, w - 0.20, 0.44,
              [("mismo query · 16 criterios", 10.5, False, GREY_T),
               ("→ se concatenan = 512", 10.5, False, GREY_T)], align=PP_ALIGN.CENTER)
    return box


def add_dim_callout(slide, x, y, w, h, lines, attach_y):
    """Callout de DIMENSIONES pegado a la derecha del bloque (estilo expansión CLAM,
    convención §5.c): entrada → operación → salida con la forma del tensor. La familia
    teal = 'forma del tensor / features in-out' de la leyenda. Conector fino al bloque."""
    box = add_box(slide, x, y, w, h, lines, DIM_F, DIM_E, lw=1.25)
    add_connector(slide, x - 0.07, attach_y, x, y + h / 2, color=DIM_E, w=1.0, arrow=False)
    return box


def build_slide2(prs):
    """Zoom CODE-ACCURATE de mammoth como GRAFO RAMIFICADO (árbol).

    Tronco (parche → query → ruteo por slots) que se ABRE a los expertos (fan-out) y se
    vuelve a CERRAR en la combinación (fan-in) → salida. Cada nodo lleva fórmula + dimensión;
    el detalle LoRA va en un callout matemático. Parámetros reales (models_mammoth/ + paquete
    pin, paper Shao et al. ICLR 2026): Wq 512→256, 16 cabezas, 30 expertos × 10 slots = 300,
    LoRA rank 8 (auto_rank), keep_slots=False.
    """
    blank = None
    for lay in prs.slide_layouts:
        if lay.name and "blank" in lay.name.lower():
            blank = lay; break
    blank = blank or prs.slide_layouts[-1]
    s = prs.slides.add_slide(blank)

    TEAL_E   = RGBColor(0x2C, 0x7A, 0x8C)
    BANNER_F = RGBColor(0xFD, 0xF1, 0xE5)
    FAN      = RGBColor(0x7A, 0x86, 0x92)

    # ---- TRONCO central (parche → query → ruteo) ----
    cx = 6.40
    tw = 3.7; tl = cx - tw / 2
    add_box(s, tl, 0.98, tw, 0.80,
            [("PARCHES DE LA SLIDE  (CONCH)", 16, True, INK),
             ("z : [N × 512]   (N parches × 512)", 14, False, GREY_T)],
            TEAL_F, TEAL_E)
    add_box(s, tl, 2.02, tw, 0.96,
            [("PROYECCIÓN A QUERY", 16, True, ORA_E),
             ("q = LN(W_q · z) : [N × 256]", 15, False, INK),
             ("W_q : 512 → 256   ·   256 = 16 cab × 16 →", 12, False, GREY_T)],
            ORA_F, ORA_E, lw=2)
    add_box(s, tl, 3.24, tw, 1.00,
            [("RUTEO POR SLOTS", 16, True, ORA_E),
             ("u_(es) = Σ_n D_n · q_n   →   [300 slots]", 14, False, INK),
             ("D = softmax_n ⟨q, S⟩   (sobre los N parches)", 12, False, GREY_T)],
            ORA_F, ORA_E, lw=2)

    # ---- conectores del tronco ----
    add_connector(s, cx, 1.78, cx, 2.02, w=2.0)
    add_connector(s, cx, 2.98, cx, 3.24, w=2.0)

    # ---- RAMIFICACIÓN a expertos (fan-out) ----
    ex_cx = [cx - 2.30, cx, cx + 2.30]      # 4.10, 6.40, 8.70
    ew, eh, ey = 1.98, 1.22, 4.62
    labels = ["EXPERTO 1", "EXPERTO 2", "EXPERTO 30"]
    add_label(s, cx - 2.3, 4.28, 4.6, 0.28,
              [("300 slots = 30 expertos × 10   ·   S : [30,16,10,16]", 12.5, True, FAN)],
              align=PP_ALIGN.CENTER)
    for xc, lab in zip(ex_cx, labels):
        add_connector(s, cx, 4.24, xc, ey, color=FAN, w=1.75)        # fan-out
        add_box(s, xc - ew / 2, ey, ew, eh,
                [(lab, 14, True, BLU_E), ("10 slots → [10 × 512]", 11.5, False, GREY_T),
                 ("o_e = (u_e·A)·B_e", 13, False, INK), ("rank 8  ·  ℝ^(512)", 11.5, False, GREY_T)],
                BLU_F, BLU_E, lw=1.75)
        add_connector(s, xc, ey + eh, cx, 6.18, color=FAN, w=1.75)   # fan-in
    # puntos suspensivos entre experto 2 y 30
    add_label(s, (ex_cx[1] + ew / 2), ey + 0.30, (ex_cx[2] - ew / 2) - (ex_cx[1] + ew / 2),
              0.5, [("· · ·", 28, True, FAN)], align=PP_ALIGN.CENTER)

    # ---- COMBINACIÓN (fan-in) → salida ----
    add_label(s, cx - 1.8, 5.84, 3.6, 0.26,
              [("C = softmax_(300) ⟨q, S⟩", 13, True, FAN)], align=PP_ALIGN.CENTER)
    add_box(s, tl, 6.18, tw, 0.92,
            [("COMBINACIÓN", 16, True, INK),
             ("h_n = Σ_(es) C_n · o_(es)   :   [N × 512]", 14.5, False, INK),
             ("C : [N,16,300]  softmax ↓ 300 slots  →  CLAM", 11, False, GREY_T)],
            TEAL_F, TEAL_E)

    # ---- callouts de DIMENSIONES por bloque (entrada → operación → salida), a la
    #      derecha del tronco, estilo expansión CLAM (convención §5.c) ----
    DXL, DCW = tl + tw + 0.07, 1.40          # x pegado al tronco · ancho del chip
    add_dim_callout(s, DXL, 2.06, DCW, 0.88,
                    [("[N, 512]", 10.5, False, GREY_T),
                     ("W_q : 512 → 256", 9.5, False, INK),
                     ("[N, 256]", 12.5, True, TEAL_E)], 2.50)
    add_dim_callout(s, DXL, 3.30, DCW, 0.88,
                    [("[N, 256]", 10.5, False, GREY_T),
                     ("⟨q, S⟩  ·  D ↓ N", 9.5, False, INK),
                     ("[30, 16, 10, 16]", 11.5, True, TEAL_E)], 3.74)
    add_dim_callout(s, DXL, 6.20, DCW, 0.88,
                    [("[300, 512]", 10.5, False, GREY_T),
                     ("C ↓ 300 slots", 9.5, False, INK),
                     ("[N, 512]", 12.5, True, TEAL_E)], 6.64)

    # ---- callout MATEMÁTICO: experto LoRA de bajo rango (izquierda, a la altura del fan) ----
    add_box(s, 0.28, 4.06, 2.74, 1.94,
            [("EXPERTO = LoRA bajo rango", 13, True, BLU_E),
             ("[30,16,10,16] → [300 × 512]", 11, True, DIM_E),
             ("A : 16 → 8   (compartido)", 12.5, False, INK),
             ("B_e : 8 → 32   (por experto)", 12.5, False, INK),
             ("16 cabezas → concat = 512", 11.5, False, GREY_T),
             ("rank r = 8 (automático)", 12, True, GREY_T)],
            RGBColor(0xF1, 0xF6, 0xFB), BLU_E, lw=1)
    add_connector(s, 3.02, 5.05, ex_cx[0] - ew / 2, 5.05, color=BLU_E, w=1.0, arrow=False)

    # ---- panel 'lentes paralelas': qué son las CABEZAS (derecha, simétrico al callout LoRA) ----
    #      (sin conector al tronco: ahora la lane derecha la ocupan los callouts de dims;
    #       el panel queda asociado por color naranja = query y por proximidad)
    add_heads_panel(s, 9.82, 2.22, 3.18, 2.20)

    # ---- banner: mismo presupuesto de parámetros (top-derecha) ----
    add_box(s, 9.55, 0.98, 3.45, 1.06,
            [("MISMO PRESUPUESTO", 14, True, ORA_E),
             ("DE PARÁMETROS", 14, True, ORA_E),
             ("#params(MoE) ≈ FC 512×512", 12.5, False, INK)],
            BANNER_F, ORA_E, lw=2.25)
    # ---- nota integración (top-izquierda) ----
    add_box(s, 0.35, 0.98, 3.45, 1.06,
            [("MAMMOTH = MoE bajo rango", 14, True, ORA_E),
             ("reemplaza la 1ª capa lineal (FC)", 13, False, INK),
             ("drop-in:  ℝ^(512) → ℝ^(512)", 12.5, False, GREY_T)],
            BANNER_F, ORA_E, lw=2.25)

    # ---- leyenda compacta (esquina inferior izquierda) ----
    add_label(s, 0.35, 6.5, 3.6, 0.7,
              [("teal = dimensiones (forma in → out)", 11, False, TEAL_E),
               ("naranjo = MAMMOTH (entrenable)", 11, False, ORA_E),
               ("azul = expertos  ·  [.] = forma del tensor", 11, False, BLU_E)], align=PP_ALIGN.LEFT)
    return s


def main():
    prs = Presentation(SRC)
    edit_slide1(prs.slides[0])
    build_slide2(prs)
    prs.save(DST)
    # verificación textual (sin imagen)
    chk = Presentation(DST)
    s1_has = any(sh.has_text_frame and "MAMMOTH" in sh.text for sh in chk.slides[0].shapes)
    print(f"OK  {os.path.relpath(DST, REPO)}  slides={len(chk.slides)}  "
          f"slide1_mammoth={s1_has}  slide2_shapes={len(chk.slides[1].shapes)}")


if __name__ == "__main__":
    main()
